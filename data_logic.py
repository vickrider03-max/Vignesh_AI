import csv
import hashlib
import html
import re
import time
import zipfile
from collections import Counter, OrderedDict
from io import BytesIO, StringIO
import xml.etree.ElementTree as ET

import docx
import openpyxl
import pandas as pd
import pdfplumber
from bs4 import BeautifulSoup
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


MAX_VECTOR_TEXT_CHARS = 250000


class CacheManager:
    """LRU cache manager for expensive operations with TTL support."""
    def __init__(self, max_size=50):
        self.cache = OrderedDict()
        self.max_size = max_size
        self.timestamps = {}

    def get(self, key, ttl_seconds=3600):
        if key not in self.cache:
            return None
        if key in self.timestamps:
            age = time.time() - self.timestamps[key]
            if age > ttl_seconds:
                del self.cache[key]
                del self.timestamps[key]
                return None
        self.cache.move_to_end(key)
        return self.cache[key]

    def set(self, key, value):
        if len(self.cache) >= self.max_size:
            oldest = next(iter(self.cache))
            del self.cache[oldest]
            self.timestamps.pop(oldest, None)
        self.cache[key] = value
        self.timestamps[key] = time.time()
        if len(self.cache) > 1:
            self.cache.move_to_end(key)

    def clear(self):
        self.cache.clear()
        self.timestamps.clear()


def get_file_hash(file_bytes):
    return hashlib.sha256(file_bytes).hexdigest()


def normalize_extracted_line(line):
    line = str(line or "").strip()
    line = re.sub(r"([a-z])([A-Z])", r"\1 \2", line)
    line = re.sub(r"([A-Za-z])(\d)", r"\1 \2", line)
    line = re.sub(r"(\d)([A-Za-z])", r"\1 \2", line)
    line = re.sub(r"\s+", " ", line)
    return line.strip()


def xml_text_content(xml_bytes):
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        soup = BeautifulSoup(xml_bytes, "xml")
        return soup.get_text("\n", strip=True)

    text_items = []
    for element in root.iter():
        tag_name = element.tag.split("}", 1)[-1].lower()
        if tag_name in {"p", "h", "span", "line-break", "tab"} and element.text:
            text_items.append(element.text.strip())
        if element.tail:
            text_items.append(element.tail.strip())
    return "\n".join(item for item in text_items if item)


def extract_text(file_name, file_bytes):
    text_parts = []
    bio = BytesIO(file_bytes)
    file_name_lower = file_name.lower()

    try:
        if file_name_lower.endswith(".pdf"):
            text_parts.extend(extract_pdf_content(bio))
        elif file_name_lower.endswith(".docx"):
            text_parts.extend(extract_docx_content(bio))
        elif file_name_lower.endswith(".doc"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy Word document"))
        elif file_name_lower.endswith(".pptx"):
            text_parts.extend(extract_pptx_content(bio))
        elif file_name_lower.endswith(".ppt"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy PowerPoint document"))
        elif file_name_lower.endswith(".xlsx"):
            text_parts.extend(extract_xlsx_content(bio))
        elif file_name_lower.endswith(".xls"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy Excel workbook"))
        elif file_name_lower.endswith(".csv"):
            text_parts.extend(extract_csv_content(bio))
        elif file_name_lower.endswith((".html", ".htm")):
            text_parts.extend(extract_html_content(bio))
        elif file_name_lower.endswith(".odt"):
            text_parts.extend(extract_odt_content(bio))
        elif file_name_lower.endswith(".rtf"):
            text_parts.extend(extract_rtf_content(bio))
        elif file_name_lower.endswith(".pages"):
            text_parts.extend(extract_pages_content(bio))
        elif file_name_lower.endswith((".txt", ".md", ".log", ".can", ".capl")):
            text_parts.append(("TEXT", bio.read().decode("utf-8", errors="ignore")))
        else:
            text_parts.append(("UNSUPPORTED", f"Unsupported file format: {file_name_lower}"))
    except Exception as e:
        text_parts.append(("ERROR", f"Error extracting content: {str(e)}"))

    combined_text = ""
    for content_type, content in text_parts:
        if content_type == "TEXT":
            combined_text += content + "\n"
        elif content_type == "TABLE":
            combined_text += f"\nTABLE:\n{content}\n"
        elif content_type == "IMAGE":
            combined_text += f"\n[IMAGE: {content}]\n"
        elif content_type == "EMBEDDED_IMAGE":
            combined_text += f"\n[EMBEDDED_IMAGE: {content}]\n"
        elif content_type == "METADATA":
            combined_text += f"\n{content}\n"
        elif content_type == "ERROR":
            combined_text += f"\nERROR: {content}\n"
        elif content_type == "UNSUPPORTED":
            combined_text += f"\n{content}\n"
    return combined_text.strip()


def extract_pdf_content(bio):
    content = []
    try:
        with pdfplumber.open(bio) as pdf:
            if pdf.metadata:
                metadata = [f"{key}: {value}" for key, value in pdf.metadata.items() if value]
                if metadata:
                    content.append(("METADATA", "PDF Metadata:\n" + "\n".join(metadata)))
            content.append(("METADATA", f"Total Pages: {len(pdf.pages)}"))
            for index, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    content.append(("TEXT", f"Page {index + 1} Text:\n{page_text}"))
    except Exception as e:
        content.append(("ERROR", f"PDF extraction failed: {str(e)}"))
    return content


def extract_docx_content(bio):
    content = []
    try:
        doc = docx.Document(bio)
        metadata = []
        if doc.core_properties.title:
            metadata.append(f"Title: {doc.core_properties.title}")
        if doc.core_properties.author:
            metadata.append(f"Author: {doc.core_properties.author}")
        if doc.core_properties.created:
            metadata.append(f"Created: {doc.core_properties.created}")
        if metadata:
            content.append(("METADATA", "Document Metadata:\n" + "\n".join(metadata)))

        current_heading = None
        current_lines = []
        table_count = 0

        def flush_section():
            nonlocal current_heading, current_lines
            if current_heading or current_lines:
                body = "\n\n".join(current_lines).strip()
                content.append(("TEXT", f"Heading: {current_heading}\n{body}" if current_heading else body))
            current_heading = None
            current_lines = []

        def is_heading(paragraph):
            text = paragraph.text.strip()
            style = (paragraph.style.name or "").lower() if paragraph.style else ""
            return bool(text) and ("heading" in style or style.startswith(("title", "subtitle")) or re.match(r"^\d+(?:\.\d+)*\s+.+", text))

        for element in doc.element.body:
            if element.tag.endswith("}p"):
                paragraph = Paragraph(element, doc)
                text = paragraph.text.strip()
                if not text:
                    continue
                if is_heading(paragraph):
                    flush_section()
                    current_heading = text
                else:
                    current_lines.append(text)
            elif element.tag.endswith("}tbl"):
                flush_section()
                table_count += 1
                table = Table(element, doc)
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                if rows:
                    content.append(("TABLE", f"Table {table_count}:\n" + "\n".join(rows)))
        flush_section()
        if table_count:
            content.append(("METADATA", f"Total Tables: {table_count}"))
    except Exception as e:
        content.append(("ERROR", f"DOCX extraction failed: {str(e)}"))
    return content


def extract_pptx_content(bio):
    content = []
    try:
        prs = Presentation(bio)
        content.append(("METADATA", f"Total Slides: {len(prs.slides)}"))
        for index, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
                if hasattr(shape, "table"):
                    rows = [
                        " | ".join(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    ]
                    if rows:
                        slide_content.append("Table:\n" + "\n".join(rows))
            if slide_content:
                content.append(("TEXT", f"Slide {index}:\n" + "\n\n".join(slide_content)))
    except Exception as e:
        content.append(("ERROR", f"PPTX extraction failed: {str(e)}"))
    return content


def extract_xlsx_content(bio):
    content = []
    try:
        wb = openpyxl.load_workbook(bio, data_only=True)
        content.append(("METADATA", f"Workbook contains {len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)}"))
        for sheet_name in wb.sheetnames:
            rows = []
            for row in wb[sheet_name].iter_rows(values_only=True):
                if any(cell for cell in row):
                    rows.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
            if rows:
                content.append(("TABLE", f"Sheet '{sheet_name}':\n" + "\n".join(rows)))
    except Exception as e:
        content.append(("ERROR", f"XLSX extraction failed: {str(e)}"))
    return content


def extract_excel_data(file_name, file_bytes):
    data = []
    try:
        if file_name.lower().endswith(".xlsx"):
            wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
            for sheet in wb:
                headers = None
                for index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if index == 0:
                        headers = list(row)
                    elif row and any(cell is not None for cell in row):
                        data.append(dict(zip(headers, row)))
    except Exception:
        data = []
    return data


def extract_csv_content(bio):
    content = []
    try:
        df = pd.read_csv(bio)
        content.append(("METADATA", f"CSV rows: {len(df)} columns: {len(df.columns)}"))
        if not df.empty:
            preview_df = df.fillna("").head(500)
            rows = [preview_df.columns.tolist()] + preview_df.values.tolist()
            table_text = "\n".join(" | ".join(map(str, row)) for row in rows)
            content.append(("TABLE", f"CSV Data:\n{table_text}"))
    except Exception as e:
        content.append(("ERROR", f"CSV extraction failed: {str(e)}"))
    return content


def extract_html_content(bio):
    content = []
    try:
        soup = BeautifulSoup(bio.read(), "html.parser")
        if soup.title:
            content.append(("METADATA", f"Title: {soup.title.string}"))
        text = soup.get_text(separator="\n")
        if text.strip():
            content.append(("TEXT", text))
    except Exception as e:
        content.append(("ERROR", f"HTML extraction failed: {str(e)}"))
    return content


def extract_odt_content(bio):
    content = []
    try:
        with zipfile.ZipFile(bio) as odt_zip:
            if "content.xml" not in odt_zip.namelist():
                return [("ERROR", "ODT content.xml was not found.")]
            soup = BeautifulSoup(odt_zip.read("content.xml"), "xml")
            text_blocks = [node.get_text(" ", strip=True) for node in soup.find_all(["text:h", "text:p"]) if node.get_text(" ", strip=True)]
            if text_blocks:
                content.append(("TEXT", "\n".join(text_blocks)))
    except Exception as e:
        content.append(("ERROR", f"ODT extraction failed: {str(e)}"))
    return content


def strip_rtf_to_text(rtf_text):
    text = re.sub(r"\\'[0-9a-fA-F]{2}", lambda m: bytes.fromhex(m.group(0)[2:]).decode("latin-1", errors="ignore"), rtf_text)
    text = re.sub(r"\\(par|line)\b", "\n", text)
    text = re.sub(r"\\tab\b", "\t", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("\\{", "{").replace("\\}", "}").replace("\\\\", "\\")
    text = re.sub(r"[{}]", "", text)
    return html.unescape(text).strip()


def extract_rtf_content(bio):
    try:
        raw = bio.read()
        text = raw.decode("utf-8", errors="ignore") or raw.decode("latin-1", errors="ignore")
        plain_text = strip_rtf_to_text(text)
        return [("TEXT", plain_text)] if plain_text else [("ERROR", "No readable text was found in the RTF file.")]
    except Exception as e:
        return [("ERROR", f"RTF extraction failed: {str(e)}")]


def extract_legacy_office_content(bio, label):
    try:
        raw = bio.read()
        decoded = raw.decode("utf-16le", errors="ignore") + "\n" + raw.decode("latin-1", errors="ignore")
        strings = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s.,;:!?()/_+\-]{3,}", decoded)
        cleaned = []
        seen = set()
        for value in strings:
            value = re.sub(r"\s+", " ", value).strip()
            if 4 <= len(value) <= 240 and value.lower() not in seen:
                seen.add(value.lower())
                cleaned.append(value)
            if len(cleaned) >= 1000:
                break
        if cleaned:
            return [("METADATA", f"{label}: recovered readable text using best-effort binary extraction."), ("TEXT", "\n".join(cleaned))]
        return [("ERROR", f"{label} text could not be recovered. Save/export as modern Office, PDF, RTF, or TXT for full analysis.")]
    except Exception as e:
        return [("ERROR", f"{label} extraction failed: {str(e)}")]


def extract_pages_content(bio):
    content = []
    try:
        with zipfile.ZipFile(bio) as pages_zip:
            readable_parts = [
                name for name in pages_zip.namelist()
                if name.lower().endswith((".xml", ".txt", ".html", ".xhtml"))
                and not name.lower().startswith(("metadata/", "quicklook/thumbnail"))
            ]
            blocks = []
            for name in readable_parts[:20]:
                part = pages_zip.read(name)
                if name.lower().endswith((".html", ".xhtml")):
                    text = BeautifulSoup(part, "html.parser").get_text("\n", strip=True)
                elif name.lower().endswith(".xml"):
                    text = xml_text_content(part)
                else:
                    text = part.decode("utf-8", errors="ignore")
                text = re.sub(r"\s+", " ", text).strip()
                if text:
                    blocks.append(f"{name}\n{text}")
            if blocks:
                content.append(("TEXT", "\n\n".join(blocks)))
            else:
                content.append(("ERROR", "No readable text preview was found in this Pages file. Export it as DOCX/PDF for full analysis."))
    except Exception as e:
        content.append(("ERROR", f"Pages extraction failed: {str(e)}"))
    return content


def get_document_asset_counts(file_name, file_bytes, extracted_text):
    file_name_lower = file_name.lower()
    page_count = 0
    table_count = 0
    image_count = 0
    if file_name_lower.endswith(".pdf"):
        page_match = re.search(r"Total Pages:\s*(\d+)", extracted_text)
        page_count = int(page_match.group(1)) if page_match else len(re.findall(r"Page \d+ Text:", extracted_text))
    elif file_name_lower.endswith(".pptx"):
        slide_match = re.search(r"Total Slides:\s*(\d+)", extracted_text)
        page_count = int(slide_match.group(1)) if slide_match else 0
    elif file_name_lower.endswith(".xlsx"):
        sheet_match = re.search(r"Workbook contains (\d+) sheets", extracted_text)
        page_count = int(sheet_match.group(1)) if sheet_match else 0
    table_count = len(re.findall(r"\b(?:Table|Sheet|CSV Data)\b", extracted_text))
    image_count = len(re.findall(r"\[(?:IMAGE|EMBEDDED_IMAGE):", extracted_text))
    return page_count, image_count, table_count


def get_column_counts(data, column):
    counts = {}
    for row in data:
        value = row.get(column)
        if value is not None:
            counts[value] = counts.get(value, 0) + 1
    return counts
