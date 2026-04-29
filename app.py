# -------------------------------
# IMPORTS
# -------------------------------
# Standard library and third-party imports for the application.
import html, re, hashlib, os, json, base64, pickle, zipfile, sqlite3
import importlib
import math
import uuid
import urllib.parse
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from datetime import datetime, timedelta
from difflib import SequenceMatcher
from io import BytesIO
from pytz import timezone
import time
from collections import OrderedDict

import docx, openpyxl, pdfplumber, streamlit as st
import streamlit.components.v1 as components
from docx.text.paragraph import Paragraph
from docx.table import Table
import pandas as pd
from openpyxl.styles import PatternFill
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image, ImageDraw, ImageFont
import plotly.express as px

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import HuggingFacePipeline
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_text_splitters import RecursiveCharacterTextSplitter

# -------------------------------
# GLOBAL VARIABLES & CONSTANTS
# -------------------------------
# Persistent storage for document previews across app reruns.
# PREVIEW_TOKENS maps tokens to metadata, PREVIEW_STORE holds file data.
PREVIEW_TOKENS = {}  # token -> {'file_name': str, 'timestamp': datetime}
PREVIEW_STORE = {}   # token -> file_dict

APP_DIR = os.path.dirname(os.path.abspath(__file__))
PREVIEW_DATA_FILE = os.path.join(APP_DIR, "preview_data.pkl")
PDF_PREVIEW_RESOLUTION = 100
PDF_PREVIEW_WINDOW = 25
PDF_ASSET_SCAN_PAGE_LIMIT = 10
MAX_VECTOR_TEXT_CHARS = 250000

# ===============================================
# CACHING SYSTEM FOR PERFORMANCE OPTIMIZATION
# ===============================================
# This caching layer dramatically improves load times by avoiding redundant processing

class CacheManager:
    """LRU cache manager for expensive operations with TTL support"""
    def __init__(self, max_size=50):
        self.cache = OrderedDict()
        self.max_size = max_size
        self.timestamps = {}
        
    def get(self, key, ttl_seconds=3600):
        if key not in self.cache:
            return None
        if key in self.timestamps:
            age = time.time() - self.timestamps[key]
            if age > ttl_seconds:
                del self.cache[key]
                del self.timestamps[key]
                return None
        # Move to end (most recently used)
        self.cache.move_to_end(key)
        return self.cache[key]
    
    def set(self, key, value):
        if len(self.cache) >= self.max_size:
            oldest = next(iter(self.cache))
            del self.cache[oldest]
            if oldest in self.timestamps:
                del self.timestamps[oldest]
        self.cache[key] = value
        self.timestamps[key] = time.time()
        if len(self.cache) > 1:
            self.cache.move_to_end(key)
    
    def clear(self):
        self.cache.clear()
        self.timestamps.clear()

# Global cache instances
FILE_TEXT_CACHE = CacheManager(max_size=100)
VECTOR_STORE_CACHE = CacheManager(max_size=20)
EXCEL_DATA_CACHE = CacheManager(max_size=50)
EMBEDDINGS_CACHE = CacheManager(max_size=200)

# Cache for file hashes to detect modifications
FILE_HASH_CACHE = {}

# ===============================================
# HASH-BASED CHANGE DETECTION
# ===============================================
def get_file_hash(file_bytes):
    """Generate SHA256 hash of file contents for change detection"""
    return hashlib.sha256(file_bytes).hexdigest()

def file_has_changed(file_name, file_bytes):
    """Check if file has been modified since last processing"""
    new_hash = get_file_hash(file_bytes)
    cache_key = f"{file_name}_hash"
    old_hash = FILE_HASH_CACHE.get(cache_key)
    FILE_HASH_CACHE[cache_key] = new_hash
    return old_hash != new_hash

# ===============================================
# PREVIEW PERSISTENCE HELPERS
# -------------------------------
# Functions to save/load preview data so document previews persist across Streamlit reruns.
# This allows users to share preview links that work even after app restart.
def load_preview_data():
    """Load preview data from file"""
    global PREVIEW_TOKENS, PREVIEW_STORE
    if os.path.exists(PREVIEW_DATA_FILE):
        try:
            with open(PREVIEW_DATA_FILE, "rb") as f:
                data = pickle.load(f)
            if not isinstance(data, dict):
                raise ValueError("preview data is not a dictionary")
            PREVIEW_TOKENS = data.get("tokens", {}) if isinstance(data.get("tokens", {}), dict) else {}
            PREVIEW_STORE = data.get("store", {}) if isinstance(data.get("store", {}), dict) else {}
        except Exception as e:
            backup_path = f"{PREVIEW_DATA_FILE}.corrupt.{datetime.now().strftime('%Y%m%d%H%M%S')}"
            try:
                os.replace(PREVIEW_DATA_FILE, backup_path)
                st.warning("Preview cache was corrupted and has been reset. Please open the document preview again from the sidebar.")
            except Exception:
                st.warning(f"Could not load preview data: {e}")
            PREVIEW_TOKENS = {}
            PREVIEW_STORE = {}

def save_preview_data():
    """Save preview data to file"""
    temp_file = None
    try:
        data = {
            "tokens": PREVIEW_TOKENS,
            "store": PREVIEW_STORE
        }
        preview_dir = os.path.dirname(PREVIEW_DATA_FILE) or "."
        os.makedirs(preview_dir, exist_ok=True)
        temp_file = f"{PREVIEW_DATA_FILE}.{os.getpid()}.{uuid.uuid4().hex}.tmp"
        with open(temp_file, "wb") as f:
            pickle.dump(data, f)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_file, PREVIEW_DATA_FILE)
    except Exception as e:
        # Streamlit reruns can overlap during upload/sidebar rendering. If an
        # atomic temp replace fails, keep the in-memory preview store working and
        # fall back to a direct write before showing a warning.
        try:
            with open(PREVIEW_DATA_FILE, "wb") as f:
                pickle.dump({
                    "tokens": PREVIEW_TOKENS,
                    "store": PREVIEW_STORE
                }, f)
        except Exception as fallback_error:
            st.warning(f"Could not save preview data: {fallback_error}")
    finally:
        if temp_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except Exception:
                pass

def cleanup_expired_preview_tokens():
    """Remove preview tokens older than 1 hour to prevent memory accumulation."""
    now = datetime.now()
    expired_tokens = []
    for token, data in PREVIEW_TOKENS.items():
        if now - data['timestamp'] > timedelta(hours=1):
            expired_tokens.append(token)
    
    for token in expired_tokens:
        del PREVIEW_TOKENS[token]
        if token in PREVIEW_STORE:
            del PREVIEW_STORE[token]
    
    if expired_tokens:
        save_preview_data()

def render_html_frame(html_content, height="content", width="stretch"):
    """Render inline HTML with Streamlit's supported components API."""
    if height == "content":
        height = 240
    if isinstance(height, int) and height < 1:
        height = 1
    component_width = None if width in (None, "stretch") else width
    components.html(str(html_content), width=component_width, height=height, scrolling=True)


# -------------------------------
# MERCEDES LOGO GENERATION
# -------------------------------
# Generates an animated GIF of the Mercedes-Benz logo using matplotlib.
# Used in the sidebar and login screen for branding.
@st.cache_data(show_spinner=False)
def get_needle_minimalist_logo():
    frames = []

    silver_grey = "#A0A0A0"
    star_light = "#DCDCDC"
    star_shadow = "#B8B8B8"
    canvas_size = 220
    center = canvas_size // 2
    radius = 86

    for angle_deg in range(360, 0, -15):
        image = Image.new("RGBA", (canvas_size, canvas_size), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image)
        raw_scale = math.cos(math.radians(angle_deg))
        flip_scale = raw_scale if abs(raw_scale) > 0.08 else (0.08 if raw_scale >= 0 else -0.08)

        ellipse_box = [
            center - int(radius * abs(flip_scale)),
            center - radius,
            center + int(radius * abs(flip_scale)),
            center + radius,
        ]
        draw.ellipse(ellipse_box, outline=silver_grey, width=4)

        for base_angle in [90, 210, 330]:
            angle = math.radians(base_angle)
            tip = (
                center + int(radius * 0.88 * math.cos(angle) * flip_scale),
                center - int(radius * 0.88 * math.sin(angle)),
            )
            side_l = (
                center + int(radius * 0.13 * math.cos(angle + 2.15) * flip_scale),
                center - int(radius * 0.13 * math.sin(angle + 2.15)),
            )
            side_r = (
                center + int(radius * 0.13 * math.cos(angle - 2.15) * flip_scale),
                center - int(radius * 0.13 * math.sin(angle - 2.15)),
            )
            c_l, c_r = (star_light, star_shadow) if flip_scale > 0 else (star_shadow, star_light)
            draw.polygon([(center, center), tip, side_l], fill=c_l)
            draw.polygon([(center, center), tip, side_r], fill=c_r)

        frames.append(image)

    gif_buf = BytesIO()
    if frames:
        frames[0].save(
            gif_buf,
            format='GIF',
            save_all=True,
            append_images=frames[1:],
            duration=80,
            loop=0,
            disposal=2
        )

    return base64.b64encode(gif_buf.getvalue()).decode('utf-8')


# -------------------------------
# STREAMLIT PAGE CONFIG
# -------------------------------
st.set_page_config(page_title="🧠 IntelliDoc AI ", layout="wide")

# Mobile viewport meta tag for proper scaling
st.markdown(
    """
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
    """,
    unsafe_allow_html=True,
)

# -------------------------------
# GLOBAL CSS STYLING
# -------------------------------
# Custom CSS to hide Streamlit branding, style buttons, and add animations.
# Applied globally to override default Streamlit appearance.
st.markdown(
    """
    <style>
        #MainMenu,
        header,
        footer,
        [data-testid="stToolbar"],
        [data-testid="stFooter"],
        [data-testid="stDecoration"],
        [data-testid="stStatusWidget"],
        .viewerBadge_container__1QSob,
        .css-1lsmgbg,
        .css-rtg1gx,
        .stAppDeployButton {
            display: none !important;
            visibility: hidden !important;
            opacity: 0 !important;
            height: 0 !important;
            width: 0 !important;
            pointer-events: none !important;
        }
    </style>
    <script>
        const hideStreamlitBranding = () => {
            const selectors = [
                '#MainMenu',
                'header',
                'footer',
                '[data-testid="stToolbar"]',
                '[data-testid="stFooter"]',
                '[data-testid="stDecoration"]',
                '[data-testid="stStatusWidget"]',
                '.viewerBadge_container__1QSob',
                '.css-1lsmgbg',
                '.css-rtg1gx',
                '.stAppDeployButton'
            ];

            selectors.forEach(selector => {
                document.querySelectorAll(selector).forEach(el => {
                    el.style.setProperty('display', 'none', 'important');
                    el.style.setProperty('visibility', 'hidden', 'important');
                    el.style.setProperty('opacity', '0', 'important');
                    el.style.setProperty('height', '0', 'important');
                    el.style.setProperty('width', '0', 'important');
                    el.style.setProperty('pointer-events', 'none', 'important');
                });
            });

            document.querySelectorAll('*').forEach(el => {
                try {
                    const text = (el.innerText || el.textContent || '').trim();
                    if (/made with streamlit/i.test(text) || /github/i.test(text)) {
                        el.style.setProperty('display', 'none', 'important');
                        el.style.setProperty('visibility', 'hidden', 'important');
                        el.style.setProperty('opacity', '0', 'important');
                        el.style.setProperty('height', '0', 'important');
                        el.style.setProperty('width', '0', 'important');
                        el.style.setProperty('pointer-events', 'none', 'important');
                    }
                } catch (err) {
                    // ignore inaccessible nodes
                }
            });

            document.querySelectorAll('img, svg').forEach(el => {
                try {
                    const src = (el.src || '') + (el.outerHTML || '');
                    if (/streamlit/i.test(src) || /github/i.test(src)) {
                        el.style.setProperty('display', 'none', 'important');
                        el.style.setProperty('visibility', 'hidden', 'important');
                        el.style.setProperty('opacity', '0', 'important');
                        el.style.setProperty('height', '0', 'important');
                        el.style.setProperty('width', '0', 'important');
                        el.style.setProperty('pointer-events', 'none', 'important');
                    }
                } catch (err) {
                    // ignore inaccessible nodes
                }
            });
        };

        const observer = new MutationObserver(hideStreamlitBranding);
        observer.observe(document.documentElement, { childList: true, subtree: true });
        hideStreamlitBranding();
        setInterval(hideStreamlitBranding, 1000);
    </script>
    """,
    unsafe_allow_html=True,
)

# Load preview data from file
load_preview_data()

# Clean up expired preview tokens on app start
cleanup_expired_preview_tokens()

try:
    logo_data = get_needle_minimalist_logo()
except Exception:
    logo_data = None

# -------------------------------
# ADDITIONAL CSS STYLING
# -------------------------------
# More CSS for dashboard grids, metric cards, loading animations, and responsive design.
st.markdown(
    """
    <style>
        :root {
            --primary: #e8f6ff;
            --secondary: #64748b;
            --background: #ffffff;
            --surface: #f8fafc;
            --text: #1e293b;
            --text-secondary: #64748b;
            --border: #e2e8f0;
            --success: #10b981;
            --warning: #f59e0b;
            --error: #ef4444;
            --button-bg: #e8f6ff;
            --button-hover: #d0e8f8;
            --button-text: #1e293b;
        }
        
        *, *::before, *::after {
            box-sizing: border-box;
        }

        html, body {
            min-width: 0;
            overflow-x: hidden;
            overflow-y: auto !important;
            height: auto !important;
            background: var(--background)!important;
            color: var(--text);
            transition: background 0.3s ease, color 0.3s ease;
        }

        div[role="main"], section.main, .stApp {
            min-width: 0;
            max-width: 1600px;
            width: 100%;
            margin: 0 auto;
        }

        .block-container {
            padding-left: clamp(1rem, 2vw, 1.75rem) !important;
            padding-right: clamp(1rem, 2vw, 1.75rem) !important;
            max-width: 1600px;
        }

        .stSidebar {
            min-width: clamp(220px, 18vw, 300px) !important;
            max-width: clamp(260px, 20vw, 380px) !important;
            width: clamp(220px, 18vw, 360px) !important;
        }

        .stSidebarNav {
            min-width: clamp(220px, 18vw, 280px) !important;
            width: 100% !important;
        }

        /* Main content positioning */
        .main .block-container {
            margin-left: 0 !important;
            margin-right: 0 !important;
            padding-left: 1rem !important;
            padding-right: 1rem !important;
            width: 100% !important;
            max-width: none !important;
        }

        /* Ensure main content is visible */
        section.main {
            display: block !important;
            visibility: visible !important;
            opacity: 1 !important;
        }

        /* Responsive spacing for smaller and larger screens */
        @media (min-width: 640px) {
            .main .block-container {
                padding-left: 1.5rem !important;
                padding-right: 1.5rem !important;
            }
        }

        @media (min-width: 900px) {
            .main .block-container {
                padding-left: 2rem !important;
                padding-right: 2rem !important;
            }
        }

        @media (min-width: 1200px) {
            .main .block-container {
                padding-left: 3rem !important;
                padding-right: 3rem !important;
            }
        }

        @media (max-width: 900px) {
            .stSidebar,
            .stSidebarNav {
                width: 100% !important;
                min-width: 0 !important;
                max-width: 100% !important;
            }
            .stSidebar {
                position: relative !important;
                transform: none !important;
            }
        }

        @media (min-width: 1024px) {
            .main .block-container {
                padding-left: 3rem !important;
                padding-right: 3rem !important;
            }
        }

        /* Streamlit page container fixes */
        .stApp,
        div[data-testid="stAppViewContainer"] {
            width: 100% !important;
            max-width: 100% !important;
            min-width: 0 !important;
            min-height: auto !important;
        }

        section.main,
        .main,
        .main .block-container,
        div[data-testid="stAppViewContainer"],
        div[data-testid="stMainContent"],
        div[data-testid="main"] {
            width: 100% !important;
            max-width: 100% !important;
            min-width: 0 !important;
            margin: 0 !important;
            padding: 0 !important;
            overflow-x: hidden !important;
            overflow-y: visible !important;
        }

        section.main {
            display: block !important;
            min-height: auto !important;
        }

        .main .block-container {
            max-width: 100% !important;
            width: 100% !important;
        }

        /* Make sure sidebar and content layout stays aligned */
        div[data-testid="stAppViewContainer"] > div {
            min-width: 0 !important;
            overflow-x: hidden !important;
        }

        .dashboard-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 1rem;
            margin: 1rem 0;
        }
        
        .metric-card {
            background: var(--surface);
            border: 1px solid var(--border);
            border-radius: 8px;
            padding: 1.5rem;
            text-align: center;
            transition: transform 0.2s ease, box-shadow 0.2s ease;
        }
        
        .metric-card:hover {
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }
        
        .card-label {
            font-size: 0.875rem;
            font-weight: 600;
            color: var(--text-secondary);
            text-transform: uppercase;
            letter-spacing: 0.05em;
            margin-bottom: 0.5rem;
        }
        
        .card-value {
            font-size: 1.5rem;
            font-weight: 700;
            color: var(--primary);
        }
        
        .spinner {
            width: 40px;
            height: 40px;
            border: 3px solid var(--border);
            border-top: 3px solid var(--primary);
            border-radius: 50%;
            animation: spin 1s linear infinite;
            margin: 20px auto;
        }
        
        .loading-overlay {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(255, 255, 255, 0.9);
            display: flex;
            justify-content: center;
            align-items: center;
            z-index: 9999;
            animation: fadeIn 0.3s ease;
        }
        
        .loading-content {
            text-align: center;
            padding: 2rem;
            background: var(--surface);
            border-radius: 12px;
            box-shadow: 0 8px 32px rgba(0,0,0,0.1);
            animation: slideUp 0.4s ease;
        }
        
        .loading-dots {
            display: inline-block;
            width: 20px;
            height: 20px;
            border-radius: 50%;
            background: var(--primary);
            animation: bounce 1.4s ease-in-out infinite both;
            margin: 0 2px;
        }
        
        .loading-dots:nth-child(1) { animation-delay: -0.32s; }
        .loading-dots:nth-child(2) { animation-delay: -0.16s; }
        
        /* Button styling - Ultra-aggressive Streamlit override */
        .stButton > button {
            background-color: #e8f6ff !important;
            color: #1e293b !important;
            border: 2px solid #c0dff0 !important;
            border-radius: 6px !important;
            padding: 0.5rem 1rem !important;
            font-weight: 600 !important;
            transition: all 0.2s ease !important;
            box-shadow: 0 2px 4px rgba(200, 230, 250, 0.2) !important;
        }
        
        .stButton > button:hover {
            background-color: #d0e8f8 !important;
            transform: translateY(-1px) !important;
            box-shadow: 0 4px 8px rgba(175, 215, 245, 0.3) !important;
            border-color: #a0c8e8 !important;
        }
        
        .stButton > button:active {
            transform: translateY(0) !important;
            box-shadow: 0 2px 4px rgba(200, 230, 250, 0.2) !important;
        }
        
        /* Override Streamlit's internal button styling */
        div.stButton > button,
        div[data-testid="stBaseButton"] button,
        button[kind="primary"],
        button[kind="secondary"],
        button[kind="tertiary"] {
            background-color: #e8f6ff !important;
            color: #1e293b !important;
            border: 2px solid #c0dff0 !important;
        }
        
        div.stButton > button:hover,
        div[data-testid="stBaseButton"] button:hover,
        button[kind="primary"]:hover,
        button[kind="secondary"]:hover,
        button[kind="tertiary"]:hover {
            background-color: #d0e8f8 !important;
        }
        
        /* Alternative Streamlit button selectors */
        [role="button"],
        [data-testid*="button"] {
            background-color: #e8f6ff !important;
            color: #1e293b !important;
        }
        
        [role="button"]:hover,
        [data-testid*="button"]:hover {
            background-color: #d0e8f8 !important;
        }
        
        /* Strip Streamlit theme blue and apply light blue */
        button {
            background-color: #e8f6ff !important;
            color: #1e293b !important;
            border-color: #c0dff0 !important;
        }
        
        button:hover {
            background-color: #d0e8f8 !important;
        }
        
        /* Smaller buttons for logout and reset */
        button[data-testid*="main_logout_btn"],
        button[data-testid*="reset_chat_selection"],
        button[data-testid*="reset_dashboard_selection"],
        button[data-testid*="reset_compare_selection"],
        button[data-testid*="reset_capl_selection"] {
            padding: 0.25rem 0.75rem !important;
            font-size: 0.875rem !important;
            min-width: auto !important;
            width: auto !important;
        }
        
        /* Animations */
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        
        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(10px); }
            to { opacity: 1; transform: translateY(0); }
        }
        
        @keyframes slideUp {
            from { 
                opacity: 0; 
                transform: translateY(20px); 
            }
            to { 
                opacity: 1; 
                transform: translateY(0); 
            }
        }
        
        @keyframes bounce {
            0%, 80%, 100% { 
                transform: scale(0);
                opacity: 0.5;
            } 
            40% { 
                transform: scale(1);
                opacity: 1;
            }
        }
        
        @keyframes pulse {
            0% { transform: scale(1); }
            50% { transform: scale(1.05); }
            100% { transform: scale(1); }
        }
        
        @keyframes slideInLeft {
            from { 
                opacity: 0; 
                transform: translateX(-30px); 
            }
            to { 
                opacity: 1; 
                transform: translateX(0); 
            }
        }
        
        @keyframes slideInRight {
            from { 
                opacity: 0; 
                transform: translateX(30px); 
            }
            to { 
                opacity: 1; 
                transform: translateX(0); 
            }
        }
        
        @keyframes glow {
            0% { box-shadow: 0 0 5px rgba(168, 216, 240, 0.3); }
            50% { box-shadow: 0 0 20px rgba(127, 197, 232, 0.6); }
            100% { box-shadow: 0 0 5px rgba(168, 216, 240, 0.3); }
        }
        
        @keyframes wiggle {
            0%, 100% { transform: rotate(0deg); }
            25% { transform: rotate(-3deg); }
            75% { transform: rotate(3deg); }
        }
        
        @keyframes float {
            0%, 100% { transform: translateY(0px); }
            50% { transform: translateY(-10px); }
        }
        
        .fade-in {
            animation: fadeIn 0.5s ease-out;
        }
        
        .slide-in-left {
            animation: slideInLeft 0.6s ease-out;
        }
        
        .slide-in-right {
            animation: slideInRight 0.6s ease-out;
        }
        
        .pulse {
            animation: pulse 2s infinite;
        }
        
        .glow {
            animation: glow 3s infinite;
        }
        
        .wiggle {
            animation: wiggle 1s ease-in-out;
        }
        
        .float {
            animation: float 3s ease-in-out infinite;
        }
        
        @media (max-width: 768px) {
            .dashboard-grid {
                grid-template-columns: 1fr;
            }
        }
        
        /* Hide Streamlit footer and GitHub icon */
        footer,
        body > footer,
        footer *,
        [data-testid="stFooter"],
        [data-testid="stFooter"] *,
        [data-testid="stDecoration"],
        [data-testid="stDecoration"] *,
        [data-testid="stToolbar"],
        [data-testid="stHeader"],
        [data-testid="stStatusWidget"],
        #MainMenu,
        #MainMenu *,
        a[href*="streamlit.io"],
        a[href*="github.com"],
        a[href*="github"],
        [title*="GitHub"],
        [aria-label*="GitHub"],
        [aria-label*="Streamlit"],
        [role="complementary"] img,
        [role="complementary"] svg,
        [style*="position: fixed"][style*="bottom: 0px"][style*="right: 0px"],
        [style*="position: fixed"][style*="bottom: 1px"][style*="right: 1px"],
        [style*="position: fixed"][style*="bottom: 10px"][style*="right: 10px"],
        [style*="position: fixed"][style*="bottom: 8px"][style*="right: 8px"],
        [data-testid="stDecoration"][style*="position: fixed"] {
            display: none !important;
            visibility: hidden !important;
            opacity: 0 !important;
            height: 0 !important;
            width: 0 !important;
            pointer-events: none !important;
        }
    </style>
    <script>
        const hideStreamlitFooter = () => {
            const selectors = [
                'footer',
                'footer *',
                '[data-testid="stFooter"]',
                '[data-testid="stDecoration"]',
                '[data-testid="stToolbar"]',
                '[data-testid="stHeader"]',
                '#MainMenu',
                'a[href*="streamlit.io"]',
                'a[href*="github.com"]',
                'a[href*="github"]',
                '[aria-label*="GitHub"]',
                '[aria-label*="Streamlit"]',
                '[title*="GitHub"]',
                '[title*="Streamlit"]'
            ];
            selectors.forEach(selector => {
                document.querySelectorAll(selector).forEach(el => {
                    el.style.setProperty('display', 'none', 'important');
                    el.style.setProperty('visibility', 'hidden', 'important');
                    el.style.setProperty('opacity', '0', 'important');
                    el.style.setProperty('height', '0', 'important');
                    el.style.setProperty('width', '0', 'important');
                    el.style.setProperty('pointer-events', 'none', 'important');
                });
            });

            document.querySelectorAll('*').forEach(el => {
                try {
                    const text = (el.innerText || el.textContent || '').trim();
                    const hasIcon = el.querySelector('img[src*="streamlit"], img[src*="github"], svg');
                    if (text.includes('Made with Streamlit') || text.includes('GitHub') || hasIcon) {
                        el.style.setProperty('display', 'none', 'important');
                        el.style.setProperty('visibility', 'hidden', 'important');
                        el.style.setProperty('opacity', '0', 'important');
                        el.style.setProperty('height', '0', 'important');
                        el.style.setProperty('width', '0', 'important');
                        el.style.setProperty('pointer-events', 'none', 'important');
                    }
                } catch (err) {
                    // ignore inaccessible nodes
                }
            });

            document.querySelectorAll('body *').forEach(el => {
                try {
                    const style = window.getComputedStyle(el);
                    if (style.position === 'fixed') {
                        const rect = el.getBoundingClientRect();
                        if (rect.bottom >= window.innerHeight - 90 && rect.right >= window.innerWidth - 90 && rect.width <= 90 && rect.height <= 90) {
                            if (el.querySelector('img, svg') || /Streamlit|GitHub/i.test(el.innerText || el.textContent || '')) {
                                el.style.setProperty('display', 'none', 'important');
                                el.style.setProperty('visibility', 'hidden', 'important');
                                el.style.setProperty('opacity', '0', 'important');
                                el.style.setProperty('height', '0', 'important');
                                el.style.setProperty('width', '0', 'important');
                                el.style.setProperty('pointer-events', 'none', 'important');
                            }
                        }
                    }
                } catch (err) {
                    // ignore inaccessible nodes
                }
            });
        };

        const footerObserver = new MutationObserver(() => hideStreamlitFooter());
        footerObserver.observe(document.body, { childList: true, subtree: true });
        hideStreamlitFooter();
        setInterval(hideStreamlitFooter, 1000);

        const applyLightButtonStyles = () => {
            const selectors = [
                'button',
                'input[type="button"]',
                'input[type="submit"]',
                '[role="button"]',
                '.stButton > button',
                'div[data-testid="stBaseButton"] button'
            ];
            document.querySelectorAll(selectors.join(',')).forEach(el => {
                el.style.setProperty('background-color', '#e8f6ff', 'important');
                el.style.setProperty('color', '#1e293b', 'important');
                el.style.setProperty('border-color', '#c0dff0', 'important');
                el.style.setProperty('border-style', 'solid', 'important');
                el.style.setProperty('border-width', '2px', 'important');
                el.style.setProperty('box-shadow', '0 2px 4px rgba(200, 230, 250, 0.2)', 'important');
            });
        };

        const buttonObserver = new MutationObserver(() => {
            applyLightButtonStyles();
        });
        buttonObserver.observe(document.body, { childList: true, subtree: true });
        applyLightButtonStyles();
    </script>
    """,
    unsafe_allow_html=True
)

# Ensure session state keys exist before rendering login status
for key, default_value in [
    ("is_authenticated", False),
    ("logged_in_username", ""),
    ("user_role", None),
    ("login_history", []),
    ("selected_files", []),
    ("file_texts", {}),
    ("vector_stores", {}),
    ("chat_file_selection", []),
    ("chat_summary_downloads", {"images": [], "tables": [], "csv": [], "diagrams": []}),
    ("chat_item_downloads", {"csv": [], "diagrams": []}),
    ("messages", []),
    ("welcome_shown", False),
    ("mobile_sidebar_visible", False),
]:
    if key not in st.session_state:
        st.session_state[key] = default_value


def render_status_strip():
    if not st.session_state.get("is_authenticated"):
        return

    # Get current elapsed time
    if 'start_time' not in st.session_state or st.session_state.start_time is None:
        st.session_state.start_time = time.time()
    
    elapsed = int(time.time() - st.session_state.start_time)
    hours, rem = divmod(elapsed, 3600)
    mins, secs = divmod(rem, 60)
    timer_str = f"{hours:02d}:{mins:02d}:{secs:02d}"

    username = st.session_state.get("logged_in_username") or "Vignesh"
    role = st.session_state.get("user_role") or "User"
    available_files = len(st.session_state.get("selected_files", []))

    # Create a placeholder for the live timer
    timer_placeholder = st.empty()
    
    # JavaScript for live timer
    live_timer_js = f"""
    <script>
        // Get the start time from Python
        var startTime = {st.session_state.start_time * 1000}; // Convert to milliseconds
        
        function updateTimer() {{
            var now = new Date().getTime();
            var elapsed = Math.floor((now - startTime) / 1000);
            
            var hours = Math.floor(elapsed / 3600);
            var minutes = Math.floor((elapsed % 3600) / 60);
            var seconds = elapsed % 60;
            
            var timerStr = 
                hours.toString().padStart(2, '0') + ':' +
                minutes.toString().padStart(2, '0') + ':' +
                seconds.toString().padStart(2, '0');
            
            // Update the timer display
            var timerElement = document.getElementById('live-timer');
            if (timerElement) {{
                timerElement.textContent = timerStr;
            }}
        }}
        
        // Update immediately and then every second
        updateTimer();
        setInterval(updateTimer, 1000);
    </script>
    """

    status_html = f"""
    <style>
        body {{
            margin: 0;
            background: transparent;
            font-family: 'Segoe UI', Tahoma, sans-serif;
        }}
        .dashboard-grid {{
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 8px;
            padding: 0 2px 4px;
        }}
        .metric-card {{
            position: relative;
            overflow: hidden;
            min-height: 58px;
            border-radius: 12px;
            padding: 8px 10px;
            border: 1px solid rgba(255, 255, 255, 0.5) !important;
            box-shadow: 0 6px 14px rgba(15, 23, 42, 0.12);
            transform: translateY(8px) scale(0.98);
            opacity: 0;
            animation: riseIn 0.65s ease-out forwards, cardFloat 4.5s ease-in-out infinite;
        }}
        .metric-card:nth-child(1) {{ animation-delay: 0.05s, 0.9s; }}
        .metric-card:nth-child(2) {{ animation-delay: 0.15s, 1.05s; }}
        .metric-card:nth-child(3) {{ animation-delay: 0.25s, 1.2s; }}
        .metric-card:nth-child(4) {{ animation-delay: 0.35s, 1.35s; }}
        .metric-card::before {{
            content: "";
            position: absolute;
            width: 100px;
            height: 100px;
            top: -40px;
            right: -20px;
            border-radius: 50%;
            background: rgba(255, 255, 255, 0.22);
            filter: blur(4px);
            animation: bubbleDrift 8s ease-in-out infinite;
        }}
        .metric-card::after {{
            content: "";
            position: absolute;
            width: 60px;
            height: 60px;
            bottom: -20px;
            left: -8px;
            border-radius: 50%;
            background: rgba(255, 255, 255, 0.16);
        }}
        .card-label, .card-value {{
            position: relative;
            z-index: 1;
            display: block;
        }}
        .card-label {{
            font-size: 0.58rem !important;
            font-weight: 700 !important;
            letter-spacing: 0.05em;
            text-transform: uppercase;
            color: rgba(15, 23, 42, 0.68) !important;
            margin-bottom: 4px;
        }}
        .card-value {{
            font-size: 1rem !important;
            font-weight: 800 !important;
            line-height: 1.2;
            word-break: break-word;
        }}
        #live-timer {{
            letter-spacing: 0.06em;
            animation: timerGlow 1.8s ease-in-out infinite;
        }}
        @keyframes riseIn {{
            from {{ opacity: 0; transform: translateY(8px) scale(0.98); }}
            to {{ opacity: 1; transform: translateY(0) scale(1); }}
        }}
        @keyframes cardFloat {{
            0%, 100% {{ transform: translateY(0); }}
            50% {{ transform: translateY(-3px); }}
        }}
        @keyframes bubbleDrift {{
            0%, 100% {{ transform: translate(0, 0); }}
            50% {{ transform: translate(-8px, 10px); }}
        }}
        @keyframes timerGlow {{
            0%, 100% {{ text-shadow: 0 0 0 rgba(46, 125, 50, 0); }}
            50% {{ text-shadow: 0 0 12px rgba(46, 125, 50, 0.25); }}
        }}
        @media (max-width: 900px) {{
            .dashboard-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }}
        }}
        @media (max-width: 560px) {{
            .dashboard-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
                gap: 8px;
                padding: 4px 2px 8px;
            }}
            .metric-card {{
                min-height: 74px;
                border-radius: 12px;
                padding: 9px 8px;
                box-shadow: 0 6px 14px rgba(15, 23, 42, 0.12);
            }}
            .card-label {{
                font-size: 0.54rem !important;
                letter-spacing: 0.04em;
                margin-bottom: 4px;
                white-space: normal;
            }}
            .card-value {{
                font-size: clamp(0.82rem, 4vw, 1rem) !important;
                line-height: 1.15;
                overflow-wrap: anywhere;
            }}
            #live-timer {{
                font-size: clamp(0.78rem, 3.8vw, 0.95rem) !important;
                letter-spacing: 0.02em;
            }}
        }}
    </style>
    <script>
        // Dynamically override button colors on page load and mutations
        function applyButtonColors() {{
            const buttons = document.querySelectorAll('button, [role="button"]');
            buttons.forEach(btn => {{
                btn.style.backgroundColor = '#e8f6ff !important';
                btn.style.color = '#1e293b !important';
                btn.style.borderColor = '#c0dff0 !important';
                btn.style.border = '2px solid #c0dff0 !important';
            }});
        }}
        
        // Apply on page load
        document.addEventListener('DOMContentLoaded', applyButtonColors);
        
        // Apply on mutation (when new buttons are added)
        const observer = new MutationObserver(applyButtonColors);
        observer.observe(document.body, {{ childList: true, subtree: true }});
        
        // Apply immediately
        applyButtonColors();
    </script>
    {live_timer_js}
    <div class="dashboard-grid">
        <div class="metric-card" style="background: linear-gradient(135deg, #e8eaf6 0%, #c5cae9 100%); color: #3c4f7e; border: 1px solid #e8eaf6;">
            <span class="card-label" style="color: #666;">👤 User</span>
            <span class="card-value" style="color: #3c4f7e;">{html.escape(username)}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #fce4ec 0%, #f8bbd9 100%); color: #7b1fa2; border: 1px solid #fce4ec;">
            <span class="card-label" style="color: #666;">🔑 Role</span>
            <span class="card-value" style="color: #7b1fa2;">{html.escape(str(role).title())}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #e3f2fd 0%, #bbdefb 100%); color: #1565c0; border: 1px solid #e3f2fd;">
            <span class="card-label" style="color: #666;">📁 Available Files</span>
            <span class="card-value" style="color: #1565c0;">{available_files}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #e8f5e8 0%, #c8e6c9 100%); color: #2e7d32; border: 1px solid #e8f5e8;">
            <span class="card-label" style="color: #666;">⏱️ Usage Time</span>
            <span class="card-value" id="live-timer" style="color: #2e7d32; font-family: 'Courier New', monospace;">{timer_str}</span>
        </div>
    </div>
    """

    render_html_frame(status_html, height=118)


def _help_state_key(tab_name):
    return f"show_help_popup_{tab_name}"


def _help_query_param_key(tab_name):
    return f"help_popup_{tab_name}"


def ensure_help_popup_state(tab_name):
    key = _help_state_key(tab_name)
    query_key = _help_query_param_key(tab_name)
    if key not in st.session_state:
        st.session_state[key] = False
    if query_key in query_params and query_params[query_key]:
        query_value = query_params[query_key]
        if isinstance(query_value, list):
            query_value = query_value[0] if query_value else ""
        st.session_state[key] = str(query_value).strip().lower() in {"1", "true", "yes", "open"}
    return key


def _set_query_params(params):
    try:
        if hasattr(st, "query_params"):
            st.query_params.clear()
            for param_key, param_value in params.items():
                if isinstance(param_value, list):
                    st.query_params[param_key] = [str(v) for v in param_value]
                else:
                    st.query_params[param_key] = str(param_value)
        elif hasattr(st, "experimental_set_query_params"):
            st.experimental_set_query_params(**params)
        elif hasattr(st, "set_query_params"):
            st.set_query_params(**params)
    except Exception:
        pass


def set_help_popup_state(tab_name, is_open):
    state_key = ensure_help_popup_state(tab_name)
    query_key = _help_query_param_key(tab_name)
    st.session_state[state_key] = is_open

    updated_params = {}
    try:
        for param_key in query_params.keys():
            param_value = query_params[param_key]
            if isinstance(param_value, list):
                updated_params[param_key] = list(param_value)
            else:
                updated_params[param_key] = param_value
    except Exception:
        updated_params = dict(query_params) if isinstance(query_params, dict) else {}

    if is_open:
        updated_params[query_key] = "1"
    else:
        updated_params.pop(query_key, None)

    _set_query_params(updated_params)

# ============================================
# SIMPLE HEADER - Moved higher for better visibility
# ============================================
if st.session_state.is_authenticated:
    st.markdown(
        """
        <style>
            section.main .block-container,
            .main .block-container,
            div[data-testid="stMain"] .block-container {
                padding-top: 4px !important;
            }
            div[data-testid="stVerticalBlock"] {
                gap: 0.25rem !important;
            }
            div[data-testid="stHorizontalBlock"]:has(.st-key-header_brain_icon) {
                margin-top: -0.5rem !important;
                margin-bottom: -0.2rem !important;
                align-items: center !important;
                min-height: 40px !important;
            }
            .app-header-title {
                transform: translateY(-2px);
                line-height: 1.1;
                margin: 0 !important;
            }
            .app-header-main {
                color: #1e293b;
                font-size: 1.18rem;
                font-weight: 700;
                margin: 0;
            }
            .app-header-subtitle {
                color: #64748b;
                font-size: 0.9rem;
                font-style: italic;
                margin-top: 0.04rem;
            }
            .st-key-header_brain_icon,
            .st-key-main_logout_btn {
                margin-top: -0.3rem !important;
            }
            .st-key-main_logout_btn {
                display: flex !important;
                align-items: center !important;
            }
            .compact-header-divider {
                height: 1px;
                background: #e2e8f0;
                margin: 2px 0 4px;
            }
        </style>
        """,
        unsafe_allow_html=True,
    )

    header_col, logout_col = st.columns([7, 1.15], vertical_alignment="center")

    with header_col:
        brain_col, title_col = st.columns([0.45, 7.55], vertical_alignment="center")
        with brain_col:
            if st.button("🧠", key="header_brain_icon", help="Click to show/hide helper tips"):
                helper_tab_map = {
                    "💬 Chat": "chat",
                    "📊 Dashboard": "dashboard",
                    "📂 Compare": "compare",
                    "📡 CAPL": "capl"
                }
                current_helper_tab = helper_tab_map.get(st.session_state.get("active_main_tab", "💬 Chat"), "chat")
                state_key = _help_state_key(current_helper_tab)
                st.session_state[state_key] = not st.session_state.get(state_key, False)
        with title_col:
            st.markdown(
                """
                <div class="app-header-title">
                    <div class="app-header-main">IntelliDoc AI</div>
                    <div class="app-header-subtitle">Smart Document Assistant</div>
                </div>
                """,
                unsafe_allow_html=True
            )

    with logout_col:
        if st.button("🚶 Logout", key="main_logout_btn"):
            now = datetime.now()
            ist_tz = timezone('Asia/Kolkata')
            ist_time = now.astimezone(ist_tz).strftime("%Y-%m-%d %H:%M:%S %Z")

            # Calculate usage time
            usage_seconds = 0
            if st.session_state.start_time is not None:
                usage_seconds = int(time.time() - st.session_state.start_time)

            hours, remainder = divmod(usage_seconds, 3600)
            minutes, seconds = divmod(remainder, 60)
            usage_time_str = f"{hours}h {minutes}m {seconds}s"

            st.session_state.login_history.append({
                "username": st.session_state.logged_in_username,
                "role": st.session_state.user_role,
                "action": "logout",
                "timestamp": ist_time,
                "usage_time": usage_time_str
            })
            active_file = "active_users.json"
            if os.path.exists(active_file):
                with open(active_file, "r") as f:
                    active_users = json.load(f)
                active_users = [u for u in active_users if u["username"] != st.session_state.logged_in_username]
                with open(active_file, "w") as f:
                    json.dump(active_users, f)
            goodbye_user = st.session_state.logged_in_username or "User"

            # Clear all user workspace state so the next login starts fresh.
            st.session_state.is_authenticated = False
            st.session_state.logged_in_username = ""
            st.session_state.user_role = None
            st.session_state.user_session_start_time = None
            st.session_state.start_time = None
            st.session_state.uploaded_files = []
            st.session_state.selected_files = []
            st.session_state.file_texts = {}
            st.session_state.excel_data_by_file = {}
            st.session_state.vector_stores = {}
            st.session_state.ask_messages = []
            st.session_state.extracted_images = {}
            st.session_state.chat_file_selection = []
            st.session_state.chat_summary_downloads = {"images": [], "tables": [], "csv": [], "diagrams": []}
            st.session_state.chat_item_downloads = {"csv": [], "diagrams": []}
            st.session_state.messages = []
            st.session_state.compare_file_selection = []
            st.session_state.compare_result_html = None
            st.session_state.compare_result_excel_bytes = None
            st.session_state.compare_result_files = []
            st.session_state.compare_semantic_summary = None
            st.session_state.file_dropdown = "--Select File--"
            st.session_state.dashboard_chart_type = "Pie Chart"
            st.session_state.dashboard_bar_orientation = "Vertical"
            st.session_state.selected_capl_file = "--Select CAPL file--"
            st.session_state.capl_last_analyzed_file = None
            st.session_state.capl_last_issues = None
            st.session_state.capl_editor_ai_fix = ""
            st.session_state.capl_autonomous_goal = ""
            st.session_state.capl_agent_result = ""
            st.session_state.agent_run_history = []
            st.session_state.last_streamed_assistant_index = None
            st.session_state.mobile_sidebar_visible = False
            st.session_state.llm_task = None
            st.session_state.welcome_shown = False
            st.session_state.behavior_tracker = {
                "chat": {"queries": 0, "actions": []},
                "dashboard": {"queries": 0, "actions": []},
                "compare": {"queries": 0, "actions": []},
                "capl": {"queries": 0, "actions": []}
            }
            for helper_tab in ["chat", "dashboard", "compare", "capl"]:
                st.session_state[_help_state_key(helper_tab)] = False
            st.session_state.workspace_memory = {
                "chat": [],
                "agent_runs": [],
                "indexed_files": [],
                "memory_events": [],
                "summary": {},
                "metadata": {},
            }
            st.session_state.workspace_memory_loaded = True
            st.session_state.file_uploader_key = int(st.session_state.get("file_uploader_key", 0)) + 1

            try:
                workspace_db_file = os.path.join(APP_DIR, "workspace_memory.db")
                conn = sqlite3.connect(workspace_db_file, check_same_thread=False)
                cursor = conn.cursor()
                cursor.execute(
                    """
                    CREATE TABLE IF NOT EXISTS workspace_meta (
                        meta_key TEXT PRIMARY KEY,
                        meta_value TEXT
                    )
                    """
                )
                cursor.execute(
                    "INSERT OR REPLACE INTO workspace_meta (meta_key, meta_value) VALUES (?, ?)",
                    ("workspace_memory", json.dumps(st.session_state.workspace_memory, default=str))
                )
                conn.commit()
                conn.close()
            except Exception:
                pass

            for cache in [FILE_TEXT_CACHE, VECTOR_STORE_CACHE, EXCEL_DATA_CACHE, EMBEDDINGS_CACHE]:
                cache.clear()
            FILE_HASH_CACHE.clear()
            PREVIEW_TOKENS.clear()
            PREVIEW_STORE.clear()
            save_preview_data()

            st.success(f"Goodbye, {goodbye_user}! You have been logged out.")
            st.rerun()

    st.markdown("<div class='compact-header-divider'></div>", unsafe_allow_html=True)

    render_html_frame(
        """
        <style>
            section.main > div:first-child {
                margin-top: -10px !important;
            }
            @keyframes brainGlow {
                0%, 100% {
                    box-shadow: 0 0 16px rgba(59, 130, 246, 0.45), 0 0 32px rgba(79, 70, 229, 0.22);
                    transform: scale(1);
                }
                25% {
                    box-shadow: 0 0 20px rgba(14, 165, 233, 0.55), 0 0 38px rgba(168, 85, 247, 0.20);
                    transform: scale(1.03);
                }
                50% {
                    box-shadow: 0 0 24px rgba(168, 85, 247, 0.65), 0 0 44px rgba(59, 130, 246, 0.28);
                    transform: scale(1.05);
                }
                75% {
                    box-shadow: 0 0 20px rgba(59, 130, 246, 0.55), 0 0 40px rgba(14, 165, 233, 0.24);
                    transform: scale(1.03);
                }
            }
            .header-brain-icon-large {
                animation: brainGlow 2.2s ease-in-out infinite;
                transition: transform 0.2s ease, box-shadow 0.2s ease, background-color 0.2s ease;
                background: rgba(219, 234, 254, 0.95) !important;
                color: #1d4ed8 !important;
                border: 1px solid rgba(59, 130, 246, 0.35) !important;
                width: 46px !important;
                height: 46px !important;
                min-width: 46px !important;
                min-height: 46px !important;
                padding: 0 !important;
                font-size: 1.55rem !important;
                line-height: 1 !important;
                display: inline-flex !important;
                align-items: center !important;
                justify-content: center !important;
                border-radius: 12px !important;
            }
            .header-brain-icon-large:hover {
                transform: scale(1.08);
                box-shadow: 0 0 30px rgba(79, 70, 229, 0.46), 0 0 52px rgba(14, 165, 233, 0.24);
            }
        </style>
        <script>
            const applyBrainIconStyles = () => {
                const root = window.parent ? window.parent.document : document;
                const buttons = Array.from(root.querySelectorAll('button'));
                buttons.forEach(btn => {
                    const title = (btn.getAttribute('title') || '').trim();
                    const text = (btn.innerText || '').trim().replace(/\\s+/g, '');
                    if (title === 'Click to show/hide helper tips' || text === '🧠') {
                        btn.classList.add('header-brain-icon-large');
                        btn.style.setProperty('font-size', '1.55rem', 'important');
                        btn.style.setProperty('padding', '0', 'important');
                        btn.style.setProperty('min-width', '46px', 'important');
                        btn.style.setProperty('min-height', '46px', 'important');
                        btn.style.setProperty('width', '46px', 'important');
                        btn.style.setProperty('height', '46px', 'important');
                        btn.style.setProperty('display', 'inline-flex', 'important');
                        btn.style.setProperty('align-items', 'center', 'important');
                        btn.style.setProperty('justify-content', 'center', 'important');
                        btn.style.setProperty('overflow', 'visible', 'important');
                        btn.style.setProperty('border-radius', '12px', 'important');
                        btn.style.setProperty('line-height', '1', 'important');
                        btn.style.setProperty('box-sizing', 'border-box', 'important');
                        btn.style.setProperty('transform', 'none', 'important');
                        btn.style.setProperty('background-color', 'rgba(219, 234, 254, 0.95)', 'important');
                        btn.style.setProperty('border', '1px solid rgba(59, 130, 246, 0.35)', 'important');
                        btn.style.setProperty('color', '#1d4ed8', 'important');
                        Array.from(btn.querySelectorAll('*')).forEach(child => {
                            child.style.setProperty('font-size', '1.55rem', 'important');
                            child.style.setProperty('line-height', '1', 'important');
                        });
                    }
                });
            };

            const brainObserver = new MutationObserver(() => applyBrainIconStyles());
            if (window.parent && window.parent.document) {
                brainObserver.observe(window.parent.document.body, { childList: true, subtree: true });
            }
            requestAnimationFrame(applyBrainIconStyles);
            setTimeout(applyBrainIconStyles, 300);
        </script>
        """,
        height=0,
    )

    if not st.session_state.get('welcome_shown', False):
        user = st.session_state.get("logged_in_username", "")
        if user:
            st.toast(f"Welcome back, {user} 🎉", icon="🙋")
        else:
            st.toast("Welcome back 🎉", icon="🙋")
        st.session_state.welcome_shown = True

    render_status_strip()


def render_mobile_workspace_controls():
    if not st.session_state.get("is_authenticated"):
        return

    show_sidebar = st.session_state.get("mobile_sidebar_visible", False)
    if show_sidebar:
        mobile_mode_css = """
        <style>
        @media (max-width: 767px) {
            section.main,
            [data-testid="stMain"],
            div[data-testid="stMain"] {
                display: none !important;
            }
            [data-testid="stSidebar"],
            .stSidebar {
                display: block !important;
                width: 100% !important;
                min-width: 0 !important;
                max-width: 100% !important;
                position: relative !important;
                transform: none !important;
                visibility: visible !important;
                opacity: 1 !important;
            }
            [data-testid="stSidebar"] > div {
                width: 100% !important;
                max-width: 100% !important;
            }
        }
        </style>
        """
    else:
        mobile_mode_css = """
        <style>
        @media (max-width: 767px) {
            [data-testid="stSidebar"],
            .stSidebar {
                display: none !important;
                visibility: hidden !important;
                width: 0 !important;
                min-width: 0 !important;
                max-width: 0 !important;
                transform: translateX(-100%) !important;
            }
            section.main,
            [data-testid="stMain"],
            div[data-testid="stMain"],
            div[data-testid="stAppViewContainer"] {
                display: block !important;
                width: 100% !important;
                max-width: 100% !important;
                margin-left: 0 !important;
                padding-left: 0 !important;
            }
        }
        </style>
        """

    st.markdown(
        mobile_mode_css
        + """
        <style>
        .st-key-mobile_show_files_btn,
        .st-key-mobile_open_workspace_btn {
            display: none;
        }
        @media (max-width: 767px) {
            .st-key-mobile_show_files_btn,
            .st-key-mobile_open_workspace_btn {
                display: block !important;
                margin-bottom: 0.75rem !important;
            }
            .st-key-mobile_show_files_btn button,
            .st-key-mobile_open_workspace_btn button {
                width: 100% !important;
                min-height: 46px !important;
                border-radius: 12px !important;
                border: 2px solid #93c5fd !important;
                background: #eff6ff !important;
                color: #1e3a8a !important;
                font-weight: 800 !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


render_mobile_workspace_controls()

if st.session_state.get("is_authenticated"):
    if st.button("📂 Files / Uploads", key="mobile_show_files_btn"):
        st.session_state.mobile_sidebar_visible = True
        st.rerun()


# -------------------------------
# SESSION STATE INITIALIZATION
# -------------------------------
for key in ["uploaded_files", "selected_files", "file_texts", "excel_data_by_file", "vector_stores", "messages",
            "ask_messages", "extracted_images"]:
    if key not in st.session_state:
        st.session_state[key] = [] if key in ["uploaded_files", "selected_files", "messages", "ask_messages"] else {}

if "capl_last_analyzed_file" not in st.session_state:
    st.session_state.capl_last_analyzed_file = None
if "capl_last_issues" not in st.session_state:
    st.session_state.capl_last_issues = None
if "capl_editor_name" not in st.session_state:
    st.session_state.capl_editor_name = "new_script.can"
if "capl_editor_code" not in st.session_state:
    st.session_state.capl_editor_code = """variables
{

}

on message *
{

}
"""
if "capl_editor_ai_fix" not in st.session_state:
    st.session_state.capl_editor_ai_fix = ""
if "is_authenticated" not in st.session_state:
    st.session_state.is_authenticated = False
if "logged_in_username" not in st.session_state:
    st.session_state.logged_in_username = ""
if "user_role" not in st.session_state:
    st.session_state.user_role = None
if "login_history" not in st.session_state:
    st.session_state.login_history = []
if "file_uploader_key" not in st.session_state:
    st.session_state.file_uploader_key = 0
if "compare_result_html" not in st.session_state:
    st.session_state.compare_result_html = None
if "compare_result_excel_bytes" not in st.session_state:
    st.session_state.compare_result_excel_bytes = None
if "compare_result_files" not in st.session_state:
    st.session_state.compare_result_files = []
if "chat_file_selection" not in st.session_state:
    st.session_state.chat_file_selection = []
if "chat_summary_downloads" not in st.session_state:
    st.session_state.chat_summary_downloads = {"images": [], "tables": [], "csv": [], "diagrams": []}
if "compare_file_selection" not in st.session_state:
    st.session_state.compare_file_selection = []
if "file_dropdown" not in st.session_state:
    st.session_state.file_dropdown = "--Select File--"
if "selected_capl_file" not in st.session_state:
    st.session_state.selected_capl_file = "--Select CAPL file--"
if "llm_task" not in st.session_state:
    st.session_state.llm_task = None
if "user_session_start_time" not in st.session_state:
    st.session_state.user_session_start_time = None
if "start_time" not in st.session_state:
    st.session_state.start_time = None
if "active_main_tab" not in st.session_state:
    st.session_state.active_main_tab = "💬 Chat"
if "workspace_memory" not in st.session_state:
    st.session_state.workspace_memory = {
        "chat": [],
        "agent_runs": [],
        "indexed_files": [],
        "memory_events": [],
        "summary": {},
        "metadata": {},
    }
if "workspace_memory_loaded" not in st.session_state:
    st.session_state.workspace_memory_loaded = False
if "capl_autonomous_goal" not in st.session_state:
    st.session_state.capl_autonomous_goal = ""
if "capl_agent_result" not in st.session_state:
    st.session_state.capl_agent_result = ""
if "agent_run_history" not in st.session_state:
    st.session_state.agent_run_history = []
if "last_streamed_assistant_index" not in st.session_state:
    st.session_state.last_streamed_assistant_index = None
if "compare_semantic_summary" not in st.session_state:
    st.session_state.compare_semantic_summary = None

# -------------------------------
# DOCUMENT PREVIEW FUNCTION
# -------------------------------
# Preview processing helpers:
# These functions support the standalone preview page and also provide extracted
# text/data reused by Chat, Dashboard, Compare, and CAPL when files are selected.
# ===============================================
# LAZY LOADING & PERFORMANCE OPTIMIZATION
# ===============================================

def lazy_load_file_section(file_name, section_id, loader_func):
    """Lazily load file sections on demand to reduce initial load time"""
    cache_key = f"{file_name}_{section_id}"
    cached_result = FILE_TEXT_CACHE.get(cache_key, ttl_seconds=7200)
    if cached_result is not None:
        return cached_result
    
    result = loader_func()
    FILE_TEXT_CACHE.set(cache_key, result)
    return result

def optimize_tab_rendering():
    """Optimize rendering by deferring non-active tab loading"""
    active_tab = st.session_state.get("active_main_tab", "💬 Chat")
    return active_tab

def ensure_file_processed(file_name):
    """Process file with caching to avoid redundant extraction"""
    file_info = get_uploaded_file_entry(file_name)
    if not file_info:
        return
    file_name_lower = file_name.lower()
    file_bytes = file_info["bytes"]
    new_hash = get_file_hash(file_bytes)
    hash_cache_key = f"{file_name}_hash"
    has_changed = FILE_HASH_CACHE.get(hash_cache_key) != new_hash
    FILE_HASH_CACHE[hash_cache_key] = new_hash
    
    # Check cache first
    cached_text = FILE_TEXT_CACHE.get(file_name)
    if cached_text is not None and not has_changed:
        st.session_state.file_texts[file_name] = cached_text
        if file_name_lower.endswith(".xlsx"):
            cached_excel = EXCEL_DATA_CACHE.get(file_name)
            if cached_excel is not None:
                st.session_state.excel_data_by_file[file_name] = cached_excel
        return
    
    # Process file if not in cache
    if file_name not in st.session_state.file_texts or has_changed:
        extracted_text = extract_text(file_name, file_bytes)
        st.session_state.file_texts[file_name] = extracted_text
        FILE_TEXT_CACHE.set(file_name, extracted_text)

    if file_name_lower.endswith(".xlsx") and (file_name not in st.session_state.excel_data_by_file or has_changed):
        excel_data = extract_excel_data(file_name, file_bytes)
        st.session_state.excel_data_by_file[file_name] = excel_data
        EXCEL_DATA_CACHE.set(file_name, excel_data)


def extract_text(file_name, file_bytes):
    """Extract comprehensive text content from various file formats including tables and metadata."""
    text_parts = []
    bio = BytesIO(file_bytes)
    file_name_lower = file_name.lower()
    
    try:
        if file_name_lower.endswith(".pdf"):
            text_parts.extend(extract_pdf_content(bio))
        elif file_name_lower.endswith(".docx"):
            text_parts.extend(extract_docx_content(bio))
        elif file_name_lower.endswith(".doc"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy Word document"))
        elif file_name_lower.endswith(".pptx"):
            text_parts.extend(extract_pptx_content(bio))
        elif file_name_lower.endswith(".ppt"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy PowerPoint document"))
        elif file_name_lower.endswith(".xlsx"):
            text_parts.extend(extract_xlsx_content(bio))
        elif file_name_lower.endswith(".xls"):
            text_parts.extend(extract_legacy_office_content(bio, "Legacy Excel workbook"))
        elif file_name_lower.endswith(".csv"):
            text_parts.extend(extract_csv_content(bio))
        elif file_name_lower.endswith((".html", ".htm")):
            text_parts.extend(extract_html_content(bio))
        elif file_name_lower.endswith(".odt"):
            text_parts.extend(extract_odt_content(bio))
        elif file_name_lower.endswith(".rtf"):
            text_parts.extend(extract_rtf_content(bio))
        elif file_name_lower.endswith(".pages"):
            text_parts.extend(extract_pages_content(bio))
        elif file_name_lower.endswith((".txt", ".md", ".log")):
            text_parts.append(("TEXT", bio.read().decode("utf-8", errors="ignore")))
        elif file_name_lower.endswith(".can"):
            text_parts.append(("TEXT", bio.read().decode("utf-8", errors="ignore")))
        else:
            text_parts.append(("UNSUPPORTED", f"Unsupported file format: {file_name_lower}"))
    
    except Exception as e:
        text_parts.append(("ERROR", f"Error extracting content: {str(e)}"))
    
    # Combine all extracted content
    combined_text = ""
    for content_type, content in text_parts:
        if content_type == "TEXT":
            combined_text += content + "\n"
        elif content_type == "TABLE":
            combined_text += f"\nTABLE:\n{content}\n"
        elif content_type == "IMAGE":
            combined_text += f"\n[IMAGE: {content}]\n"
        elif content_type == "EMBEDDED_IMAGE":
            combined_text += f"\n[EMBEDDED_IMAGE: {content}]\n"
        elif content_type == "METADATA":
            combined_text += f"\n{content}\n"
        elif content_type == "ERROR":
            combined_text += f"\nERROR: {content}\n"
        elif content_type == "UNSUPPORTED":
            combined_text += f"\n{content}\n"
    
    return combined_text.strip()


def extract_pdf_content(bio):
    """Extract searchable PDF text quickly.

    Full table/image extraction is intentionally deferred to preview/download
    helpers because scanning every page of a large manual is slow.
    """
    content = []
    try:
        with pdfplumber.open(bio) as pdf:
            # Add metadata
            if pdf.metadata:
                metadata = []
                for key, value in pdf.metadata.items():
                    if value:
                        metadata.append(f"{key}: {value}")
                if metadata:
                    content.append(("METADATA", "PDF Metadata:\n" + "\n".join(metadata)))
            
            content.append(("METADATA", f"Total Pages: {len(pdf.pages)}"))
            
            # Extract text from each page. Avoid table/image scans here so Chat
            # can load large manuals without blocking on expensive page parsing.
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    content.append(("TEXT", f"Page {i+1} Text:\n{page_text}"))
    
    except Exception as e:
        content.append(("ERROR", f"PDF extraction failed: {str(e)}"))
    
    return content


def table_to_png_bytes(table_data, title=None):
    """Render table rows as a PNG image and return the bytes."""
    try:
        font = ImageFont.load_default()
    except Exception:
        font = None

    padding = 10
    row_height = 22
    col_padding = 18

    # Normalize table data
    normalized_table = [[str(cell) for cell in row] for row in table_data]
    col_widths = []
    for col_idx in range(len(normalized_table[0])):
        col_width = max(len(row[col_idx]) for row in normalized_table) * 7 + col_padding
        col_widths.append(col_width)

    width = sum(col_widths) + padding * 2
    height = row_height * len(normalized_table) + padding * 2
    if title:
        height += row_height

    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    y = padding

    if title:
        draw.text((padding, y), title, fill="black", font=font)
        y += row_height

    for row in normalized_table:
        x = padding
        for col_idx, cell in enumerate(row):
            draw.text((x, y), cell, fill="black", font=font)
            x += col_widths[col_idx]
        y += row_height

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def image_bytes_to_png_bytes(image_bytes):
    """Convert an uploaded image to PNG bytes."""
    with Image.open(BytesIO(image_bytes)) as image:
        png_buffer = BytesIO()
        image.save(png_buffer, format="PNG")
        return png_buffer.getvalue()


def dataframe_to_table_rows(df):
    """Convert a dataframe to table rows suitable for PNG rendering."""
    safe_df = df.fillna("")
    rows = [list(map(str, safe_df.columns.tolist()))]
    rows.extend([list(map(str, row)) for row in safe_df.values.tolist()])
    return rows


def crop_pdf_region_to_png(page, bbox, resolution=150):
    """Crop a rectangular region from a PDF page and return it as PNG bytes."""
    cropped_page = page.crop(bbox)
    cropped_image = cropped_page.to_image(resolution=resolution)
    output = BytesIO()
    cropped_image.original.save(output, format="PNG")
    return output.getvalue()


def extract_docx_content(bio):
    """Extract text, tables, and images from DOCX."""
    content = []
    try:
        doc = docx.Document(bio)
        
        # Extract metadata
        core_props = doc.core_properties
        metadata = []
        if core_props.title:
            metadata.append(f"Title: {core_props.title}")
        if core_props.author:
            metadata.append(f"Author: {core_props.author}")
        if core_props.created:
            metadata.append(f"Created: {core_props.created}")
        if metadata:
            content.append(("METADATA", "Document Metadata:\n" + "\n".join(metadata)))
        
        # Extract embedded images for preview and downloads
        image_count = 0
        if 'extracted_images' not in st.session_state:
            st.session_state.extracted_images = {}
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    image_ext = image_part.content_type.split('/')[-1]
                    if image_ext == 'jpeg':
                        image_ext = 'jpg'
                    image_key = f"docx_image_{image_count}"
                    st.session_state.extracted_images[image_key] = {
                        'bytes': image_bytes,
                        'ext': image_ext,
                        'filename': f"image_{image_count}.{image_ext}"
                    }
                    content.append(("EMBEDDED_IMAGE", f"Embedded Image {image_count + 1}: {image_key}"))
                    image_count += 1
                except Exception as e:
                    content.append(("ERROR", f"Could not extract image {image_count + 1}: {e}"))
        if image_count > 0:
            content.append(("METADATA", f"Total Images: {image_count}"))

        # Walk the document body in order and preserve headings, tables, and text.
        current_section_title = None
        current_section_lines = []
        table_count = 0

        def flush_current_section():
            nonlocal current_section_title, current_section_lines
            if current_section_title or current_section_lines:
                if current_section_title:
                    section_text = "\n\n".join(current_section_lines).strip()
                    if section_text:
                        content.append(("TEXT", f"Heading: {current_section_title}\n{section_text}"))
                    else:
                        content.append(("TEXT", f"Heading: {current_section_title}"))
                else:
                    section_text = "\n\n".join(current_section_lines).strip()
                    if section_text:
                        content.append(("TEXT", section_text))
                current_section_title = None
                current_section_lines = []

        def is_docx_heading(para):
            text = para.text.strip()
            if not text:
                return False
            try:
                style_name = (para.style.name or "").lower()
            except Exception:
                style_name = ""
            if "heading" in style_name or style_name.startswith("title") or style_name.startswith("subtitle"):
                return True
            if re.match(r'^\d+(?:\.\d+)*\s+.+', text) and len(text) <= 120:
                return True
            return False

        for element in doc.element.body:
            if element.tag.endswith('}p'):
                paragraph = Paragraph(element, doc)
                paragraph_text = paragraph.text.strip()
                if not paragraph_text:
                    continue
                if is_docx_heading(paragraph):
                    flush_current_section()
                    current_section_title = paragraph_text
                else:
                    current_section_lines.append(paragraph_text)
            elif element.tag.endswith('}tbl'):
                flush_current_section()
                table_count += 1
                table = Table(element, doc)
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    content.append(("TABLE", f"Table {table_count}:\n{table_text}"))

        flush_current_section()

        if table_count > 0:
            content.append(("METADATA", f"Total Tables: {table_count}"))

        if not any(item[0] in ("TEXT", "TABLE", "EMBEDDED_IMAGE", "IMAGE", "METADATA") for item in content):
            paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            if paragraphs:
                content.append(("TEXT", "\n\n".join(paragraphs)))
    except Exception as e:
        content.append(("ERROR", f"DOCX extraction failed: {str(e)}"))
    return content


def extract_pptx_content(bio):
    """Extract text, tables, and images from PPTX."""
    content = []
    try:
        prs = Presentation(bio)
        
        content.append(("METADATA", f"Total Slides: {len(prs.slides)}"))
        
        # Extract and display images
        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_ext = shape.image.content_type.split('/')[-1]
                        if image_ext == 'jpeg':
                            image_ext = 'jpg'
                        
                        # Create a unique key for the image
                        image_key = f"pptx_image_{image_count}"
                        
                        # Store image data for display
                        if 'extracted_images' not in st.session_state:
                            st.session_state.extracted_images = {}
                        st.session_state.extracted_images[image_key] = {
                            'bytes': image_bytes,
                            'ext': image_ext,
                            'filename': f"slide_image_{image_count}.{image_ext}"
                        }
                        
                        content.append(("EMBEDDED_IMAGE", f"Slide Image {image_count + 1}: {image_key}"))
                        image_count += 1
                    except Exception as e:
                        content.append(("ERROR", f"Could not extract slide image {image_count + 1}: {e}"))
        
        if image_count > 0:
            content.append(("METADATA", f"Total Images: {image_count}"))
        
        # Extract text and tables from each slide
        for i, slide in enumerate(prs.slides):
            slide_content = []
            
            # Extract text shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            
            # Extract tables
            for shape in slide.shapes:
                if hasattr(shape, 'table'):
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    
                    if table_data:
                        table_text = "\n".join([" | ".join(row) for row in table_data])
                        slide_content.append(f"Table:\n{table_text}")
            
            if slide_content:
                content.append(("TEXT", f"Slide {i+1}:\n" + "\n\n".join(slide_content)))
    
    except Exception as e:
        content.append(("ERROR", f"PPTX extraction failed: {str(e)}"))
    
    return content


def extract_xlsx_content(bio):
    """Extract data from all sheets in XLSX."""
    content = []
    try:
        wb = openpyxl.load_workbook(bio, data_only=True)
        
        content.append(("METADATA", f"Workbook contains {len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)}"))
        
        # Extract data from each sheet
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_data = []
            
            # Get all rows with data
            for row in sheet.iter_rows(values_only=True):
                if any(cell for cell in row):  # Skip empty rows
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    sheet_data.append(row_data)
            
            if sheet_data:
                table_text = "\n".join([" | ".join(row) for row in sheet_data])
                content.append(("TABLE", f"Sheet '{sheet_name}':\n{table_text}"))
    
    except Exception as e:
        content.append(("ERROR", f"XLSX extraction failed: {str(e)}"))
    
    return content


def extract_csv_content(bio):
    """Extract CSV rows as table text for Chat, Compare, and Preview."""
    content = []
    try:
        df = pd.read_csv(bio)
        content.append(("METADATA", f"CSV rows: {len(df)} columns: {len(df.columns)}"))
        if not df.empty:
            preview_df = df.fillna("").head(500)
            table_text = "\n".join(
                " | ".join(map(str, row))
                for row in [preview_df.columns.tolist()] + preview_df.values.tolist()
            )
            content.append(("TABLE", f"CSV Data:\n{table_text}"))
    except Exception as e:
        content.append(("ERROR", f"CSV extraction failed: {str(e)}"))
    return content


def extract_html_content(bio):
    """Extract text and metadata from HTML."""
    content = []
    try:
        html_content = bio.read()
        soup = BeautifulSoup(html_content, "html.parser")
        
        # Extract title
        title = soup.title.string if soup.title else "No title"
        content.append(("METADATA", f"Title: {title}"))
        
        # Extract meta tags
        meta_info = []
        for meta in soup.find_all('meta'):
            if meta.get('name') and meta.get('content'):
                meta_info.append(f"{meta['name']}: {meta['content']}")
        if meta_info:
            content.append(("METADATA", "Meta Tags:\n" + "\n".join(meta_info)))
        
        # Count images
        images = soup.find_all('img')
        if images:
            content.append(("IMAGE", f"{len(images)} images found in HTML"))
        
        # Extract text content
        text = soup.get_text(separator="\n")
        if text.strip():
            content.append(("TEXT", text))
    
    except Exception as e:
        content.append(("ERROR", f"HTML extraction failed: {str(e)}"))
    
    return content


def xml_text_content(xml_bytes):
    """Extract readable text from XML-based office documents."""
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        soup = BeautifulSoup(xml_bytes, "xml")
        return soup.get_text("\n", strip=True)

    text_items = []
    for element in root.iter():
        tag_name = element.tag.split("}", 1)[-1].lower()
        if tag_name in {"p", "h", "span", "line-break", "tab"} and element.text:
            text_items.append(element.text.strip())
        if element.tail:
            text_items.append(element.tail.strip())
    return "\n".join(item for item in text_items if item)


def extract_odt_content(bio):
    """Extract text and simple tables from OpenDocument Text files."""
    content = []
    try:
        with zipfile.ZipFile(bio) as odt_zip:
            if "meta.xml" in odt_zip.namelist():
                meta_text = xml_text_content(odt_zip.read("meta.xml"))
                if meta_text.strip():
                    content.append(("METADATA", "ODT Metadata:\n" + meta_text[:2000]))

            if "content.xml" not in odt_zip.namelist():
                content.append(("ERROR", "ODT content.xml was not found."))
                return content

            content_xml = odt_zip.read("content.xml")
            soup = BeautifulSoup(content_xml, "xml")
            text_blocks = []
            for node in soup.find_all(["text:h", "text:p"]):
                text_value = node.get_text(" ", strip=True)
                if text_value:
                    text_blocks.append(text_value)
            if not text_blocks:
                fallback_text = xml_text_content(content_xml)
                if fallback_text.strip():
                    text_blocks.append(fallback_text)
            if text_blocks:
                content.append(("TEXT", "\n".join(text_blocks)))

            for table_index, table in enumerate(soup.find_all("table:table"), start=1):
                rows = []
                for row in table.find_all("table:table-row"):
                    cells = [cell.get_text(" ", strip=True) for cell in row.find_all("table:table-cell")]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    content.append(("TABLE", f"Table {table_index}:\n" + "\n".join(rows)))
    except Exception as e:
        content.append(("ERROR", f"ODT extraction failed: {str(e)}"))
    return content


def strip_rtf_to_text(rtf_text):
    """Best-effort RTF to plain text conversion without external dependencies."""
    text = rtf_text
    text = re.sub(r"\\'[0-9a-fA-F]{2}", lambda m: bytes.fromhex(m.group(0)[2:]).decode("latin-1", errors="ignore"), text)
    text = re.sub(r"\\(par|line)\b", "\n", text)
    text = re.sub(r"\\tab\b", "\t", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("\\{", "{").replace("\\}", "}").replace("\\\\", "\\")
    text = re.sub(r"[{}]", "", text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return html.unescape(text).strip()


def extract_rtf_content(bio):
    """Extract text from Rich Text Format documents."""
    content = []
    try:
        raw = bio.read()
        rtf_text = raw.decode("utf-8", errors="ignore")
        if not rtf_text.strip():
            rtf_text = raw.decode("latin-1", errors="ignore")
        plain_text = strip_rtf_to_text(rtf_text)
        if plain_text:
            content.append(("TEXT", plain_text))
        else:
            content.append(("ERROR", "No readable text was found in the RTF file."))
    except Exception as e:
        content.append(("ERROR", f"RTF extraction failed: {str(e)}"))
    return content


def extract_legacy_office_content(bio, label):
    """Best-effort text recovery for legacy binary Office files."""
    content = []
    try:
        raw = bio.read()
        decoded = raw.decode("utf-16le", errors="ignore") + "\n" + raw.decode("latin-1", errors="ignore")
        strings = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s.,;:!?()/_+\-]{3,}", decoded)
        cleaned = []
        seen = set()
        for value in strings:
            value = re.sub(r"\s+", " ", value).strip()
            if len(value) < 4 or len(value) > 240:
                continue
            if value.lower() in seen:
                continue
            seen.add(value.lower())
            cleaned.append(value)
            if len(cleaned) >= 1000:
                break

        if cleaned:
            content.append(("METADATA", f"{label}: recovered readable text using best-effort binary extraction."))
            content.append(("TEXT", "\n".join(cleaned)))
        else:
            content.append(("ERROR", f"{label} text could not be recovered. Save/export as DOCX, PPTX, XLSX, PDF, RTF, or TXT for full analysis."))
    except Exception as e:
        content.append(("ERROR", f"{label} extraction failed: {str(e)}"))
    return content


def extract_pages_content(bio):
    """Extract readable text from Apple Pages files when XML/text previews exist."""
    content = []
    try:
        with zipfile.ZipFile(bio) as pages_zip:
            names = pages_zip.namelist()
            readable_parts = [
                name for name in names
                if name.lower().endswith((".xml", ".txt", ".html", ".xhtml"))
                and not name.lower().startswith(("metadata/", "quicklook/thumbnail"))
            ]

            extracted_blocks = []
            for name in readable_parts[:20]:
                try:
                    part_bytes = pages_zip.read(name)
                    if name.lower().endswith((".html", ".xhtml")):
                        text_value = BeautifulSoup(part_bytes, "html.parser").get_text("\n", strip=True)
                    elif name.lower().endswith(".xml"):
                        text_value = xml_text_content(part_bytes)
                    else:
                        text_value = part_bytes.decode("utf-8", errors="ignore")
                    text_value = re.sub(r"\s+", " ", text_value).strip()
                    if text_value:
                        extracted_blocks.append(f"{name}\n{text_value}")
                except Exception:
                    pass

            if extracted_blocks:
                content.append(("TEXT", "\n\n".join(extracted_blocks)))
            else:
                content.append(("ERROR", "No readable text preview was found in this Pages file. Export it as DOCX/PDF for full analysis."))
    except zipfile.BadZipFile:
        content.append(("ERROR", "Pages extraction failed because this file is not a readable Pages ZIP package."))
    except Exception as e:
        content.append(("ERROR", f"Pages extraction failed: {str(e)}"))
    return content


def get_uploaded_file_entry(file_name):
    for file_info in st.session_state.uploaded_files:
        if file_info["name"] == file_name:
            return file_info
    return None


def create_preview_link(file_name, highlight_term=None, page_num=None):
    file_entry = get_uploaded_file_entry(file_name)
    if not file_entry:
        return None

    token = None
    for existing_token, token_data in list(PREVIEW_TOKENS.items()):
        if token_data.get("file_name") == file_name and existing_token in PREVIEW_STORE:
            token = existing_token
            token_data["timestamp"] = datetime.now()
            PREVIEW_STORE[existing_token] = file_entry
            break

    if token is None:
        token = str(uuid.uuid4())
        PREVIEW_TOKENS[token] = {'file_name': file_name, 'timestamp': datetime.now()}
        PREVIEW_STORE[token] = file_entry
        save_preview_data()

    params = [f"preview_token={token}"]
    if highlight_term:
        params.append(f"highlight={urllib.parse.quote_plus(highlight_term)}")
    if page_num is not None:
        params.append(f"page={urllib.parse.quote_plus(str(page_num))}")
    return "?" + "&".join(params)


def create_heading_anchor(text):
    anchor_text = str(text or "").strip().lower()
    anchor_text = re.sub(r'[^a-z0-9]+', '-', anchor_text)
    anchor_text = re.sub(r'-{2,}', '-', anchor_text).strip('-')
    if not anchor_text:
        anchor_text = 'preview'
    return f"heading-{anchor_text}"


def highlight_for_preview(text, highlight_term=None):
    if not highlight_term:
        return html.escape(text)
    escaped = html.escape(text)
    pattern = re.compile(re.escape(highlight_term), re.IGNORECASE)
    highlighted = pattern.sub(
        lambda m: f"<mark style='background:#fff3a3; padding:0 2px;'>{html.escape(m.group(0))}</mark>",
        escaped
    )
    return highlighted


def render_text_block(text, highlight_term=None, anchor_id=None):
    if not highlight_term:
        return None

    escaped = html.escape(text)
    pattern = re.compile(re.escape(highlight_term), re.IGNORECASE)
    first_match = True

    def replace_match(match):
        nonlocal first_match
        content = html.escape(match.group(0))
        if first_match and anchor_id:
            first_match = False
            return f"<span id='{anchor_id}'></span><mark style='background:#fff3a3; padding:0 2px;'>{content}</mark>"
        first_match = False
        return f"<mark style='background:#fff3a3; padding:0 2px;'>{content}</mark>"

    highlighted = pattern.sub(replace_match, escaped)
    return f"<pre style='white-space: pre-wrap; word-break: break-word; background:#f4f7fb; padding:12px; border-radius:8px; font-family: inherit;'>{highlighted}</pre>"


def render_document_preview(file_name, file_entry=None, highlight_term=None, highlight_page=None):
    """Render document preview with caching and error handling"""
    st.markdown(f"**Preview: {file_name}**")
    file_name_lower = file_name.lower()
    
    # Ensure file entry exists
    if file_entry is None:
        if not file_name_lower.endswith(".pdf"):
            ensure_file_processed(file_name)
        file_entry = get_uploaded_file_entry(file_name)
    else:
        # Keep PDF rendering lazy. Full PDF text extraction happens only when an
        # analysis panel explicitly asks for it.
        if file_name not in st.session_state.file_texts and not file_name_lower.endswith(".pdf"):
            extracted = extract_text(file_name, file_entry["bytes"])
            st.session_state.file_texts[file_name] = extracted
            FILE_TEXT_CACHE.set(file_name, extracted)
        if file_name.lower().endswith(".xlsx") and file_name not in st.session_state.excel_data_by_file:
            excel_data = extract_excel_data(file_name, file_entry["bytes"])
            st.session_state.excel_data_by_file[file_name] = excel_data
            EXCEL_DATA_CACHE.set(file_name, excel_data)
    
    if not file_entry:
        st.error("❌ File preview unavailable - file could not be loaded.")
        return

    image_download_items = []
    table_download_items = []

    # Special handling for PDF files: render actual page images for true preview
    if file_name_lower.endswith(".pdf"):
        try:
            pdf_bio = BytesIO(file_entry["bytes"])
            with pdfplumber.open(pdf_bio) as pdf:
                total_pages = len(pdf.pages)
                if total_pages == 0:
                    st.warning("⚠️ PDF has no readable pages")
                    return
                    
                st.markdown(f"**PDF Pages: {total_pages}**")
                
                # Determine which page to show
                selected_page_index = None
                if highlight_page is not None:
                    if 1 <= highlight_page <= total_pages:
                        selected_page_index = highlight_page - 1
                    else:
                        st.warning(f"⚠️ Page {highlight_page} not found. Showing all pages.")
                else:
                    # Large PDFs are rendered through the navigation controls below.
                    if total_pages > 10:
                        st.info(f"This PDF has {total_pages} pages. Use the controls below to preview any page.")
                        pages_to_show = []
                    else:
                        pages_to_show = list(range(total_pages))
                
                preview_key_base = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
                start_key = f"pdf_preview_start_{preview_key_base}"
                count_key = f"pdf_preview_count_{preview_key_base}"
                default_page = highlight_page if highlight_page and 1 <= highlight_page <= total_pages else 1
                if start_key not in st.session_state:
                    st.session_state[start_key] = default_page
                if count_key not in st.session_state:
                    st.session_state[count_key] = 1

                nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
                with nav_cols[0]:
                    if st.button(
                        "←",
                        key=f"pdf_prev_{preview_key_base}",
                        use_container_width=True,
                        help="Previous page",
                    ):
                        current_count = int(st.session_state.get(count_key, 1))
                        st.session_state[start_key] = max(1, int(st.session_state.get(start_key, 1)) - current_count)
                        st.rerun()
                with nav_cols[1]:
                    if st.button(
                        "→",
                        key=f"pdf_next_{preview_key_base}",
                        use_container_width=True,
                        help="Next page",
                    ):
                        current_count = int(st.session_state.get(count_key, 1))
                        next_start = int(st.session_state.get(start_key, 1)) + current_count
                        st.session_state[start_key] = min(total_pages, next_start)
                        st.rerun()
                with nav_cols[2]:
                    st.caption("Use page navigation for large PDFs. Only the selected window is rendered.")

                page_start = st.number_input(
                    "Start page",
                    min_value=1,
                    max_value=total_pages,
                    value=int(st.session_state.get(start_key, default_page)),
                    step=1,
                    key=start_key,
                )
                pages_to_render_count = st.slider(
                    "Pages to render on demand",
                    min_value=1,
                    max_value=min(PDF_PREVIEW_WINDOW, total_pages),
                    value=int(st.session_state.get(count_key, 1)),
                    key=count_key,
                    help="Default is one page for lazy loading. Increase only when you explicitly want a multi-page preview window.",
                )
                render_tables = st.checkbox(
                    "Detect tables on previewed pages",
                    value=False,
                    key=f"pdf_preview_tables_{preview_key_base}",
                    help="Table detection is slower, so it only runs when enabled.",
                )
                page_end = min(total_pages, int(page_start) + int(pages_to_render_count) - 1)
                pages_to_show = range(int(page_start) - 1, page_end)
                st.caption(f"Showing page {int(page_start)} to {page_end} of {total_pages}. Change the start page to preview any page.")

                # Render pages
                highlight_found = False
                for i in pages_to_show:
                    page = pdf.pages[i]
                    
                    try:
                        # Extract page text for highlighting
                        page_text = page.extract_text() or ""
                        page_anchor_id = None
                        
                        if highlight_term and highlight_term.lower() in page_text.lower():
                            page_anchor_id = create_heading_anchor(highlight_term)
                            highlight_found = True
                            if page_anchor_id:
                                st.markdown(f"<div id='{page_anchor_id}'></div>", unsafe_allow_html=True)
                        
                        # Render page image with caching
                        page_cache_key = f"{file_name}_page_{i}_image_{PDF_PREVIEW_RESOLUTION}"
                        cached_image = FILE_TEXT_CACHE.get(page_cache_key)
                        
                        if cached_image is None:
                            page_image = page.to_image(resolution=PDF_PREVIEW_RESOLUTION)
                            image_bytes_io = BytesIO()
                            page_image.original.save(image_bytes_io, format="PNG")
                            image_bytes = image_bytes_io.getvalue()
                            FILE_TEXT_CACHE.set(page_cache_key, image_bytes)
                        else:
                            image_bytes = cached_image
                        
                        st.image(image_bytes, caption=f"Page {i+1}", use_container_width=True)
                        
                        # Show highlighted text if match found
                        if page_anchor_id and highlight_term:
                            st.markdown("### Highlighted Text", unsafe_allow_html=True)
                            st.markdown(render_text_block(page_text, highlight_term, anchor_id=None), unsafe_allow_html=True)
                        
                        # Extract tables only when requested; this is slow on large manuals.
                        tables = page.extract_tables() if render_tables else []
                        if tables:
                            for j, table in enumerate(tables):
                                if table and any(any(cell for cell in row) for row in table):
                                    table_cache_key = f"{file_name}_page_{i}_table_{j}"
                                    cached_table = FILE_TEXT_CACHE.get(table_cache_key)
                                    
                                    if cached_table is None:
                                        table_png = table_to_png_bytes(table, title=f"Page {i+1} Table {j+1}")
                                        FILE_TEXT_CACHE.set(table_cache_key, table_png)
                                    else:
                                        table_png = cached_table
                                    
                                    st.image(table_png, caption=f"Page {i+1} Table {j+1}", use_container_width=True)
                                    table_download_items.append({
                                        "label": f"📥 Download Table {j+1} as PNG",
                                        "data": table_png,
                                        "file_name": f"{os.path.splitext(file_name)[0]}_page_{i+1}_table_{j+1}.png",
                                        "mime": "image/png",
                                        "key": f"download_pdf_table_{file_name}_{i}_{j}"
                                    })
                        
                        image_download_items.append({
                            "label": f"📥 Download Page {i+1} as PNG",
                            "data": image_bytes,
                            "file_name": f"{os.path.splitext(file_name)[0]}_page_{i+1}.png",
                            "mime": "image/png",
                            "key": f"download_pdf_page_{file_name}_{i}"
                        })
                    
                    except Exception as page_err:
                        st.warning(f"⚠️ Could not render page {i+1} as image: {str(page_err)[:100]}")
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            st.markdown(f"#### Page {i+1} Text")
                            st.code(page_text[:1000], language="text")
                
                # Download sections
                if image_download_items:
                    with st.expander("🖼️ Image Downloads", expanded=False):
                        for item in image_download_items[:10]:  # Limit to first 10
                            st.download_button(
                                label=item["label"],
                                data=item["data"],
                                file_name=item["file_name"],
                                mime=item["mime"],
                                key=item["key"]
                            )
                
                if table_download_items:
                    with st.expander("📊 Table Downloads", expanded=False):
                        for item in table_download_items[:10]:
                            st.download_button(
                                label=item["label"],
                                data=item["data"],
                                file_name=item["file_name"],
                                mime=item["mime"],
                                key=item["key"]
                            )
            
            return
        
        except Exception as pdf_err:
            st.error(f"❌ PDF rendering error: {str(pdf_err)[:200]}")
            st.info("Falling back to text-based document preview...")

    # Special handling for images
    if file_name_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
        try:
            st.image(file_entry["bytes"], caption=file_name, use_container_width=True)
            png_bytes = image_bytes_to_png_bytes(file_entry["bytes"])
            ext = file_name_lower.split('.')[-1]
            mime_type = f"image/{ext}"
            if ext == "jpg":
                mime_type = "image/jpeg"
            st.download_button(
                "📥 Download Image",
                png_bytes,
                file_name=file_name,
                mime=mime_type
            )
            return
        except Exception as img_err:
            st.error(f"❌ Could not display image: {str(img_err)[:100]}")

    # Fallback: text-based preview
    try:
        file_text = st.session_state.file_texts.get(file_name, "")
        if not file_text.strip():
            file_text = extract_text(file_name, file_entry["bytes"])
        
        if file_text.strip():
            preview_length = 2000
            preview_text = file_text[:preview_length]
            if len(file_text) > preview_length:
                preview_text += f"\n\n... ({len(file_text) - preview_length} more characters)"
            
            if highlight_term:
                st.markdown(render_text_block(preview_text, highlight_term), unsafe_allow_html=True)
            else:
                st.code(preview_text, language="text")
        else:
            st.info("📄 No readable content found in document.")
    
    except Exception as fallback_err:
        st.error(f"❌ Preview error: {str(fallback_err)[:100]}")
    # Special handling for Excel files (show as table)
    if file_name_lower.endswith(".xlsx"):
        with st.spinner("Loading preview..."):
            data = st.session_state.excel_data_by_file.get(file_name, [])
            if data:
                st.markdown("### Excel Preview")
                preview_df = pd.DataFrame(data).head(20)
                st.dataframe(preview_df, use_container_width=True, hide_index=True)
                st.download_button(
                    label="Download table as CSV",
                    data=preview_df.to_csv(index=False),
                    file_name=f"{os.path.splitext(file_name)[0]}_preview.csv",
                    mime="text/csv",
                    key=f"download_excel_csv_{file_name}"
                )
                try:
                    table_png = table_to_png_bytes(
                        dataframe_to_table_rows(preview_df),
                        title=f"{file_name} Preview"
                    )
                    st.download_button(
                        label="Download table as PNG",
                        data=table_png,
                        file_name=f"{os.path.splitext(file_name)[0]}_preview.png",
                        mime="image/png",
                        key=f"download_excel_png_{file_name}"
                    )
                except Exception:
                    pass
            else:
                st.info("No preview data available for this spreadsheet.")
        return

    # For all other files, show comprehensive extracted content
    with st.spinner("Loading preview..."):
        # Ensure text extraction has run
        if file_name not in st.session_state.file_texts:
            st.session_state.file_texts[file_name] = extract_text(file_name, file_entry["bytes"])
        
        full_content = st.session_state.file_texts.get(file_name, "")
        
        if not full_content.strip():
            st.warning("No content could be extracted from this file. This might indicate an issue with the file format or content extraction.")
            # Try to show basic file info
            file_size = len(file_entry["bytes"])
            st.info(f"File size: {file_size} bytes")
            st.info(f"File type: {file_name.split('.')[-1].upper()}")
            return

        # Parse the content into sections
        sections = parse_extracted_content(full_content)
        
        # For PDF files, if parsing doesn't produce good content, show raw text
        if file_name_lower.endswith('.pdf') and sections:
            # Check if we have meaningful text content
            has_meaningful_text = any(
                section_type == "TEXT" and section_content.strip()
                for section_type, _, section_content in sections
            )
            if not has_meaningful_text:
                # Fall back to showing raw content
                st.markdown("### Document Content")
                anchor_id = create_heading_anchor(highlight_term) if highlight_term else None
                with st.expander("Show extracted text", expanded=True):
                    if highlight_term:
                        st.markdown(render_text_block(full_content.strip(), highlight_term, anchor_id=anchor_id), unsafe_allow_html=True)
                    else:
                        st.code(full_content.strip(), language="text")
                return
        
        if not sections:
            st.warning("Content was extracted but could not be parsed into displayable sections.")
            anchor_id = create_heading_anchor(highlight_term) if highlight_term else None
            if highlight_term:
                st.markdown(render_text_block(full_content[:1000] + ("..." if len(full_content) > 1000 else ""), highlight_term, anchor_id=anchor_id), unsafe_allow_html=True)
            else:
                st.code(full_content[:1000] + ("..." if len(full_content) > 1000 else ""), language="text")
            return
        
        # Display each section
        for section_type, section_title, section_content in sections:
            if section_type == "METADATA":
                with st.expander(f"📋 {section_title}", expanded=False):
                    if section_content.strip():
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
                    else:
                        st.info("No metadata available")
            elif section_type == "TEXT":
                section_anchor_id = create_heading_anchor(section_title)
                st.markdown(f"<div id='{section_anchor_id}'></div>", unsafe_allow_html=True)
                st.markdown(f"### {section_title}")
                if section_content.strip():
                    if len(section_content) > 2000:
                        with st.expander("Show text content", expanded=False):
                            if highlight_term:
                                st.markdown(render_text_block(section_content, highlight_term, anchor_id=None), unsafe_allow_html=True)
                            else:
                                st.code(section_content, language="text")
                    else:
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term, anchor_id=None), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
                else:
                    st.info("No text content available for this section.")
            elif section_type == "TABLE":
                with st.expander(f"📊 {section_title}", expanded=True):
                    # Try to parse table and display as dataframe
                    try:
                        lines = section_content.strip().split('\n')
                        if lines:
                            # Parse table data
                            table_data = []
                            for line in lines:
                                if ' | ' in line:
                                    row = [cell.strip() for cell in line.split(' | ')]
                                    table_data.append(row)
                            
                            if table_data:
                                df = pd.DataFrame(table_data[1:] if len(table_data) > 1 else table_data, 
                                                columns=table_data[0] if len(table_data) > 1 else None)
                                st.dataframe(df, use_container_width=True, hide_index=True)
                                
                                # Add download button for table as CSV
                                csv_data = df.to_csv(index=False)
                                st.download_button(
                                    label="Download as CSV",
                                    data=csv_data,
                                    file_name=f"{section_title.replace(' ', '_')}.csv",
                                    mime="text/csv",
                                    key=f"download_table_{file_name}_{section_title}"
                                )
                                try:
                                    table_png = table_to_png_bytes(table_data, title=section_title)
                                    table_download_items.append({
                                        "label": f"Download {section_title} as PNG",
                                        "data": table_png,
                                        "file_name": f"{section_title.replace(' ', '_')}.png",
                                        "mime": "image/png",
                                        "key": f"download_table_png_{file_name}_{section_title}"
                                    })
                                except Exception:
                                    pass
                            else:
                                if highlight_term:
                                    st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                                else:
                                    st.code(section_content, language="text")
                    except Exception:
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
            elif section_type == "EMBEDDED_IMAGE":
                # Display embedded image
                image_key = section_content.split(": ")[-1] if ": " in section_content else section_content
                if 'extracted_images' in st.session_state and image_key in st.session_state.extracted_images:
                    image_data = st.session_state.extracted_images[image_key]
                    try:
                        st.image(image_data['bytes'], caption=image_data['filename'], use_container_width=True)
                        # Add download button
                        mime_type = f"image/{image_data['ext']}"
                        if image_data['ext'] == "jpg":
                            mime_type = "image/jpeg"
                        st.download_button(
                            label="Download Image",
                            data=image_data['bytes'],
                            file_name=image_data['filename'],
                            mime=mime_type,
                            key=f"download_embedded_{image_key}"
                        )
                        try:
                            png_bytes = image_bytes_to_png_bytes(image_data['bytes'])
                            st.download_button(
                                label="Download as PNG",
                                data=png_bytes,
                                file_name=f"{os.path.splitext(image_data['filename'])[0]}.png",
                                mime="image/png",
                                key=f"download_embedded_png_{image_key}"
                            )
                        except Exception:
                            pass
                    except Exception as e:
                        st.error(f"Could not display image: {e}")
                else:
                    st.info(f"🖼️ {section_content}")
            elif section_type == "IMAGE":
                st.info(f"🖼️ {section_content}")
            elif section_type == "ERROR":
                st.error(f"❌ {section_content}")
            elif section_type == "UNSUPPORTED":
                st.warning(f"⚠️ {section_content}")

        if image_download_items:
            with st.expander("🖼️ Image Downloads", expanded=False):
                for item in image_download_items:
                    st.download_button(
                        label=item["label"],
                        data=item["data"],
                        file_name=item["file_name"],
                        mime=item["mime"],
                        key=item["key"]
                    )
        if table_download_items:
            with st.expander("📊 Table Downloads", expanded=False):
                for item in table_download_items:
                    st.download_button(
                        label=item["label"],
                        data=item["data"],
                        file_name=item["file_name"],
                        mime=item["mime"],
                        key=item["key"]
                    )


def parse_extracted_content(content):
    """Parse the extracted content into sections for display."""
    sections = []
    lines = content.split('\n')
    current_section = None
    current_content = []

    def flush_section():
        nonlocal current_section, current_content
        if not current_section:
            return

        section_type, section_title, section_value = current_section
        if section_type in ("TEXT", "TABLE"):
            sections.append((section_type, section_title, '\n'.join(current_content).strip()))
        else:
            final_value = section_value
            if current_content:
                final_value = '\n'.join(current_content).strip()
            sections.append((section_type, section_title, final_value))

        current_section = None
        current_content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check for section markers
        if line.startswith('TABLE:'):
            # Save previous section
            flush_section()
            
            # Start new table section
            current_section = ("TABLE", "Table Content", "")
            current_content = []
            
        elif line.startswith('[IMAGE:'):
            # Save previous section
            flush_section()
            
            # Add image info
            sections.append(("IMAGE", "Images", line))
            current_section = None
            current_content = []
            
        elif line.startswith('[EMBEDDED_IMAGE:'):
            # Save previous section
            flush_section()
            
            # Add embedded image info
            image_info = line.replace('[EMBEDDED_IMAGE: ', '').replace(']', '')
            sections.append(("EMBEDDED_IMAGE", "Embedded Image", image_info))
            current_section = None
            current_content = []
            
        elif line.startswith('PDF Metadata:') or line.startswith('Document Metadata:') or line.startswith('Meta Tags:') or line.startswith('Title:') or 'Pages:' in line or 'Slides:' in line or 'sheets:' in line:
            # Save previous section
            flush_section()
            
            # Start metadata section
            if 'Metadata:' in line or 'Tags:' in line:
                section_title = line.split(':')[0] + " Information"
            else:
                section_title = "Document Information"
            
            current_section = ("METADATA", section_title, line)
            current_content = [line]
            
        elif line.startswith('Heading:'):
            # Save previous section
            flush_section()
            
            heading_title = line.replace('Heading:', '', 1).strip()
            current_section = ("TEXT", heading_title or "Heading", "")
            current_content = []

        elif line.startswith('Page ') and 'Text:' in line:
            # Save previous section
            flush_section()
            
            # Start new text section
            current_section = ("TEXT", f"Page {line.split()[1]} Content", "")
            current_content = []
            
        elif line.startswith('Slide ') and ':' in line:
            # Save previous section
            flush_section()
            
            # Start new slide section
            slide_num = line.split(':')[0]
            current_section = ("TEXT", f"{slide_num} Content", "")
            current_content = []
            
        elif line.startswith('Sheet ') and ':' in line:
            # Save previous section
            flush_section()
            
            # Start new sheet section
            sheet_name = line.split(':')[0].replace("'", "")
            current_section = ("TABLE", f"{sheet_name} Data", "")
            current_content = []
            
        elif current_section:
            # Add to current section
            if current_section[0] == "METADATA":
                current_content.append(line)
                current_section = (current_section[0], current_section[1], '\n'.join(current_content))
            else:
                current_content.append(line)
        else:
            # Start default text section
            if not current_section:
                current_section = ("TEXT", "Document Content", "")
                current_content = []
            current_content.append(line)
    
    # Save final section
    flush_section()
    
    # If no sections were created but we have content, create a default text section
    if not sections and content.strip():
        sections.append(("TEXT", "Document Content", content.strip()))
    
    return sections


@st.cache_data(show_spinner=False)
def build_summary_download_assets(file_name, file_bytes):
    file_name_lower = file_name.lower()
    image_items = []
    table_items = []

    try:
        if file_name_lower.endswith(".pdf"):
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for page_index, page in enumerate(pdf.pages[:PDF_ASSET_SCAN_PAGE_LIMIT], start=1):
                    for image_index, image_info in enumerate(page.images or [], start=1):
                        try:
                            bbox = (
                                image_info.get("x0", 0),
                                image_info.get("top", 0),
                                image_info.get("x1", 0),
                                image_info.get("bottom", 0)
                            )
                            if bbox[0] < bbox[2] and bbox[1] < bbox[3]:
                                image_items.append({
                                    "label": f"{file_name} - Page {page_index} Image {image_index}",
                                    "data": crop_pdf_region_to_png(page, bbox),
                                    "file_name": f"{os.path.splitext(file_name)[0]}_page_{page_index}_image_{image_index}.png",
                                    "mime": "image/png"
                                })
                        except Exception:
                            pass

                    for table_index, table in enumerate(page.extract_tables() or [], start=1):
                        if table and any(any(cell for cell in row) for row in table):
                            try:
                                table_items.append({
                                    "label": f"{file_name} - Page {page_index} Table {table_index}",
                                    "data": table_to_png_bytes(table, title=f"Page {page_index} Table {table_index}"),
                                    "file_name": f"{os.path.splitext(file_name)[0]}_page_{page_index}_table_{table_index}.png",
                                    "mime": "image/png"
                                })
                            except Exception:
                                pass

        elif file_name_lower.endswith(".docx"):
            doc = docx.Document(BytesIO(file_bytes))
            image_index = 1
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_bytes = rel.target_part.blob
                        image_items.append({
                            "label": f"{file_name} - Image {image_index}",
                            "data": image_bytes_to_png_bytes(image_bytes),
                            "file_name": f"{os.path.splitext(file_name)[0]}_image_{image_index}.png",
                            "mime": "image/png"
                        })
                        image_index += 1
                    except Exception:
                        pass

            for table_index, table in enumerate(doc.tables, start=1):
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if table_data:
                    try:
                        table_items.append({
                            "label": f"{file_name} - Table {table_index}",
                            "data": table_to_png_bytes(table_data, title=f"Table {table_index}"),
                            "file_name": f"{os.path.splitext(file_name)[0]}_table_{table_index}.png",
                            "mime": "image/png"
                        })
                    except Exception:
                        pass

        elif file_name_lower.endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            image_index = 1
            table_index = 1
            for slide_index, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_items.append({
                                "label": f"{file_name} - Slide {slide_index} Image {image_index}",
                                "data": image_bytes_to_png_bytes(shape.image.blob),
                                "file_name": f"{os.path.splitext(file_name)[0]}_slide_{slide_index}_image_{image_index}.png",
                                "mime": "image/png"
                            })
                            image_index += 1
                        except Exception:
                            pass
                    if hasattr(shape, "table"):
                        try:
                            table_data = [
                                [cell.text.strip() for cell in row.cells]
                                for row in shape.table.rows
                            ]
                            table_items.append({
                                "label": f"{file_name} - Slide {slide_index} Table {table_index}",
                                "data": table_to_png_bytes(table_data, title=f"Slide {slide_index} Table {table_index}"),
                                "file_name": f"{os.path.splitext(file_name)[0]}_slide_{slide_index}_table_{table_index}.png",
                                "mime": "image/png"
                            })
                            table_index += 1
                        except Exception:
                            pass

        elif file_name_lower.endswith(".xlsx"):
            data = extract_excel_data(file_name, file_bytes)
            if data:
                preview_df = pd.DataFrame(data).head(20)
                table_items.append({
                    "label": f"{file_name} - Preview Table",
                    "data": table_to_png_bytes(dataframe_to_table_rows(preview_df), title=f"{file_name} Preview"),
                    "file_name": f"{os.path.splitext(file_name)[0]}_preview.png",
                    "mime": "image/png"
                })

        elif file_name_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
            image_items.append({
                "label": f"{file_name} - Image",
                "data": image_bytes_to_png_bytes(file_bytes),
                "file_name": f"{os.path.splitext(file_name)[0]}.png",
                "mime": "image/png"
            })
    except Exception:
        return {"images": [], "tables": []}

    return {"images": image_items, "tables": table_items}


def render_extracted_assets_preview(file_name, file_entry):
    st.markdown(f"**Preview: {file_name}**")
    assets = build_summary_download_assets(file_name, file_entry["bytes"])
    image_items = assets.get("images", [])
    table_items = assets.get("tables", [])

    if not image_items and not table_items:
        st.info("No extractable images or tables were found in this file.")
        return

    if image_items:
        st.markdown("### Extracted Images")
        for index, item in enumerate(image_items):
            st.image(item["data"], caption=item["label"], use_container_width=True)
            st.download_button(
                label=f"Download {item['label']} as PNG",
                data=item["data"],
                file_name=item["file_name"],
                mime=item["mime"],
                key=f"preview_image_download_{index}_{item['file_name']}"
            )

    if table_items:
        st.markdown("### Extracted Tables")
        for index, item in enumerate(table_items):
            st.image(item["data"], caption=item["label"], use_container_width=True)
            st.download_button(
                label=f"Download {item['label']} as PNG",
                data=item["data"],
                file_name=item["file_name"],
                mime=item["mime"],
                key=f"preview_table_download_{index}_{item['file_name']}"
            )


def format_file_size(byte_count):
    """Return a readable file size label."""
    try:
        size = float(byte_count or 0)
        for unit in ["B", "KB", "MB", "GB"]:
            if size < 1024 or unit == "GB":
                return f"{size:.1f} {unit}" if unit != "B" else f"{int(size)} {unit}"
            size /= 1024
    except Exception:
        return "Unknown"


@st.cache_data(show_spinner=False)
def get_preview_metadata(file_name, file_bytes, extracted_text):
    """Collect lightweight metadata without rendering the whole document."""
    metadata = {
        "File name": file_name,
        "File type": os.path.splitext(file_name)[1].lower().lstrip(".").upper() or "Unknown",
        "File size": format_file_size(len(file_bytes or b"")),
        "Extracted text": f"{len(str(extracted_text or '')):,} characters",
        "Pages / Slides / Sheets": "Not available",
        "Tables": "0",
        "Images": "0",
    }

    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".pdf"):
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                metadata["Pages / Slides / Sheets"] = f"{len(pdf.pages)} pages"
        elif file_name_lower.endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            metadata["Pages / Slides / Sheets"] = f"{len(prs.slides)} slides"
        elif file_name_lower.endswith(".xlsx"):
            wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            metadata["Pages / Slides / Sheets"] = f"{len(wb.sheetnames)} sheets"
        elif file_name_lower.endswith(".docx"):
            doc = docx.Document(BytesIO(file_bytes))
            metadata["Tables"] = str(len(doc.tables))

        page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, extracted_text)
        if page_count and metadata["Pages / Slides / Sheets"] == "Not available":
            metadata["Pages / Slides / Sheets"] = str(page_count)
        metadata["Tables"] = str(max(int(metadata.get("Tables", "0") or 0), table_count))
        metadata["Images"] = str(image_count)
    except Exception:
        pass

    return metadata


def render_preview_metadata_cards(metadata):
    """Render compact metadata cards for the document viewer."""
    try:
        labels = list(metadata.items())
        cols = st.columns(3)
        for index, (label, value) in enumerate(labels):
            with cols[index % 3]:
                st.metric(label, value)
    except Exception as e:
        st.warning(f"Could not render preview metadata: {e}")


def chunk_preview_text(text, chunk_size=1200, overlap=180):
    """Split extracted text into chunks for search and Q&A."""
    try:
        clean_text = re.sub(r"\s+", " ", str(text or "")).strip()
        if not clean_text:
            return []
        chunks = []
        start = 0
        while start < len(clean_text):
            end = min(len(clean_text), start + chunk_size)
            chunk = clean_text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(clean_text):
                break
            start = max(0, end - overlap)
        return chunks
    except Exception:
        return []


@st.cache_data(show_spinner=False)
def keyword_search_preview_chunks(text, query, limit=8):
    """Fast fallback search for preview Q&A and Search tabs."""
    try:
        query_terms = [
            term.lower()
            for term in re.findall(r"[A-Za-z0-9_+\-/]{2,}", str(query or ""))
        ]
        if not query_terms:
            return []

        scored = []
        for chunk in chunk_preview_text(text):
            chunk_lower = chunk.lower()
            score = sum(chunk_lower.count(term) for term in query_terms)
            if score:
                scored.append((score, chunk))
        scored.sort(key=lambda item: item[0], reverse=True)
        return [chunk for _, chunk in scored[:limit]]
    except Exception:
        return []


def semantic_search_preview_chunks(file_name, text, query, limit=5):
    """Retrieve relevant chunks with FAISS when available, then fall back to keyword search."""
    try:
        if not str(query or "").strip() or not str(text or "").strip():
            return []
        vector_key = f"preview_vector::{file_name}"
        if vector_key not in st.session_state.vector_stores:
            st.session_state.vector_stores[vector_key] = create_vector_store(str(text)[:MAX_VECTOR_TEXT_CHARS])
        docs = st.session_state.vector_stores[vector_key].similarity_search(query, k=limit)
        return [doc.page_content for doc in docs if getattr(doc, "page_content", "")]
    except Exception:
        return keyword_search_preview_chunks(text, query, limit=limit)


def build_preview_answer(file_name, text, question):
    """Create an extractive answer from retrieved chunks."""
    try:
        chunks = semantic_search_preview_chunks(file_name, text, question, limit=5)
        if not chunks:
            return "No relevant information was found in the extracted text."
        answer_parts = []
        for chunk in chunks[:3]:
            sentences = re.split(r"(?<=[.!?])\s+", chunk)
            useful = [normalize_extracted_line(sentence) for sentence in sentences if len(sentence.strip()) > 30]
            answer_parts.extend(useful[:2])
        answer = "\n".join(f"- {part}" for part in answer_parts[:6])
        return answer or "\n\n".join(chunks[:2])
    except Exception as e:
        return f"Could not answer the question: {e}"


def extract_preview_tables(file_name, file_bytes, extracted_text):
    """Extract interactive tables from spreadsheets or pipe-delimited extracted text."""
    tables = []
    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".csv"):
            tables.append(("CSV Data", pd.read_csv(BytesIO(file_bytes))))
        elif file_name_lower.endswith(".xlsx"):
            workbook = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            for sheet_name, df in workbook.items():
                tables.append((sheet_name, df))
        elif file_name_lower.endswith(".xls"):
            try:
                workbook = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
                for sheet_name, df in workbook.items():
                    tables.append((sheet_name, df))
            except Exception:
                pass
        else:
            current_rows = []
            table_index = 1
            for line in str(extracted_text or "").splitlines():
                if " | " in line:
                    current_rows.append([cell.strip() for cell in line.split(" | ")])
                elif current_rows:
                    width = max(len(row) for row in current_rows)
                    normalized = [row + [""] * (width - len(row)) for row in current_rows]
                    tables.append((f"Extracted Table {table_index}", pd.DataFrame(normalized)))
                    table_index += 1
                    current_rows = []
            if current_rows:
                width = max(len(row) for row in current_rows)
                normalized = [row + [""] * (width - len(row)) for row in current_rows]
                tables.append((f"Extracted Table {table_index}", pd.DataFrame(normalized)))
    except Exception:
        pass
    return tables


def dataframe_to_xlsx_bytes(df):
    """Convert a dataframe to an XLSX download."""
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Data")
    return output.getvalue()


def ocr_image_best_effort(file_bytes):
    """Run optional OCR for image files when pytesseract is installed."""
    try:
        pytesseract = importlib.import_module("pytesseract")
        with Image.open(BytesIO(file_bytes)) as image:
            return pytesseract.image_to_string(image).strip()
    except Exception as e:
        return f"OCR is unavailable or failed: {e}"


def build_preview_summary_markdown(file_name, file_bytes, extracted_text):
    """Build a concise markdown summary for downloads."""
    try:
        plain_lines = [
            normalize_extracted_line(line)
            for line in str(extracted_text or "").splitlines()
            if 20 <= len(line.strip()) <= 220
        ]
        words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{2,}", str(extracted_text or ""))
        keyword_counts = Counter(
            word.lower()
            for word in words
            if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS
        )
        keywords = [word.title() for word, _ in keyword_counts.most_common(8)]

        sections = [
            f"# Document Intelligence Summary: {file_name}",
            "## Overview",
            f"- File type: {os.path.splitext(file_name)[1].lower() or 'unknown'}",
            f"- Extracted text length: {len(str(extracted_text or '')):,} characters",
            f"- Main themes: {', '.join(keywords) if keywords else 'Not enough readable text detected'}",
            "## Key Insights",
        ]
        key_lines = plain_lines[:5] or ["No reliable readable text was extracted."]
        sections.extend(f"- {line}" for line in key_lines)
        sections.extend([
            "## Practical Use",
            "- Use the Viewer tab for visual inspection.",
            "- Use Search and Q&A for targeted analysis over extracted chunks.",
            "- Use Tables, Images, and Downloads to export reusable assets.",
        ])
        return "\n".join(sections)
    except Exception as e:
        return f"# Summary unavailable\n\n{e}"


def render_paginated_text_document_viewer(file_name, extracted_text, page_size=9000):
    """Render long text-like documents with explicit pagination or full-load mode."""
    try:
        text = str(extracted_text or "").strip()
        if not text:
            st.info("No readable text was extracted for this document.")
            return

        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
        full_key = f"text_full_view_{base_key}"
        st.checkbox("Load full document", key=full_key, help="Shows all extracted text. For very large files this may be slower.")

        if st.session_state.get(full_key):
            st.info("Full document mode is active. No content is truncated.")
            st.text_area("Full extracted document", value=text, height=720, key=f"text_full_area_{base_key}")
            return

        total_sections = max(1, (len(text) + page_size - 1) // page_size)
        st.info("Preview mode: showing one section at a time. Use Load full document to view everything at once.")

        nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
        page_key = f"text_section_{base_key}"
        if page_key not in st.session_state:
            st.session_state[page_key] = 1
        with nav_cols[0]:
            if st.button("Previous section", key=f"text_prev_{base_key}", use_container_width=True):
                st.session_state[page_key] = max(1, int(st.session_state[page_key]) - 1)
                st.rerun()
        with nav_cols[1]:
            if st.button("Next section", key=f"text_next_{base_key}", use_container_width=True):
                st.session_state[page_key] = min(total_sections, int(st.session_state[page_key]) + 1)
                st.rerun()
        with nav_cols[2]:
            section_number = st.number_input(
                "Jump to section",
                min_value=1,
                max_value=total_sections,
                value=int(st.session_state[page_key]),
                key=f"text_jump_{base_key}",
            )
            st.session_state[page_key] = int(section_number)

        start = (int(st.session_state[page_key]) - 1) * page_size
        end = min(len(text), start + page_size)
        st.caption(f"Showing section {st.session_state[page_key]} of {total_sections}. Characters {start + 1:,}-{end:,} of {len(text):,}.")
        st.text_area("Document section", value=text[start:end], height=620, key=f"text_section_area_{base_key}_{st.session_state[page_key]}")
    except Exception as e:
        st.error(f"Could not render text document viewer: {e}")


def render_spreadsheet_document_viewer(file_name, file_bytes):
    """Render CSV/XLS/XLSX with sheet selection and table pagination."""
    try:
        file_name_lower = file_name.lower()
        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]

        if file_name_lower.endswith(".csv"):
            sheets = {"CSV Data": pd.read_csv(BytesIO(file_bytes))}
        else:
            try:
                sheets = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            except Exception as e:
                st.warning(f"Spreadsheet preview is best-effort for this format: {e}")
                rows = extract_excel_data(file_name, file_bytes)
                sheets = {"Recovered Data": pd.DataFrame(rows)}

        if not sheets:
            st.info("No sheets or rows were found.")
            return

        sheet_name = st.selectbox("Sheet", list(sheets.keys()), key=f"sheet_select_{base_key}")
        df = sheets[sheet_name]
        total_rows = len(df)
        rows_per_page = st.selectbox("Rows per page", [25, 50, 100, 250, 500, 1000], index=1, key=f"rows_page_{base_key}")
        total_pages = max(1, (total_rows + rows_per_page - 1) // rows_per_page)

        page_key = f"sheet_page_{base_key}_{sheet_name}"
        if page_key not in st.session_state:
            st.session_state[page_key] = 1

        nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
        with nav_cols[0]:
            if st.button("Previous rows", key=f"sheet_prev_{base_key}", use_container_width=True):
                st.session_state[page_key] = max(1, int(st.session_state[page_key]) - 1)
                st.rerun()
        with nav_cols[1]:
            if st.button("Next rows", key=f"sheet_next_{base_key}", use_container_width=True):
                st.session_state[page_key] = min(total_pages, int(st.session_state[page_key]) + 1)
                st.rerun()
        with nav_cols[2]:
            jump_page = st.number_input(
                "Jump to row page",
                min_value=1,
                max_value=total_pages,
                value=int(st.session_state[page_key]),
                key=f"sheet_jump_{base_key}",
            )
            st.session_state[page_key] = int(jump_page)

        start = (int(st.session_state[page_key]) - 1) * rows_per_page
        end = min(total_rows, start + rows_per_page)
        st.caption(f"Sheet '{sheet_name}': showing rows {start + 1:,}-{end:,} of {total_rows:,}. All sheets are selectable above.")
        st.dataframe(df.iloc[start:end], use_container_width=True, hide_index=True)
    except Exception as e:
        st.error(f"Could not render spreadsheet viewer: {e}")


def render_presentation_document_viewer(file_name, file_bytes, extracted_text):
    """Render PPTX slides on demand, with best-effort text pagination for PPT."""
    try:
        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
        if file_name.lower().endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            total_slides = len(prs.slides)
            if total_slides == 0:
                st.info("No slides were found.")
                return

            slide_key = f"slide_number_{base_key}"
            if slide_key not in st.session_state:
                st.session_state[slide_key] = 1

            nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
            with nav_cols[0]:
                if st.button("Previous slide", key=f"slide_prev_{base_key}", use_container_width=True):
                    st.session_state[slide_key] = max(1, int(st.session_state[slide_key]) - 1)
                    st.rerun()
            with nav_cols[1]:
                if st.button("Next slide", key=f"slide_next_{base_key}", use_container_width=True):
                    st.session_state[slide_key] = min(total_slides, int(st.session_state[slide_key]) + 1)
                    st.rerun()
            with nav_cols[2]:
                slide_number = st.number_input(
                    "Jump to slide",
                    min_value=1,
                    max_value=total_slides,
                    value=int(st.session_state[slide_key]),
                    key=f"slide_jump_{base_key}",
                )
                st.session_state[slide_key] = int(slide_number)

            slide = prs.slides[int(st.session_state[slide_key]) - 1]
            st.caption(f"Showing slide {st.session_state[slide_key]} of {total_slides}. Slides are rendered on demand.")
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
                if hasattr(shape, "table"):
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text.strip() for cell in row.cells])
                    if table_rows:
                        st.dataframe(pd.DataFrame(table_rows), use_container_width=True, hide_index=True)
            if slide_lines:
                st.markdown("### Slide Content")
                st.markdown("\n\n".join(html.escape(line) for line in slide_lines), unsafe_allow_html=True)
            else:
                st.info("This slide has no readable text. Native PPTX image rendering is not available without an external converter.")
        else:
            st.warning("Legacy PPT preview uses best-effort extracted text. Convert to PPTX for slide-level rendering.")
            render_paginated_text_document_viewer(file_name, extracted_text, page_size=7000)
    except Exception as e:
        st.error(f"Could not render presentation viewer: {e}")


def render_universal_document_viewer(file_name, file_entry, extracted_text, highlight_term=None, highlight_page=None):
    """Route the Viewer tab to a full-navigation renderer by file type."""
    try:
        file_bytes = file_entry["bytes"]
        file_name_lower = file_name.lower()
        if file_name_lower.endswith(".pdf"):
            render_document_preview(file_name, file_entry=file_entry, highlight_term=highlight_term, highlight_page=highlight_page)
        elif file_name_lower.endswith((".pptx", ".ppt")):
            render_presentation_document_viewer(file_name, file_bytes, extracted_text)
        elif file_name_lower.endswith((".xlsx", ".xls", ".csv")):
            render_spreadsheet_document_viewer(file_name, file_bytes)
        elif file_name_lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp")):
            st.image(file_bytes, caption=file_name, use_container_width=True)
        else:
            render_paginated_text_document_viewer(file_name, extracted_text)
    except Exception as e:
        st.error(f"Could not render universal document viewer: {e}")


def get_preview_extracted_text(file_name, file_bytes):
    """Extract full text only when an analysis panel actually needs it."""
    try:
        if file_name not in st.session_state.file_texts:
            with st.spinner("Extracting full document text for this analysis panel..."):
                extracted = extract_text(file_name, file_bytes)
                st.session_state.file_texts[file_name] = extracted
                FILE_TEXT_CACHE.set(file_name, extracted)
        return st.session_state.file_texts.get(file_name, "")
    except Exception as e:
        st.warning(f"Could not extract full document text: {e}")
        return ""


def render_professional_document_preview(file_name, file_entry=None, highlight_term=None, highlight_page=None):
    """Render a professional multi-tab document intelligence preview."""
    global PDF_PREVIEW_RESOLUTION

    try:
        if file_entry is None:
            file_entry = get_uploaded_file_entry(file_name)
        if not file_entry:
            st.error("File preview unavailable - file could not be loaded.")
            return

        file_bytes = file_entry["bytes"]
        file_name_lower = file_name.lower()
        reliable_formats = (".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".odt", ".rtf", ".txt", ".md", ".html", ".htm", ".png", ".jpg", ".jpeg", ".webp")
        best_effort_formats = (".doc", ".ppt", ".xls", ".pages")

        extracted_text = st.session_state.file_texts.get(file_name, "")

        metadata = get_preview_metadata(file_name, file_bytes, extracted_text)

        st.markdown(
            """
            <style>
                .preview-shell {
                    padding: 12px 0 2px;
                }
                .preview-note {
                    border: 1px solid #dbeafe;
                    background: #f8fbff;
                    border-radius: 10px;
                    padding: 12px 14px;
                    color: #173152;
                    margin: 8px 0 16px;
                }
            </style>
            """,
            unsafe_allow_html=True,
        )

        st.sidebar.markdown("### Preview Controls")
        zoom_percent = st.sidebar.slider("PDF render zoom", 80, 180, 100, 10)
        PDF_PREVIEW_RESOLUTION = max(72, int(zoom_percent))
        quick_search = st.sidebar.text_input("Search in this document", value=highlight_term or "")
        st.sidebar.caption("Large files are rendered in pages, chunks, and cached extraction layers.")

        st.markdown(f"## {html.escape(file_name)}")
        render_preview_metadata_cards(metadata)

        if file_name_lower.endswith(best_effort_formats):
            st.warning(
                "This is a binary or proprietary format. Preview and extraction use best-effort recovery. "
                "For reliable formatting and analysis, convert to DOCX, PPTX, XLSX, PDF, ODT, RTF, TXT, or HTML."
            )
        elif not file_name_lower.endswith(reliable_formats):
            st.warning("This file type is not fully supported. The app will attempt best-effort extraction.")

        active_preview_panel = st.radio(
            "Preview panel",
            ["Viewer", "Summary", "Search", "Q&A", "Tables", "Images", "Downloads"],
            horizontal=True,
            key=f"preview_panel_{hashlib.md5(file_name.encode('utf-8')).hexdigest()[:12]}",
        )

        if active_preview_panel == "Viewer":
            st.markdown("<div class='preview-note'>Viewer mode supports full-document navigation. PDFs, slides, sheets, and long text documents are rendered on demand instead of silently truncating content.</div>", unsafe_allow_html=True)
            viewer_text = extracted_text
            if file_name_lower.endswith((".doc", ".docx", ".odt", ".rtf", ".txt", ".md", ".html", ".htm", ".pages", ".ppt")):
                viewer_text = get_preview_extracted_text(file_name, file_bytes)
            render_universal_document_viewer(
                file_name,
                file_entry,
                viewer_text,
                highlight_term=quick_search or highlight_term,
                highlight_page=highlight_page,
            )

        elif active_preview_panel == "Summary":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            summary_md = build_preview_summary_markdown(file_name, file_bytes, extracted_text)
            st.markdown(summary_md)
            st.download_button(
                "Download summary as Markdown",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.md",
                mime="text/markdown",
                key=f"preview_summary_md_{file_name}",
            )
            st.download_button(
                "Download summary as TXT",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.txt",
                mime="text/plain",
                key=f"preview_summary_txt_{file_name}",
            )

        elif active_preview_panel == "Search":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            search_query = st.text_input("Search extracted text", value=quick_search, key=f"preview_search_{file_name}")
            if search_query:
                chunks = keyword_search_preview_chunks(extracted_text, search_query, limit=20)
                st.caption(f"{len(chunks)} relevant chunks found.")
                for index, chunk in enumerate(chunks, start=1):
                    with st.expander(f"Match {index}", expanded=index == 1):
                        st.write(chunk)
            else:
                st.info("Enter a search term to search extracted text chunks.")

        elif active_preview_panel == "Q&A":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            question = st.text_input("Ask a question about this document", key=f"preview_qa_{file_name}")
            if question:
                answer = build_preview_answer(file_name, extracted_text, question)
                st.markdown("### Answer")
                st.markdown(answer)
                st.download_button(
                    "Download Q&A result",
                    data=f"Question: {question}\n\nAnswer:\n{answer}".encode("utf-8"),
                    file_name=f"{os.path.splitext(file_name)[0]}_qa.txt",
                    mime="text/plain",
                    key=f"preview_qa_download_{file_name}",
                )
            else:
                st.info("Ask a focused question. The app retrieves only relevant chunks before answering.")

        elif active_preview_panel == "Tables":
            if not file_name_lower.endswith((".csv", ".xlsx", ".xls")):
                extracted_text = get_preview_extracted_text(file_name, file_bytes)
            tables = extract_preview_tables(file_name, file_bytes, extracted_text)
            if not tables:
                st.info("No tables were detected. For scanned PDFs, enable page-level table detection in the Viewer tab.")
            for index, (table_name, df) in enumerate(tables, start=1):
                with st.expander(table_name, expanded=index == 1):
                    st.dataframe(df, use_container_width=True, hide_index=True)
                    st.download_button(
                        "Download CSV",
                        data=df.to_csv(index=False).encode("utf-8"),
                        file_name=f"{os.path.splitext(file_name)[0]}_{index}.csv",
                        mime="text/csv",
                        key=f"preview_table_csv_{file_name}_{index}",
                    )
                    st.download_button(
                        "Download XLSX",
                        data=dataframe_to_xlsx_bytes(df),
                        file_name=f"{os.path.splitext(file_name)[0]}_{index}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"preview_table_xlsx_{file_name}_{index}",
                    )

        elif active_preview_panel == "Images":
            if file_name_lower.endswith((".png", ".jpg", ".jpeg", ".webp")):
                st.image(file_bytes, caption=file_name, use_container_width=True)
                if st.checkbox("Run OCR if available", key=f"preview_ocr_{file_name}"):
                    st.text_area("OCR text", value=ocr_image_best_effort(file_bytes), height=220)
            render_extracted_assets_preview(file_name, file_entry)

        elif active_preview_panel == "Downloads":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            summary_md = build_preview_summary_markdown(file_name, file_bytes, extracted_text)
            st.download_button(
                "Summary - Markdown",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.md",
                mime="text/markdown",
                key=f"preview_download_summary_md_{file_name}",
            )
            st.download_button(
                "Extracted text - TXT",
                data=str(extracted_text or "").encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_extracted_text.txt",
                mime="text/plain",
                key=f"preview_download_text_{file_name}",
            )
            report_bytes = None
            try:
                report_doc = docx.Document()
                report_doc.add_heading(f"Document Intelligence Report: {file_name}", level=1)
                report_doc.add_paragraph(summary_md)
                report_output = BytesIO()
                report_doc.save(report_output)
                report_bytes = report_output.getvalue()
            except Exception:
                report_bytes = None
            if report_bytes:
                st.download_button(
                    "Generated report - DOCX",
                    data=report_bytes,
                    file_name=f"{os.path.splitext(file_name)[0]}_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key=f"preview_download_report_{file_name}",
                )
            with st.expander("Recommended install commands"):
                st.code(
                    "pip install streamlit pandas openpyxl python-docx python-pptx pdfplumber beautifulsoup4 pillow plotly faiss-cpu langchain-community langchain-text-splitters sentence-transformers pytesseract",
                    language="bash",
                )
    except Exception as e:
        st.error(f"Error rendering professional document preview: {e}")


# -------------------------------
# PREVIEW ROUTE HANDLING
# -------------------------------
# If a preview token is present in the query params, the app short-circuits into
# a dedicated document preview screen instead of rendering the main multi-tab UI.
preview_file_from_url = None
query_params = {}

# Try different methods to get query params (for compatibility across Streamlit versions)
try:
    if hasattr(st, "query_params"):
        query_params = st.query_params
    elif hasattr(st, "experimental_get_query_params"):
        query_params = st.experimental_get_query_params() or {}
    elif hasattr(st, "get_query_params"):
        query_params = st.get_query_params() or {}
    else:
        query_params = {}
except Exception:
    query_params = {}

highlight_term = None
preview_page = None
preview_token = None
if "preview_token" in query_params and query_params["preview_token"]:
    preview_value = query_params["preview_token"]
    if isinstance(preview_value, list):
        preview_value = preview_value[0] if preview_value else ""
    preview_token = str(preview_value)
    token_data = PREVIEW_TOKENS.get(preview_token)
    preview_file_from_url = token_data['file_name'] if token_data else None
    if "highlight" in query_params and query_params["highlight"]:
        highlight_value = query_params["highlight"]
        if isinstance(highlight_value, list):
            highlight_value = highlight_value[0] if highlight_value else ""
        highlight_term = urllib.parse.unquote_plus(str(highlight_value))
    if "page" in query_params and query_params["page"]:
        page_value = query_params["page"]
        if isinstance(page_value, list):
            page_value = page_value[0] if page_value else ""
        try:
            preview_page = int(str(page_value))
        except ValueError:
            preview_page = None

    # If preview_token is present but not found, show error and stop
    if not token_data:
        st.title("Document Preview")
        st.warning("This preview link is no longer available.")
        st.info("Return to the main app and click the eye preview button next to the uploaded file again.")
        st.stop()

if preview_file_from_url:
    preview_entry = PREVIEW_STORE.get(preview_token)
    st.title("📄 Document Preview")
    if preview_entry is not None:
        st.markdown(f"### {preview_entry['name']}")
        st.markdown("---")
        render_professional_document_preview(
            preview_entry['name'],
            file_entry=preview_entry,
            highlight_term=highlight_term,
            highlight_page=preview_page,
        )
    else:
        st.error("Preview file not found in the preview store. Please return to the app and click preview again.")
    st.stop()


# -------------------------------
# -------------------------------
# FILE UPLOAD & MANAGEMENT (SIDEBAR)
# -------------------------------
# Sidebar area:
# This block manages login state, file upload/delete, global file selection, and
# preview launch links. Files selected here become available to the individual tabs.
with st.sidebar:
    if st.session_state.is_authenticated:
        if st.button("Open Workspace", key="mobile_open_workspace_btn", use_container_width=True):
            st.session_state.mobile_sidebar_visible = False
            st.rerun()

        creator_timestamp = None
        if st.session_state.user_role == "creator" and st.session_state.login_history:
            creator_entries = [
                entry for entry in st.session_state.login_history
                if entry.get("username") == st.session_state.logged_in_username and entry.get("role") == "creator" and entry.get("action") == "login"
            ]
            if creator_entries:
                creator_timestamp = creator_entries[-1].get("timestamp")

        st.markdown(
            """
            <style>
                [data-testid="stSidebar"] div[data-testid="column"] {
                    display: flex;
                    align-items: center;
                }
                [data-testid="stSidebar"] div[data-testid="column"] > div {
                    width: 100%;
                }
                .file-box {
                    background: #f8fbff;
                    border: 1px solid #d7e3f4;
                    border-radius: 12px;
                    margin-bottom: 8px;
                    color: #173152;
                    font-size: 14px;
                    overflow-wrap: anywhere;
                }
                [data-testid="stSidebar"] [class*="st-key-select_file_"] button[kind="secondary"],
                [data-testid="stSidebar"] [class*="st-key-select_file_"] button[kind="primary"] {
                    width: 100%;
                    min-height: 72px;
                    border-radius: 12px;
                    padding: 12px 14px;
                    font-size: 14px;
                    text-align: left;
                    justify-content: flex-start;
                    white-space: normal;
                    line-height: 1.4;
                }
                [data-testid="stSidebar"] [class*="st-key-select_file_"] button[kind="secondary"] {
                    background: #f5fbff;
                    border: 1px solid #d0e8f8;
                    color: #1e293b;
                }
                [data-testid="stSidebar"] [class*="st-key-select_file_"] button[kind="primary"] {
                    background: #ffe7d6 !important;
                    border: 2px solid #ffbea3 !important;
                    color: #5f351c !important;
                    box-shadow: inset 0 0 0 1px rgba(255, 190, 163, 0.55);
                }
                [data-testid="stSidebar"] [class*="st-key-select_file_"] button[kind="primary"]:hover {
                    background: #ffd3b5 !important;
                    border-color: #ffb18f !important;
                }
                [data-testid="stSidebar"] div[data-testid="stLinkButton"] a {
                    min-height: 38px;
                    border-radius: 12px;
                    padding: 0.35rem 0.5rem;
                    white-space: nowrap;
                    justify-content: center;
                    background: #f5fbff !important;
                    border: 1px solid #d0e8f8 !important;
                    color: #173152 !important;
                }
                [data-testid="stSidebar"] div[data-testid="stLinkButton"] a:hover {
                    background: #e8f6ff !important;
                    border-color: #a0c8e8 !important;
                }
                [data-testid="stSidebar"] [class*="st-key-del_file_"] button[kind="tertiary"] {
                    background: transparent;
                    border: none;
                    box-shadow: none;
                    color: #6c7280;
                    min-height: 38px;
                }
                [data-testid="stSidebar"] [class*="st-key-del_file_"] button[kind="tertiary"]:hover {
                    background: #f3f6fb;
                    color: #b42318;
                }
            </style>
            """,
            unsafe_allow_html=True
        )

        # Mercedes Logo in Sidebar
        if logo_data:
            st.markdown(
                f'''
                <div style="text-align: center; margin-bottom: 20px;">
                    <img src="data:image/gif;base64,{logo_data}" style="width: 40px; height: 40px;">
                    <div style="font-size: 15px; color: var(--text-secondary); margin-top: 4px; font-weight: 500;">Mercedes-Benz</div>
                </div>
                ''',
                unsafe_allow_html=True,
            )
        
        st.header("Upload Documents")
        st.info("1) Upload files." \
        " 2) Click the file cards you need. " \
        "3) Switch tabs and work with selected files.")
        new_files = st.file_uploader(
            "Upload PDF, Word, PPT, Excel, CSV, TXT, HTML, ODT, RTF, Pages, CAPL, Images",
            type=["pdf", "doc", "docx", "txt", "md", "log", "ppt", "pptx", "xls", "xlsx", "csv", "html", "htm", "odt", "rtf", "pages", "capl", "can", "png", "jpg", "jpeg", "gif", "bmp", "webp"],
            accept_multiple_files=True,
            key=f"file_uploader_{st.session_state.file_uploader_key}"
        )

        if new_files:
            with st.spinner("Loading uploaded files..."):
                existing_names = {f["name"] for f in st.session_state.uploaded_files}
                uploaded_new_file = False
                for file in new_files:
                    if file.name not in existing_names:
                        file_bytes = file.read()
                        st.session_state.uploaded_files.append({"name": file.name, "bytes": file_bytes})
                        uploaded_new_file = True
                if uploaded_new_file:
                    new_file_names = [file.name for file in new_files if file.name not in existing_names]
                    ensure_files_processed(new_file_names)
                    st.session_state.workspace_memory["indexed_files"] = sorted(set(st.session_state.workspace_memory.get("indexed_files", []) + new_file_names))
                    record_workspace_memory_event(
                        "upload",
                        "Documents indexed into shared memory",
                        "Uploaded and indexed: " + ", ".join(new_file_names),
                        source="Upload",
                    )
                    get_unified_workspace_vector_store(new_file_names)
                    save_workspace_memory()
                    save_memory_log("upload", f"Uploaded {len(new_file_names)} new file(s)", {"files": new_file_names})
                    st.session_state.messages = []
                    st.session_state.chat_summary_downloads = {"images": [], "tables": [], "csv": [], "diagrams": []}
                    st.session_state.chat_file_selection = []
                    st.success("✅ New files uploaded. Chat history has been cleared.")

        st.markdown("---")
        st.markdown("### Uploaded files")
        
        for idx, file_dict in enumerate(st.session_state.uploaded_files[:]):
            cols = st.columns([0.56, 0.27, 0.17], vertical_alignment="center")
            with cols[0]:
                file_name = file_dict["name"]
                is_selected = file_name in st.session_state.selected_files
                button_label = file_name if not is_selected else f"Selected: {file_name}"
                if st.button(
                    button_label,
                    key=f"select_file_{idx}",
                    help=f"Click to {'remove' if is_selected else 'add'} {file_name}",
                    use_container_width=True,
                    type="primary" if is_selected else "secondary",
                ):
                    if is_selected:
                        st.session_state.selected_files.remove(file_name)
                    else:
                        st.session_state.selected_files.append(file_name)
                    st.rerun()
            with cols[1]:
                preview_url = create_preview_link(file_name)
                st.link_button(
                    "Open",
                    preview_url or "#",
                    help=f"Open preview for {file_name}",
                    icon="👁️",
                    type="secondary",
                    use_container_width=True,
                    disabled=preview_url is None,
                )
            with cols[2]:
                if st.button("🗑️", key=f"del_file_{idx}", help=f"Delete {file_name}", use_container_width=True, type="tertiary"):
                    deleted_name = file_name
                    st.session_state.uploaded_files = [f for f in st.session_state.uploaded_files if
                                                       f["name"] != deleted_name]

                    if deleted_name in st.session_state.selected_files:
                        st.session_state.selected_files.remove(deleted_name)

                    for key in ["file_texts", "excel_data_by_file", "vector_stores"]:
                        st.session_state.get(key, {}).pop(deleted_name, None)

                    if st.session_state.capl_last_analyzed_file == deleted_name:
                        st.session_state.capl_last_analyzed_file = None
                        st.session_state.capl_last_issues = None
                    st.rerun()
        st.markdown("*Selected files above are available across all tabs.*")
        st.markdown("---")
        if st.button("Clear All Files"):
            for key in ["uploaded_files", "selected_files", "file_texts", "excel_data_by_file", "vector_stores",
                        "messages"]:
                st.session_state[key].clear()
            st.session_state.chat_summary_downloads = {"images": [], "tables": [], "csv": [], "diagrams": []}
            st.session_state.chat_file_selection = []
            st.session_state.capl_last_analyzed_file = None
            st.session_state.capl_last_issues = None
            st.session_state.file_uploader_key += 1
            st.rerun()

        if st.session_state.user_role == "creator":
            st.markdown("---")
            st.subheader("Creator Login / Logout History")
            if st.session_state.login_history:
                st.table(pd.DataFrame(st.session_state.login_history))
            st.markdown("---")
            st.subheader("User Statistics")
            total_opened = len(set(h["username"] for h in st.session_state.login_history))
            active_file = "active_users.json"
            if os.path.exists(active_file):
                with open(active_file, "r") as f:
                    active_users = json.load(f)
                now = datetime.now()
                current_active = len(set(u["username"] for u in active_users if
                                         datetime.fromisoformat(u["timestamp"]) > now - timedelta(minutes=30)))
            else:
                current_active = 0
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Total Users Opened", total_opened)
            with col2:
                st.metric("Currently Active Users", current_active)
            st.markdown("---")
            if st.button("Clean Old Active Users"):
                active_file = "active_users.json"
                if os.path.exists(active_file):
                    with open(active_file, "r") as f:
                        active_users = json.load(f)
                    now = datetime.now()
                    active_users = [u for u in active_users if
                                    datetime.fromisoformat(u["timestamp"]) > now - timedelta(hours=1)]
                    with open(active_file, "w") as f:
                        json.dump(active_users, f)
                    st.success("Cleaned old active users.")
                else:
                    st.info("No active users file found.")


# -------------------------------
# TEXT EXTRACTION
# -------------------------------
# Excel extraction helper:
# Primarily used by Dashboard previews/charts, but cached so other tabs can reuse
# the parsed spreadsheet rows without re-reading the uploaded file each rerun.
@st.cache_data(show_spinner=False)
def extract_excel_data(file_name, file_bytes):
    data = []
    bio = BytesIO(file_bytes)
    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".xlsx"):
            wb = openpyxl.load_workbook(bio, data_only=True)
            for sheet in wb:
                headers = None
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i == 0:
                        headers = list(row)
                    else:
                        if row and any(cell is not None for cell in row):
                            data.append(dict(zip(headers, row)))
    except Exception:
        data = []
    return data


# -------------------------------
# PROCESS FILES & BUILD VECTOR STORES
# -------------------------------
# AI/vector helpers:
# These are mainly used by the Chat tab for semantic retrieval and LLM answers.
# They also centralize file preprocessing so each tab can rely on the same cache.
@st.cache_data(show_spinner=False)
def create_vector_store(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    chunks = splitter.split_text(text[:MAX_VECTOR_TEXT_CHARS])
    emb = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(chunks, emb)


@st.cache_resource(show_spinner=False)
def load_llm():
    try:
        from transformers import pipeline
    except Exception:
        st.warning(
            "LLM is unavailable because transformers could not be imported (torchvision unmet).\nInstall torch & torchvision if you want AI features.")
        return None

    candidate_tasks = ["text2text-generation", "text-generation", "image-text-to-text", "table-question-answering"]
    task_errors = []

    for task in candidate_tasks:
        try:
            pipe = pipeline(task, model="google/flan-t5-small", max_new_tokens=128, return_full_text=False)
            st.session_state.llm_task = task
            return HuggingFacePipeline(pipeline=pipe)
        except Exception as e:
            task_errors.append((task, str(e)))

    # Only show a single warning if no task could be initialized
    if task_errors:
        st.warning("LLM initialization failed for all candidate tasks; AI features are unavailable.")
        st.session_state.llm_task = None

    return None


def ensure_files_processed(file_names):
    for file_name in file_names:
        ensure_file_processed(file_name)


def process_selected_files():
    """Fully extract all currently selected files so every tab uses the same data."""
    ensure_files_processed(st.session_state.selected_files)


WORKSPACE_DB_FILE = os.path.join(APP_DIR, "workspace_memory.db")
WORKSPACE_MEMORY_KEY = "workspace_memory"


def init_workspace_db():
    """Initialize persistent workspace storage for memory and logs."""
    os.makedirs(APP_DIR, exist_ok=True)
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workspace_meta (
            meta_key TEXT PRIMARY KEY,
            meta_value TEXT
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workspace_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            timestamp TEXT,
            log_type TEXT,
            message TEXT,
            details TEXT
        )
        """
    )
    conn.commit()
    conn.close()


def default_workspace_memory():
    return {
        "chat": [],
        "agent_runs": [],
        "indexed_files": [],
        "memory_events": [],
        "summary": {},
        "metadata": {},
    }


def normalize_workspace_memory(memory):
    """Keep older saved memory compatible with the autonomous workspace schema."""
    normalized = default_workspace_memory()
    if isinstance(memory, dict):
        for key, value in memory.items():
            normalized[key] = value
    for list_key in ["chat", "agent_runs", "indexed_files", "memory_events"]:
        if not isinstance(normalized.get(list_key), list):
            normalized[list_key] = []
    for dict_key in ["summary", "metadata"]:
        if not isinstance(normalized.get(dict_key), dict):
            normalized[dict_key] = {}
    return normalized


def load_workspace_memory():
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "SELECT meta_value FROM workspace_meta WHERE meta_key = ?",
        (WORKSPACE_MEMORY_KEY,)
    )
    row = cursor.fetchone()
    conn.close()
    if row:
        try:
            return normalize_workspace_memory(json.loads(row[0]))
        except Exception:
            pass
    return default_workspace_memory()


def save_workspace_memory():
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "INSERT OR REPLACE INTO workspace_meta (meta_key, meta_value) VALUES (?, ?)",
        (WORKSPACE_MEMORY_KEY, json.dumps(st.session_state.workspace_memory, default=str))
    )
    conn.commit()
    conn.close()


def save_memory_log(log_type, message, details=None):
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    details_json = json.dumps(details, default=str) if details is not None else None
    cursor.execute(
        "INSERT INTO workspace_logs (timestamp, log_type, message, details) VALUES (?, ?, ?, ?)",
        (datetime.now().isoformat(), log_type, message, details_json)
    )
    conn.commit()
    conn.close()


def get_memory_logs(limit=50):
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "SELECT timestamp, log_type, message, details FROM workspace_logs ORDER BY id DESC LIMIT ?",
        (limit,)
    )
    rows = cursor.fetchall()
    conn.close()
    results = []
    for timestamp, log_type, message, details in rows:
        try:
            details = json.loads(details) if details else None
        except Exception:
            details = details
        results.append({
            "timestamp": timestamp,
            "type": log_type,
            "message": message,
            "details": details,
        })
    return results


def record_workspace_memory_event(event_type, title, content, source=None):
    """Append a compact memory event that every module can retrieve later."""
    event = {
        "type": event_type,
        "title": title,
        "content": str(content or "")[:4000],
        "source": source or "workspace",
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    st.session_state.workspace_memory["memory_events"].append(event)
    st.session_state.workspace_memory["memory_events"] = st.session_state.workspace_memory["memory_events"][-200:]
    return event


def append_chat_to_workspace_memory(user_input, assistant_response, file_names):
    """Store chat as durable workspace memory instead of isolated chat history."""
    chat_entry = {
        "user": user_input,
        "assistant": assistant_response,
        "files": list(file_names or []),
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    st.session_state.workspace_memory["chat"].append(chat_entry)
    st.session_state.workspace_memory["chat"] = st.session_state.workspace_memory["chat"][-200:]
    record_workspace_memory_event(
        "chat",
        "Conversation memory",
        f"User: {user_input}\nAssistant: {assistant_response}",
        source=", ".join(file_names or []) or "chat",
    )
    return chat_entry


def build_unified_memory_text(file_names=None, include_chat=True, include_agents=True, max_chars=MAX_VECTOR_TEXT_CHARS):
    """Compose documents, conversations, and agent runs into one AI brain text."""
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    sections = []

    candidate_files = file_names if file_names is not None else [f["name"] for f in st.session_state.get("uploaded_files", [])]
    candidate_files = [file_name for file_name in candidate_files if file_name in st.session_state.file_texts]
    for file_name in candidate_files:
        text = str(st.session_state.file_texts.get(file_name, "")).strip()
        if text:
            sections.append(f"[DOCUMENT: {file_name}]\n{text[:60000]}")

    if include_chat:
        for entry in st.session_state.workspace_memory.get("chat", [])[-80:]:
            sections.append(
                "[CHAT MEMORY]\n"
                f"User: {entry.get('user', '')}\n"
                f"Assistant: {entry.get('assistant', '')}\n"
                f"Files: {', '.join(entry.get('files', []) or [])}"
            )

    for event in st.session_state.workspace_memory.get("memory_events", [])[-100:]:
        sections.append(
            "[MEMORY EVENT]\n"
            f"Type: {event.get('type', '')}\n"
            f"Title: {event.get('title', '')}\n"
            f"Source: {event.get('source', '')}\n"
            f"{event.get('content', '')}"
        )

    if include_agents:
        for run in st.session_state.workspace_memory.get("agent_runs", [])[-40:]:
            sections.append(
                "[CAPL AGENT RUN]\n"
                f"Goal: {run.get('goal', '')}\n"
                f"Plan: {', '.join(run.get('plan', []) or [])}\n"
                f"Final: {run.get('final_response', '')[:2500]}"
            )

    memory_text = "\n\n".join(sections).strip()
    return memory_text[:max_chars]


def get_unified_workspace_vector_store(file_names=None):
    memory_text = build_unified_memory_text(file_names=file_names)
    if not memory_text.strip():
        return None

    digest = hashlib.md5(memory_text.encode("utf-8", errors="ignore")).hexdigest()
    selection_key = f"unified_memory::{digest}"
    cached_vs = VECTOR_STORE_CACHE.get(selection_key)
    if cached_vs is not None:
        st.session_state.vector_stores[selection_key] = cached_vs
        return cached_vs

    try:
        vs = create_vector_store(memory_text)
        st.session_state.vector_stores[selection_key] = vs
        VECTOR_STORE_CACHE.set(selection_key, vs)
        return vs
    except Exception:
        return None


def get_workspace_vector_store(file_names=None):
    return get_unified_workspace_vector_store(file_names)


def search_workspace_memory(query, limit=4):
    vector_store = get_workspace_vector_store()
    if vector_store is None:
        return []
    try:
        docs = vector_store.similarity_search(query, k=limit)
        return [getattr(doc, "page_content", str(doc)) for doc in docs]
    except Exception:
        return []


def extract_risk_signals(text, limit=12):
    risk_terms = [
        "risk", "failure", "hazard", "issue", "problem", "danger", "alert",
        "warning", "fault", "breach", "leak", "vulnerability", "impact",
        "delay", "downtime", "non-compliance",
    ]
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    results = []
    for line in lines:
        lower_line = line.lower()
        if any(term in lower_line for term in risk_terms):
            results.append(line)
        if len(results) >= limit:
            break
    return results or ["No explicit risk signals were found in the selected documents."]


def extract_entities(text, limit=20):
    raw_text = str(text or "")
    candidates = re.findall(r"\b[A-Z][A-Za-z0-9]{2,}(?: [A-Z][A-Za-z0-9]{2,})*\b", raw_text)
    counts = Counter(candidates)
    return [entity for entity, _ in counts.most_common(limit)]


def extract_key_themes(text, limit=8):
    words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{3,}", str(text or ""))
    counts = Counter(
        word.lower()
        for word in words
        if word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    return [word.title() for word, _ in counts.most_common(limit)]


def build_workspace_intelligence_summary(file_names=None):
    """Produce the live AI Insight Engine view from shared memory."""
    if file_names is None:
        file_names = st.session_state.get("selected_files", []) or [
            f["name"] for f in st.session_state.get("uploaded_files", [])
        ]
    ensure_files_processed(file_names)
    memory_text = build_unified_memory_text(file_names=file_names, max_chars=180000)
    chat_count = len(st.session_state.workspace_memory.get("chat", []))
    agent_count = len(st.session_state.workspace_memory.get("agent_runs", []))
    indexed_files = sorted(set(st.session_state.workspace_memory.get("indexed_files", []) + list(file_names or [])))

    themes = extract_key_themes(memory_text, limit=10)
    entities = extract_entities(memory_text, limit=16)
    risks = extract_risk_signals(memory_text, limit=8)
    recent_logs = get_memory_logs(limit=6)

    insights = []
    if themes:
        insights.append(f"Dominant knowledge themes: {', '.join(themes[:5])}.")
    if chat_count:
        insights.append(f"Conversation memory is active with {chat_count} stored exchange(s).")
    if indexed_files:
        insights.append(f"{len(indexed_files)} document(s) are connected to the shared AI memory.")
    if agent_count:
        insights.append(f"CAPL agents have completed {agent_count} autonomous run(s).")
    if not insights:
        insights.append("Upload documents or start a conversation to grow the workspace memory.")

    return {
        "themes": themes,
        "insights": insights,
        "entities": entities,
        "risks": risks,
        "state": {
            "indexed_files": len(indexed_files),
            "chat_entries": chat_count,
            "agent_runs": agent_count,
            "memory_events": len(st.session_state.workspace_memory.get("memory_events", [])),
            "memory_chars": len(memory_text),
        },
        "logs": recent_logs,
    }


def render_workspace_intelligence_panel(file_names=None):
    summary = build_workspace_intelligence_summary(file_names=file_names)
    state = summary["state"]

    st.markdown("### AI Insight Engine")
    metric_cols = st.columns(4)
    metric_cols[0].metric("Indexed Files", state["indexed_files"])
    metric_cols[1].metric("Chat Memory", state["chat_entries"])
    metric_cols[2].metric("CAPL Runs", state["agent_runs"])
    metric_cols[3].metric("Memory Events", state["memory_events"])

    cols = st.columns(2)
    with cols[0]:
        st.markdown("#### Key Themes")
        st.markdown(", ".join(summary["themes"]) if summary["themes"] else "No themes detected yet.")
        st.markdown("#### Entities")
        st.markdown(", ".join(summary["entities"][:12]) if summary["entities"] else "No entities detected yet.")
    with cols[1]:
        st.markdown("#### Insights")
        for item in summary["insights"]:
            st.markdown(f"- {html.escape(item)}")
        st.markdown("#### Risks / Signals")
        for item in summary["risks"][:6]:
            st.markdown(f"- {html.escape(str(item)[:240])}")

    if summary["logs"]:
        with st.expander("Live Memory Loop", expanded=False):
            for log in summary["logs"]:
                st.markdown(f"**{html.escape(log.get('type', 'log'))}** - {html.escape(log.get('message', ''))}")


def plan_autonomous_task(goal):
    lower_goal = str(goal or "").lower()
    tasks = ["retrieve_relevant_memory"]
    if any(term in lower_goal for term in ["compare", "difference", "diff", "semantic diff"]):
        tasks.append("compare_documents")
    if any(term in lower_goal for term in ["risk", "issue", "signal", "warning", "hazard"]):
        tasks.append("extract_risks")
    if any(term in lower_goal for term in ["entity", "entities", "extract entities", "parts", "components"]):
        tasks.append("extract_entities")
    if any(term in lower_goal for term in ["summarize", "summary", "overview", "insight", "analyze"]):
        tasks.append("summarize_findings")
    if any(term in lower_goal for term in ["analyze", "review", "inspect"]) and "summarize_findings" not in tasks:
        tasks.append("analyze_documents")
    if "summarize_findings" not in tasks and "analyze_documents" not in tasks:
        tasks.append("summarize_findings")
    return list(dict.fromkeys(tasks))


def planning_agent(goal):
    steps = plan_autonomous_task(goal)
    return {
        "agent": "Planning Agent",
        "role": "Brain / Orchestrator",
        "strategy": "Break the goal into memory retrieval, tool execution, reasoning, and coordination steps.",
        "steps": steps,
    }


def retrieve_autonomous_context(goal, file_names):
    if not goal:
        return []
    vector_store = get_workspace_vector_store(file_names)
    if vector_store is None:
        return []
    try:
        docs = vector_store.similarity_search(goal, k=4)
        return [getattr(doc, "page_content", str(doc)) for doc in docs]
    except Exception:
        return []


def retrieval_agent(goal, file_names):
    context = retrieve_autonomous_context(goal, file_names)
    return {
        "agent": "Retrieval Agent",
        "role": "Memory Brain",
        "context": context,
        "summary": f"Retrieved {len(context)} relevant memory fragment(s) from shared FAISS memory.",
    }


def execute_autonomous_tool(task, file_names, context):
    file_names = file_names or [f["name"] for f in st.session_state.uploaded_files]
    if task == "retrieve_relevant_memory":
        return {"context": context}
    if task == "analyze_documents":
        results = []
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            file_entry = get_uploaded_file_entry(file_name)
            if file_text.strip() and file_entry is not None:
                results.append(build_detailed_document_summary(file_name, file_entry["bytes"], file_text))
        return {"analysis": results or ["No document content available for analysis."]}
    if task == "summarize_findings":
        context_text = "\n\n".join(context[:3]) if context else "No relevant memory context found."
        return {"summary": context_text}
    if task == "extract_risks":
        results = {}
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            if file_text.strip():
                results[file_name] = extract_risk_signals(file_text)
        return {"risks": results or {"message": "No risk-related content found."}}
    if task == "extract_entities":
        results = {}
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            if file_text.strip():
                results[file_name] = extract_entities(file_text)
        return {"entities": results or {"message": "No entities could be extracted."}}
    if task == "compare_documents":
        if len(file_names) < 2:
            return {"compare": "At least two files are required for document comparison."}
        file_texts = {f: st.session_state.file_texts.get(f, "") for f in file_names}
        return {"compare": highlight_multi_file_differences(file_texts)}
    return {"result": f"No tool implementation exists for '{task}'."}


def execution_agent(steps, file_names, context):
    outputs = {}
    for task in steps:
        if task == "retrieve_relevant_memory":
            outputs[task] = {"context": context}
        else:
            outputs[task] = execute_autonomous_tool(task, file_names, context)
    return {
        "agent": "Execution Agent",
        "role": "Tool Runner",
        "outputs": outputs,
        "summary": f"Executed {len(outputs)} autonomous tool step(s).",
    }


def reason_over_results(agent_outputs):
    summary_parts = []
    for step, output in agent_outputs.items():
        if isinstance(output, dict):
            if "summary" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** {output['summary']}")
            elif "context" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Retrieved {len(output['context'])} memory fragments.")
            elif "risks" in output:
                risk_lines = sum(len(v) for v in output['risks'].values() if isinstance(v, list))
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Extracted {risk_lines} risk lines.")
            elif "entities" in output:
                entity_count = sum(len(v) for v in output['entities'].values() if isinstance(v, list))
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Extracted {entity_count} entities.")
            elif "compare" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Comparison results are available in the agent output." )
            elif "analysis" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Document analysis completed for {len(output['analysis'])} files.")
            else:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Output produced.")
        else:
            summary_parts.append(f"**{step.replace('_', ' ').title()}:** {str(output)[:200]}")
    return "\n\n".join(summary_parts)


def reasoning_agent(goal, file_names, context, tool_outputs):
    memory_text = build_unified_memory_text(file_names=file_names, max_chars=80000)
    combined_reasoning_text = "\n\n".join(context or []) + "\n\n" + memory_text[:40000]
    themes = extract_key_themes(combined_reasoning_text, limit=8)
    entities = extract_entities(combined_reasoning_text, limit=12)
    risks = extract_risk_signals(combined_reasoning_text, limit=8)
    interpretation = [
        f"Goal interpreted as: {goal}",
        f"Key themes: {', '.join(themes) if themes else 'None detected'}",
        f"Important entities: {', '.join(entities[:8]) if entities else 'None detected'}",
        f"Risk/signal count: {len(risks)}",
        reason_over_results(tool_outputs),
    ]
    return {
        "agent": "Reasoning Agent",
        "role": "Analyst",
        "themes": themes,
        "entities": entities,
        "risks": risks,
        "interpretation": "\n\n".join(interpretation),
    }


def coordinate_agents(goal, steps, outputs):
    sections = [f"### Autonomous CAPL Agent Result\n**Goal:** {html.escape(goal)}\n"]
    sections.append(reason_over_results(outputs))
    for step in steps:
        output = outputs.get(step)
        if isinstance(output, dict) and step == "compare_documents":
            sections.append("### Comparison Output")
            sections.append(output.get("compare", "No comparison result."))
        elif isinstance(output, dict) and step == "analyze_documents":
            analysis = output.get("analysis", [])
            sections.append("### Analysis Output")
            sections.extend(analysis[:2] if isinstance(analysis, list) else [str(analysis)])
    return "\n\n".join(sections)


def coordination_agent(goal, planning, retrieval, execution, reasoning):
    tool_outputs = execution.get("outputs", {})
    sections = [
        "### Autonomous CAPL Agent Result",
        f"**Goal:** {html.escape(goal)}",
        "#### Agent Trace",
        f"- **Planning Agent:** {html.escape(planning.get('strategy', ''))}",
        f"- **Retrieval Agent:** {html.escape(retrieval.get('summary', ''))}",
        f"- **Execution Agent:** {html.escape(execution.get('summary', ''))}",
        "- **Reasoning Agent:** Interpreted tool outputs against shared memory.",
        "- **Coordination Agent:** Merged agent outputs into this final response.",
        "#### Execution Plan",
    ]
    sections.extend(f"- {html.escape(step.replace('_', ' ').title())}" for step in planning.get("steps", []))
    sections.extend(["#### Reasoned Findings", reasoning.get("interpretation", "")])

    if reasoning.get("risks"):
        sections.append("#### Risks / Signals")
        sections.extend(f"- {html.escape(str(item)[:240])}" for item in reasoning.get("risks", [])[:8])

    for step in planning.get("steps", []):
        output = tool_outputs.get(step)
        if isinstance(output, dict) and step == "compare_documents":
            sections.append("#### Comparison Output")
            sections.append(output.get("compare", "No comparison result."))
        elif isinstance(output, dict) and step == "analyze_documents":
            analysis = output.get("analysis", [])
            sections.append("#### Analysis Output")
            sections.extend(analysis[:2] if isinstance(analysis, list) else [str(analysis)])

    return "\n\n".join(sections)


def run_capl_agent(goal, file_names):
    if not goal or not str(goal).strip():
        return "Provide a goal for the autonomous CAPL agents."
    target_files = file_names or [f["name"] for f in st.session_state.uploaded_files]
    target_files = [f for f in target_files if get_uploaded_file_entry(f)]
    if not target_files:
        return "No processed files are available. Upload files and allow the system to extract them first."

    ensure_files_processed(target_files)
    planning = planning_agent(goal)
    plan = planning["steps"]
    retrieval = retrieval_agent(goal, target_files)
    execution = execution_agent(plan, target_files, retrieval.get("context", []))
    reasoning = reasoning_agent(goal, target_files, retrieval.get("context", []), execution.get("outputs", {}))
    outputs = {
        "planning": planning,
        "retrieval": retrieval,
        "execution": execution,
        "reasoning": reasoning,
    }
    final_response = coordination_agent(goal, planning, retrieval, execution, reasoning)

    run_entry = {
        "goal": goal,
        "files": target_files,
        "plan": plan,
        "outputs": outputs,
        "final_response": final_response,
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.agent_run_history.append(run_entry)
    st.session_state.workspace_memory["agent_runs"].append(run_entry)
    record_workspace_memory_event("capl_agent", f"Autonomous goal: {goal}", final_response, source="CAPL")
    st.session_state.workspace_memory["indexed_files"] = sorted(set(st.session_state.workspace_memory.get("indexed_files", []) + target_files))
    save_workspace_memory()
    save_memory_log("capl_agent", f"Ran autonomous CAPL goal: {goal}", {"files": target_files, "plan": plan})

    return final_response


def update_workspace_memory_selection(file_names):
    file_names = sorted(set(file_names or []))
    st.session_state.workspace_memory["indexed_files"] = file_names
    save_workspace_memory()


def ensure_workspace_memory_loaded():
    if not st.session_state.workspace_memory_loaded:
        st.session_state.workspace_memory = normalize_workspace_memory(load_workspace_memory())
        st.session_state.workspace_memory_loaded = True
    else:
        st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)


ensure_workspace_memory_loaded()


def get_selection_signature(file_names):
    digest = hashlib.md5()
    for file_name in sorted(file_names):
        digest.update(file_name.encode("utf-8"))
        digest.update(st.session_state.file_texts.get(file_name, "").encode("utf-8"))
    return f"combined::{digest.hexdigest()}"


SUMMARY_STOPWORDS = {
    "the", "and", "for", "with", "that", "this", "from", "are", "was", "were", "into", "your", "have",
    "has", "had", "not", "but", "you", "all", "can", "will", "use", "using", "used", "how", "what",
    "when", "where", "which", "while", "into", "more", "than", "their", "there", "about", "after",
    "before", "within", "without", "each", "page", "pages", "table", "tables", "image", "images",
    "document", "content", "metadata", "information", "product", "file", "text"
}


@st.cache_data(show_spinner=False)
def get_document_asset_counts(file_name, file_bytes, extracted_text):
    file_name_lower = file_name.lower()
    page_count = 0
    table_count = 0
    image_count = 0

    if file_name_lower.endswith(".pdf"):
        page_match = re.search(r"Total Pages:\s*(\d+)", extracted_text)
        page_count = int(page_match.group(1)) if page_match else len(re.findall(r"Page \d+ Text:", extracted_text))
        table_count = len(re.findall(r"Page \d+ Table \d+:", extracted_text))
        image_count = len(re.findall(r"\[IMAGE:", extracted_text))
    elif file_name_lower.endswith(".pptx"):
        slide_match = re.search(r"Total Slides:\s*(\d+)", extracted_text)
        page_count = int(slide_match.group(1)) if slide_match else 0
        table_count = len(re.findall(r"\bTable:\n", extracted_text))
        image_count = len(re.findall(r"\[EMBEDDED_IMAGE:", extracted_text))
    elif file_name_lower.endswith(".docx"):
        table_count = len(re.findall(r"Table \d+:", extracted_text))
        image_count = len(re.findall(r"\[EMBEDDED_IMAGE:", extracted_text))
    elif file_name_lower.endswith(".xlsx"):
        sheet_match = re.search(r"Workbook contains (\d+) sheets", extracted_text)
        page_count = int(sheet_match.group(1)) if sheet_match else 0
        table_count = len(re.findall(r"Sheet '.*?':", extracted_text))
    elif file_name_lower.endswith((".html", ".htm")):
        image_match = re.search(r"(\d+) images found in HTML", extracted_text)
        image_count = int(image_match.group(1)) if image_match else 0

    return page_count, image_count, table_count


def empty_chat_summary_downloads():
    return {"images": [], "tables": [], "csv": [], "diagrams": []}


# Chat summary helper:
# Used only in the Chat tab after summarize/analyze actions to expose extracted
# images and tables as downloads for the files currently selected in chat.
def render_chat_summary_downloads():
    downloads = st.session_state.get("chat_summary_downloads", empty_chat_summary_downloads())
    image_items = downloads.get("images", [])
    table_items = downloads.get("tables", [])
    csv_items = downloads.get("csv", [])
    diagram_items = downloads.get("diagrams", [])

    if not image_items and not table_items and not csv_items and not diagram_items:
        return

    st.markdown("### Summary Downloads")

    if image_items:
        with st.expander("🖼️ Image PNG Downloads", expanded=False):
            for index, item in enumerate(image_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_image_{index}_{item['file_name']}"
                )

    if table_items:
        with st.expander("📊 Table PNG Downloads", expanded=False):
            for index, item in enumerate(table_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_table_{index}_{item['file_name']}"
                )

    if csv_items:
        with st.expander("Pin Table CSV Downloads", expanded=False):
            for index, item in enumerate(csv_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_csv_{index}_{item['file_name']}"
                )

    if diagram_items:
        with st.expander("ASCII Diagram Downloads", expanded=False):
            for index, item in enumerate(diagram_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_diagram_{index}_{item['file_name']}"
                )


def build_adaptive_document_analysis(file_name, file_bytes, text):
    raw_text = str(text or "")
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    words = re.findall(r"\w+", raw_text)
    title_match = re.search(r"Title:\s*(.+)", raw_text)
    title = title_match.group(1).strip() if title_match and title_match.group(1).strip() else file_name

    keyword_counts = Counter(
        word.lower()
        for word in words
        if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    keywords = [word.title() for word, _ in keyword_counts.most_common(8)]
    keyword_text = ", ".join(keywords) if keywords else "Not available"
    page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, raw_text)

    ignored_prefixes = (
        "pdf metadata:", "document metadata:", "meta tags:", "total pages:", "total slides:",
        "workbook contains", "error:", "[image:", "[embedded_image:", "table:"
    )
    metadata_prefixes = (
        "producer:", "creationdate:", "moddate:", "author:", "creator:", "title:",
        "subject:", "keywords:", "trapped:", "pdfversion:"
    )

    def prettify_extracted_text(value):
        value = str(value or "").strip()
        if not value:
            return value
        value = re.sub(r"([a-z])([A-Z])", r"\1 \2", value)
        value = re.sub(r"([A-Za-z])(\d)", r"\1 \2", value)
        value = re.sub(r"(\d)([A-Za-z])", r"\1 \2", value)
        value = re.sub(r"\s+", " ", value)
        value = value.replace("e. g.", "e.g.").replace("i. e.", "i.e.")
        return value.strip()

    keywords = [prettify_extracted_text(keyword) for keyword in keywords]
    keyword_text = ", ".join(keywords) if keywords else "Not available"

    def clean_content_lines(max_items=12):
        cleaned = []
        seen = set()
        for line in lines:
            line_lower = line.lower()
            if line_lower.startswith(ignored_prefixes):
                continue
            if line_lower.startswith(metadata_prefixes):
                continue
            if len(line) < 8 or len(line) > 240:
                continue
            if re.fullmatch(r"[\W_]+", line):
                continue
            line = prettify_extracted_text(line)
            line_lower = line.lower()
            if line in seen:
                continue
            seen.add(line)
            cleaned.append(line)
            if len(cleaned) >= max_items:
                break
        return cleaned

    key_lines = clean_content_lines(12)
    headings = extract_document_headings(raw_text)
    toc_entries = extract_toc_with_page_numbers(raw_text)
    lower_text = raw_text.lower()
    file_name_lower = file_name.lower()
    type_scores = {
        "technical": sum(1 for term in [
            "architecture", "system", "module", "component", "workflow", "api", "interface",
            "configuration", "requirement", "software", "hardware", "capl", "diagnostic", "test"
        ] if term in lower_text),
        "business": sum(1 for term in [
            "strategy", "market", "customer", "revenue", "business", "goal", "objective",
            "stakeholder", "risk", "cost", "benefit", "performance", "operation"
        ] if term in lower_text),
        "research": sum(1 for term in [
            "abstract", "methodology", "experiment", "hypothesis", "dataset", "findings",
            "results", "conclusion", "references", "study", "analysis"
        ] if term in lower_text),
    }
    if file_name_lower.endswith((".can", ".capl")):
        type_scores["technical"] += 4
    if file_name_lower.endswith((".xlsx", ".html", ".htm")):
        type_scores["business"] += 1
    document_type = max(type_scores, key=type_scores.get) if max(type_scores.values() or [0]) > 0 else "general"

    def pick_lines(patterns, limit=5):
        selected = []
        seen = set()
        for line in key_lines + lines:
            line = prettify_extracted_text(line)
            if len(line) < 8 or len(line) > 260:
                continue
            line_lower = line.lower()
            if line_lower.startswith(ignored_prefixes) or line_lower.startswith(metadata_prefixes):
                continue
            if any(pattern in line_lower for pattern in patterns) and line not in seen:
                selected.append(line)
                seen.add(line)
            if len(selected) >= limit:
                break
        return selected

    feature_lines = pick_lines(["feature", "component", "module", "function", "capability", "system", "interface", "configuration"])
    workflow_lines = pick_lines(["step", "process", "workflow", "flow", "first", "then", "after", "before", "execute", "upload", "select"])
    use_case_lines = pick_lines(["use case", "application", "used for", "used to", "can be used", "supports", "helps", "enables"])
    important_note_lines = pick_lines(["warning", "caution", "note", "limit", "constraint", "assumption", "must", "shall", "required", "error"])

    context_bits = []
    if page_count:
        context_bits.append(f"{page_count} pages/sections")
    if image_count:
        context_bits.append(f"{image_count} images")
    if table_count:
        context_bits.append(f"{table_count} tables")
    context_text = ", ".join(context_bits) if context_bits else f"{len(lines)} content lines"

    def bullet_list(items, fallback):
        usable_items = [item for item in items if item] or fallback
        return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in usable_items) + "</ul>"

    def section(title_text, body_html):
        if not body_html:
            return ""
        return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title_text)}</h4>{body_html}"

    structure_items = []
    if toc_entries:
        structure_items = [
            f"{num + ' ' if num else ''}{prettify_extracted_text(heading)}" + (f" - page {page_num}" if page_num else "")
            for num, heading, page_num in toc_entries[:6]
        ]
    elif headings:
        structure_items = [f"{num + ' ' if num else ''}{prettify_extracted_text(heading)}" for num, heading in headings[:6]]
    else:
        structure_items = [
            "Content is presented as extracted text rather than clearly labeled sections.",
            f"Detected document assets: {context_text}.",
        ]

    purpose_by_type = {
        "technical": "to describe a system, process, implementation, test, or technical capability",
        "business": "to communicate objectives, operational context, metrics, or decision-oriented information",
        "research": "to explain a problem, method, evidence, findings, and conclusions",
        "general": "to present information in a readable and referenceable form",
    }
    summary_focus = ", ".join(keywords[:4]) if keywords else "the extracted document content"
    key_point_items = key_lines[:5] or ["No detailed content lines could be extracted, but document metadata was detected."]
    insight_items = [
        f"The document appears to focus on {summary_focus}.",
        f"It should be read as a {document_type} document.",
    ]
    if image_count or table_count:
        insight_items.append("Visual or tabular assets may contain supporting details that complement the extracted text.")
    if important_note_lines:
        insight_items.append("Several lines contain requirements, constraints, warnings, or operational notes.")

    simplified_items = [
        f"In simple terms, this document is about {summary_focus}.",
        "It collects the main information a reader needs to understand the topic, context, and next actions.",
    ]
    takeaway_items = []
    if keywords:
        takeaway_items.append(f"Primary themes: {', '.join(keywords[:5])}.")
    takeaway_items.extend(key_point_items[:4])
    takeaway_items = takeaway_items[:5]

    summary_html = (
        "<div>"
        f"<div><b>What the document is about:</b> {html.escape(title)}</div>"
        f"<div><b>Main purpose:</b> {html.escape(purpose_by_type[document_type])}.</div>"
        f"<div><b>Key context:</b> {html.escape(context_text)}. Detected type: {html.escape(document_type.title())}.</div>"
        f"<div><b>Important themes:</b> {html.escape(keyword_text)}</div>"
        "</div>"
    )

    optional_sections = ""
    if feature_lines or document_type == "technical":
        optional_sections += section("Features / Concepts / Components", bullet_list(
            feature_lines,
            [f"Relevant concepts include {keyword_text}.", "No explicit feature list was detected in the extracted text."]
        ))
    if workflow_lines or document_type == "technical":
        optional_sections += section("Workflow / Process", bullet_list(
            workflow_lines,
            ["No clear step-by-step workflow was detected in the extracted text."]
        ))
    if use_case_lines or document_type in ("technical", "business"):
        optional_sections += section("Use Cases / Applications", bullet_list(
            use_case_lines,
            ["Use this document as a reference for understanding the topic, validating details, or planning related work."]
        ))
    if important_note_lines:
        optional_sections += section("Important Notes", bullet_list(important_note_lines, []))

    analysis_sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Document Analysis: {html.escape(file_name)}</h3>",
        section("Summary", summary_html),
        section("Key Points", bullet_list(key_point_items, [])),
        section("Structure Breakdown", bullet_list(structure_items, [])),
        section("Key Insights / Core Insights", bullet_list(insight_items, [])),
        optional_sections,
        section("Simplified Explanation", bullet_list(simplified_items, [])),
        section("Key Takeaways", bullet_list(takeaway_items, [])),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(analysis_sections) + "</div>"


def build_product_documentation_analysis(file_name, file_bytes, text):
    raw_text = str(text or "")
    lines = [normalize_extracted_line(line) for line in raw_text.splitlines() if line.strip()]
    lower_text = raw_text.lower()
    words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{2,}", raw_text)

    ignored_prefixes = (
        "pdf metadata:", "document metadata:", "odt metadata:", "meta tags:", "total pages:",
        "total slides:", "workbook contains", "csv rows:", "error:", "[image:", "[embedded_image:",
        "page ", "slide ", "sheet ", "table:"
    )
    metadata_prefixes = (
        "producer:", "creationdate:", "moddate:", "author:", "creator:", "title:",
        "subject:", "keywords:", "trapped:", "pdfversion:"
    )

    def clean_sentence(value):
        value = normalize_extracted_line(value)
        value = re.sub(r"^(?:page|slide|sheet)\s+\d+\s*(?:text|content)?\s*:?", "", value, flags=re.IGNORECASE).strip()
        value = re.sub(r"\s+", " ", value)
        return value.strip(" -:")

    def meaningful_lines(max_items=180):
        selected = []
        seen = set()
        for line in lines:
            cleaned = clean_sentence(line)
            if not cleaned:
                continue
            lowered = cleaned.lower()
            if lowered.startswith(ignored_prefixes) or lowered.startswith(metadata_prefixes):
                continue
            if len(cleaned) < 18 or len(cleaned) > 220:
                continue
            if re.fullmatch(r"[\W_]+", cleaned):
                continue
            if re.fullmatch(r"\d+(?:\.\d+)*\s+.+", cleaned) and len(cleaned.split()) <= 7:
                continue
            key = lowered
            if key in seen:
                continue
            seen.add(key)
            selected.append(cleaned)
            if len(selected) >= max_items:
                break
        return selected

    clean_lines = meaningful_lines()
    keyword_counts = Counter(
        word.lower()
        for word in words
        if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    keywords = [normalize_extracted_line(word).title() for word, _ in keyword_counts.most_common(10)]

    type_scores = {
        "technical system": sum(1 for term in [
            "interface", "module", "device", "hardware", "software", "configuration", "channel",
            "connector", "signal", "protocol", "diagnostic", "architecture", "firmware", "driver"
        ] if term in lower_text),
        "process or workflow": sum(1 for term in [
            "process", "workflow", "procedure", "step", "approval", "operation", "execute", "setup",
            "install", "configure", "use", "report"
        ] if term in lower_text),
        "data or report": sum(1 for term in [
            "metric", "statistics", "result", "dashboard", "table", "test case", "passed", "failed",
            "executed", "summary", "analysis"
        ] if term in lower_text),
        "business document": sum(1 for term in [
            "customer", "market", "objective", "stakeholder", "cost", "risk", "benefit", "strategy",
            "requirement", "decision"
        ] if term in lower_text),
    }
    document_kind = max(type_scores, key=type_scores.get) if max(type_scores.values() or [0]) > 0 else "reference document"

    page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, raw_text)
    title_match = re.search(r"Title:\s*(.+)", raw_text)
    detected_title = clean_sentence(title_match.group(1)) if title_match else ""
    display_name = detected_title if detected_title and len(detected_title) < 120 else os.path.splitext(file_name)[0]
    topic_terms = keywords[:5] or [display_name]
    topic_phrase = ", ".join(topic_terms[:4])

    def collect_by_terms(terms, limit=6):
        selected = []
        seen = set()
        for line in clean_lines:
            lowered = line.lower()
            if any(term in lowered for term in terms):
                simplified = synthesize_line(line)
                if simplified and simplified.lower() not in seen:
                    selected.append(simplified)
                    seen.add(simplified.lower())
            if len(selected) >= limit:
                break
        return selected

    def synthesize_line(line):
        line = clean_sentence(line)
        if not line:
            return ""
        line = re.sub(r"\b(?:note|warning|caution)\s*[:\-]\s*", "", line, flags=re.IGNORECASE)
        line = re.sub(r"\s*\.+\s*\d+\s*$", "", line)
        if len(line) > 150:
            line = line[:147].rsplit(" ", 1)[0] + "..."
        return line[0].upper() + line[1:] if line else line

    capability_terms = [
        "support", "supports", "feature", "function", "capability", "enable", "allows", "provide",
        "communication", "measurement", "analysis", "diagnostic", "configuration", "export"
    ]
    architecture_terms = [
        "component", "module", "interface", "channel", "connector", "port", "device", "unit",
        "software", "hardware", "driver", "network", "table", "sheet", "slide"
    ]
    workflow_terms = [
        "install", "configure", "connect", "select", "upload", "execute", "start", "use",
        "create", "open", "set", "download", "export", "analyze"
    ]
    use_case_terms = [
        "application", "used for", "used to", "use case", "measurement", "testing", "diagnostic",
        "report", "automation", "monitoring", "analysis", "configuration"
    ]

    capabilities = collect_by_terms(capability_terms, 7)
    architecture_evidence = collect_by_terms(architecture_terms, 7)
    workflow_evidence = collect_by_terms(workflow_terms, 6)
    use_case_evidence = collect_by_terms(use_case_terms, 6)

    def has_any(*terms):
        return any(term in lower_text for term in terms)

    generated_capabilities = []
    if has_any("can ", "can-fd", "can fd", "lin", "flexray", "ethernet", "protocol", "interface", "communication"):
        generated_capabilities.append("Supports communication-oriented work through documented protocols, interfaces, or channels.")
    if has_any("configuration", "configure", "setup", "install", "driver", "software"):
        generated_capabilities.append("Provides configuration and setup guidance so the system can be prepared for practical use.")
    if has_any("measurement", "test", "diagnostic", "analysis", "monitor", "report"):
        generated_capabilities.append("Supports analysis, measurement, diagnostics, reporting, or validation activities.")
    if has_any("connector", "pin", "port", "socket", "plug", "channel"):
        generated_capabilities.append("Documents physical or logical connectivity details needed for integration.")
    if has_any("table", "sheet", "csv", "dashboard", "statistics", "result"):
        generated_capabilities.append("Contains structured data or results that can be reviewed, summarized, or exported.")
    if has_any("image", "figure", "diagram", "visual", "illustration"):
        generated_capabilities.append("Includes visual or diagram-like information that can support engineering reference work.")
    if generated_capabilities:
        capabilities = generated_capabilities

    generated_architecture = []
    if has_any("hardware", "device", "unit", "module", "component"):
        generated_architecture.append("Hardware or device layer: the physical units, modules, or components described by the source.")
    if has_any("software", "driver", "application", "tool", "configuration"):
        generated_architecture.append("Software and configuration layer: tools, drivers, settings, and setup behavior around the system.")
    if has_any("interface", "protocol", "channel", "network", "communication"):
        generated_architecture.append("Communication layer: interfaces, protocols, and channels that connect the system to other tools or networks.")
    if has_any("connector", "pin", "port", "socket", "plug"):
        generated_architecture.append("Connectivity layer: ports, connectors, pin assignments, or wiring-related details.")
    if has_any("table", "sheet", "report", "result", "metadata"):
        generated_architecture.append("Information layer: tables, results, metadata, and reference data used to interpret the document.")
    if generated_architecture:
        architecture_evidence = generated_architecture

    generated_workflow = []
    if has_any("upload", "select", "open", "choose"):
        generated_workflow.append("Select the relevant file, section, component, or dataset.")
    if has_any("install", "setup", "driver", "connect"):
        generated_workflow.append("Prepare the environment by installing, connecting, or setting up the required parts.")
    if has_any("configure", "configuration", "setting", "parameter"):
        generated_workflow.append("Configure the required options, channels, interfaces, or parameters.")
    if has_any("execute", "run", "start", "measurement", "test", "analysis"):
        generated_workflow.append("Run the intended operation such as measurement, testing, communication, analysis, or review.")
    if has_any("result", "report", "export", "download", "table"):
        generated_workflow.append("Review outputs, results, tables, or reports and export anything needed for reference.")
    if generated_workflow:
        workflow_evidence = generated_workflow

    generated_use_cases = []
    if has_any("measurement", "canalyzer", "canoe", "diagnostic", "test"):
        generated_use_cases.append("Vehicle/network measurement, diagnostics, testing, and validation workflows.")
    if has_any("configuration", "install", "setup", "driver"):
        generated_use_cases.append("Setup and configuration reference for engineers or technicians.")
    if has_any("interface", "connector", "pin", "channel", "protocol"):
        generated_use_cases.append("Integration reference for ports, channels, protocols, connectors, or pin mappings.")
    if has_any("report", "dashboard", "statistics", "result", "table"):
        generated_use_cases.append("Report review, structured data analysis, and documentation support.")
    if has_any("warning", "caution", "note", "safety", "required"):
        generated_use_cases.append("Operational guidance where constraints, warnings, or required practices matter.")
    if generated_use_cases:
        use_case_evidence = generated_use_cases

    components = []
    component_candidates = []
    for pattern in [
        r"\b[A-Z]{2,}[A-Za-z0-9_+\-/]*\b",
        r"\b[A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z0-9]+){0,2}\b",
    ]:
        component_candidates.extend(re.findall(pattern, raw_text))
    component_counts = Counter(
        normalize_extracted_line(candidate).strip()
        for candidate in component_candidates
        if 3 <= len(normalize_extracted_line(candidate).strip()) <= 45
        and normalize_extracted_line(candidate).lower() not in SUMMARY_STOPWORDS
    )
    for candidate, _ in component_counts.most_common(8):
        lowered = candidate.lower()
        if lowered in {"pdf", "metadata", "page", "text", "table", "figure"}:
            continue
        components.append(f"{candidate}: appears to be a major referenced part, concept, interface, or artifact in the document.")
        if len(components) >= 5:
            break

    if not components and architecture_evidence:
        components = architecture_evidence[:5]

    assets = []
    if page_count:
        assets.append(f"about {page_count} pages or sections")
    if table_count:
        assets.append(f"{table_count} table-like data areas")
    if image_count:
        assets.append(f"{image_count} visual assets")
    asset_phrase = ", ".join(assets) if assets else "the available extracted content"

    overview_items = [
        f"This is a {document_kind} centered on {topic_phrase}.",
        f"It serves as a practical reference for understanding the subject, its purpose, and how the relevant pieces fit together.",
    ]
    if assets:
        overview_items.append(f"The source contains {asset_phrase}, but the summary below reorganizes the content by meaning rather than document order.")

    core_concept_items = [
        f"In simple terms, the document explains how {topic_terms[0] if topic_terms else 'the subject'} is used, configured, or understood in context.",
        "The important ideas are grouped into purpose, structure, capabilities, usage flow, and practical value so a reader can act on them quickly.",
    ]
    if capabilities:
        core_concept_items.append(f"The central behavior is reflected in capabilities such as {', '.join(keywords[:4])}.")

    if not architecture_evidence:
        architecture_evidence = [
            "The document content is best understood as a set of related concepts, interfaces, configuration details, and operational notes.",
            "Related elements are grouped logically instead of following the original document layout."
        ]

    if not capabilities:
        capabilities = [
            "Provides reference information needed to understand and apply the documented subject.",
            "Combines functional context with technical details where the source provides them."
        ]

    if not workflow_evidence:
        workflow_evidence = [
            "Identify the relevant subject or component.",
            "Review its purpose, interfaces, configuration needs, and constraints.",
            "Apply the information in implementation, testing, documentation, or troubleshooting work."
        ]

    if not use_case_evidence:
        use_case_evidence = [
            "Engineering reference and onboarding.",
            "Configuration or implementation planning.",
            "Troubleshooting, validation, and documentation support."
        ]

    takeaway_items = []
    if keywords:
        takeaway_items.append(f"The main focus areas are {', '.join(keywords[:5])}.")
    takeaway_items.append(f"The document is most useful as a {document_kind} rather than as a narrative document.")
    if capabilities:
        takeaway_items.append("The key value is translating scattered technical or functional details into usable reference knowledge.")
    if architecture_evidence:
        takeaway_items.append("Understanding the relationships between components, interfaces, and usage flow is more important than memorizing the original section order.")
    takeaway_items = takeaway_items[:5]

    def bullet_list(items):
        clean_items = [item for item in items if item]
        return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in clean_items) + "</ul>"

    def section(title_text, items):
        if not items:
            return ""
        return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title_text)}</h4>{bullet_list(items)}"

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Document Analysis: {html.escape(file_name)}</h3>",
        section("Overview", overview_items),
        section("Core Concept", core_concept_items),
        section("Architecture / Structure", architecture_evidence[:7]),
        section("Key Capabilities", capabilities[:7]),
        section("Components / Modules", components[:6]),
        section("Workflow / How It Is Used", workflow_evidence[:6]),
        section("Practical Use Cases", use_case_evidence[:6]),
        section("Key Takeaways", takeaway_items[:5]),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(part for part in sections if part) + "</div>"


def build_detailed_document_summary(file_name, file_bytes, text):
    return build_product_documentation_analysis(file_name, file_bytes, text)


def extract_quoted_item_name(user_input):
    match = re.search(r"'(.*?)'|\"(.*?)\"", str(user_input or ""))
    if match:
        return (match.group(1) or match.group(2) or "").strip()

    patterns = [
        r"\b(?:item|about|for|related to)\s+([A-Za-z0-9][A-Za-z0-9 _./+\-]{1,80})",
        r"\b(?:pin(?:s)?|diagram|connector|visual)\s+([A-Za-z0-9][A-Za-z0-9 _./+\-]{1,80})",
    ]
    for pattern in patterns:
        match = re.search(pattern, str(user_input or ""), re.IGNORECASE)
        if match:
            item = re.split(r"\b(?:from|in|with|please|and|details?|info|information)\b", match.group(1), 1, flags=re.IGNORECASE)[0]
            return item.strip(" :-")
    return ""


def normalize_extracted_line(line):
    line = str(line or "").strip()
    line = re.sub(r"([a-z])([A-Z])", r"\1 \2", line)
    line = re.sub(r"([A-Za-z])(\d)", r"\1 \2", line)
    line = re.sub(r"(\d)([A-Za-z])", r"\1 \2", line)
    line = re.sub(r"\s+", " ", line)
    return line.strip()


def collect_item_context_lines(text, item_name, window=4, limit=80):
    item_name = str(item_name or "").strip()
    if not item_name:
        return []

    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    item_tokens = [token.lower() for token in re.findall(r"[A-Za-z0-9]+", item_name) if len(token) > 1]
    if not item_tokens:
        return []

    selected = []
    seen = set()
    for index, line in enumerate(lines):
        line_lower = line.lower()
        compact_line = re.sub(r"\s+", "", line_lower)
        compact_item = re.sub(r"\s+", "", item_name.lower())
        has_match = compact_item in compact_line or all(token in line_lower for token in item_tokens)
        if not has_match:
            continue

        start = max(0, index - window)
        end = min(len(lines), index + window + 1)
        for context_line in lines[start:end]:
            pretty_line = normalize_extracted_line(context_line)
            if len(pretty_line) < 3 or len(pretty_line) > 300:
                continue
            key = pretty_line.lower()
            if key in seen:
                continue
            seen.add(key)
            selected.append(pretty_line)
            if len(selected) >= limit:
                return selected
    return selected


def select_relevant_lines(context_lines, patterns, limit=8):
    selected = []
    seen = set()
    for line in context_lines:
        line_lower = line.lower()
        if any(pattern in line_lower for pattern in patterns) and line not in seen:
            selected.append(line)
            seen.add(line)
        if len(selected) >= limit:
            break
    return selected


def html_bullet_list(items):
    if not items:
        return ""
    return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in items) + "</ul>"


def html_section(title, items):
    if not items:
        return ""
    return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title)}</h4>{html_bullet_list(items)}"


def build_item_information_response(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=5, limit=100)
    if not context_lines:
        return f"<div><h3>Item Information: {html.escape(item_name)}</h3><p>No relevant information for this item was found in {html.escape(file_name)}.</p></div>"

    overview = context_lines[:5]
    features = select_relevant_lines(context_lines, ["feature", "support", "capability", "function", "operation", "application"])
    technical = select_relevant_lines(context_lines, ["mbit", "kbit", "volt", "channel", "standard", "protocol", "can", "lin", "flexray", "interface", "specification", "iso"])
    architecture = select_relevant_lines(context_lines, ["structure", "component", "module", "piggy", "channel", "internal", "family", "device"])
    configuration = select_relevant_lines(context_lines, ["configure", "configuration", "install", "insert", "setup", "use", "driver", "software", "hardware"])
    connectivity = select_relevant_lines(context_lines, ["connector", "port", "pin", "d-sub", "usb", "channel", "plug", "socket", "interface"])
    special = select_relevant_lines(context_lines, ["special", "unique", "only", "limitation", "difference", "optional", "available", "not supported"])
    notes = select_relevant_lines(context_lines, ["note", "warning", "caution", "must", "shall", "important", "avoid", "required"])

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Item Information: {html.escape(item_name)}</h3>",
        f"<p><b>Source:</b> {html.escape(file_name)}</p>",
        html_section("Overview", overview),
        html_section("Key Features", features),
        html_section("Technical Details", technical),
        html_section("Architecture / Structure", architecture),
        html_section("Configuration / Usage", configuration),
        html_section("Interfaces & Connectivity", connectivity),
        html_section("Special Capabilities", special),
        html_section("Important Notes", notes),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(section for section in sections if section) + "</div>"


def extract_pin_rows(context_lines):
    rows = []
    seen = set()
    pin_patterns = [
        r"\bpin\s*(\d+)\b\s*[:\-]?\s*([A-Za-z0-9_+/.\- ]{0,40})\s*(.*)",
        r"^\s*(\d{1,2})\s+([A-Za-z][A-Za-z0-9_+/.\-]*)\s*(.*)",
    ]
    for line in context_lines:
        line_lower = line.lower()
        if not any(term in line_lower for term in ["pin", "signal", "d-sub", "connector", "ground", "shield", "can", "lin", "vbat"]):
            continue
        for pattern in pin_patterns:
            match = re.search(pattern, line, re.IGNORECASE)
            if not match:
                continue
            pin_no = match.group(1).strip()
            signal = (match.group(2) or "").strip(" :-") or "Not specified"
            description = (match.group(3) or "").strip(" :-") or line
            key = (pin_no, signal.lower(), description.lower())
            if key in seen:
                break
            seen.add(key)
            rows.append({
                "pin": pin_no,
                "signal": signal,
                "description": description,
                "notes": ""
            })
            break
    return rows[:40]


def build_pin_csv(pin_rows):
    lines = ["Pin Number,Signal Name,Description,Notes"]
    for row in pin_rows:
        values = [row["pin"], row["signal"], row["description"], row.get("notes", "")]
        escaped_values = ['"' + str(value).replace('"', '""') + '"' for value in values]
        lines.append(",".join(escaped_values))
    return "\n".join(lines)


def build_ascii_pin_diagram(pin_rows, item_name):
    if not pin_rows:
        return f"+------------------------------+\n| {item_name[:28]:<28} |\n| Pin diagram not available    |\n+------------------------------+"
    left = pin_rows[::2]
    right = pin_rows[1::2]
    width = 34
    lines = [f"+{'-' * width}+", f"| {item_name[:width-4]:<{width-4}} |", f"+{'-' * width}+"]
    max_len = max(len(left), len(right))
    for index in range(max_len):
        left_text = ""
        right_text = ""
        if index < len(left):
            left_text = f"{left[index]['pin']}:{left[index]['signal']}"[:15]
        if index < len(right):
            right_text = f"{right[index]['pin']}:{right[index]['signal']}"[:15]
        lines.append(f"| {left_text:<15}  {right_text:>15} |")
    lines.append(f"+{'-' * width}+")
    return "\n".join(lines)


def build_item_visual_assets(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=8, limit=140)
    pin_rows = extract_pin_rows(context_lines)
    if not pin_rows:
        return {"csv": [], "diagrams": []}

    safe_item_name = re.sub(r"[^A-Za-z0-9_-]+", "_", str(item_name)).strip("_") or "item"
    file_base = re.sub(r"[^A-Za-z0-9_-]+", "_", os.path.splitext(file_name)[0]).strip("_") or "document"
    csv_text = build_pin_csv(pin_rows)
    ascii_diagram = build_ascii_pin_diagram(pin_rows, item_name)

    return {
        "csv": [{
            "label": f"{file_name} - {item_name} pin table CSV",
            "data": csv_text.encode("utf-8"),
            "file_name": f"{file_base}_{safe_item_name}_pin_table.csv",
            "mime": "text/csv",
        }],
        "diagrams": [{
            "label": f"{file_name} - {item_name} ASCII diagram",
            "data": ascii_diagram.encode("utf-8"),
            "file_name": f"{file_base}_{safe_item_name}_diagram.txt",
            "mime": "text/plain",
        }],
    }


def html_table(headers, rows):
    if not rows:
        return ""
    head_html = "".join(f"<th>{html.escape(header)}</th>" for header in headers)
    body_html = ""
    for row in rows:
        body_html += "<tr>" + "".join(f"<td>{html.escape(str(cell))}</td>" for cell in row) + "</tr>"
    return f"<table style='border-collapse:collapse; width:100%; margin:8px 0;'><thead><tr>{head_html}</tr></thead><tbody>{body_html}</tbody></table>"


def build_item_visual_response(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=8, limit=140)
    if not context_lines:
        return f"<div><h3>Visual / Pin Reference: {html.escape(item_name)}</h3><p>No relevant visual or structural information for this item was found in {html.escape(file_name)}.</p></div>"

    pin_rows = extract_pin_rows(context_lines)
    connector_lines = select_relevant_lines(context_lines, ["connector", "port", "d-sub", "usb", "channel", "plug", "socket", "interface"], limit=12)
    image_lines = select_relevant_lines(context_lines, ["figure", "image", "diagram", "pin assignment", "illustration"], limit=10)
    table_lines = select_relevant_lines(context_lines, ["table", "specification", "signal", "configuration", "pin"], limit=12)
    csv_text = build_pin_csv(pin_rows) if pin_rows else "Pin Number,Signal Name,Description,Notes\n"
    ascii_diagram = build_ascii_pin_diagram(pin_rows, item_name)
    pin_table_rows = [[row["pin"], row["signal"], row["description"], row.get("notes", "")] for row in pin_rows]

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Visual / Pin Reference: {html.escape(item_name)}</h3>",
        f"<p><b>Source:</b> {html.escape(file_name)}</p>",
        html_section("Pin Diagrams", ["Recreated below from extracted pin/signal lines." if pin_rows else "No explicit pin diagram was found in the extracted text."]),
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>",
        f"<h4 style='margin:16px 0 6px 0; color:#173152;'>Pin Configuration Table</h4>",
        html_table(["Pin Number", "Signal Name", "Description", "Notes"], pin_table_rows) if pin_rows else "<p>No pin table data was found.</p>",
        html_section("Connector Details", connector_lines),
        html_section("Images & Visuals", image_lines),
        html_section("Technical Tables", table_lines),
        f"<h4 style='margin:16px 0 6px 0; color:#173152;'>Downloadable Outputs</h4>",
        "<p><b>a) Pin table as CSV</b></p>",
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(csv_text)}</pre>",
        "<p><b>b) Diagram as ASCII / structured format</b></p>",
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>",
        "<p><b>c) Image references or recreated diagrams</b></p>",
        html_bullet_list(image_lines or ["No direct image reference was found in extracted text; use the recreated ASCII diagram above when pin rows are available."]),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(section for section in sections if section) + "</div>"


# Document structure helpers:
# Used by the Chat "overview" flow and preview links to identify headings,
# table-of-contents entries, and likely page numbers inside uploaded documents.
def extract_page_text(text, page_number=1):
    text = str(text)
    pattern = rf"Page {page_number}\s+Text:\s*(.*?)(?=Page \d+\s+Text:|\Z)"
    match = re.search(pattern, text, re.S | re.IGNORECASE)
    if match:
        return match.group(1).strip()

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines[:80])


def find_heading_page_number(text, heading):
    text = str(text)
    lines = [line for line in text.splitlines()]
    heading_pattern = re.escape(str(heading).strip())
    for index, line in enumerate(lines):
        if re.search(rf"\b{heading_pattern}\b", line, re.IGNORECASE):
            for j in range(index, -1, -1):
                page_match = re.search(r'Page\s+(\d+)\s+Text:', lines[j], re.IGNORECASE)
                if page_match:
                    return int(page_match.group(1))
    return None


def resolve_heading_page_number(text, heading, toc_entries=None):
    if not heading:
        return None
    heading_text = str(heading).strip()
    if toc_entries is None:
        toc_entries = extract_toc_with_page_numbers(text)
    for num, title, page_num in toc_entries:
        if title.strip().lower() == heading_text.lower():
            return page_num
        if heading_text.lower() in title.strip().lower() or title.strip().lower() in heading_text.lower():
            return page_num
    return find_heading_page_number(text, heading_text)


def extract_document_headings(text):
    """Extract numbered headings and explicit DOCX headings from extracted text."""
    headings = []
    text = str(text)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    
    for line in lines:
        # Skip lines that are too long
        if len(line) > 120:
            continue
        
        # Skip metadata, page markers, and special content
        if (line.isupper() or line.endswith(":") or 
            "Page" in line or "PDF Metadata" in line or 
            "Total Pages" in line or "TABLE:" in line):
            continue

        # Match explicit heading markers from DOCX extraction
        if line.startswith("Heading:"):
            heading_text = line.replace("Heading:", "", 1).strip()
            if 3 <= len(heading_text) <= 120:
                headings.append(("", heading_text))
            continue
        
        # Match numbered headings at start: "1 Overview", "1.1 Introduction", etc.
        match = re.match(r'^(\d+(?:\.\d+)*)\s+([A-Za-z\s][^.]*?)(?:\s*\.+\s*\d+)?\s*$', line)
        if match:
            num = match.group(1)
            title = match.group(2).strip()
            
            # Clean up any trailing dots or page numbers
            title = re.sub(r'\s*\.+\s*\d*\s*$', '', title).strip()
            
            if 3 <= len(title) <= 120:
                headings.append((num, title))
    
    # Remove duplicates while preserving order
    seen = set()
    deduped = []
    for num, title in headings:
        key = f"{num}:{title}"
        if key not in seen:
            seen.add(key)
            deduped.append((num, title))
    
    return deduped


def extract_toc_with_page_numbers(text):
    """Extract table of contents entries with page numbers from document."""
    toc_entries = []
    text = str(text)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    
    # First, try explicit TOC patterns on full text
    for regex in [
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s+\.{2,}\s*(\d+)\s*$',
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s{3,}(\d+)\s*$',
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s+(\d+)\s*$'
    ]:
        for match in re.finditer(regex, text):
            num = match.group(1)
            title = match.group(2).strip()
            page_num = match.group(3)
            if 3 <= len(title) <= 120 and len(re.findall(r'\d+', title)) <= 2:
                toc_entries.append((num, title, page_num))
        if toc_entries:
            return toc_entries

    # Fallback: build TOC from detected headings and page markers
    headings = extract_document_headings(text)
    if headings:
        for num, title in headings:
            page_num = None
            search_pattern = re.escape(title)
            for i, line in enumerate(lines):
                if title in line or re.search(search_pattern, line, re.IGNORECASE):
                    for j in range(i, max(0, i - 20), -1):
                        page_match = re.search(r'Page\s+(\d+)\s+Text:', lines[j])
                        if page_match:
                            page_num = page_match.group(1)
                            break
                    if page_num:
                        break
            toc_entries.append((num, title, page_num))
    return toc_entries


def build_file_overview(file_name, text):
    text = str(text)
    toc_entries = extract_toc_with_page_numbers(text)
    all_headings = extract_document_headings(text)

    overview_parts = [f"📄 **{file_name}**"]
    
    # Table of Contents section
    overview_parts.append("### Table of Contents")
    if toc_entries:
        overview_parts.append("| Contents | Page No |")
        overview_parts.append("|----------|---------|")
        for num, title, page_num in toc_entries:
            content_str = f"{num} {title}" if num else title
            display_text = f"{content_str} (Page {page_num})" if page_num else content_str
            preview_link = create_preview_link(file_name, highlight_term=title, page_num=page_num)
            anchor_id = create_heading_anchor(title)
            if preview_link:
                page_display = page_num if page_num else "-"
                overview_parts.append(f"| <a href='{preview_link}#{anchor_id}' target='_blank'>{html.escape(display_text)}</a> | {page_display} |")
            else:
                page_display = page_num if page_num else "-"
                overview_parts.append(f"| {html.escape(display_text)} | {page_display} |")
    else:
        overview_parts.append("- No table of contents found with page numbers.")

    # Document Headings section
    overview_parts.append("### Document Headings")
    if all_headings:
        for num, title in all_headings:
            content_str = f"{num} {title}" if num else title
            anchor_id = create_heading_anchor(title)
            page_num = resolve_heading_page_number(text, title, toc_entries)
            preview_link = create_preview_link(file_name, highlight_term=title, page_num=page_num)
            if preview_link:
                overview_parts.append(f"- <a href='{preview_link}#{anchor_id}' target='_blank'>{html.escape(content_str)}</a>")
            else:
                overview_parts.append(f"- {content_str}")
    else:
        overview_parts.append("- No document headings were detected.")

    return "\n".join(overview_parts)


@st.cache_data(show_spinner=False)
def build_highlighted_search_results(file_name, text, query):
    if not query:
        return ""

    pattern = re.compile(re.escape(query), re.IGNORECASE)
    matches = []

    for line_no, raw_line in enumerate(str(text).splitlines(), 1):
        if pattern.search(raw_line):
            escaped_line = html.escape(raw_line)
            highlighted_line = pattern.sub(
                lambda match: f"<mark style='background:#fff3a3; padding:0 2px;'>{html.escape(match.group(0))}</mark>",
                escaped_line
            )
            matches.append(
                f"<div style='margin:0 0 8px 0;'><b>Line {line_no}</b>: {highlighted_line}</div>"
            )

    if not matches:
        return f"<div><b>{html.escape(file_name)}</b><br>No matches found for <code>{html.escape(query)}</code>.</div>"

    return (
        f"<div style='margin-bottom:14px;'>"
        f"<h4 style='margin:0 0 8px 0; color:#a8d8f0;'>{html.escape(file_name)} ({len(matches)} matches)</h4>"
        f"{''.join(matches)}"
        f"</div>"
    )


@st.cache_data(show_spinner=False)
def extract_login_name_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    text = soup.get_text(" ", strip=True)
    match = re.search(r'login name[:\s]+(.+?)(version|$)', text, re.IGNORECASE)
    if match:
        name = match.group(1).strip()
        parts = name.split()
        return " ".join(parts[:1])
    return "Not found"


@st.cache_data(show_spinner=False)
def extract_statistics_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    stats = {
        "Executed": 0,
        "Passed": 0,
        "Failed": 0,
        "Inconclusive": 0,
        "Error": 0
    }

    text = soup.get_text(" ", strip=True).lower()
    patterns = {
        "Executed": r'executed test cases[:\s]+(\d+)',
        "Passed": r'passed[:\s]+(\d+)',
        "Failed": r'failed[:\s]+(\d+)',
        "Inconclusive": r'inconclusive[:\s]+(\d+)',
        "Error": r'error[:\s]+(\d+)'
    }

    for key, pattern in patterns.items():
        match = re.search(pattern, text)
        if match:
            stats[key] = int(match.group(1))

    return stats


@st.cache_data(show_spinner=False)
def extract_test_results_grouped_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    results = {}

    group_tables = soup.find_all('table', class_='GroupHeadingTable')

    for group_table in group_tables:
        try:
            rows = group_table.find_all('tr')
            if len(rows) >= 2:
                first_row = rows[0]
                heading = first_row.find('big', class_='Heading3')

                if heading:
                    heading_text = heading.get_text(strip=True)
                    fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', heading_text, re.IGNORECASE)

                    if fixture_match:
                        fixture_name = fixture_match.group(1).strip()
                        second_row = rows[1]
                        overview_table = second_row.find('table', class_='OverviewResultTable')

                        if overview_table:
                            count_cell = overview_table.find('td')
                            if count_cell:
                                try:
                                    count = int(count_cell.get_text(strip=True))

                                    if fixture_name not in results:
                                        results[fixture_name] = {
                                            "name": fixture_name,
                                            "test_cases": [],
                                            "pass": count,
                                            "fail": 0,
                                            "error": 0,
                                            "not executed": 0,
                                            "inconclusive": 0,
                                            "total": count,
                                            "count_cell_class": count_cell.get('class', [''])[0]
                                        }
                                except ValueError:
                                    pass
        except Exception:
            pass

    full_text = soup.get_text("\n", strip=True)
    lines = [l.strip() for l in full_text.split("\n") if l.strip()]

    current_fixture = None

    for i, line in enumerate(lines):
        line_lower = line.lower()

        if "test fixture:" in line_lower:
            fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', line, re.IGNORECASE)
            if fixture_match:
                current_fixture = fixture_match.group(1).strip()
                if current_fixture not in results:
                    results[current_fixture] = {
                        "name": current_fixture,
                        "test_cases": [],
                        "pass": 0,
                        "fail": 0,
                        "error": 0,
                        "not executed": 0,
                        "inconclusive": 0,
                        "total": 0
                    }

        elif re.match(r'^\d+\.\d+', line) and current_fixture:
            verdict_match = re.search(r':\s*(Passed|Failed|Pass|Fail|Error|Not Executed|Inconclusive)\s*$', line,
                                      re.IGNORECASE)

            if verdict_match:
                verdict_word = verdict_match.group(1).lower()

                if "pass" in verdict_word:
                    verdict_type = "Pass"
                    results[current_fixture]["pass"] += 1
                elif "fail" in verdict_word:
                    verdict_type = "Failed"
                    results[current_fixture]["fail"] += 1
                elif "error" in verdict_word:
                    verdict_type = "Error"
                    results[current_fixture]["error"] += 1
                elif "not executed" in verdict_word:
                    verdict_type = "Not Executed"
                    results[current_fixture]["not executed"] += 1
                elif "inconclusive" in verdict_word:
                    verdict_type = "Inconclusive"
                    results[current_fixture]["inconclusive"] += 1
                else:
                    continue

                timestamp = None
                test_step = "Step"
                failure_step_id = ""

                def score_timestamp(candidate):
                    if not candidate:
                        return -1
                    parts = candidate.split('.')
                    if len(parts) != 2 or not parts[0].isdigit() or not parts[1].isdigit():
                        return -1
                    leading_num = int(parts[0])
                    decimal_places = len(parts[1])
                    decimal_bonus = 10000 if decimal_places >= 3 else (100 if decimal_places == 2 else 0)
                    return decimal_bonus + leading_num

                def find_best_timestamp(text):
                    matches = re.findall(r'\b(\d+\.\d+)\b', text)
                    if not matches:
                        return None
                    return max(matches, key=score_timestamp)

                def find_first_relevant_timestamp(text):
                    for m in re.findall(r'\b(\d+\.\d+)\b', text):
                        if len(m.split('.')[1]) >= 3:
                            return m
                    for m in re.findall(r'\b(\d+\.\d+)\b', text):
                        if len(m.split('.')[1]) >= 2:
                            return m
                    return None

                def consider_timestamp(candidate):
                    nonlocal timestamp
                    if not candidate:
                        return
                    if not timestamp:
                        timestamp = candidate
                        return
                    if len(timestamp.split('.')[1]) >= 3:
                        return
                    if len(candidate.split('.')[1]) > len(timestamp.split('.')[1]):
                        timestamp = candidate
                        return
                    if score_timestamp(candidate) > score_timestamp(timestamp):
                        timestamp = candidate

                same_line_step = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', line,
                                           re.IGNORECASE)
                if same_line_step:
                    failure_step_id = same_line_step.group(1)
                    action_text = same_line_step.group(2).strip()
                    test_step = action_text
                    consider_timestamp(find_first_relevant_timestamp(line) or find_best_timestamp(line))

                for k in range(i + 1, min(i + 150, len(lines))):
                    next_line = lines[k]

                    if re.match(r'^\d+\.\d+(?:\s|$)', next_line) and k > i + 5:
                        break

                    consider_timestamp(find_first_relevant_timestamp(next_line) or find_best_timestamp(next_line))

                    if verdict_type in ["Failed", "Error"] and not failure_step_id:
                        next_line_lower = next_line.lower()

                        step_match = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', next_line,
                                               re.IGNORECASE)
                        if step_match:
                            failure_step_id = step_match.group(1)
                            action_text = step_match.group(2).strip()
                            test_step = action_text
                            consider_timestamp(find_best_timestamp(next_line))
                        else:
                            if any(keyword in next_line_lower for keyword in
                                   ["condition", "value", "expected", "actual", "mismatch", "not found",
                                    "exception", "error", "failed to", "failed"]):
                                if not re.match(r'^\d+\.\d+', next_line):
                                    step_num_match = re.match(r'^(\d+(?:\.\d+)*)', next_line.strip())
                                    if step_num_match:
                                        failure_step_id = step_num_match.group(1)
                                        test_step = next_line[:80]

                    if verdict_type == "Pass":
                        next_line_lower = next_line.lower()

                        if "execute" in next_line_lower:
                            match = re.search(r'execute\s+(\w+)', next_line_lower)
                            if match:
                                test_step = match.group(1).capitalize()
                        elif "question" in next_line_lower and "text" in next_line_lower:
                            test_step = "Question/Text"
                        elif "await" in next_line_lower or "wait" in next_line_lower:
                            test_step = "Await Value Match"
                        elif "resume" in next_line_lower:
                            test_step = "Resume"
                        elif "set" in next_line_lower:
                            test_step = "Set"
                        elif "tester" in next_line_lower and "confirmed" in next_line_lower:
                            test_step = "Tester Confirmation"

                if timestamp:
                    results[current_fixture]["test_cases"].append({
                        "timestamp": timestamp,
                        "verdict": verdict_type,
                        "details": test_step
                    })

    for fixture_name in results:
        parsed_count = len(results[fixture_name]["test_cases"])
        initial_count = results[fixture_name].get("total", 0)
        results[fixture_name]["total"] = max(parsed_count, initial_count)

    return results


CREATOR_USERNAME = "Vignesh"
CREATOR_PASSWORD = "Rider@100"

# Login gate:
# This runs before the main app tabs are shown. It keeps the creator/user access
# flow in one place so authentication checks do not have to be repeated per tab.
# ================================
# PREMIUM LOGIN EXPERIENCE
# Replace your current login gate block with this
# ================================
if not st.session_state.is_authenticated and "preview_token" not in query_params:
    # Custom CSS for the login page
    st.markdown("""
    <style>
        html, body { margin: 0; padding: 0; height: 100%; }
        .block-container {
            max-width: none !important;
            padding: clamp(20px, 3vw, 36px) !important;
            margin: 0 !important;
            width: 100vw !important;
            min-height: 100vh !important;
        }
        .main {
            width: 100% !important;
            padding: 0 !important;
            min-height: 100vh !important;
        }

        /* Full viewport container */
        [data-testid="stAppViewContainer"] {
            background:
                radial-gradient(circle at top left, rgba(176, 224, 230, 0.35) 0%, rgba(176, 224, 230, 0) 40%),
                radial-gradient(circle at top right, rgba(230, 245, 250, 0.3) 0%, rgba(230, 245, 250, 0) 45%),
                radial-gradient(circle at bottom left, rgba(176, 196, 222, 0.25) 0%, rgba(176, 196, 222, 0) 40%),
                linear-gradient(135deg, #ffffff 0%, #f8fafc 40%, #eef2f7 100%) !important;  # login background
            min-height: 100vh !important;
            display: flex !important;
            align-items: stretch !important;
        }

        /* Override Streamlit's default column behavior */
        [data-testid="column"]:first-child {
            flex: 1 !important;
            padding: clamp(36px, 4vw, 64px) clamp(32px, 5vw, 84px) !important;
            display: flex !important;
            flex-direction: column !important;
            justify-content: center !important;
            color: #3B5E7F !important;
            min-height: calc(100vh - 2 * clamp(20px, 3vw, 36px)) !important;
        }
        [data-testid="column"]:nth-child(2) {
            width: min(460px, 100%) !important;
            display: flex !important;
            align-items: center !important;
            justify-content: center !important;
            padding: clamp(20px, 2.6vw, 32px) !important;
            flex-shrink: 0 !important;
            min-height: calc(100vh - 2 * clamp(20px, 3vw, 36px)) !important;
        }

        /* Branding elements */
        .brand-strip {
            display: flex !important;
            align-items: center !important;
            gap: 12px !important;
            #margin-bottom: 32px !important;
           perspective: 800px;      
        }
        .login-panel {
            width: min(440px, 100%);
            padding: 36px 32px 30px;
            border-radius: 24px;
            background: linear-gradient(180deg, rgba(255, 253, 250, 0.95) 0%, rgba(245, 250, 252, 0.92) 100%);
            border: 1.5px solid rgba(176, 224, 230, 0.45);
            box-shadow: 0 20px 50px rgba(135, 206, 235, 0.1), 0 0 30px rgba(176, 196, 222, 0.08);
            backdrop-filter: blur(12px);
        }
        .brand-logo-3d {
            width: 60px !important;
            height: 60 px !important;
            perspective: 800px;    
            display: flex !important;
            align-items: center !important;
            justify-content: center !important;
            animation: float 3s ease-in-out infinite;       
        }  
        @keyframes float {
           0%, 100% { transform: translateY(0px); }
          50% { transform: translateY(-4px); }
        }  
        .logo-inner{
                width : 100%;
                height: 100%;
                border-radius:50%;
                background: radial-gradient(circle at 30% 30%, #dcdcdc, #a0a0a0);
                display: flex;
                align-items: center;
                justify-content: center;
                transform-style: preserve-3d;
                transition: transform 0.2s ease;
                box-shadow: 0 8px 20px rgba(0,0,0,0.15);
            }   
        .star{
              font-size: 28px;
              color: white;
              transform: translateZ(10px); 
            }                      
        .brand-label {
            font-size: 0.82rem !important;
            color: rgba(91, 127, 166, 0.9) !important;
            letter-spacing: 0.18em !important;
            text-transform: uppercase !important;
        }
        .ai-branding {
            display : flex;
            gap : 8px;
            align-items : baseline;
            margin-bottom: 28px!important;
        }
        
        .ai-title {
           font-size: clamp(2rem, 3.5vw, 2.8rem);
           font-weight: 800;
           color: #2C5F7F;
        }

        .ai-subtitle {
         font-size: 14px;
         font-style: italic;
         font-weight: 400;
         color: rgba(44, 95, 127, 0.6);
        }
        .ai-tagline {
            font-size: clamp(2.2rem, 4vw, 3.2rem) !important;
            font-weight: 800 !important;
            line-height: 1.15 !important;
            color: #3B5E7F !important;
            margin-bottom: 32px !important;
            letter-spacing: -0.03em !important;
        }
        .ai-description {
            font-size: 1.02rem !important;
            line-height: 1.7 !important;
            color: rgba(91, 127, 166, 0.9) !important;
            margin-bottom: 48px !important;
        }
        .trust-row {
            color: rgba(91, 127, 166, 0.8) !important;
            font-size: 0.95rem !important;
            margin-bottom: 56px !important;
            line-height: 1.6 !important;
        }
        /* Feature cards */
        .feature-grid {
            display: grid;
            grid-template-columns: repeat(2, minmax(280px, 1fr));
            gap: 16px;
            align-items: stretch;
            margin-top: 8px;
        }
        .feature-card {
            background: rgba(176, 224, 230, 0.15);
            border: 1px solid rgba(176, 196, 222, 0.3);
            border-radius: 12px;
            color: #3B5E7F;
            padding: 18px 18px 16px;
            min-height: 150px;
            display: flex;
            flex-direction: column;
            justify-content: flex-start;
            backdrop-filter: blur(4px);
        }
        .feature-card h4 {
            margin: 0 0 14px 0;
            color: #0d5aa7;
            font-size: 1.08rem;
            font-weight: 700;
        }
        .feature-card ul {
            margin: 0;
            padding-left: 20px;
        }
        .feature-card li {
            margin: 0 0 8px 0;
            line-height: 1.45;
            color: #0d5aa7;
        }
        .login-panel .stButton button {
            background: linear-gradient(135deg, #87CEEB 0%, #6BA3C5 100%) !important;
            color: #ffffff !important;
            border: 0 !important;
            border-radius: 8px !important;
            font-weight: 600 !important;
            width: 100% !important;
            padding: 12px 16px !important;
            margin-top: 8px !important;
        }

        .login-panel .stButton button:hover {
            background: linear-gradient(135deg, #ADD8E6 0%, #87CEEB 100%) !important;
            transform: translateY(-1px) !important;
        }

        .login-note {
            color: rgba(100, 100, 140, 0.7) !important;
            font-size: 0.85rem !important;
            text-align: center !important;
            margin-top: 16px !important;
        }

        /* Responsive design */
        @media (max-width: 768px) {
            [data-testid="column"]:first-child {
                padding: 24px 12px 16px !important;
                order: 2 !important;
                min-height: auto !important;
            }
            [data-testid="column"]:nth-child(2) {
                width: 100% !important;
                padding: 12px !important;
                order: 1 !important;
                min-height: auto !important;
            }
            .ai-tagline {
                font-size: 2rem !important;
            }
            .login-panel {
                width: 100% !important;
                padding: 24px 20px !important;
            }
            .feature-grid {
                grid-template-columns: 1fr;
            }
        }
    </style>
    """, unsafe_allow_html=True)

    # Create the new flexbox layout using Streamlit columns - Login page layout
    left_col, right_col = st.columns([3, 1.3])

    # Apply additional form styling for login elements
    st.markdown("""
    <style>
        /* Login heading and subheading */
        .login-panel .login-heading {
            margin-top: 0 !important;
        }
        .login-heading {
            font-size: clamp(1.8rem, 4vw, 2.8rem) !important;
            color: #2C5F7F !important;
            font-weight: 700 !important;
            margin-bottom: 10px !important;
            line-height: 1.1 !important;
        }

        .login-subheading {
            font-size: clamp(0.95rem, 2vw, 1.05rem) !important;
            color: #5B7FA6 !important;
            font-weight: 500 !important;
            margin-bottom: 20px !important;
        }

        /* Form elements */
        .login-panel [data-testid="stTextInput"] {
            margin-bottom: 8px !important;
        }
        .login-panel [data-testid="stTextInput"] label,
        .login-panel [data-testid="stTextInput"] label p,
        .login-panel [data-testid="stTextInput"] label span,
        .login-panel [data-testid="stTextInput"] p {
            color: #3B5E7F !important;
            -webkit-text-fill-color: #3B5E7F !important;
            opacity: 1 !important;
        }
        .login-panel [data-baseweb="base-input"],
        .login-panel [data-baseweb="input"],
        .login-panel [data-testid="stTextInput"] > div > div {
            background: rgba(230, 244, 248, 0.85) !important;
            border: 1.5px solid rgba(176, 224, 230, 0.6) !important;
            border-radius: 12px !important;
            box-shadow: none !important;
        }
        .login-panel [data-baseweb="base-input"]:focus-within,
        .login-panel [data-baseweb="input"]:focus-within,
        .login-panel [data-testid="stTextInput"] > div > div:focus-within {
            border: 1.5px solid #87CEEB !important;
            box-shadow: 0 0 0 2px rgba(135, 206, 235, 0.3) !important;
                outline : none !important;
        }
                    
         /*Remove streamlit red focus */
        .login-panel input:focus,
        .login-panel input:invalid {
                outline:none !important;
                box-shadow:none ! important;
                border-color: #87CEEB !important;
            }   
                                        
        .login-panel input[type="text"],
        .login-panel input[type="password"] {
            background: transparent !important;
            color: #2C5F7F !important;
            caret-color: #4D94B9 !important;
            -webkit-text-fill-color: #2C5F7F !important;
            border: none !important;
            box-shadow: none !important;
            font-weight: 600 !important;
            letter-spacing: 0.2px !important;
            font-size: 1rem !important;
            padding: 12px 16px !important;
        }
        .login-panel input[type="text"]::placeholder,
        .login-panel input[type="password"]::placeholder {
            color: rgba(91, 127, 166, 0.65) !important;
            opacity: 1 !important;
            -webkit-text-fill-color: rgba(91, 127, 166, 0.65) !important;
        }
        .login-panel input[type="text"]:focus,
        .login-panel input[type="password"]:focus {
            background: transparent !important;
            color: #2C5F7F !important;
            -webkit-text-fill-color: #2C5F7F !important;
            outline: none !important;
        }

        .login-note {
            color: rgba(91, 127, 166, 0.75) !important;
            font-size: 0.85rem !important;
            text-align: center !important;
            margin-top: 16px !important;
        }
        .login-panel [data-testid="stCaptionContainer"],
        .login-panel [data-testid="stCaptionContainer"] p,
        .login-panel .stCaption,
        .login-panel .stCaption p {
            color: rgba(91, 127, 166, 0.85) !important;
            -webkit-text-fill-color: rgba(91, 127, 166, 0.85) !important;
            opacity: 1 !important;
            font-size: 0.84rem !important;
        }
        .login-panel .stButton > button,
        .login-panel div.stButton > button {
            width: 100% !important;
            min-height: 42px !important;
            padding: 0.45rem 0.9rem !important;
            font-size: 0.98rem !important;
            border-radius: 10px !important;
        }

        /* Responsive design */
        @media (max-width: 768px) {
            [data-testid="column"]:first-child {
                padding: 24px 12px 16px !important;
                order: 2 !important;
            }
            [data-testid="column"]:nth-child(2) {
                width: 100% !important;
                padding: 12px !important;
                order: 1 !important;
            }
            .ai-tagline {
                font-size: 2rem !important;
            }
            .login-panel {
                width: 100% !important;
                padding: 24px 20px !important;
            }
        }
         /* 🔥 HARD OVERRIDE - kills ALL red focus from Streamlit/BaseWeb */

            .login-panel *:focus {
             outline: none !important;
            }  
            .login-panel input,
            .login-panel textarea {
               box-shadow: none !important;
            }

           /* Target BaseWeb internal input container */
            .login-panel div[data-baseweb="base-input"] {
              border: 1.5px solid rgba(176, 224, 230, 0.6) !important;
            }   

            /* Focus state */
           .login-panel div[data-baseweb="base-input"]:focus-within {
               border: 1.5px solid #87CEEB !important;
              box-shadow: 0 0 0 2px rgba(135, 206, 235, 0.3) !important;
            }    

            /* Remove error (red) state completely */
            .login-panel div[data-baseweb="base-input"][aria-invalid="true"] {
                border: 1.5px solid #87CEEB !important;
                box-shadow: none !important;
            }

            /* Also kill any red glow from deeper layers */
            .login-panel div[data-baseweb="input"] {
               box-shadow: none !important;
            }   

         /* 🚨 NUCLEAR OVERRIDE — removes ALL red states */

            .login-panel *[aria-invalid="true"],
            .login-panel *[data-baseweb="base-input"][aria-invalid="true"],
            .login-panel *[data-baseweb="input"][aria-invalid="true"] {
                border-color: #87CEEB !important;
                box-shadow: none !important;
            }

         /* Remove any red focus ring from BaseWeb */
            .login-panel *:focus-visible {
                outline: none !important;
                box-shadow: none !important;
            }

            /* Force our focus style ONLY */
            .login-panel div[data-baseweb="base-input"]:focus-within {
                border: 1.5px solid #87CEEB !important;
              box-shadow: 0 0 0 2px rgba(135, 206, 235, 0.3) !important;
            }

            /* Kill ANY red borders globally inside login */
            .login-panel input,
            .login-panel div,
            .login-panel textarea {
             border-color: rgba(176, 224, 230, 0.6) !important;
            } 
            /* 🚨 HARD KILL: BaseWeb invalid (red) state */

            .login-panel div[data-baseweb="base-input"][aria-invalid="true"],
            .login-panel div[data-baseweb="input"][aria-invalid="true"],
            .login-panel [data-baseweb="base-input"] {
               border: 1.5px solid rgba(176, 224, 230, 0.6) !important;
               box-shadow: none !important;
            }

            /* Override inner input focus ring */
            .login-panel input {
               outline: none !important;
               box-shadow: none !important;
            }

            /* Force consistent focus (blue theme only) */
            .login-panel div[data-baseweb="base-input"]:focus-within {
              border: 1.5px solid #87CEEB !important;
             box-shadow: 0 0 0 2px rgba(135, 206, 235, 0.25) !important;
            }         

    </style>
    """, unsafe_allow_html=True)

    with left_col:
        if logo_data:
            logo_display = f'''
            <img src = "data:image/gif;base64,{logo_data}"
            style = "width: 36px; height: 36px; object-fit;contain;">'''
        else:
             logo_display = '<div class = "star">★</div>'
        st.markdown(f"""
                    <div class = "brand-strip">
                        <div class = " brand-logo-3d">
                           {logo_display}
                        </div>
                    </div>
                    <div class = "brand-label">Mercedes_Benz</div>
                </div>
                """,unsafe_allow_html=True)
            
        
        
        st.markdown("""<div class="ai-branding"> 
        <span class= "ai-title"> IntelliDoc AI </span> 
        <span class="ai-subtitle">-Smart Document Assistant </span>
        </div>""", unsafe_allow_html=True)
        st.markdown('<h1 class="ai-tagline">Where Documents Become Intelligence</h1>', unsafe_allow_html=True)
        st.markdown('<p class="ai-description">for secure document insight, comparison, dashboards, and automation.</p>', unsafe_allow_html=True)
        st.markdown(
            """
            <div class="feature-grid">
                <div class="feature-card">
                    <h4>💬 Chat</h4>
                    <ul>
                        <li>Ask questions about uploaded files</li>
                        <li>Context-aware responses</li>
                        <li>Multi-file semantic understanding</li>
                    </ul>
                </div>
                <div class="feature-card">
                    <h4>📊 Dashboard</h4>
                    <ul>
                        <li>Excel/CSV visualization</li>
                        <li>Export insights</li>
                    </ul>
                </div>
                <div class="feature-card">
                    <h4>🔄 Compare</h4>
                    <ul>
                        <li>Compare 2+ files</li>
                        <li>Word-level diff</li>
                        <li>Inline visual comparison</li>
                        <li>Export results to Excel</li>
                    </ul>
                </div>
                <div class="feature-card">
                    <h4>📡 CAPL</h4>
                    <ul>
                        <li>Upload or create <code>.can</code> files</li>
                        <li>Built-in CAPL editor</li>
                        <li>Code analysis &amp; issue detection</li>
                        <li>Suggestions &amp; improvements</li>
                    </ul>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )

    with right_col:  # type: ignore  # right_col is defined at line 3699
        #st.markdown('<div class="login-panel">', unsafe_allow_html=True)
        #st.markdown('<div class="login-heading">Welcome back</div>', unsafe_allow_html=True)
        st.markdown('<div class="login-subheading">Sign in to IntelliDoc AI</div>', unsafe_allow_html=True)

        login_username = st.text_input("👤 Username", placeholder="Username", key="username")
        login_password = st.text_input("🔒 Password", type="password", placeholder="Password", key="password")

        st.caption("Standard users can leave password empty")

        access_clicked = st.button("Access", use_container_width=False, key="signin")

        st.markdown('</div>', unsafe_allow_html=True)

    if access_clicked:
        cleaned_username = (login_username or "").strip()
        cleaned_password = (login_password or "").strip()

        if cleaned_username == CREATOR_USERNAME and cleaned_password == CREATOR_PASSWORD:
            st.session_state.is_authenticated = True
            st.session_state.logged_in_username = cleaned_username
            st.session_state.user_role = "creator"
            st.session_state.user_session_start_time = datetime.now().isoformat()
            st.session_state.start_time = time.time()

            ist_tz = timezone('Asia/Kolkata')
            ist_time = datetime.now(ist_tz).strftime("%Y-%m-%d %H:%M:%S %Z")
            st.session_state.login_history.append({
                "username": cleaned_username,
                "role": "creator",
                "action": "login",
                "timestamp": ist_time,
                "usage_time": "-"
            })

            active_file = "active_users.json"
            now = datetime.now()
            if os.path.exists(active_file):
                with open(active_file, "r") as f:
                    active_users = json.load(f)
            else:
                active_users = []

            active_users = [u for u in active_users if u.get("username") != cleaned_username]
            active_users.append({"username": cleaned_username, "timestamp": now.isoformat()})

            with open(active_file, "w") as f:
                json.dump(active_users, f)

            st.success("✅ Creator access granted")
            st.rerun()

        elif cleaned_username and len(cleaned_username) > 3 and cleaned_password == "":
            st.session_state.is_authenticated = True
            st.session_state.logged_in_username = cleaned_username
            st.session_state.user_role = "user"
            st.session_state.user_session_start_time = datetime.now().isoformat()
            st.session_state.start_time = time.time()

            ist_tz = timezone('Asia/Kolkata')
            ist_time = datetime.now(ist_tz).strftime("%Y-%m-%d %H:%M:%S %Z")
            st.session_state.login_history.append({
                "username": cleaned_username,
                "role": "user",
                "action": "login",
                "timestamp": ist_time,
                "usage_time": "-"
            })

            active_file = "active_users.json"
            now = datetime.now()
            if os.path.exists(active_file):
                with open(active_file, "r") as f:
                    active_users = json.load(f)
            else:
                active_users = []

            active_users = [u for u in active_users if u.get("username") != cleaned_username]
            active_users.append({"username": cleaned_username, "timestamp": now.isoformat()})

            with open(active_file, "w") as f:
                json.dump(active_users, f)

            st.success(f"✅ Welcome, {cleaned_username}!")
            st.rerun()

        else:
            st.error("❌ Invalid credentials. Creator needs password. Users need username >3 chars with empty password.")

    st.stop()


# Files will be processed on-demand per tab when selected
# -------------------------------
# HELPER CHART FUNCTIONS
# -------------------------------
# Dashboard chart helpers:
# These are used by the Dashboard tab when an XLSX/HTML file is selected and the
# user wants counts shown as bar or pie charts.
def get_column_counts(data, column):
    counts = defaultdict(int)
    for row in data:
        val = row.get(column)
        if val is not None:
            counts[val] += 1
    return dict(counts)


def plot_pie_chart(counts, title):
    fig = px.pie(
        names=list(counts.keys()),
        values=list(counts.values()),
        title=title,
        hole=0.3,
    )
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(margin=dict(t=50, b=20, l=20, r=20))
    return fig


def plot_bar_chart(counts, title, horizontal=False):
    labels = list(counts.keys())
    values = list(counts.values())
    if horizontal:
        fig = px.bar(x=values, y=labels, orientation="h", title=title)
    else:
        fig = px.bar(x=labels, y=values, title=title)
    fig.update_layout(margin=dict(t=50, b=80, l=40, r=20))
    return fig


# -------------------------------
# INLINE MULTI-FILE DIFF (HTML)
# -------------------------------
# Compare tab HTML diff helper:
# Generates the inline visual comparison shown in the Compare tab and also reused
# from Chat when the user asks to compare multiple selected documents.
@st.cache_data(show_spinner=False)
def highlight_multi_file_differences_cached(file_items, comparison_mode="Exact inline word diff", reference_file=None):
    if len(file_items) < 2:
        return "Select at least two files to compare."

    files = [fname for fname, _ in file_items]
    if reference_file is None or reference_file not in files:
        reference_file = files[0]

    css = """
    <style>
        body { font-family: Arial; margin: 20px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid black; padding: 4px; vertical-align: top; white-space: pre-wrap; }
        th { background-color: #f0f0f0; }
        td.line-number { background-color: #f0f0f0; font-weight: bold; text-align: center; }
        .match { background-color: #ccffcc; }
        .mismatch { background-color: #ffcccc; }
        .scrollable { overflow:auto; max-height:800px; }
        p.legend span { display:inline-block; width:20px; height:20px; margin-right:5px; vertical-align:middle; }
    </style>
    """
    html_parts = [
        "<html><head>", css, "</head><body><div class='scrollable'>",
        "<p class='legend'><b>Legend:</b> <span class='match'></span> Matched word, <span class='mismatch'></span> Different or missing word</p>",
        "<table><tr><th>Line #</th>",
        "".join(f"<th>{html.escape(fname)}</th>" for fname in files),
        "</tr>",
    ]

    file_lines = {fname: text.splitlines() for fname, text in file_items}
    max_lines = max(len(lines) for lines in file_lines.values())

    for i in range(max_lines):
        html_parts.append(f"<tr><td class='line-number'>{i + 1}</td>")

        line_word_lists = {}
        ordered_words = []
        word_presence = defaultdict(int)

        for fname in files:
            raw_line = file_lines[fname][i] if i < len(file_lines[fname]) else ""
            words = raw_line.split()
            line_word_lists[fname] = words
            for word in words:
                if word not in ordered_words:
                    ordered_words.append(word)
            for word in set(words):
                word_presence[word] += 1

        reference_words = line_word_lists.get(reference_file, [])

        for fname in files:
            words = line_word_lists[fname]
            if comparison_mode == "Word presence summary":
                highlighted = []
                word_set = set(words)
                for word in ordered_words:
                    escaped_word = html.escape(word)
                    if word in word_set and word_presence[word] == len(files):
                        highlighted.append(f"<span class='match'>{escaped_word}</span>")
                    else:
                        highlighted.append(f"<span class='mismatch'>{escaped_word}</span>")
                cell_html = ' '.join(highlighted) if highlighted else '&nbsp;'
            else:
                highlighted = []
                matcher = SequenceMatcher(None, reference_words, words)
                for tag, i1, i2, j1, j2 in matcher.get_opcodes():
                    if tag == "equal":
                        highlighted.extend(f"<span class='match'>{html.escape(w)}</span>" for w in words[j1:j2])
                    else:
                        highlighted.extend(f"<span class='mismatch'>{html.escape(w)}</span>" for w in words[j1:j2])
                cell_html = ' '.join(highlighted) if highlighted else '&nbsp;'
            html_parts.append(f"<td>{cell_html}</td>")

        html_parts.append("</tr>")

    html_parts.append("</table></div></body></html>")
    return "".join(html_parts)


def highlight_side_by_side_differences_cached(file_items, reference_file=None):
    files = [fname for fname, _ in file_items]
    if len(files) < 2:
        return "Select at least two files to compare."
    if reference_file is None or reference_file not in files:
        reference_file = files[0]

    file_lines = {fname: text.splitlines() for fname, text in file_items}
    max_lines = max(len(lines) for lines in file_lines.values())

    css = """
    <style>
        body { font-family: Arial; margin: 20px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid black; padding: 4px; vertical-align: top; white-space: pre-wrap; }
        th { background-color: #f0f0f0; }
        td.line-number { background-color: #f0f0f0; font-weight: bold; text-align: center; }
        .line-match { background-color: #ccffcc; display: block; width: 100%; }
        .line-mismatch { background-color: #ffcccc; display: block; width: 100%; }
        .scrollable { overflow:auto; max-height:800px; }
        p.legend span { display:inline-block; width:20px; height:20px; margin-right:5px; vertical-align:middle; }
    </style>
    """
    html_parts = [
        "<html><head>", css, "</head><body><div class='scrollable'>",
        "<p class='legend'><b>Legend:</b> <span class='line-match'></span> Same as reference line, <span class='line-mismatch'></span> Different from reference or missing line</p>",
        "<p><b>Reference file:</b> " + html.escape(reference_file) + "</p>",
        "<table><tr><th>Line #</th>",
        "".join(f"<th>{html.escape(fname)}</th>" for fname in files),
        "</tr>",
    ]

    for i in range(max_lines):
        html_parts.append(f"<tr><td class='line-number'>{i + 1}</td>")
        reference_line = file_lines[reference_file][i] if i < len(file_lines[reference_file]) else ""
        for fname in files:
            line_text = file_lines[fname][i] if i < len(file_lines[fname]) else ""
            if line_text == reference_line and line_text != "":
                cell_html = f"<span class='line-match'>{html.escape(line_text)}</span>"
            elif line_text == reference_line == "":
                cell_html = "&nbsp;"
            else:
                cell_html = f"<span class='line-mismatch'>{html.escape(line_text)}</span>"
            html_parts.append(f"<td>{cell_html}</td>")
        html_parts.append("</tr>")

    html_parts.append("</table></div></body></html>")
    return "".join(html_parts)


def highlight_multi_file_differences(file_texts, comparison_mode="Exact inline word diff", reference_file=None):
    if comparison_mode == "Side-by-side line diff":
        return highlight_side_by_side_differences_cached(
            tuple((fname, str(text)) for fname, text in file_texts.items()),
            reference_file=reference_file
        )
    return highlight_multi_file_differences_cached(
        tuple((fname, str(text)) for fname, text in file_texts.items()),
        comparison_mode=comparison_mode,
        reference_file=reference_file
    )


def build_semantic_diff_explanation(file_texts):
    """Explain meaning-level changes beside the visual word/line diff."""
    if not file_texts or len(file_texts) < 2:
        return "Select at least two files to generate a semantic difference explanation."

    file_names = list(file_texts.keys())
    per_file = {}
    for file_name, text in file_texts.items():
        text = str(text or "")
        per_file[file_name] = {
            "themes": set(extract_key_themes(text, limit=14)),
            "entities": set(extract_entities(text, limit=20)),
            "risks": set(extract_risk_signals(text, limit=10)),
            "length": len(text),
        }

    baseline = file_names[0]
    base = per_file[baseline]
    sections = [f"### Semantic Diff Explanation\nBaseline: **{html.escape(baseline)}**"]
    for file_name in file_names[1:]:
        current = per_file[file_name]
        added_themes = sorted(current["themes"] - base["themes"])[:8]
        removed_themes = sorted(base["themes"] - current["themes"])[:8]
        added_entities = sorted(current["entities"] - base["entities"])[:8]
        removed_entities = sorted(base["entities"] - current["entities"])[:8]
        added_risks = sorted(current["risks"] - base["risks"])[:5]
        delta = current["length"] - base["length"]
        delta_label = "expanded" if delta > 0 else "contracted" if delta < 0 else "unchanged in size"

        sections.append(f"#### {html.escape(file_name)}")
        sections.append(f"- Structural signal: content {delta_label} by {abs(delta):,} extracted characters.")
        sections.append(f"- New themes: {html.escape(', '.join(added_themes) if added_themes else 'None detected')}.")
        sections.append(f"- Missing themes: {html.escape(', '.join(removed_themes) if removed_themes else 'None detected')}.")
        sections.append(f"- New entities: {html.escape(', '.join(added_entities) if added_entities else 'None detected')}.")
        sections.append(f"- Missing entities: {html.escape(', '.join(removed_entities) if removed_entities else 'None detected')}.")
        if added_risks:
            sections.append("- Risk/signals introduced:")
            sections.extend(f"  - {html.escape(str(risk)[:220])}" for risk in added_risks)
        else:
            sections.append("- Risk/signals introduced: None detected.")

    return "\n".join(sections)


# -------------------------------
# COMPARE EXCEL HIGHLIGHT
# -------------------------------
# Compare tab Excel export helper:
# Builds the downloadable workbook used in the Compare tab so users can inspect
# mismatches outside the Streamlit UI.
@st.cache_data(show_spinner=False)
def generate_word_level_comparison_excel_cached(file_items):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Comparison"

    files = [fname for fname, _ in file_items]
    file_texts = {fname: text for fname, text in file_items}
    ws.append(["Line #"] + files)
    file_lines = {f: [l.split() for l in t.splitlines()] for f, t in file_texts.items()}

    red_fill = PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid")
    green_fill = PatternFill(start_color="CCFFCC", end_color="CCFFCC", fill_type="solid")

    max_lines = max(len(l) for l in file_lines.values())

    for i in range(max_lines):
        max_words = max(len(file_lines[f][i]) if i < len(file_lines[f]) else 0 for f in files)
        for w_idx in range(max_words):
            row_values = [i + 1 if w_idx == 0 else ""]
            for f in files:
                line_words = file_lines[f][i] if i < len(file_lines[f]) else []
                word = line_words[w_idx] if w_idx < len(line_words) else ""
                row_values.append(word)
            ws.append(row_values)

            # Highlight exact matches in green and missing/mismatched content in red
            all_words_set = set()
            for f in files:
                if i < len(file_lines[f]):
                    all_words_set.update(file_lines[f][i])
            for col_idx, f in enumerate(files, start=2):
                cell = ws.cell(row=ws.max_row, column=col_idx)
                line_words = file_lines[f][i] if i < len(file_lines[f]) else []
                if w_idx >= len(line_words):
                    cell.fill = red_fill
                elif all(word == line_words[w_idx] for other_file in files
                         for word in ([file_lines[other_file][i][w_idx]]
                                      if i < len(file_lines[other_file]) and w_idx < len(file_lines[other_file][i])
                                      else ["__missing__"])):
                    cell.fill = green_fill
                else:
                    cell.fill = red_fill

    excel_io = BytesIO()
    wb.save(excel_io)
    return excel_io.getvalue()


def generate_word_level_comparison_excel(file_texts):
    excel_io = BytesIO(generate_word_level_comparison_excel_cached(tuple((fname, str(text)) for fname, text in file_texts.items())))
    excel_io.seek(0)
    return excel_io


# -------------------------------
# CAPL Complier
# -------------------------------
# CAPL analyzer helpers:
# These functions are used only by the CAPL tab for syntax checking, issue
# listing, and highlighted code rendering inside the CAPL editor/viewer panels.
@st.cache_data(show_spinner=False)
def analyze_capl_code_with_suggestions_cached(code):
    issues = []

    brace_stack = []
    declared_vars = []
    used_vars = []

    lines = code.splitlines()

    for i, line in enumerate(lines, 1):
        stripped = line.strip()

        if not stripped or stripped.startswith("//"):
            continue  # Skip empty lines/comments

        # Track braces
        for c in stripped:
            if c == "{":
                brace_stack.append(i)
            elif c == "}":
                if brace_stack:
                    brace_stack.pop()
                else:
                    issues.append({
                        "line": i,
                        "error": "Unmatched closing brace",
                        "suggestion": "Remove or match with an opening '{'"
                    })

        # Detect variable declarations
        var_match = re.match(r'\b(int|float|byte|char|mstimer|timer|enum)\b\s+(\w+)', stripped)
        if var_match:
            declared_vars.append(var_match.group(2))

        # Track all used variable names
        used_vars += re.findall(r'\b([a-zA-Z_]\w*)\b', stripped)

        # Check for case sensitivity in keywords
        if re.search(r'\b(If|Else|For|While|Switch|Case|Return|On|Variables|Includes|Enum|Mstimer|Timer)\b', stripped):
            issues.append({
                "line": i,
                "error": "CAPL keywords should be lowercase",
                "suggestion": "Use lowercase keywords like 'if', 'else', 'on', etc."
            })

        # Check for incomplete if conditions
        if re.match(r'^\s*(if|else if)\s*\(', stripped) and not re.search(r'\)\s*(\{)?\s*$', stripped):
            issues.append({
                "line": i,
                "error": "Incomplete if condition",
                "suggestion": "Add closing parenthesis ')' and possibly opening brace '{'"
            })

        # Check for missing opening brace after control statements
        if re.match(r'^\s*(if|else if|else|for|while|switch)\b', stripped) and not stripped.endswith(
                '{') and not re.search(r'\)\s*\{', stripped):
            # Check if next line starts with '{'
            if i < len(lines) and not lines[i].strip().startswith('{'):
                issues.append({
                    "line": i,
                    "error": "Missing opening brace after control statement",
                    "suggestion": "Add '{' after the condition or on the next line"
                })

        # Detect missing semicolon
        if not stripped.endswith(";") and not stripped.endswith("{") and not stripped.endswith("}"):
            if not re.match(r'^(on|variables|includes|enum|mstimer|timer|if|else|switch|case|for|while|return)\b',
                            stripped):
                issues.append({
                    "line": i,
                    "error": "Missing semicolon",
                    "suggestion": "Add ';' at the end of this line"
                })

    # Check unmatched opening braces
    for open_line in brace_stack:
        issues.append({
            "line": open_line,
            "error": "Unmatched opening brace",
            "suggestion": "Add closing '}' to match this '{'"
        })

    # Check for 'on message' presence
    if "on message" not in code.lower():
        issues.append({
            "line": None,
            "error": "No 'on message' handler found",
            "suggestion": "Add an 'on message' event handler as required"
        })

    # Check for unused declared variables
    for var in declared_vars:
        if var not in used_vars:
            issues.append({
                "line": None,
                "error": f"Unused variable: {var}",
                "suggestion": "Consider removing this variable or using it in the code"
            })

    # Detect undeclared variables starting with PT4_ or $PT4_ used in code
    for i, line in enumerate(lines, 1):
        pt4_vars = re.findall(r'\b(PT4_[a-zA-Z_]\w*|\$PT4_[a-zA-Z_]\w*)\b', line)
        for var in pt4_vars:
            if var not in declared_vars and not var.startswith("$"):
                issues.append({
                    "line": i,
                    "error": f"Undeclared variable used: {var}",
                    "suggestion": f"Declare '{var}' in the variables section before using it"
                })

    return issues


def analyze_capl_code_with_suggestions(code):
    return analyze_capl_code_with_suggestions_cached(code)


# -------------------------------
# CAPL CODE DETECTION
# -------------------------------
def is_capl_code(text):
    """Check if the given text contains CAPL-specific keywords or syntax."""
    capl_keywords = [
        "on message", "variables", "includes", "mstimer", "timer", "byte", "char", "int", "float",
        "enum", "if", "else", "switch", "case", "for", "while", "return", "write", "output",
        "setTimer", "cancelTimer", "getTimer", "putValue", "getValue", "testcase", "teststep"
    ]
    text_lower = text.lower()
    return any(keyword in text_lower for keyword in capl_keywords)


@st.cache_data(show_spinner=False)
def render_capl_code_with_highlights_cached(code, issues_key):
    """Render CAPL code with IDE-like line highlighting for detected issues."""
    issues = [
        {"line": line, "error": error, "suggestion": suggestion}
        for line, error, suggestion in issues_key
    ]
    issue_lines = defaultdict(list)

    for issue in issues:
        line_no = issue.get("line")
        if isinstance(line_no, int):
            issue_lines[line_no].append(issue.get("error", "Issue detected"))

    code_lines = code.splitlines() or [""]
    rendered_lines = []

    for line_no, line in enumerate(code_lines, 1):
        escaped_line = html.escape(line) if line else "&nbsp;"
        line_classes = ["capl-line"]
        if line_no in issue_lines:
            line_classes.append("capl-line-error")

        tooltip = " | ".join(issue_lines[line_no]) if line_no in issue_lines else ""
        title_attr = f' title="{html.escape(tooltip)}"' if tooltip else ""

        rendered_lines.append(
            f"<div class=\"{' '.join(line_classes)}\"{title_attr}>"
            f"<span class=\"capl-gutter\">{line_no:>4}</span>"
            f"<span class=\"capl-code-text\">{escaped_line}</span>"
            f"</div>"
        )

    code_html = """
    <style>
        .capl-code-block {
            background: #0f172a;
            border: 1px solid #cbd5e1;
            border-radius: 10px;
            font-family: Consolas, "Courier New", monospace;
            font-size: 14px;
            line-height: 1.5;
            max-height: 420px;
            overflow: auto;
            padding: 12px 0;
        }
        .capl-line {
            color: #e2e8f0;
            display: flex;
            white-space: pre;
        }
        .capl-line-error {
            background: rgba(239, 68, 68, 0.22);
            border-left: 4px solid #ef4444;
        }
        .capl-gutter {
            color: #94a3b8;
            display: inline-block;
            min-width: 52px;
            padding: 0 12px;
            text-align: right;
            user-select: none;
        }
        .capl-code-text {
            display: inline-block;
            padding: 0 16px 0 0;
            width: 100%;
        }
    </style>
    """
    return code_html + f"<div class='capl-code-block'>{''.join(rendered_lines)}</div>"


def render_capl_code_with_highlights(code, issues=None):
    issues_key = tuple(
        (
            issue.get("line"),
            issue.get("error", "Issue detected"),
            issue.get("suggestion", "")
        )
        for issue in (issues or [])
    )
    return render_capl_code_with_highlights_cached(code, issues_key)


def render_capl_issue_table(issues):
    if not issues:
        st.success("✅ No issues detected!")
        return

    df_issues = pd.DataFrame(issues).fillna("-")
    st.dataframe(df_issues, use_container_width=True, hide_index=True)


def get_combined_vector_store(file_names):
    """Get vector store with intelligent caching to avoid redundant processing"""
    ensure_files_processed(file_names)
    selection_key = get_selection_signature(file_names)
    
    # Check cache first
    cached_vs = VECTOR_STORE_CACHE.get(selection_key)
    if cached_vs is not None:
        st.session_state.vector_stores[selection_key] = cached_vs
        return cached_vs
    
    # Create vector store if not cached
    if selection_key not in st.session_state.vector_stores:
        combined_text = "\n".join(st.session_state.file_texts.get(file_name, "") for file_name in file_names)
        vs = create_vector_store(combined_text)
        st.session_state.vector_stores[selection_key] = vs
        VECTOR_STORE_CACHE.set(selection_key, vs)
    return st.session_state.vector_stores[selection_key]


# Shared UI helpers:
# These small functions are reused across multiple tabs to show the current
# sidebar selection, tab-level file context, and floating helper popups.
def show_current_sidebar_selection():
    selected = st.session_state.get("selected_files", [])
    if selected:
        st.info("Sidebar selected files: " + ", ".join(selected))
    else:
        st.info("No sidebar files selected yet. Upload and select files from the sidebar first.")


def render_file_context_card(title, available_files, active_files=None):
    active_files = active_files or []
    chips_html = "".join(
        f"<span class='file-chip'>{html.escape(file_name)}</span>"
        for file_name in active_files[:12]
    )
    if len(active_files) > 12:
        chips_html += f"<span class='file-chip'>+{len(active_files) - 12} more</span>"

    st.markdown(
        f"""
        <div class="app-card">
            <h4>{html.escape(title)}</h4>
            <p class="app-muted">Available from sidebar: {len(available_files)} file(s)</p>
            <p class="app-muted">Selected in this tab: {len(active_files)} file(s)</p>
            <div class="file-chip-wrap">
                {chips_html if chips_html else "<span class='file-chip'>No tab files selected yet</span>"}
            </div>
        </div>
        """,
        unsafe_allow_html=True
    )




# Floating icon function removed - helper is now triggered by header 🧠 icon


def render_autonomous_workspace_shell():
    """Render the connected AI operating system layer above module controls."""
    memory = normalize_workspace_memory(st.session_state.workspace_memory)
    indexed_count = len(memory.get("indexed_files", []))
    chat_count = len(memory.get("chat", []))
    agent_count = len(memory.get("agent_runs", []))
    event_count = len(memory.get("memory_events", []))
    st.markdown(
        f"""
        <div class="ai-os-shell">
            <div class="ai-os-kicker">Autonomous AI Operating System</div>
            <div class="ai-os-title">One shared AI brain across Chat, Upload, Dashboard, Compare, and CAPL.</div>
            <div class="ai-os-loop">
                <span>Chat -> Memory</span>
                <span>Upload -> FAISS</span>
                <span>Dashboard -> Insights</span>
                <span>Compare -> Semantic Diff</span>
                <span>CAPL -> Agents</span>
            </div>
            <div class="ai-os-metrics">
                <span>{indexed_count} indexed files</span>
                <span>{chat_count} chat memories</span>
                <span>{agent_count} agent runs</span>
                <span>{event_count} memory events</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


# Define keywords for each tab
tab_keywords = {
    "chat": ["memory", "overview", "summary", "count", "find", "analyze", "details", "downloads"],
    "dashboard": ["insights", "memory", "themes", "entities", "risks", "charts", "metrics", "reports"],
    "compare": ["semantic", "differences", "compare", "changes", "diff", "side-by-side", "inline", "excel"],
    "capl": ["agents", "syntax", "variables", "errors", "debug", "code", "fix", "run history"]
}


def track_user_behavior(tab_name):
    """Tracks user actions to detect skill level progression."""
    if "behavior_tracker" not in st.session_state:
        st.session_state.behavior_tracker = {
            "chat": {"queries": 0, "actions": []},
            "dashboard": {"queries": 0, "actions": []},
            "compare": {"queries": 0, "actions": []},
            "capl": {"queries": 0, "actions": []}
        }
    
    if tab_name not in st.session_state.behavior_tracker:
        st.session_state.behavior_tracker[tab_name] = {"queries": 0, "actions": []}
    
    st.session_state.behavior_tracker[tab_name]["queries"] += 1
    tracker = st.session_state.behavior_tracker[tab_name]
    return tracker


def infer_user_workflow():
    """Auto-detects user skill level from query patterns."""
    if "behavior_tracker" not in st.session_state:
        return "beginner"
    
    total_queries = sum(t.get("queries", 0) for t in st.session_state.behavior_tracker.values())
    
    if total_queries > 15:
        return "advanced"
    elif total_queries > 5:
        return "intermediate"
    else:
        return "beginner"


def get_dynamic_suggestions(tab_name, skill_level):
    """Returns context-aware suggestions based on skill level."""
    suggestions_by_skill = {
        "chat": {
            "beginner": ["Summarize selected files", "Find a keyword", "Count a phrase"],
            "intermediate": ["Ask from workspace memory", "Extract item details", "Build summary downloads"],
            "advanced": ["Compare answers across files", "Generate engineering reference", "Use chat memory as context", "Extract workflow"]
        },
        "dashboard": {
            "beginner": ["Review memory snapshot", "Show key themes", "Check indexed files"],
            "intermediate": ["Review entities and risks", "Analyze report metrics", "Create charts from reports"],
            "advanced": ["Cross-file insight review", "Memory log inspection", "Risk/theme triage", "Structured report analysis"]
        },
        "compare": {
            "beginner": ["Exact inline diff", "Side-by-side line diff", "Select two files"],
            "intermediate": ["Word presence summary", "Download Excel diff", "Review semantic summary"],
            "advanced": ["Multi-file comparison", "Change impact analysis", "Store comparison memory", "Validate changed sections"]
        },
        "capl": {
            "beginner": ["Analyze CAPL syntax", "Review issue table", "Select a .can file"],
            "intermediate": ["Generate AI fix", "Run CAPL agents", "Inspect unused variables"],
            "advanced": ["Goal-driven agent run", "Retrieve FAISS memory", "Review agent history", "Coordinate final output"]
        }
    }
    
    return suggestions_by_skill.get(tab_name, {}).get(skill_level, [])


def get_next_best_action(tab_name, skill_level):
    """Intelligently recommends the next workflow step."""
    workflow_paths = {
        "chat": {
            "beginner": "Pro Tip: Select files, then ask 'summarize', 'overview', 'find keyword', or 'count phrase'.",
            "intermediate": "Next: Ask targeted questions that use both selected documents and stored chat memory.",
            "advanced": "Next: Combine item extraction, direct commands, and prior chat context to validate details across files."
        },
        "dashboard": {
            "beginner": "Pro Tip: Start with the workspace memory snapshot to confirm what the app has indexed.",
            "intermediate": "Next: Review themes, entities, risks, and report charts together for faster triage.",
            "advanced": "Next: Use the dashboard as a cross-module intelligence view over uploads, chat, compare, and CAPL runs."
        },
        "compare": {
            "beginner": "Pro Tip: Select at least two files, then start with exact inline word diff.",
            "intermediate": "Next: Switch modes to line diff or word presence summary, then export the Excel workbook.",
            "advanced": "Next: Use the semantic summary to capture comparison findings into shared memory."
        },
        "capl": {
            "beginner": "Pro Tip: Select a CAPL file, run analysis, and review line-level issues first.",
            "intermediate": "Next: Use AI fix suggestions or run a focused autonomous CAPL goal.",
            "advanced": "Next: Let the planning, retrieval, execution, reasoning, and coordination agents work across selected files."
        }
    }
    
    return workflow_paths.get(tab_name, {}).get(skill_level, "Keep exploring the features available.")


def show_help_popup(tab_name, selected_files):
    state_key = ensure_help_popup_state(tab_name)

    if not st.session_state[state_key]:
        return

    tracker = track_user_behavior(tab_name)
    skill_level = infer_user_workflow()
    selected_types = {os.path.splitext(f)[1].lower() for f in (selected_files or [])}
    selected_types_text = ", ".join(sorted(selected_types)) if selected_types else "No files selected"

    helper_defs = {
        "chat": {
            "title": "Chat Helper",
            "text": "Chat uses selected uploads plus shared workspace memory. It can summarize, answer questions, retrieve prior context, and create downloadable extracted assets.",
            "hint": "Try: summarize, overview, find \"keyword\", count \"signal\", item details \"VN1630A\", visual reference \"VN1640A\", or compare."
        },
        "dashboard": {
            "title": "Dashboard Helper",
            "text": "Dashboard shows the hidden shared AI brain as useful insights: indexed files, chat memory, agent runs, themes, entities, risks, logs, and report charts.",
            "hint": "Use this tab to inspect workspace intelligence, then open structured HTML/XLSX reports for metrics, verdicts, fixtures, and visual summaries."
        },
        "compare": {
            "title": "Compare Helper",
            "text": "Compare supports exact inline word diff, side-by-side line diff, word presence summaries, semantic explanations, and Excel export.",
            "hint": "Select two or more files, choose a comparison mode, run the comparison, then review the semantic summary and download workbook."
        },
        "capl": {
            "title": "CAPL Helper",
            "text": "CAPL combines code analysis with autonomous agents. It checks syntax, reports issues, suggests fixes, and can run goal-driven workflows over shared memory.",
            "hint": "Select a .can or .txt file for direct analysis, or enter an autonomous goal to run the planning, retrieval, execution, reasoning, and coordination agents."
        }
    }

    helper_def = helper_defs.get(tab_name, helper_defs["chat"])
    suggestions = get_dynamic_suggestions(tab_name, skill_level)[:4]
    suggestion_tags = "".join(f"<span>{html.escape(s)}</span>" for s in suggestions)
    next_action = get_next_best_action(tab_name, skill_level)

    st.markdown(
        f"""
        <style>
        .helper-popup-overlay {{
            position: fixed;
            bottom: 18px;
            right: 18px;
            top: auto;
            width: min(380px, 88vw);
            max-height: 72vh;
            padding: 18px 20px;
            background: rgba(255, 250, 245, 0.98);
            border: 1px solid #ffd5c5;
            border-radius: 22px;
            box-shadow: 0 28px 60px rgba(0, 0, 0, 0.14);
            overflow-y: auto;
            z-index: 99999;
            color: #2f2a26;
            font-family: system-ui, sans-serif;
        }}
        .helper-popup-overlay h3 {{
            margin: 0 0 10px;
            font-size: 1.05rem;
            color: #3c2e2a;
        }}
        .helper-popup-overlay .helper-close-btn {{
            position: absolute;
            top: 12px;
            right: 14px;
            width: 34px;
            height: 34px;
            border: none;
            border-radius: 50%;
            background: rgba(255, 255, 255, 0.95);
            color: #7a4f3a;
            font-size: 1.1rem;
            line-height: 1;
            cursor: pointer;
            box-shadow: 0 2px 10px rgba(0, 0, 0, 0.08);
        }}
        .helper-popup-overlay .helper-close-btn:hover {{
            background: #fee7d0;
        }}
        .helper-popup-overlay .helper-meta {{
            margin-bottom: 12px;
            color: #5f3a31;
            font-size: 0.92rem;
        }}
        .helper-popup-overlay p {{ margin: 10px 0; line-height: 1.55; }}
        .helper-popup-overlay .helper-hint {{
            margin-top: 10px;
            color: #5e3b33;
            font-size: 0.93rem;
        }}
        .helper-popup-overlay .helper-next-action {{
            margin-top: 10px;
            padding: 10px 12px;
            border-radius: 12px;
            background: #fff4eb;
            border: 1px solid rgba(161, 106, 85, 0.22);
            color: #4a2f28;
            font-size: 0.92rem;
            line-height: 1.45;
        }}
        .helper-popup-overlay .helper-pill {{
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 10px;
        }}
        .helper-popup-overlay .helper-pill span {{
            background: #ffe7d6;
            color: #8a3f1b;
            border-radius: 999px;
            padding: 6px 10px;
            font-size: 0.84rem;
            line-height: 1.4;
        }}
        .helper-popup-overlay .helper-footer {{
            margin-top: 14px;
            padding-top: 12px;
            border-top: 1px solid rgba(161, 106, 85, 0.18);
            color: #5f3d34;
            font-size: 0.9rem;
        }}
        </style>
        <div class="helper-popup-overlay">
            <button class="helper-close-btn" onclick="this.closest('.helper-popup-overlay').style.display='none'" aria-label="Close helper">×</button>
            <h3>{html.escape(helper_def['title'])}</h3>
            <div class="helper-meta">Skill: <strong>{skill_level.title()}</strong> · Queries: <strong>{tracker.get('queries', 0)}</strong></div>
            <p>{html.escape(helper_def['text'])}</p>
            <p class="helper-hint">{html.escape(helper_def['hint'])}</p>
            <div class="helper-next-action">{html.escape(next_action)}</div>
            <div class="helper-pill">{suggestion_tags}</div>
            <div class="helper-footer">Selected file types: <strong>{html.escape(selected_types_text)}</strong><br>Click the header 🧠 to close this helper.</div>
        </div>
        """,
        unsafe_allow_html=True
    )


# -------------------------------
# TABS
# -------------------------------
# Main application tabs:
# Each block below owns one visible area of the app. If you want to change a
# feature, start in the matching tab block and then follow the helper comments above.

# Session-backed main navigation:
# 1. Premium "Soft-Glow" Navigation CSS
st.markdown("""
    <style>
    .st-key-active_main_tab {
        margin-top: 16px !important;
        margin-bottom: 8px !important;
        position: relative;
    }

    /* Text-only tabs with glass sliding indicator */
    div[role="radiogroup"] > label > div:first-child { display: none !important; }
    div[role="radiogroup"] {
        gap: 24px;
        display: flex;
        margin: 0 !important;
        padding: 0 !important;
        position: relative;
        justify-content: center;
        border-bottom: 1px solid rgba(59, 130, 246, 0.1);
        padding-bottom: 8px;
    }

    /* Text-only tab styling */
    div[role="radiogroup"] > label {
        background: none !important;
        border: none !important;
        padding: 8px 16px !important;
        border-radius: 0 !important;
        display: flex !important;
        align-items: center !important;
        height: auto;
        min-height: auto !important;
        font-weight: 500;
        line-height: 1.2 !important;
        transition: all 0.3s ease;
        position: relative;
        cursor: pointer;
        color: #64748b;
    }

    /* Active tab text color */
    div[role="radiogroup"] > label[data-checked="true"] {
        color: #3b82f6 !important;
        background: none !important;
        box-shadow: none !important;
    }

    /* Hover effects */
    div[role="radiogroup"] > label:hover {
        color: #3b82f6 !important;
        transform: translateY(-1px);
    }

    /* Glass sliding indicator under active tab */
    div[role="radiogroup"] > label[data-checked="true"]::after {
        content: '';
        position: absolute;
        bottom: -8px;
        left: 50%;
        width: 60px;
        height: 2px;
        background: linear-gradient(90deg, rgba(59, 130, 246, 0.3), rgba(59, 130, 246, 0.6), rgba(59, 130, 246, 0.3));
        border-radius: 1px;
        transform: translateX(-50%);
        box-shadow: 0 0 8px rgba(59, 130, 246, 0.4);
        animation: indicatorBreath 3s ease-in-out infinite;
    }

    /* Breathing animation for indicator */
    @keyframes indicatorBreath {
        0%, 100% { opacity: 0.6; transform: translateX(-50%) scaleY(1); }
        50% { opacity: 1; transform: translateX(-50%) scaleY(1.2); }
    }

    /* Remove old styling */
    div[role="radiogroup"] > label::after {
        display: none !important;
    }

    /* --- RESPONSIVE DESIGN --- */
    /* Mobile First: Base styles are mobile-friendly */

    /* Tablet and up */
    @media (min-width: 768px) {
        div[role="radiogroup"] {
            flex-direction: row !important;
            justify-content: center;
        }
        div[role="radiogroup"] > label {
            flex: 0 0 auto;
        }
    }

    /* Desktop */
    @media (min-width: 1024px) {
        div[role="radiogroup"] > label {
            padding: 8px 20px !important;
        }
    }

    /* Mobile: Stack tabs vertically */
    @media (max-width: 767px) {
        div[role="radiogroup"] {
            flex-direction: column !important;
            gap: 16px;
            align-items: center;
            padding-bottom: 16px;
        }
        div[role="radiogroup"] > label {
            width: auto !important;
            min-height: auto;
            padding: 12px 16px !important;
            font-size: 16px !important; /* Prevent zoom on iOS */
            text-align: center;
        }
        div[role="radiogroup"] > label[data-checked="true"]::after {
            width: 40px;
            bottom: -12px;
        }
    }

    /* Scroll-sync navigation */
    </style>
    <script>
    document.addEventListener('DOMContentLoaded', function() {
        const sections = ['chat-section', 'dashboard-section', 'compare-section', 'capl-section'];
        const tabLabels = document.querySelectorAll('div[role="radiogroup"] > label');
        
        const observerOptions = {
            root: null,
            rootMargin: '-50% 0px -50% 0px',
            threshold: 0
        };
        
        const observer = new IntersectionObserver((entries) => {
            entries.forEach(entry => {
                if (entry.isIntersecting) {
                    const sectionId = entry.target.id;
                    const tabIndex = sections.indexOf(sectionId);
                    if (tabIndex !== -1 && tabLabels[tabIndex]) {
                        // Remove active state from all tabs
                        tabLabels.forEach(label => label.removeAttribute('data-checked'));
                        // Add active state to current tab
                        tabLabels[tabIndex].setAttribute('data-checked', 'true');
                    }
                }
            });
        }, observerOptions);
        
        // Observe sections
        sections.forEach(sectionId => {
            const section = document.getElementById(sectionId);
            if (section) observer.observe(section);
        });
    });
    </script>
    <style>
    @media (max-width: 767px) {
        .stButton > button {
            min-height: 48px !important;
            padding: 12px 16px !important;
            font-size: 16px !important;
        }
        
        /* Larger file cards for touch */
        [data-testid="stSidebar"] [class*="st-key-select_file_"] button {
            min-height: 80px !important;
            padding: 16px 12px !important;
        }
        
        /* Sidebar improvements for mobile */
        [data-testid="stSidebar"] {
            width: 280px !important;
        }
    }

    /* Prevent horizontal scroll */
    body {
        overflow-x: hidden;
    }

    /* Responsive text sizes */
    .login-heading {
        font-size: clamp(1.8rem, 5vw, 3.4rem) !important;
        color: #F8FAFC !important;
        font-weight: 700 !important;
        margin-bottom: 12px !important;
    }

    .login-subheading {
        font-size: clamp(0.9rem, 2.5vw, 1.1rem) !important;
        color: rgba(248, 250, 252, 0.85) !important;
        font-weight: 500 !important;
        margin-bottom: 24px !important;
    }

    .login-tagline {
        font-size: clamp(0.9rem, 2.5vw, 1rem);
    }

    /* Content responsiveness */
    .app-card {
        padding: clamp(1rem, 2vw, 1.5rem);
    }

    .metric-card {
        padding: clamp(1rem, 2vw, 1.5rem);
    }

    /* Table responsiveness */
    .scrollable {
        overflow-x: auto;
    }

    /* Ensure images are responsive */
    img {
        max-width: 100%;
        height: auto;
    }

    /* Form inputs on mobile */
    @media (max-width: 767px) {
        .glass-card .stTextInput input {
            font-size: 16px !important; /* Prevent zoom */
            padding: 16px !important;
            min-height: 48px !important;
        }
        
        .glass-card .stButton > button {
            min-height: 48px !important;
            font-size: 16px !important;
        }
    }

    /* Dashboard grid responsiveness */
    @media (max-width: 767px) {
        .dashboard-grid {
            grid-template-columns: 1fr;
            gap: 1rem;
        }
    }

    /* Dashboard grid responsiveness */
    @media (max-width: 767px) {
        .dashboard-grid {
            grid-template-columns: 1fr;
            gap: 1rem;
        }
    }

    /* Chat and other content */
    @media (max-width: 767px) {
        .stMarkdown {
            font-size: 14px;
        }
        
        .stDataFrame {
            font-size: 12px;
        }
    }

    /* ============================================ */
    /* ENHANCED RESPONSIVE DESIGN FOR ALL DEVICES */
    /* ============================================ */
    
    /* Ultra-responsive design for all screen sizes */
    @media (max-width: 320px) {
        /* Smartphones - extra small */
        .stApp {
            padding: 0.25rem;
        }
        .stSidebar {
            width: 100% !important;
        }
        h1 {
            font-size: 1.2rem !important;
        }
        h2 {
            font-size: 1rem !important;
        }
        .stButton > button {
            padding: 0.5rem !important;
            font-size: 0.75rem !important;
            min-width: 40px !important;
        }
    }

    @media (min-width: 321px) and (max-width: 480px) {
        /* Smartphones */
        .block-container {
            padding-left: 0.5rem !important;
            padding-right: 0.5rem !important;
        }
        .stColumn {
            padding: 0.5rem 0.25rem;
        }
        .metric-card {
            padding: 0.75rem !important;
        }
        .stMetric {
            font-size: 0.9rem;
        }
    }

    @media (min-width: 481px) and (max-width: 768px) {
        /* Tablets (portrait) */
        .block-container {
            padding-left: 1rem !important;
            padding-right: 1rem !important;
        }
        .dashboard-grid {
            grid-template-columns: repeat(2, 1fr);
        }
    }

    @media (min-width: 769px) and (max-width: 1024px) {
        /* Tablets (landscape) & small laptops */
        .dashboard-grid {
            grid-template-columns: repeat(3, 1fr);
        }
    }

    @media (min-width: 1025px) and (max-width: 1440px) {
        /* Laptops & desktops */
        .dashboard-grid {
            grid-template-columns: repeat(4, 1fr);
        }
    }

    @media (min-width: 1441px) {
        /* Large screens & projectors */
        .dashboard-grid {
            grid-template-columns: repeat(5, 1fr);
        }
        .block-container {
            max-width: 1920px;
        }
    }

    /* Projector & presentation mode */
    @media screen and (min-height: 1080px) {
        .stApp {
            background: #ffffff;
            font-size: 18px;
        }
        .stMarkdown {
            font-size: 18px;
        }
        h1 { font-size: 2.5rem; }
        h2 { font-size: 2rem; }
        h3 { font-size: 1.5rem; }
        .stButton > button {
            min-height: 60px;
            font-size: 18px;
            padding: 1rem;
        }
    }

    /* Touch-friendly interface (mobile & tablet) */
    @media (hover: none) {
        .stButton > button,
        [role="button"] {
            min-height: 48px;
            padding: 12px 16px;
            font-size: 16px;
        }
        input, select, textarea {
            min-height: 44px;
            padding: 12px;
        }
    }

    /* High DPI screens (retina) */
    @media (-webkit-min-device-pixel-ratio: 2), (min-resolution: 192dpi) {
        body {
            -webkit-font-smoothing: antialiased;
            -moz-osx-font-smoothing: grayscale;
        }
    }

    /* Landscape mode adjustments */
    @media (orientation: landscape) and (max-height: 500px) {
        .stApp {
            margin: 0;
            padding: 0;
        }
        .block-container {
            padding: 0.5rem;
        }
    }

    /* Print friendly */
    @media print {
        .stButton, [role="button"], .stSidebar,
        [data-testid="stHeader"],
        [data-testid="stToolbar"],
        footer {
            display: none !important;
        }
        body {
            background: white;
            color: black;
        }
        img {
            max-width: 100%;
            page-break-inside: avoid;
        }
    }

    /* Dark mode optimization */
    @media (prefers-color-scheme: dark) {
        :root {
            --background: #1a1a1a;
            --surface: #2d2d2d;
            --text: #e0e0e0;
        }
    }

    @keyframes selectedPillGlow {
        0%, 100% {
            box-shadow: 0 0 0 3px rgba(96, 165, 250, 0.18), 0 0 14px rgba(59, 130, 246, 0.28);
        }
        50% {
            box-shadow: 0 0 0 4px rgba(96, 165, 250, 0.24), 0 0 26px rgba(59, 130, 246, 0.48);
        }
    }

    .st-key-active_main_tab div[role="radiogroup"] > label[data-checked="true"] {
        background: linear-gradient(135deg, #dbeafe 0%, #e0f2fe 48%, #ede9fe 100%) !important;
        border: 2px solid #60a5fa !important;
        color: #0f172a !important;
        font-weight: 800 !important;
        transform: translateY(-1px);
        animation: selectedPillGlow 2.2s ease-in-out infinite;
    }

    .st-key-active_main_tab div[role="radiogroup"] > label[data-checked="true"] p,
    .st-key-dashboard_chart_type div[role="radiogroup"] > label[data-checked="true"] p,
    .st-key-dashboard_bar_orientation div[role="radiogroup"] > label[data-checked="true"] p {
        color: #0f172a !important;
        font-weight: 800 !important;
    }

    .st-key-dashboard_chart_type div[role="radiogroup"],
    .st-key-dashboard_bar_orientation div[role="radiogroup"] {
        justify-content: flex-start !important;
        margin-bottom: 0.5rem;
    }

    .st-key-dashboard_chart_type div[role="radiogroup"] > label,
    .st-key-dashboard_bar_orientation div[role="radiogroup"] > label {
        min-width: 132px !important;
        border-radius: 14px !important;
        border: 1px solid #d7e3f4 !important;
        background: #f8fbff !important;
    }

    .st-key-dashboard_chart_type div[role="radiogroup"] > label[data-checked="true"],
    .st-key-dashboard_bar_orientation div[role="radiogroup"] > label[data-checked="true"] {
        background: linear-gradient(135deg, #ecfeff 0%, #dbeafe 100%) !important;
        border: 2px solid #38bdf8 !important;
        color: #0f172a !important;
        animation: selectedPillGlow 2.2s ease-in-out infinite;
    }

    /* Reduced motion for accessibility */
    @media (prefers-reduced-motion: reduce) {
        *,
        *::before,
        *::after {
            animation-duration: 0.01ms !important;
            animation-iteration-count: 1 !important;
            transition-duration: 0.01ms !important;
        }
    }
    </style>
""", unsafe_allow_html=True)


st.markdown(
    """
    <style>
    /* Header sizing fixes */
    [data-testid="stHorizontalBlock"]:has(.st-key-header_brain_icon) {
        align-items: center !important;
        gap: 0.75rem !important;
    }
    .st-key-header_brain_icon button {
        width: 46px !important;
        min-width: 46px !important;
        height: 46px !important;
        min-height: 46px !important;
        padding: 0 !important;
        font-size: 1.55rem !important;
        line-height: 1 !important;
        border-radius: 12px !important;
    }
    .st-key-main_logout_btn button {
        min-height: 42px !important;
        padding: 0.55rem 0.75rem !important;
        white-space: nowrap !important;
    }

    /* Mobile layout repair: keep sidebar and main content from overlapping. */
    @media (max-width: 767px) {
        html, body, .stApp {
            width: 100% !important;
            max-width: 100% !important;
            overflow-x: hidden !important;
        }
        div[data-testid="stAppViewContainer"] {
            width: 100% !important;
            max-width: 100% !important;
            overflow-x: hidden !important;
        }
        .block-container,
        .main .block-container,
        section.main .block-container,
        div[data-testid="stMain"] .block-container {
            width: 100% !important;
            max-width: 100% !important;
            padding: 0.75rem !important;
            margin: 0 !important;
        }
        [data-testid="stSidebar"],
        .stSidebar {
            width: 100% !important;
            min-width: 0 !important;
            max-width: 100% !important;
            position: relative !important;
            transform: none !important;
            overflow-x: hidden !important;
        }
        [data-testid="stSidebar"] > div {
            width: 100% !important;
            max-width: 100% !important;
            padding-left: 0.75rem !important;
            padding-right: 0.75rem !important;
        }
        [data-testid="stHorizontalBlock"] {
            flex-wrap: wrap !important;
            gap: 0.5rem !important;
        }
        [data-testid="stHorizontalBlock"] > div {
            min-width: 0 !important;
        }
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_chat_selection),
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_dashboard_selection),
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_compare_selection),
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_capl_selection) {
            display: grid !important;
            grid-template-columns: minmax(0, 1fr) auto !important;
            align-items: center !important;
            column-gap: 0.75rem !important;
        }
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_chat_selection) > div,
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_dashboard_selection) > div,
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_compare_selection) > div,
        [data-testid="stHorizontalBlock"]:has(.st-key-reset_capl_selection) > div {
            width: auto !important;
            min-width: 0 !important;
        }
        .st-key-reset_chat_selection button,
        .st-key-reset_dashboard_selection button,
        .st-key-reset_compare_selection button,
        .st-key-reset_capl_selection button {
            width: auto !important;
            min-width: 92px !important;
            white-space: nowrap !important;
            padding-left: 0.75rem !important;
            padding-right: 0.75rem !important;
        }
        [data-testid="stHorizontalBlock"]:has(.st-key-header_brain_icon) {
            display: grid !important;
            grid-template-columns: 52px minmax(0, 1fr) !important;
        }
        [data-testid="stHorizontalBlock"]:has(.st-key-main_logout_btn) {
            display: flex !important;
            align-items: center !important;
        }
        .st-key-main_logout_btn {
            margin-left: auto !important;
        }
        div[role="radiogroup"] {
            flex-direction: row !important;
            flex-wrap: wrap !important;
            justify-content: stretch !important;
            gap: 8px !important;
        }
        div[role="radiogroup"] > label {
            flex: 1 1 calc(50% - 8px) !important;
            min-width: 135px !important;
            width: auto !important;
        }
        .stSelectbox,
        .stMultiSelect,
        .stTextInput,
        .stTextArea,
        .stNumberInput,
        .stSlider,
        .stCheckbox,
        .stRadio,
        .stDownloadButton,
        .stFileUploader,
        .stDataFrame,
        [data-testid="stDataFrame"],
        [data-testid="stMetric"],
        [data-testid="stPlotlyChart"],
        iframe {
            width: 100% !important;
            max-width: 100% !important;
            min-width: 0 !important;
        }
        .stDataFrame,
        [data-testid="stDataFrame"],
        [data-testid="stTable"],
        [data-testid="stPlotlyChart"] {
            overflow-x: auto !important;
        }
        [data-testid="column"],
        [data-testid="stColumn"] {
            min-width: min(100%, 260px) !important;
        }
        [data-testid="stMetric"] {
            overflow-wrap: anywhere !important;
        }
        .app-card,
        .file-chip-wrap,
        .file-chip {
            max-width: 100% !important;
            overflow-wrap: anywhere !important;
        }
        .dashboard-grid {
            grid-template-columns: 1fr !important;
        }
    }

    @media (min-width: 768px) and (max-width: 1180px) {
        [data-testid="stHorizontalBlock"] {
            gap: 0.75rem !important;
        }
        [data-testid="column"],
        [data-testid="stColumn"] {
            min-width: 0 !important;
        }
        .stDataFrame,
        [data-testid="stDataFrame"],
        [data-testid="stTable"],
        [data-testid="stPlotlyChart"],
        iframe {
            max-width: 100% !important;
            overflow-x: auto !important;
        }
        div[role="radiogroup"] {
            flex-wrap: wrap !important;
        }
    }
    </style>
    """,
    unsafe_allow_html=True,
)

st.markdown(
    """
    <style>
    .ai-os-shell {
        margin: 10px 0 14px;
        padding: 16px 18px;
        border: 1px solid rgba(148, 163, 184, 0.24);
        border-radius: 18px;
        background: rgba(248, 251, 255, 0.72);
        box-shadow: 0 18px 50px rgba(15, 23, 42, 0.08);
        backdrop-filter: blur(14px);
        animation: aiOsRise 0.45s ease-out both;
    }
    .ai-os-kicker {
        color: #2563eb;
        font-size: 0.78rem;
        font-weight: 800;
        text-transform: uppercase;
        letter-spacing: 0.08em;
    }
    .ai-os-title {
        margin-top: 4px;
        color: #0f172a;
        font-size: 1.08rem;
        font-weight: 800;
    }
    .ai-os-loop,
    .ai-os-metrics {
        display: flex;
        flex-wrap: wrap;
        gap: 8px;
        margin-top: 10px;
    }
    .ai-os-loop span,
    .ai-os-metrics span {
        padding: 6px 10px;
        border: 1px solid rgba(147, 197, 253, 0.45);
        border-radius: 999px;
        background: rgba(255, 255, 255, 0.62);
        color: #173152;
        font-size: 0.84rem;
        transition: transform 0.18s ease, background-color 0.18s ease;
    }
    .ai-os-loop span:hover,
    .ai-os-metrics span:hover {
        transform: translateY(-1px) scale(1.01);
        background: rgba(239, 246, 255, 0.92);
    }
    @keyframes aiOsRise {
        from { opacity: 0; transform: translateY(8px); }
        to { opacity: 1; transform: translateY(0); }
    }
    </style>
    """,
    unsafe_allow_html=True,
)


# -------------------------------
# MAIN TAB NAVIGATION
# -------------------------------
# Creates the horizontal tab navigation with custom styling.
# Each tab corresponds to a major feature area of the application.
main_tab_options = ["💬 Chat", "📊 Dashboard", "📂 Compare", "📡 CAPL"]
active_main_tab = st.radio("Open Section", main_tab_options, horizontal=True, key="active_main_tab", label_visibility="collapsed")

# -------------------------------
# TAB 1: CHAT
# -------------------------------
# Chat tab:
# Handles per-tab file selection, direct commands like summarize/find/count,
# semantic Q&A via vector search, and chat-specific download assets.
if active_main_tab == "💬 Chat":
    st.markdown('<div id="chat-section">', unsafe_allow_html=True)
    chat_header_col, chat_reset_col = st.columns([8, 1])
    with chat_header_col:
        st.subheader("Chat with Selected Documents")
    with chat_reset_col:
        if st.button(" 🧼 Reset", key="reset_chat_selection", help="Reset chat selection"):
            st.session_state.chat_file_selection = []
            st.session_state.chat_summary_downloads = empty_chat_summary_downloads()
            st.session_state.messages = []
            st.success("✅ Chat reset!")
            st.rerun()

    st.info(
        "Choose files in the sidebar to make them available here. Then select only the files you want for Chat in this tab.")
    show_current_sidebar_selection()
    render_file_context_card("Chat File Context", st.session_state.selected_files, st.session_state.chat_file_selection)

    show_help_popup('chat', st.session_state.selected_files)

    if st.session_state.selected_files:
        st.session_state.chat_file_selection = [
            file_name for file_name in st.session_state.chat_file_selection
            if file_name in st.session_state.selected_files
        ]
        chat_files = st.multiselect("Choose file(s) for Chat", options=st.session_state.selected_files,
                                    default=st.session_state.chat_file_selection, key="chat_file_selection")
        if not chat_files:
            st.info("Choose one or more files in this tab to start chatting.")
        else:
            with st.spinner("Loading selected files..."):
                ensure_files_processed(chat_files)
            combined_text = "\n".join([st.session_state.file_texts.get(f, "") for f in chat_files])
    
            user_input = st.chat_input("Ask something... Try: analyze, item details \"VN1630A\", pin diagram \"D-SUB9\", find \"keyword\", count \"signal\", overview")
            if user_input:
                if user_input.strip().lower() == "clear":
                    st.session_state.messages = []
                    st.session_state.chat_summary_downloads = empty_chat_summary_downloads()
                    st.success("✅ Chat cleared!")
                else:
                    st.session_state.messages.append({"role": "user", "content": user_input})
                    with st.spinner("Processing your request..."):
                        st.session_state.chat_summary_downloads = empty_chat_summary_downloads()
                        user_input_lower = user_input.lower()
                        # Word count queries
                        if any(t in user_input_lower for t in ["how many", "count", "number of", "occurrences"]):
                            match = re.search(r"'(.*?)'|\"(.*?)\"", user_input)
                            if match:
                                word = match.group(1) or match.group(2)
                                count = len(
                                    re.findall(rf'(?<![\w-]){re.escape(word)}(?![\w-])', combined_text, re.IGNORECASE))
                                response = f"🔢 The word/phrase '{word}' appears {count} times in the selected documents."
                            else:
                                response = "⚠️ Specify the word/phrase in quotes. Example: count('keyword') or count(\"keyword\")"
                        elif any(term in user_input_lower for term in ["find", "search", "locate"]) or "highlight" in user_input_lower:
                            match = re.search(r"'(.*?)'|\"(.*?)\"", user_input)
                            if match:
                                query = match.group(1) or match.group(2)
                                response_blocks = []
                                for f in chat_files:
                                    file_text = st.session_state.file_texts.get(f, "")
                                    response_blocks.append(build_highlighted_search_results(f, file_text, query))
                                response = "".join(response_blocks)
                            else:
                                response = "⚠️ Specify the search word or phrase in quotes. Example: find('keyword') or search(\"keyword\")"
                        elif "overview" in user_input_lower:
                            response_lines = []
                            for f in chat_files:
                                file_text = st.session_state.file_texts.get(f, "")
                                if file_text.strip():
                                    toc_entries = extract_toc_with_page_numbers(file_text)
                                    all_headings = extract_document_headings(file_text)
                                    if all_headings:
                                        response_lines.append(f"📄 **{f}**")
                                        response_lines.append("### Document Headings")
                                        for num, title in all_headings:
                                            content_str = f"{num} {title}" if num else title
                                            page_num = resolve_heading_page_number(file_text, title, toc_entries)
                                            display_text = f"{content_str} (Page {page_num})" if page_num else content_str
                                            preview_link = create_preview_link(f, highlight_term=title, page_num=page_num)
                                            if preview_link:
                                                anchor_id = create_heading_anchor(title)
                                                response_lines.append(f"- <a href='{preview_link}#{anchor_id}' target='_blank'>{html.escape(display_text)}</a>")
                                            else:
                                                response_lines.append(f"- {html.escape(display_text)}")
                                    else:
                                        response_lines.append(f"📄 **{f}**\n\nNo document headings were detected.")
                                else:
                                    response_lines.append(f"📄 **{f}**\n\nNo readable content found in this document.")
                            response = "\n\n".join(response_lines)
                        elif any(term in user_input_lower for term in ["pin diagram", "pin table", "pin configuration", "visual reference", "visual and structural", "connector details", "technical tables"]):
                            item_name = extract_quoted_item_name(user_input)
                            if not item_name:
                                response = "⚠️ Specify the item name in quotes. Example: pin diagram \"D-SUB9\" or visual reference \"VN1640A\""
                            else:
                                response_blocks = []
                                pin_csv_downloads = []
                                ascii_diagram_downloads = []
                                for f in chat_files:
                                    file_text = st.session_state.file_texts.get(f, "")
                                    if file_text.strip():
                                        response_blocks.append(build_item_visual_response(f, file_text, item_name))
                                        visual_assets = build_item_visual_assets(f, file_text, item_name)
                                        pin_csv_downloads.extend(visual_assets.get("csv", []))
                                        ascii_diagram_downloads.extend(visual_assets.get("diagrams", []))
                                    else:
                                        response_blocks.append(f"📄 **{f}**\n\nNo readable content found in this document.")
                                st.session_state.chat_summary_downloads = {
                                    "images": [],
                                    "tables": [],
                                    "csv": pin_csv_downloads,
                                    "diagrams": ascii_diagram_downloads,
                                }
                                response = "\n\n---\n\n".join(response_blocks)
                        elif any(term in user_input_lower for term in ["item details", "item information", "extract item", "about item", "information about", "details about"]):
                            item_name = extract_quoted_item_name(user_input)
                            if not item_name:
                                response = "⚠️ Specify the item name in quotes. Example: item details \"VN1630A\" or information about \"D-SUB9\""
                            else:
                                response_blocks = []
                                for f in chat_files:
                                    file_text = st.session_state.file_texts.get(f, "")
                                    if file_text.strip():
                                        response_blocks.append(build_item_information_response(f, file_text, item_name))
                                    else:
                                        response_blocks.append(f"📄 **{f}**\n\nNo readable content found in this document.")
                                response = "\n\n---\n\n".join(response_blocks)
                        elif any(term in user_input_lower for term in ["analyze", "summary", "summarize", "summarise"]):
                            result = []
                            summary_image_downloads = []
                            summary_table_downloads = []
                            for f in chat_files:
                                file_text = st.session_state.file_texts.get(f, "")
                                file_entry = get_uploaded_file_entry(f)
                                if file_text.strip():
                                    file_bytes = file_entry["bytes"] if file_entry else b""
                                    result.append(build_detailed_document_summary(f, file_bytes, file_text))
                                    page_match = re.search(r"Total Pages:\s*(\d+)", file_text)
                                    page_count = int(page_match.group(1)) if page_match else 0
                                    is_large_pdf = f.lower().endswith(".pdf") and page_count > PDF_ASSET_SCAN_PAGE_LIMIT
                                    if not is_large_pdf:
                                        summary_assets = build_summary_download_assets(f, file_bytes)
                                        summary_image_downloads.extend(summary_assets.get("images", []))
                                        summary_table_downloads.extend(summary_assets.get("tables", []))
                                else:
                                    result.append(f"📄 **{f}**\n\nNo readable content found in this document.")

                            st.session_state.chat_summary_downloads = {
                                "images": summary_image_downloads,
                                "tables": summary_table_downloads,
                                "csv": [],
                                "diagrams": [],
                            }
                            response = "\n\n---\n\n".join(result)
                        elif "compare" in user_input_lower:
                            if len(chat_files) >= 2:
                                selected_texts = {f: st.session_state.file_texts[f] for f in chat_files}
                                response = highlight_multi_file_differences(selected_texts)
                            else:
                                response = "⚠️ Please select at least 2 files to compare documents."
                        else:
                            combined_vs = get_workspace_vector_store(chat_files) or get_combined_vector_store(chat_files)
                            retriever = combined_vs.as_retriever(search_kwargs={"k": 3})
                            llm = load_llm()
                            prompt = ChatPromptTemplate.from_messages([
                                ("system",
                                 "You are the conversational engine inside a unified AI workspace. Use uploaded document context, prior chat memory, agent memory, and workspace logs when relevant.\nIf information is not found in shared memory, say 'This information is not available in the uploaded documents or workspace memory.'\nContext:\n{context}"),
                                ("human", "{question}")
                            ])
                            chain = None
                            if llm is not None:
                                try:
                                    chain = ({"context": retriever | (lambda docs: '\n'.join(getattr(doc, "page_content", str(doc)) for doc in docs)),
                                              "question": RunnablePassthrough()} | prompt | llm)
                                except Exception as e:
                                    st.warning(f"Could not create LLM chain: {e}")
                                    chain = None

                            if chain is not None:
                                response = str(chain.invoke(user_input))
                            else:
                                memory_hits = search_workspace_memory(user_input, limit=4)
                                if memory_hits:
                                    response = "AI model is unavailable, so I retrieved the closest workspace memory:\n\n" + "\n\n---\n\n".join(memory_hits)
                                else:
                                    response = "⚠️ AI model is unavailable. Use direct extraction questions such as 'count(\"keyword\")', 'find(\"phrase\")', 'summarize', or 'overview'."
                        st.session_state.messages.append({"role": "assistant", "content": response})
                        st.session_state.last_streamed_assistant_index = len(st.session_state.messages) - 1
                        append_chat_to_workspace_memory(user_input, response, chat_files)
                        save_workspace_memory()
                        save_memory_log("chat", "Chat interaction stored in workspace memory.", {
                            "user": user_input,
                            "files": chat_files,
                            "assistant_preview": response[:300],
                        })
                        if "⚠️" in response or "not found" in response.lower() or "please select" in response.lower() or "ai model is unavailable" in response.lower():
                            set_help_popup_state("chat", True)

        for msg_index, msg in enumerate(st.session_state.messages):
            role = "🧑" if msg["role"] == "user" else "🤖"
            st.markdown(f"**{role}**", unsafe_allow_html=True)
            if msg["role"] == "assistant" and msg_index == st.session_state.get("last_streamed_assistant_index"):
                placeholder = st.empty()
                content = str(msg["content"])
                tokens = re.split(r"(\s+)", content)
                streamed = ""
                for token_index, token in enumerate(tokens):
                    streamed += token
                    if token_index < 240:
                        placeholder.markdown(streamed + "▌", unsafe_allow_html=True)
                        time.sleep(0.006)
                placeholder.markdown(content, unsafe_allow_html=True)
                st.session_state.last_streamed_assistant_index = None
            else:
                st.markdown(msg["content"], unsafe_allow_html=True)

        render_chat_summary_downloads()
    else:
        st.info("Select files from the sidebar to start chatting.")

    st.markdown('</div>', unsafe_allow_html=True)

# -------------------------------
# TAB 2: DASHBOARD
# -------------------------------
# Dashboard tab:
# Focused on structured HTML/XLSX analysis, charts, login/stat extraction, and
# grouped test fixture reporting for uploaded report files.
if active_main_tab == "📊 Dashboard":
    st.markdown('<div id="dashboard-section">', unsafe_allow_html=True)
    dashboard_header_col, dashboard_reset_col = st.columns([8, 1])
    with dashboard_header_col:
        st.subheader("Dashboard")
    with dashboard_reset_col:
        if st.button("🧼 Reset", key="reset_dashboard_selection"):
            st.session_state.file_dropdown = "--Select File--"
            st.session_state.dashboard_chart_type = "Pie Chart"
            st.session_state.dashboard_bar_orientation = "Vertical"
            st.rerun()

    show_current_sidebar_selection()
    show_help_popup('dashboard', [
        f for f in st.session_state.selected_files
        if f.lower().endswith((".html", ".htm", ".xlsx"))
    ])

    st.markdown("### Workspace Memory Snapshot")
    st.markdown(f"- Indexed files: **{len(st.session_state.workspace_memory.get('indexed_files', []))}**")
    st.markdown(f"- Chat history entries: **{len(st.session_state.workspace_memory.get('chat', []))}**")
    st.markdown(f"- Autonomous CAPL runs: **{len(st.session_state.workspace_memory.get('agent_runs', []))}**")
    render_workspace_intelligence_panel(st.session_state.selected_files)

    # Filter selected files for dashboard-compatible formats
    dashboard_files = [
        f for f in st.session_state.selected_files
        if f.lower().endswith((".html", ".htm", ".xlsx"))
    ]
    active_dashboard_files = [] if st.session_state.file_dropdown == "--Select File--" else [st.session_state.file_dropdown]
    render_file_context_card("Dashboard File Context", dashboard_files, active_dashboard_files)


    def clean_text(x):
        return re.sub(r"\s+", " ", x).strip().lower()


    def safe_int(x):
        nums = re.findall(r"\d+", str(x))
        return int(nums[0]) if nums else 0


    def sort_key(x):
        return [int(i) for i in x.split(".")]


    def format_label(x):
        level = x.count(".")
        return "   " * level + x


    def extract_timestamp_from_line(line):
        """Extract timestamp or identifier from a verdict line"""
        # Try to extract decimal number first (e.g., 35.877473)
        decimal_match = re.search(r'\b(\d+\.\d+)\b', line)
        if decimal_match:
            return decimal_match.group(1)

        # Try to extract time pattern (HH:MM:SS or HH:MM)
        time_match = re.search(r'\b(\d{1,2}:\d{2}:\d{2})\b', line)
        if time_match:
            return time_match.group(1)

        # Try to extract date-time pattern
        datetime_match = re.search(r'\d{4}-\d{2}-\d{2}\s+\d{1,2}:\d{2}:\d{2}', line)
        if datetime_match:
            return datetime_match.group(0)

        # Try to extract test case ID (e.g., "1.", "Test Case 1", "TC_001")
        id_match = re.search(r'(?:test\s+case|tc)[\s_-]*(\d+)', line, re.IGNORECASE)
        if id_match:
            return f"TC {id_match.group(1)}"

        # Extract first number or identifier
        num_match = re.search(r'^\d+', line)
        if num_match:
            return num_match.group(0)

        # Extract the first meaningful text before colon or special char
        first_word = re.match(r'^([a-zA-Z0-9_\-\.]+)', line)
        if first_word:
            return first_word.group(1)

        # Fallback to full line (truncated)
        return line[:30]


    def extract_test_results_grouped(soup):
        """
        Extract test results from HTML grouped by test fixtures.
        Dynamically discovers fixtures from HTML structure and their numerical counts.
        Only counts the main test case verdict, not sub-step verdicts.
        Captures detailed failure information for failed test cases (step ID, action, timestamp).
        """
        results = {}

        # First, extract fixture names and counts from GroupHeadingTable structures
        group_tables = soup.find_all('table', class_='GroupHeadingTable')

        for group_table in group_tables:
            try:
                rows = group_table.find_all('tr')
                if len(rows) >= 2:
                    # First row contains fixture name in Heading3
                    first_row = rows[0]
                    heading = first_row.find('big', class_='Heading3')

                    if heading:
                        heading_text = heading.get_text(strip=True)
                        fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', heading_text, re.IGNORECASE)

                        if fixture_match:
                            fixture_name = fixture_match.group(1).strip()

                            # Second row contains the count in OverviewResultTable
                            second_row = rows[1]
                            overview_table = second_row.find('table', class_='OverviewResultTable')

                            if overview_table:
                                count_cell = overview_table.find('td')
                                if count_cell:
                                    try:
                                        count = int(count_cell.get_text(strip=True))

                                        if fixture_name not in results:
                                            results[fixture_name] = {
                                                "name": fixture_name,
                                                "test_cases": [],
                                                "pass": count,  # All are passed by default from green cells
                                                "fail": 0,
                                                "error": 0,
                                                "not executed": 0,
                                                "inconclusive": 0,
                                                "total": count,
                                                "count_cell_class": count_cell.get('class', [''])[0]
                                            }
                                    except ValueError:
                                        pass
            except Exception as e:
                pass

        # Parse full text for verdict distribution
        full_text = soup.get_text("\n", strip=True)
        lines = [l.strip() for l in full_text.split("\n") if l.strip()]

        current_fixture = None

        for i, line in enumerate(lines):
            line_lower = line.lower()

            # Check for fixture marker
            if "test fixture:" in line_lower:
                fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', line, re.IGNORECASE)
                if fixture_match:
                    current_fixture = fixture_match.group(1).strip()
                    if current_fixture not in results:
                        results[current_fixture] = {
                            "name": current_fixture,
                            "test_cases": [],
                            "pass": 0,
                            "fail": 0,
                            "error": 0,
                            "not executed": 0,
                            "inconclusive": 0,
                            "total": 0
                        }

            # Check for test case number with verdict on SAME line (e.g., "1.2.1 ...description...: Passed")
            # Only match if line contains both test number and a verdict
            elif re.match(r'^\d+\.\d+', line) and current_fixture:
                verdict_match = re.search(r':\s*(Passed|Failed|Pass|Fail|Error|Not Executed|Inconclusive)\s*$', line,
                                          re.IGNORECASE)

                if verdict_match:
                    verdict_word = verdict_match.group(1).lower()

                    # Normalize verdict and increment counter
                    if "pass" in verdict_word:
                        verdict_type = "Pass"
                        results[current_fixture]["pass"] += 1
                    elif "fail" in verdict_word:
                        verdict_type = "Failed"
                        results[current_fixture]["fail"] += 1
                    elif "error" in verdict_word:
                        verdict_type = "Error"
                        results[current_fixture]["error"] += 1
                    elif "not executed" in verdict_word:
                        verdict_type = "Not Executed"
                        results[current_fixture]["not executed"] += 1
                    elif "inconclusive" in verdict_word:
                        verdict_type = "Inconclusive"
                        results[current_fixture]["inconclusive"] += 1
                    else:
                        continue

                    # Now look forward for timestamp and failure details
                    timestamp = None
                    test_step = "Step"
                    failure_step_id = ""

                    def score_timestamp(candidate):
                        if not candidate:
                            return -1
                        parts = candidate.split('.')
                        if len(parts) != 2 or not parts[0].isdigit() or not parts[1].isdigit():
                            return -1
                        leading_num = int(parts[0])
                        decimal_places = len(parts[1])
                        decimal_bonus = 10000 if decimal_places >= 3 else (100 if decimal_places == 2 else 0)
                        return decimal_bonus + leading_num

                    def find_best_timestamp(text):
                        # Find all decimal numbers
                        matches = re.findall(r'\b(\d+\.\d+)\b', text)
                        if not matches:
                            return None

                        return max(matches, key=score_timestamp)

                    def find_first_relevant_timestamp(text):
                        # Prefer first high-precision timestamp (3+ decimals) in text
                        for m in re.findall(r'\b(\d+\.\d+)\b', text):
                            if len(m.split('.')[1]) >= 3:
                                return m
                        # fallback to lower precision (2+ decimals)
                        for m in re.findall(r'\b(\d+\.\d+)\b', text):
                            if len(m.split('.')[1]) >= 2:
                                return m
                        return None

                    def consider_timestamp(candidate):
                        nonlocal timestamp
                        if not candidate:
                            return
                        # Prefer the first nearby relevant timestamp (prefer earlier block context)
                        if not timestamp:
                            timestamp = candidate
                            return
                        # if already have a high precision timestamp, keep it
                        if len(timestamp.split('.')[1]) >= 3:
                            return
                        # if candidate is higher precision than existing timestamp, replace
                        if len(candidate.split('.')[1]) > len(timestamp.split('.')[1]):
                            timestamp = candidate
                            return
                        # as fallback, use original scoring strategy
                        if score_timestamp(candidate) > score_timestamp(timestamp):
                            timestamp = candidate

                    same_line_step = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', line,
                                               re.IGNORECASE)
                    if same_line_step:
                        failure_step_id = same_line_step.group(1)
                        action_text = same_line_step.group(2).strip()
                        test_step = action_text  # prefer just action text as details
                        # timestamp candidate in same line: prefer first relevant and then best
                        consider_timestamp(find_first_relevant_timestamp(line) or find_best_timestamp(line))

                    for k in range(i + 1, min(i + 150, len(lines))):
                        next_line = lines[k]

                        # Stop if we hit next test case
                        if re.match(r'^\d+\.\d+(?:\s|$)', next_line) and k > i + 5:
                            break

                        # Look for timestamp (decimal number) - general case
                        consider_timestamp(find_first_relevant_timestamp(next_line) or find_best_timestamp(next_line))

                        # For Failed/Error verdicts, look for failed step details
                        if verdict_type in ["Failed", "Error"] and not failure_step_id:
                            next_line_lower = next_line.lower()

                            # Look for step identifier with action (e.g., "10.1.6.9.4. Await Value Match: Failed")
                            step_match = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', next_line,
                                                   re.IGNORECASE)
                            if step_match:
                                failure_step_id = step_match.group(1)
                                action_text = step_match.group(2).strip()
                                test_step = action_text  # details should not carry step ID prefix
                                # Extract timestamp from the failure step line if present using best candidate selection
                                consider_timestamp(find_best_timestamp(next_line))
                            else:
                                # Look for common failure indicators in other formats
                                if any(keyword in next_line_lower for keyword in
                                       ["condition", "value", "expected", "actual", "mismatch", "not found",
                                        "exception", "error", "failed to", "failed"]):
                                    if not re.match(r'^\d+\.\d+', next_line):
                                        # Extract step number if it starts with numbers
                                        step_num_match = re.match(r'^(\d+(?:\.\d+)*)', next_line.strip())
                                        if step_num_match:
                                            failure_step_id = step_num_match.group(1)
                                            test_step = next_line[:80]

                        # Look for action keywords for passed cases
                        if verdict_type == "Pass":
                            next_line_lower = next_line.lower()

                            if "execute" in next_line_lower:
                                match = re.search(r'execute\s+(\w+)', next_line_lower)
                                if match:
                                    test_step = match.group(1).capitalize()
                            elif "question" in next_line_lower and "text" in next_line_lower:
                                test_step = "Question/Text"
                            elif "await" in next_line_lower or "wait" in next_line_lower:
                                test_step = "Await Value Match"
                            elif "resume" in next_line_lower:
                                test_step = "Resume"
                            elif "set" in next_line_lower:
                                test_step = "Set"
                            elif "tester" in next_line_lower and "confirmed" in next_line_lower:
                                test_step = "Tester Confirmation"

                    # Add test case if we found a timestamp
                    if timestamp:
                        results[current_fixture]["test_cases"].append({
                            "timestamp": timestamp,
                            "verdict": verdict_type,
                            "details": test_step
                        })

        # Calculate totals
        for fixture_name in results:
            # Use the maximum of: parsed test cases OR initial summary count
            # This preserves the fixture summary count if detailed parsing didn't find individual cases
            parsed_count = len(results[fixture_name]["test_cases"])
            initial_count = results[fixture_name].get("total", 0)
            results[fixture_name]["total"] = max(parsed_count, initial_count)

        return results


    if not st.session_state.selected_files:
        st.info("Select files from the sidebar to show dashboard.")
    elif not dashboard_files:
        st.info("No dashboard-friendly files selected. Choose HTML/HTM/XLSX in sidebar for dashboard details.")
    else:
        st.info(
            "Files selected in the sidebar are available here. Choose only the dashboard file you want to inspect in this tab.")
        dashboard_options = ["--Select File--"] + dashboard_files
        if st.session_state.file_dropdown not in dashboard_options:
            st.session_state.file_dropdown = "--Select File--"
        file_dropdown = st.selectbox("Select a dashboard file", dashboard_options, key="file_dropdown")

        if file_dropdown != "--Select File--":
            with st.spinner("Loading dashboard file..."):
                ensure_file_processed(file_dropdown)
            file_entry = get_uploaded_file_entry(file_dropdown)
            file_bytes = file_entry["bytes"] if file_entry else b""
            chart_type = st.radio(
                "Chart type",
                ["Pie Chart", "Bar Chart"],
                index=0,
                horizontal=True,
                key="dashboard_chart_type",
            )
            bar_orientation = "Vertical"
            if chart_type == "Bar Chart":
                bar_orientation = st.radio(
                    "Bar orientation",
                    ["Vertical", "Horizontal"],
                    index=0,
                    horizontal=True,
                    key="dashboard_bar_orientation",
                )
            horizontal_bar = (bar_orientation == "Horizontal")

            if file_dropdown.lower().endswith(".xlsx"):
                data = st.session_state.excel_data_by_file.get(file_dropdown, [])

                if data:
                    columns = ["--Select Column--"] + list(data[0].keys())
                    col = st.selectbox("Column to analyze", columns, index=0)

                    if col != "--Select Column--":
                        counts = get_column_counts(data, col)
                        if counts:
                            fig = plot_pie_chart(counts,
                                                 f"{col} Distribution") if chart_type == "Pie Chart" else plot_bar_chart(
                                counts, f"{col} Distribution", horizontal=horizontal_bar
                            )
                            st.plotly_chart(fig, use_container_width=True)
                        else:
                            st.warning("No data found for the selected column.")
                else:
                    st.warning("No Excel data available for analysis.")

            elif file_dropdown.lower().endswith((".html", ".htm")):
                st.markdown("### 🔐 Login Info")
                login = extract_login_name_from_html(file_bytes)
                st.write("Login Name:", login)

                st.markdown("### 📊 Statistics")
                stats = extract_statistics_from_html(file_bytes)

                if any(v > 0 for v in stats.values()):
                    c1, c2, c3 = st.columns(3)
                    c1.metric("Executed", stats["Executed"])
                    c2.metric("Passed", stats["Passed"])
                    c3.metric("Failed", stats["Failed"])

                    fig = plot_pie_chart(stats, "Statistics") if chart_type == "Pie Chart" else plot_bar_chart(
                        stats, "Statistics", horizontal=horizontal_bar
                    )
                    st.plotly_chart(fig, use_container_width=True)
                else:
                    st.warning("⚠️ Could not extract statistics")

                st.markdown("### 📋 Test Results")
                grouped_results = extract_test_results_grouped_from_html(file_bytes)

                if grouped_results:
                    st.markdown("#### 👉 Executed Test Cases Summary")
                    all_fixtures = sorted(grouped_results.keys())
                    max_cols = 5
                    num_cols = min(len(all_fixtures), max_cols)

                    if num_cols > 0:
                        cols = st.columns(num_cols)
                        for idx, fixture_name in enumerate(all_fixtures[:max_cols]):
                            if fixture_name in grouped_results:
                                total_cases = grouped_results[fixture_name].get("total", 0)
                                pass_count = grouped_results[fixture_name].get("pass", 0)
                                with cols[idx % num_cols]:
                                    st.metric(
                                        label=fixture_name,
                                        value=total_cases,
                                        delta=f"{pass_count} Passed" if total_cases > 0 else None,
                                        delta_color="inverse"
                                    )

                        if len(all_fixtures) > max_cols:
                            st.markdown(f"**And {len(all_fixtures) - max_cols} more fixtures...**")

                    st.markdown("---")
                    st.markdown("#### Full Test Fixtures Overview")

                    fixture_data = []
                    for fixture_name in sorted(grouped_results.keys()):
                        data = grouped_results[fixture_name]
                        fixture_data.append({
                            "Test Fixture": fixture_name,
                            "✅ Pass": data.get("pass", 0),
                            "❌ Fail": data.get("fail", 0),
                            "⚠️ Error": data.get("error", 0),
                            "⏭️ Not Executed": data.get("not executed", 0),
                            "❓ Inconclusive": data.get("inconclusive", 0),
                            "📊 Total": data.get("total", 0)
                        })

                    df = pd.DataFrame(fixture_data)


                    def style_fixture_table(row):
                        styles = [''] * len(row)
                        if '✅ Pass' in df.columns:
                            pass_col_idx = df.columns.get_loc('✅ Pass')
                            if row['✅ Pass'] > 0:
                                styles[pass_col_idx] = 'background-color: #90EE90; color: black; font-weight: bold;'
                        if '❌ Fail' in df.columns:
                            fail_col_idx = df.columns.get_loc('❌ Fail')
                            if row['❌ Fail'] > 0:
                                styles[fail_col_idx] = 'background-color: #FF6B6B; color: white; font-weight: bold;'
                        if '⚠️ Error' in df.columns:
                            error_col_idx = df.columns.get_loc('⚠️ Error')
                            if row['⚠️ Error'] > 0:
                                styles[error_col_idx] = 'background-color: #FF8C8C; color: white; font-weight: bold;'
                        if '⏭️ Not Executed' in df.columns:
                            notexec_col_idx = df.columns.get_loc('⏭️ Not Executed')
                            if row['⏭️ Not Executed'] > 0:
                                styles[notexec_col_idx] = 'background-color: #FFD700; color: black;'
                        if '❓ Inconclusive' in df.columns:
                            inconc_col_idx = df.columns.get_loc('❓ Inconclusive')
                            if row['❓ Inconclusive'] > 0:
                                styles[inconc_col_idx] = 'background-color: #FFA500; color: black;'
                        return styles


                    styled_fixture_df = df.style.apply(style_fixture_table, axis=1)
                    st.dataframe(styled_fixture_df, use_container_width=True)

                    st.markdown("#### View Test Cases by Fixture")
                    fixture_options = sorted(grouped_results.keys())
                    selected_fixture = st.selectbox("Select Test Fixture:", fixture_options, key="fixture_select")

                    if selected_fixture:
                        fixture_info = grouped_results[selected_fixture]
                        st.markdown(f"### 📋 Test Fixture: **{selected_fixture}**")

                        col1, col2, col3, col4, col5 = st.columns(5)
                        col1.metric("✅ Pass", fixture_info.get("pass", 0))
                        col2.metric("❌ Fail", fixture_info.get("fail", 0))
                        col3.metric("⚠️ Error", fixture_info.get("error", 0))
                        col4.metric("⏭️ Not Executed", fixture_info.get("not executed", 0))
                        col5.metric("❓ Inconclusive", fixture_info.get("inconclusive", 0))

                        mode = st.radio("Show Test Cases", ["All", "Passed only", "Failed/Error only"], index=2,
                                        key="test_case_mode")

                        if fixture_info["test_cases"]:
                            test_cases_df = pd.DataFrame()
                            if mode == "All":
                                test_cases_to_show = fixture_info["test_cases"]
                                heading = "Test Cases (All)"
                            elif mode == "Passed only":
                                test_cases_to_show = [
                                    tc for tc in fixture_info["test_cases"]
                                    if tc.get("verdict", "").lower() in ["pass", "passed"]
                                ]
                                heading = "Test Cases (Passed Only)"
                            else:
                                test_cases_to_show = [
                                    tc for tc in fixture_info["test_cases"]
                                    if tc.get("verdict", "").lower() in ["failed", "fail", "error"]
                                ]
                                heading = "Test Cases (Failed/Error Only)"

                            st.markdown(f"#### {heading}")
                            if not test_cases_to_show:
                                st.info("No test cases to display for selected filter.")
                            else:
                                test_cases_df = pd.DataFrame(test_cases_to_show)


                            def color_verdict(val):
                                if val == "Pass":
                                    return "background-color: #90EE90; color: black; font-weight: bold;"
                                if val in ["Fail", "Failed", "Error"]:
                                    return "background-color: #FF6B6B; color: white; font-weight: bold;"
                                if val == "Not Executed":
                                    return "background-color: #FFD700; color: black;"
                                if val == "Inconclusive":
                                    return "background-color: #FFA500; color: black;"
                                return ""


                            if not test_cases_df.empty and "verdict" in test_cases_df.columns:
                                styled_df = test_cases_df.style.map(
                                    lambda x: color_verdict(x) if isinstance(x, str) else "",
                                    subset=["verdict"]
                                )
                                st.dataframe(styled_df, use_container_width=True, hide_index=True)
                            elif not test_cases_df.empty:
                                st.dataframe(test_cases_df, use_container_width=True, hide_index=True)

                        verdict_counts = {
                            "Pass": fixture_info.get("pass", 0),
                            "Fail": fixture_info.get("fail", 0),
                            "Error": fixture_info.get("error", 0),
                            "Not Executed": fixture_info.get("not executed", 0),
                            "Inconclusive": fixture_info.get("inconclusive", 0)
                        }
                        verdict_counts = {k: v for k, v in verdict_counts.items() if v > 0}

                        if verdict_counts:
                            fig = plot_pie_chart(verdict_counts,
                                                 f"Verdict Distribution - {selected_fixture}") if chart_type == "Pie Chart" else plot_bar_chart(
                                verdict_counts, f"Verdict Distribution - {selected_fixture}", horizontal=horizontal_bar
                            )
                            st.plotly_chart(fig, use_container_width=True)
                else:
                    st.warning("No structured test results found.")

    st.markdown('</div>', unsafe_allow_html=True)

# -------------------------------
# TAB 3: COMPARE
# -------------------------------
# Compare tab:
# Lets users pick 2+ selected files, generate inline word-level differences,
# and download the comparison results as an Excel file.
if active_main_tab == "📂 Compare":
    st.markdown('<div id="compare-section">', unsafe_allow_html=True)
    compare_header_col, compare_reset_col = st.columns([8, 1])
    with compare_header_col:
        st.subheader("Compare Files")
    with compare_reset_col:
        if st.button("🧼 Reset", key="reset_compare_selection"):
            st.session_state.compare_file_selection = []
            st.session_state.compare_result_html = None
            st.session_state.compare_result_excel_bytes = None
            st.session_state.compare_result_files = []
            st.session_state.compare_semantic_summary = None
            st.rerun()

    st.info("Select files in the sidebar to make them available here, then choose only the files you want to compare in this tab.")
    show_current_sidebar_selection()
    show_help_popup('compare', st.session_state.selected_files)
    render_file_context_card("Compare File Context", st.session_state.selected_files, st.session_state.compare_file_selection)

    st.markdown("**Comparison options:**")
    st.markdown(
        "- Inline word-level diff (sequence-aware highlighting)\n"
        "- Side-by-side line diff for direct file comparison\n"
        "- Word presence summary across selected files\n"
        "- Downloadable Excel comparison workbook\n"
        "- Supports PDF, DOCX, PPTX, XLSX, HTML, TXT, CAN/CAPL formats"
    )

    # Use selected files in multiselect (user must choose from selected_files independently)
    st.session_state.compare_file_selection = [
        file_name for file_name in st.session_state.compare_file_selection
        if file_name in st.session_state.selected_files
    ]
    selected_files_for_comparison = st.multiselect(
        "Choose files to compare",
        options=st.session_state.selected_files,
        default=st.session_state.compare_file_selection,
        key="compare_file_selection"
    )

    compare_mode = st.selectbox(
        "Comparison mode",
        ["Exact inline word diff", "Side-by-side line diff", "Word presence summary"],
        index=0,
        key="compare_mode"
    )

    reference_file = None
    # Reference file selection is no longer shown; the first selected file is used automatically for comparison baselines.
    if len(selected_files_for_comparison) == 1:
        st.info("Select at least two files to compare.")

    compare_clicked = st.button("Compare Selected Files", key="run_compare_button", use_container_width=True)

    if compare_clicked:
        if len(selected_files_for_comparison) >= 2:
            with st.spinner("Loading files for comparison..."):
                ensure_files_processed(selected_files_for_comparison)

            selected_texts = {}
            for f in selected_files_for_comparison:
                raw_text = st.session_state.file_texts.get(f, "")
                selected_texts[f] = raw_text if isinstance(raw_text, str) else str(raw_text)

            with st.spinner("Loading comparison results..."):
                html_diff = highlight_multi_file_differences(
                    selected_texts,
                    comparison_mode=compare_mode,
                    reference_file=reference_file
                )
                excel_io = generate_word_level_comparison_excel(selected_texts)
                semantic_summary = build_semantic_diff_explanation(selected_texts)

            st.session_state.compare_result_html = html_diff
            st.session_state.compare_result_excel_bytes = excel_io.getvalue()
            st.session_state.compare_result_files = selected_files_for_comparison.copy()
            st.session_state.compare_semantic_summary = semantic_summary
            record_workspace_memory_event(
                "compare",
                "Semantic comparison completed",
                semantic_summary,
                source=", ".join(selected_files_for_comparison),
            )
            save_workspace_memory()
            save_memory_log("compare", "Semantic comparison stored in workspace memory.", {
                "files": selected_files_for_comparison,
            })
        else:
            st.warning("Select at least two files to compare.")

    if st.session_state.compare_result_html and st.session_state.compare_result_files:
        st.info("Compared files: " + ", ".join(st.session_state.compare_result_files))
        st.markdown(f"### Comparison Results ({len(st.session_state.compare_result_files)} files)")
        if st.session_state.get("compare_semantic_summary"):
            st.markdown(st.session_state.compare_semantic_summary, unsafe_allow_html=True)
        render_html_frame(st.session_state.compare_result_html, height=800)

        st.markdown("### Download Excel Comparison")
        st.download_button(
            "Download Comparison Excel",
            st.session_state.compare_result_excel_bytes,
            file_name="comparison.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    elif len(selected_files_for_comparison) < 2:
        st.info("Select at least two files to compare.")

    st.markdown('</div>', unsafe_allow_html=True)

# -------------------------------
# TAB 4: CAPL
# -------------------------------
# CAPL tab:
# Dedicated to CAPL file selection, live editing, compile/analyze checks, issue
# reporting, and optional AI-assisted fix generation for CAPL scripts.
if active_main_tab == "📡 CAPL":
    st.markdown('<div id="capl-section">', unsafe_allow_html=True)
    capl_header_col, capl_reset_col = st.columns([8, 1])
    with capl_header_col:
        st.subheader("⚙️ CAPL Compiler & Analyzer")
    with capl_reset_col:
        if st.button("🧼 Reset", key="reset_capl_selection"):
            st.session_state.selected_capl_file = "--Select CAPL file--"
            st.session_state.capl_last_analyzed_file = None
            st.session_state.capl_last_issues = None
            st.rerun()

    st.info(
        "Use sidebar selection to make CAPL source files available here. Then choose the CAPL file you want only in this tab.")
    show_current_sidebar_selection()
    show_help_popup('capl', [f for f in st.session_state.selected_files if f.lower().endswith((".can", ".txt"))])

    st.markdown("### Autonomous CAPL Agent Workspace")
    st.info("Use the CAPL agent system to run goal-driven workflows across uploaded documents, shared memory, and CAPL analysis.")
    agent_cols = st.columns(5)
    agent_defs = [
        ("Planning", "Breaks goals into steps"),
        ("Retrieval", "Queries FAISS memory"),
        ("Execution", "Runs tools/actions"),
        ("Reasoning", "Builds insights"),
        ("Coordination", "Merges final output"),
    ]
    for index, (agent_name, agent_role) in enumerate(agent_defs):
        with agent_cols[index]:
            st.markdown(f"**{agent_name} Agent**")
            st.caption(agent_role)
    st.text_input("Autonomous CAPL Goal", key="capl_autonomous_goal", placeholder="Analyze uploaded documents and give insights about risks")
    if st.button("🚀 Run Autonomous CAPL Agents", key="run_capl_agents"):
        with st.spinner("Running autonomous CAPL agents..."):
            st.session_state.capl_agent_result = run_capl_agent(
                st.session_state.capl_autonomous_goal,
                st.session_state.selected_files
            )
        st.success("Autonomous CAPL task completed.")

    if st.session_state.capl_agent_result:
        st.markdown("### Agent Output")
        st.markdown(st.session_state.capl_agent_result, unsafe_allow_html=True)

    if st.session_state.agent_run_history:
        with st.expander("Autonomous CAPL run history", expanded=False):
            for run in st.session_state.agent_run_history[-5:][::-1]:
                st.markdown(f"**{run['timestamp']}** — {html.escape(run['goal'])}")
                if run.get('plan'):
                    st.markdown(f"- Plan: {html.escape(', '.join(run['plan']))}")
                st.markdown("---")

    # Filter selected files for CAPL analysis
    capl_selectable_files = [f for f in st.session_state.selected_files if f.lower().endswith((".can", ".txt"))]
    active_capl_files = [] if st.session_state.selected_capl_file == "--Select CAPL file--" else [st.session_state.selected_capl_file]
    render_file_context_card("CAPL File Context", capl_selectable_files, active_capl_files)

    with st.expander("✍️ Create New CAPL Script", expanded=False):
        st.text_input("CAPL file name", key="capl_editor_name", help="Enter a file name ending with .can")
        st.text_area("Write CAPL code", key="capl_editor_code", height=260)
        live_editor_issues = analyze_capl_code_with_suggestions(st.session_state.capl_editor_code)

        st.markdown("### Live CAPL Preview")
        st.markdown(
            render_capl_code_with_highlights(st.session_state.capl_editor_code, live_editor_issues),
            unsafe_allow_html=True
        )

        st.markdown("### Live Issues & Suggestions")
        render_capl_issue_table(live_editor_issues)

        editor_ai_cols = st.columns([1, 4])
        with editor_ai_cols[0]:
            editor_ai_fix_clicked = st.button("🤖 Suggest Fix", key="capl_editor_ai_fix_button")

        if editor_ai_fix_clicked:
            llm = load_llm()
            editor_prompt = f"""
            You are a CAPL expert. Here is some CAPL code with errors. Please provide the corrected version of the code that fixes all syntax and logical errors. Only output the corrected CAPL code, nothing else.

            Code:
            {st.session_state.capl_editor_code}
            """
            if llm is None:
                st.error("AI fix feature is unavailable because model backend could not be initialized.")
            else:
                with st.spinner("Generating CAPL fix suggestion..."):
                    try:
                        st.session_state.capl_editor_ai_fix = llm.invoke(editor_prompt).strip()
                    except Exception as exc:
                        st.error(f"AI suggestion failed: {exc}")
                        st.session_state.capl_editor_ai_fix = ""

        if st.session_state.capl_editor_ai_fix:
            st.markdown("### Suggested Corrected CAPL Code")
            fixed_editor_issues = analyze_capl_code_with_suggestions(st.session_state.capl_editor_ai_fix)
            st.markdown(
                render_capl_code_with_highlights(st.session_state.capl_editor_ai_fix, fixed_editor_issues),
                unsafe_allow_html=True
            )
            if st.button("Use Suggested Fix", key="capl_editor_use_ai_fix"):
                st.session_state.capl_editor_code = st.session_state.capl_editor_ai_fix
                st.rerun()

        if st.button("💾 Save New CAPL Script"):
            new_file_name = st.session_state.capl_editor_name.strip()
            if not new_file_name:
                st.warning("Enter a file name for the CAPL script.")
            else:
                if not new_file_name.lower().endswith(".can"):
                    new_file_name += ".can"

                file_bytes = st.session_state.capl_editor_code.encode("utf-8")
                existing_index = next(
                    (idx for idx, file_info in enumerate(st.session_state.uploaded_files) if
                     file_info["name"] == new_file_name),
                    None
                )

                if existing_index is None:
                    st.session_state.uploaded_files.append({"name": new_file_name, "bytes": file_bytes})
                else:
                    st.session_state.uploaded_files[existing_index] = {"name": new_file_name, "bytes": file_bytes}

                st.session_state.file_texts[new_file_name] = st.session_state.capl_editor_code
                if new_file_name not in st.session_state.selected_files:
                    st.session_state.selected_files.append(new_file_name)
                st.session_state.capl_last_analyzed_file = None
                st.session_state.capl_last_issues = None
                st.session_state.capl_editor_ai_fix = ""
                st.success(f"Saved {new_file_name} and added it to the selected files.")

    use_all_txt = st.checkbox("Include all .txt files as CAPL")
    with st.spinner("Loading CAPL files..."):
        ensure_files_processed(capl_selectable_files)

    if use_all_txt:
        capl_files = [
            f for f in capl_selectable_files
            if f.lower().endswith((".can", ".txt"))
        ]
    else:
        capl_files = [
            f for f in capl_selectable_files
            if f.lower().endswith((".can", ".txt")) and
               is_capl_code(st.session_state.file_texts.get(f, ""))
        ]

    if not capl_files:
        st.warning("Upload/select CAPL (.can/.capl) files")
    else:
        capl_options = ["--Select CAPL file--"] + capl_files
        if st.session_state.selected_capl_file not in capl_options:
            st.session_state.selected_capl_file = "--Select CAPL file--"
        selected_capl = st.selectbox("Select CAPL file", capl_options, key="selected_capl_file")
        capl_selected = selected_capl != "--Select CAPL file--"
        action_cols = st.columns([1, 1, 4])
        with action_cols[0]:
            analyze_clicked = st.button("🚀 Compile & Analyze", disabled=not capl_selected)
        with action_cols[1]:
            clear_analysis = st.button("🧹 Clear Analysis", disabled=not capl_selected)

        if not capl_selected:
            st.info("Choose a CAPL file to view and analyze.")
        else:
            if clear_analysis:
                st.session_state.capl_last_analyzed_file = None
                st.session_state.capl_last_issues = None

            code = st.session_state.file_texts[selected_capl]
            issues = []

            if analyze_clicked:
                issues = analyze_capl_code_with_suggestions(code)
                st.session_state.capl_last_analyzed_file = selected_capl
                st.session_state.capl_last_issues = issues

            if st.session_state.capl_last_analyzed_file == selected_capl and st.session_state.capl_last_issues is not None:
                issues = st.session_state.capl_last_issues

            st.markdown("### 📄 CAPL Code")
            st.markdown(render_capl_code_with_highlights(code, issues), unsafe_allow_html=True)

            if st.session_state.capl_last_analyzed_file == selected_capl and st.session_state.capl_last_issues is not None:
                st.markdown("### 🛠️ Issues Found")
                render_capl_issue_table(issues)

                # 🤖 AI Suggestions
                if st.checkbox("🤖 AI Suggestions"):
                    llm = load_llm()
                    if llm is None:
                        st.error("AI fix feature is unavailable because model backend could not be initialized.")
                    else:
                        prompt = f"""
                        You are a CAPL expert. Here is some CAPL code with errors. Please provide the corrected version of the code that fixes all syntax and logical errors. Only output the corrected CAPL code, nothing else.

                        Code:
                        {code}
                        """
                        with st.spinner("Analyzing and fixing with AI..."):
                            try:
                                response = llm.invoke(prompt)
                                corrected_code = response.strip()
                                # Update the code in session state
                                st.session_state.file_texts[selected_capl] = corrected_code
                                code = corrected_code
                                issues = analyze_capl_code_with_suggestions(code)
                                st.session_state.capl_last_analyzed_file = selected_capl
                                st.session_state.capl_last_issues = issues
                                st.success("✅ Code corrected by AI!")
                                st.markdown("### 📄 Corrected CAPL Code")
                                st.markdown(render_capl_code_with_highlights(code, issues), unsafe_allow_html=True)
                                if issues:
                                    st.warning("⚠️ Some issues remain:")
                                render_capl_issue_table(issues)
                            except Exception as exc:
                                st.error(f"AI suggestion failed: {exc}")

    st.markdown('</div>', unsafe_allow_html=True)
