# ==============================
# BACKEND FUNCTIONS AND SHARED ENGINES
# Extracted from legacy_app.py. This module keeps Streamlit cache decorators
# and session_state-aware helpers because the original business logic depends
# on them. Tab bodies are moved to tab_*.py.
# ==============================
import html, re, hashlib, os, json, base64, pickle, zipfile, sqlite3
import importlib
import math
import random
import uuid
import urllib.parse
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timedelta
from difflib import SequenceMatcher
from io import BytesIO
from pytz import timezone
import time
from collections import OrderedDict
import docx, openpyxl, pdfplumber, streamlit as st
import streamlit.components.v1 as components
from docx.text.paragraph import Paragraph
from docx.table import Table
import pandas as pd
from openpyxl.styles import PatternFill
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image, ImageDraw, ImageFont
import plotly.express as px
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import HuggingFacePipeline
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ==============================
# GLOBAL CONSTANTS AND RUNTIME STORES
# Originally declared near the top of legacy_app.py.
# ==============================
class CacheManager:
    """LRU cache manager for expensive operations with TTL support"""
    def __init__(self, max_size=50):
        self.cache = OrderedDict()
        self.max_size = max_size
        self.timestamps = {}
        
    def get(self, key, ttl_seconds=3600):
        if key not in self.cache:
            return None
        if key in self.timestamps:
            age = time.time() - self.timestamps[key]
            if age > ttl_seconds:
                del self.cache[key]
                del self.timestamps[key]
                return None
        # Move to end (most recently used)
        self.cache.move_to_end(key)
        return self.cache[key]
    
    def set(self, key, value):
        if len(self.cache) >= self.max_size:
            oldest = next(iter(self.cache))
            del self.cache[oldest]
            if oldest in self.timestamps:
                del self.timestamps[oldest]
        self.cache[key] = value
        self.timestamps[key] = time.time()
        if len(self.cache) > 1:
            self.cache.move_to_end(key)
    
    def clear(self):
        self.cache.clear()
        self.timestamps.clear()


PREVIEW_TOKENS = {}  # token -> {'file_name': str, 'timestamp': datetime}
PREVIEW_STORE = {}   # token -> file_dict
APP_DIR = os.path.dirname(os.path.abspath(__file__))
PREVIEW_DATA_FILE = os.path.join(APP_DIR, "preview_data.pkl")
WORKSPACE_DB_FILE = os.path.join(APP_DIR, "workspace_memory.db")
WORKSPACE_MEMORY_KEY = "workspace_memory"
PDF_PREVIEW_RESOLUTION = 100
PDF_PREVIEW_WINDOW = 25
PDF_ASSET_SCAN_PAGE_LIMIT = 10
MAX_VECTOR_TEXT_CHARS = 250000
TEXT_EXTRACTION_SCHEMA_VERSION = "extract-v3"
CHATPDF_SCHEMA_VERSION = "chatpdf-schema-v3"
FILE_BRAIN_SCHEMA_VERSION = "file-brain-v1"
FILE_BRAIN_PREVIEW_CHARS = 700
FILE_BRAIN_PAGE_CONTEXT_CHARS = 3500
FILE_BRAIN_SELECTED_CHUNK_CHARS = 1400
FILE_BRAIN_MAX_TABLE_ROWS = 2000
FILE_TEXT_CACHE = CacheManager(max_size=100)
VECTOR_STORE_CACHE = CacheManager(max_size=20)
EXCEL_DATA_CACHE = CacheManager(max_size=50)
EMBEDDINGS_CACHE = CacheManager(max_size=200)
FILE_HASH_CACHE = {}
WORKSPACE_DB_FILE = os.path.join(APP_DIR, "workspace_memory.db")
WORKSPACE_MEMORY_KEY = "workspace_memory"

SUMMARY_STOPWORDS = {
    "the", "and", "for", "with", "that", "this", "from", "are", "was", "were", "into", "your", "have",
    "has", "had", "not", "but", "you", "all", "can", "will", "use", "using", "used", "how", "what",
    "when", "where", "which", "while", "into", "more", "than", "their", "there", "about", "after",
    "before", "within", "without", "each", "page", "pages", "table", "tables", "image", "images",
    "document", "content", "metadata", "information", "product", "file", "text"
}

DOCUMENT_INTELLIGENCE_BROAD_INTENTS = {
    "analysis_request",
    "summarization_request",
    "overview_request",
    "technical_overview",
    "themes_request",
}

DOCUMENT_INTELLIGENCE_DOMAIN_TERMS = {
    "automotive": {"can", "capl", "canoe", "canalyzer", "ecu", "vehicle", "lin", "flexray"},
    "software": {"api", "code", "function", "class", "module", "database", "server", "client"},
    "networking": {"ethernet", "tcp", "ip", "network", "gateway", "protocol", "bus"},
    "electronics": {"pin", "connector", "signal", "voltage", "current", "channel", "wiring"},
    "business": {"revenue", "customer", "market", "sales", "cost", "risk", "strategy"},
    "testing": {"test", "validation", "requirement", "result", "pass", "fail", "coverage"},
}

DOCUMENT_INTELLIGENCE_PROMPT_BY_INTENT = {
    "factual_query": "Answer the question directly, then briefly explain the supporting evidence in plain language.",
    "analysis_request": "Write an expert analysis with deeper reasoning, architecture insights, design patterns, strengths, risks, implications, and takeaways.",
    "summarization_request": "Write a concise executive summary that explains the document naturally.",
    "overview_request": "Give a high-level explanation suitable for quick understanding by a new reader.",
    "comparison_request": "Compare only supported items or files. Explain what the differences mean, not just what words differ.",
    "followup_question": "Use conversation memory to resolve references, then answer naturally from the document context.",
    "metadata_request": "Answer only the user-visible facts requested. Do not expose internal metadata objects or extraction details.",
    "technical_overview": "Explain the technical structure, components, interfaces, workflow, constraints, and engineering implications.",
    "themes_request": "Explain recurring themes and how they connect across the document.",
}

NATURAL_DOCUMENT_RESPONSE_PROMPT = """Your job is to produce natural, human-like, insightful document answers.

Do not expose internal extraction systems or diagnostics. Never show OCR artifacts, extracted heading dumps,
keyword lists, semantic-signal language, metadata objects, retrieval mechanics, chunking details, pipeline names,
raw extraction artifacts, page formatting noise, repeated labels, low-level parsing details, or raw context blocks.
The user cares about the document, not the machinery behind the answer.
Never turn extracted entities, figure labels, module labels, table fields, captions, or repeated section titles
into insights. Avoid claims that a term is central unless the surrounding document meaning clearly supports that conclusion.

Use retrieved content only as supporting evidence for reasoning and synthesis. Explain the document clearly,
connect ideas across sections, and write like an expert analyst.

Before answering, reason about the document's real meaning:
1. What is it fundamentally about?
2. What problem does it solve?
3. Which systems or components are central?
4. What technical concepts, workflows, applications, architecture, and takeaways matter most?

Prioritize document purpose and semantic meaning over keyword frequency. Do not assume repeated table terms
such as temperature range, power consumption, weight, voltage, dimensions, min/max values, or other isolated
specifications are the document's main themes unless the document is actually about those topics.

Intent behavior:
- Summarize: provide a concise executive summary with purpose, key systems/components, major capabilities, and practical significance.
- Analyze: provide deeper technical reasoning, architecture insights, design philosophy, strengths, limitations, relationships, engineering implications, and practical applications.
- Overview: provide a broad, accessible explanation of what the system/document is, how it is organized, what it covers, and who it is for.
- Factual questions: answer directly and precisely from the document.

Never concatenate headings into a summary. Never make keyword lists the primary answer. Never dump raw extracted
text unless the user explicitly asks for exact text, search matches, or table rows. If evidence is partial, state
the limitation briefly and still provide the best grounded synthesis."""

LOW_SIGNAL_THEME_TERMS = {
    "temperature", "range", "temperature range", "power", "power consumption", "consumption",
    "weight", "approx", "approximately", "voltage", "current", "dimension", "dimensions",
    "length", "width", "height", "minimum", "maximum", "min", "max", "typical", "value",
    "values", "unit", "units", "spec", "specs", "specification", "specifications", "table",
    "row", "column", "page", "pages", "number", "numbers", "total", "text", "content",
    "figure", "figures", "caption", "captions", "label", "labels", "field", "fields",
    "module", "modules", "channel", "channels", "technical", "data", "technical data",
    "section", "sections", "heading", "headings", "title", "titles", "diagram", "diagrams",
    "signal", "signals", "connector", "connectors", "pin", "pins", "system", "systems",
    "product", "products", "information", "overview",
}


def normalize_document_query(text):
    return re.sub(r"\s+", " ", str(text or "").strip())


def classify_query_intent(query, previous_messages=None):
    """Classify user intent for document-intelligence routing."""
    del previous_messages
    q = normalize_document_query(query).lower()
    compact = re.sub(r"[^a-z0-9]+", " ", q).strip()

    if compact in {"analyze", "analyse", "analysis"} or any(
        term in q for term in ["deep analysis", "full analysis", "analyze document", "analyse document", "key insights"]
    ):
        return "analysis_request"
    if compact in {"summary", "summarize", "summarise"} or any(
        term in q for term in ["summarize", "summarise", "short summary", "main points", "recap"]
    ):
        return "summarization_request"
    if compact == "overview" or any(term in q for term in ["give overview", "high level overview", "what is this document about"]):
        return "overview_request"
    if any(term in q for term in ["technical overview", "explain the architecture", "architecture", "system design"]):
        return "technical_overview"
    if any(term in q for term in ["main themes", "themes", "topics", "key topics"]):
        return "themes_request"
    if any(term in q for term in ["compare", "difference", "differences", " versus ", " vs "]):
        return "comparison_request"
    if any(term in q for term in ["metadata", "author", "created", "file type", "document type"]):
        return "metadata_request"
    if any(term in q for term in ["previous", "that", "it", "those", "follow up", "again", "same"]):
        return "followup_question"
    return "factual_query"


def to_technical_intent(document_intent):
    """Map document-intelligence intents onto existing response intent names."""
    return {
        "analysis_request": "FULL_DOCUMENT_ANALYSIS",
        "summarization_request": "SHORT_SUMMARY",
        "overview_request": "OVERVIEW",
        "technical_overview": "OVERVIEW",
        "themes_request": "OVERVIEW",
        "comparison_request": "COMPARISON",
        "metadata_request": "OVERVIEW",
        "followup_question": "QUESTION_ANSWERING",
        "factual_query": "QUESTION_ANSWERING",
    }.get(document_intent, "QUESTION_ANSWERING")


def requires_document_scope(document_intent):
    return document_intent in DOCUMENT_INTELLIGENCE_BROAD_INTENTS or document_intent == "metadata_request"


def document_intelligence_tokenize(text):
    return [
        token
        for token in re.findall(r"[A-Za-z][A-Za-z0-9_-]{2,}", str(text or "").lower())
        if token not in SUMMARY_STOPWORDS
    ]


def document_intelligence_clean_line(text):
    return re.sub(r"\s+", " ", str(text or "")).strip()


def document_intelligence_meaningful_sentences(text, limit=8):
    sentences = []
    for sentence in re.split(r"(?<=[.!?])\s+|\n+", str(text or "")):
        sentence = document_intelligence_clean_line(sentence)
        if 45 <= len(sentence) <= 320:
            sentences.append(sentence)
        if len(sentences) >= limit:
            break
    return sentences


def summarize_semantic_page(page_record, max_sentences=3):
    text = page_record.get("text", "")
    sentences = document_intelligence_meaningful_sentences(text, limit=max_sentences)
    summary = " ".join(sentences) if sentences else document_intelligence_clean_line(text)[:420]
    return {
        "page": page_record.get("page"),
        "section": page_record.get("section") or f"Page {page_record.get('page')}",
        "summary": summary,
        "keywords": [word for word, _ in Counter(document_intelligence_tokenize(text)).most_common(12)],
    }


def extract_semantic_topics(pages, tables=None, limit=14):
    counter = Counter()
    for page in pages or []:
        counter.update(document_intelligence_tokenize(page.get("text", "")))
    for table in tables or []:
        counter.update(document_intelligence_tokenize(" ".join(str(h) for h in table.get("headers", []))))
    return [word for word, _ in counter.most_common(limit)]


def detect_semantic_domains(topics, pages):
    topic_set = set(topics or [])
    combined_sample = " ".join(str(page.get("text", ""))[:1200].lower() for page in (pages or [])[:12])
    domains = []
    for domain, terms in DOCUMENT_INTELLIGENCE_DOMAIN_TERMS.items():
        score = len(topic_set.intersection(terms)) + sum(1 for term in terms if term in combined_sample)
        if score:
            domains.append({"domain": domain, "score": score})
    domains.sort(key=lambda item: item["score"], reverse=True)
    return domains[:5]


def build_semantic_master_summary(section_summaries, topics, tables=None, diagrams=None):
    useful = [item.get("summary", "") for item in section_summaries or [] if item.get("summary")][:6]
    overview = " ".join(useful) if useful else "No meaningful document narrative was extracted."
    supporting_assets = []
    if tables:
        supporting_assets.append("structured data")
    if diagrams:
        supporting_assets.append("visual references")
    support_sentence = ""
    if supporting_assets:
        support_sentence = " It also contains " + " and ".join(supporting_assets) + " that can support more detailed questions."
    return (
        f"{overview[:1800]}"
        f"{support_sentence}"
    )


def build_semantic_suggested_questions(topics, domains, has_tables=False, has_diagrams=False):
    topic = topics[0] if topics else "this document"
    suggestions = [
        "Summarize the document",
        "Give key insights",
        f"Explain {topic}",
        "Provide a technical overview",
    ]
    if domains:
        suggestions.append(f"Explain the {domains[0]['domain']} context")
    if has_tables:
        suggestions.append("Extract the relevant tables")
    if has_diagrams:
        suggestions.append("Explain the diagrams or figures")
    return suggestions[:6]


def build_semantic_hierarchy(section_summaries):
    return [
        {
            "level": "section",
            "page": item.get("page"),
            "title": item.get("section"),
            "keywords": item.get("keywords", []),
        }
        for item in section_summaries or []
    ]


def extract_semantic_key_concepts(pages, tables=None, limit=12):
    """Extract meaningful concept phrases instead of raw one-word keyword dumps."""
    counter = Counter()
    phrase_pattern = re.compile(
        r"\b(?:[A-Z][A-Za-z0-9_/-]{2,}|[a-z][a-z0-9_/-]{3,})"
        r"(?:\s+(?:[A-Z][A-Za-z0-9_/-]{2,}|[a-z][a-z0-9_/-]{3,})){1,4}\b"
    )
    noise = {
        "table contents", "page text", "document level", "file type", "main features",
        "important notes", "all rights", "copyright notice",
    }
    for page in pages or []:
        text = str(page.get("text", ""))
        for phrase in phrase_pattern.findall(text):
            clean_phrase = document_intelligence_clean_line(phrase).strip(" .,:;")
            lower_phrase = clean_phrase.lower()
            tokens = document_intelligence_tokenize(clean_phrase)
            if len(tokens) < 2 or lower_phrase in noise:
                continue
            if any(token in SUMMARY_STOPWORDS for token in tokens[:1]):
                continue
            counter[clean_phrase[:90]] += 1
    for table in tables or []:
        headers = " ".join(str(header) for header in table.get("headers", []))
        for phrase in phrase_pattern.findall(headers):
            clean_phrase = document_intelligence_clean_line(phrase).strip(" .,:;")
            if len(document_intelligence_tokenize(clean_phrase)) >= 2:
                counter[clean_phrase[:90]] += 2
    return [phrase for phrase, _ in counter.most_common(limit)]


def extract_architecture_components(pages, entities=None, limit=12):
    """Find component/interface/process signals for technical overview and analysis."""
    component_terms = {
        "architecture", "component", "module", "interface", "connector", "service", "engine",
        "pipeline", "workflow", "layer", "database", "server", "client", "api", "gateway",
        "controller", "device", "channel", "signal", "bus", "model", "adapter",
    }
    candidates = Counter()
    for page in pages or []:
        text = str(page.get("text", ""))
        sentences = document_intelligence_meaningful_sentences(text, limit=16)
        for sentence in sentences:
            lower_sentence = sentence.lower()
            if not any(term in lower_sentence for term in component_terms):
                continue
            for token in re.findall(r"\b[A-Za-z][A-Za-z0-9_/-]{2,}\b", sentence):
                clean_token = token.strip(".,:;()[]")
                lower_token = clean_token.lower()
                if lower_token in SUMMARY_STOPWORDS:
                    continue
                if lower_token in component_terms or clean_token.isupper() or re.search(r"[A-Z].*[a-z]|[0-9]", clean_token):
                    candidates[clean_token[:70]] += 1
    for entity in entities or []:
        if re.search(r"[A-Z0-9_-]{3,}", str(entity)):
            candidates[str(entity)[:70]] += 1
    return [name for name, _ in candidates.most_common(limit)]


def extract_semantic_relationships(pages, concepts=None, components=None, limit=10):
    """Capture lightweight semantic relationships from meaningful sentences."""
    relation_terms = [
        "supports", "provides", "connects", "uses", "requires", "includes", "contains",
        "enables", "controls", "communicates", "configured", "consists", "depends",
        "integrates", "receives", "transmits", "generates",
    ]
    relationships = []
    concept_terms = [str(item).lower() for item in (concepts or [])[:20]]
    component_terms = [str(item).lower() for item in (components or [])[:20]]
    for page in pages or []:
        for sentence in document_intelligence_meaningful_sentences(page.get("text", ""), limit=20):
            lower_sentence = sentence.lower()
            if not any(term in lower_sentence for term in relation_terms):
                continue
            if concept_terms and not any(term in lower_sentence for term in concept_terms + component_terms):
                continue
            relationships.append({
                "page": page.get("page"),
                "text": sentence[:320],
            })
            if len(relationships) >= limit:
                return relationships
    return relationships


def build_semantic_metadata(file_name, file_type, pages, tables=None, diagrams=None, entities=None):
    """Create document-level understanding artifacts at upload/file-brain time."""
    section_summaries = [summarize_semantic_page(page) for page in pages or []]
    topics = extract_semantic_topics(pages, tables=tables)
    domains = detect_semantic_domains(topics, pages)
    key_concepts = extract_semantic_key_concepts(pages, tables=tables)
    architecture_components = extract_architecture_components(pages, entities=entities)
    relationships = extract_semantic_relationships(
        pages,
        concepts=key_concepts,
        components=architecture_components,
    )
    executive_summary = build_semantic_master_summary(section_summaries[:5], topics, tables=tables, diagrams=diagrams)
    technical_parts = []
    if architecture_components:
        technical_parts.append("Core components/interfaces: " + ", ".join(architecture_components[:8]) + ".")
    if relationships:
        technical_parts.append("Key relationships: " + " ".join(item["text"] for item in relationships[:3]))
    if not technical_parts:
        technical_parts.append(executive_summary)
    return {
        "metadata": {
            "file_name": file_name,
            "file_type": file_type,
            "page_or_section_count": len(pages or []),
            "table_count": len(tables or []),
            "diagram_count": len(diagrams or []),
        },
        "section_summaries": section_summaries,
        "document_summary": executive_summary,
        "executive_summary": executive_summary,
        "technical_summary": " ".join(technical_parts)[:1800],
        "topics": topics,
        "key_concepts": key_concepts,
        "entities": list(entities or [])[:500],
        "technical_domains": domains,
        "architecture_components": architecture_components,
        "semantic_relationships": relationships,
        "suggested_questions": build_semantic_suggested_questions(topics, domains, bool(tables), bool(diagrams)),
        "hierarchy": build_semantic_hierarchy(section_summaries),
    }


def build_document_intelligence_prompt(query, document_intent, context, memory="", cross_file_hints=""):
    intent_instruction = DOCUMENT_INTELLIGENCE_PROMPT_BY_INTENT.get(
        document_intent,
        DOCUMENT_INTELLIGENCE_PROMPT_BY_INTENT["factual_query"],
    )
    format_rules = {
        "summarization_request": (
            "Output format: Executive Summary, Why It Matters, Key Insights, Key Takeaways. "
            "Keep it concise and synthesized."
        ),
        "overview_request": (
            "Output format: What this is, Who/what it is for, Main idea, What it covers. "
            "Use plain human-friendly language."
        ),
        "analysis_request": (
            "Output format: Executive readout, Architecture/structure, Key concepts and relationships, "
            "Strengths, Limitations or risks, Actionable insights, Takeaways."
        ),
        "technical_overview": (
            "Output format: Technical purpose, Components/interfaces, Workflow/data flow, Tables/diagrams, "
            "Constraints, Engineering takeaways."
        ),
        "themes_request": (
            "Output format: Main themes table, How themes connect, Evidence by section, Takeaways."
        ),
        "comparison_request": (
            "Output format: Comparison table, Similarities, Differences, Best-fit interpretation, Takeaway."
        ),
        "factual_query": "Output format: Direct answer first, then supporting evidence and sources.",
    }.get(document_intent, "Output format: structured professional answer with sources.")
    return f"""You are an advanced document intelligence assistant.

{NATURAL_DOCUMENT_RESPONSE_PROMPT}

Intent: {document_intent}
Task: {intent_instruction}
{format_rules}

Rules:
- Use only the supplied document context and conversation memory.
- Broad requests require synthesis across the available document understanding and representative passages.
- Do not dump metadata, keyword lists, OCR fragments, raw headings, source headers, or internal context labels.
- Do not treat repeated words, table fields, or isolated specifications as primary themes by default.
- For repeated specs such as temperature, weight, power, voltage, dimensions, and min/max fields, use them as supporting details unless the document is truly about those measurements.
- Connect ideas across sections and explain meaning in natural language.
- Prefer concise analyst prose over robotic templates; use headings only when they improve readability.
- Avoid generic rejection when partial evidence exists; provide best-effort, uncertainty-aware analysis.
- Clearly say "Not specified in the provided context" only for missing details.
- Cite relevant pages, sections, sheets, tables, or diagram references.
- End with a confidence label.

Conversation memory:
{memory or "No previous conversation."}

Cross-file hints:
{cross_file_hints or "No cross-file hints."}

Context:
{context}

User question:
{query}
"""


def collect_synthesis_inputs(brains, docs=None):
    """Collect clean semantic inputs for deterministic synthesis fallback."""
    files = []
    for file_name, brain in (brains or {}).items():
        semantic = brain.get("semantic_metadata", {}) or {}
        files.append({
            "file_name": file_name,
            "summary": semantic.get("executive_summary") or semantic.get("document_summary") or "",
            "technical_summary": semantic.get("technical_summary") or "",
            "topics": semantic.get("topics", [])[:10],
            "concepts": semantic.get("key_concepts", [])[:10],
            "domains": [item.get("domain", "") for item in semantic.get("technical_domains", [])[:5] if item.get("domain")],
            "components": semantic.get("architecture_components", [])[:12],
            "relationships": semantic.get("semantic_relationships", [])[:8],
            "sections": semantic.get("section_summaries", [])[:8],
            "tables": len(brain.get("tables", [])),
            "diagrams": len(brain.get("diagrams", [])),
        })

    evidence = []
    for doc in docs or []:
        meta = getattr(doc, "metadata", {}) or {}
        snippet = document_intelligence_clean_line(getattr(doc, "page_content", ""))[:420]
        if snippet:
            evidence.append({
                "source": build_chatpdf_citation_label(meta),
                "section": meta.get("section", ""),
                "snippet": snippet,
            })
    return {"files": files, "evidence": evidence[:8]}


def _is_low_signal_theme(item):
    clean_item = normalize_synthesis_text(item).lower()
    if not clean_item:
        return True
    tokens = set(re.findall(r"[a-z][a-z0-9_-]*", clean_item))
    if clean_item in LOW_SIGNAL_THEME_TERMS:
        return True
    if tokens and tokens.issubset(LOW_SIGNAL_THEME_TERMS):
        return True
    if re.fullmatch(r"[\d\s.,:+\-/%()]+", clean_item):
        return True
    return False


def _synthesis_bullets(items, limit=5, filter_low_signal=False):
    lines = []
    for item in items or []:
        clean_item = normalize_synthesis_text(item)
        if filter_low_signal and _is_low_signal_theme(clean_item):
            continue
        if clean_item and clean_item not in lines:
            lines.append(clean_item)
        if len(lines) >= limit:
            break
    return lines


def _human_join(items, limit=4, filter_low_signal=False):
    clean_items = [normalize_synthesis_text(item) for item in (items or [])]
    if filter_low_signal:
        clean_items = [item for item in clean_items if not _is_low_signal_theme(item)]
    clean_items = [item for item in clean_items if item][:limit]
    if not clean_items:
        return ""
    if len(clean_items) == 1:
        return clean_items[0]
    return ", ".join(clean_items[:-1]) + f", and {clean_items[-1]}"


def _natural_focus_sentence(items, fallback="the available document content"):
    focus = _human_join(items, limit=4, filter_low_signal=True)
    if not focus:
        focus = fallback
    return f"The document is best understood around {focus}, with the useful details organized by purpose, structure, workflow, and practical value."


def _natural_insight_lines(items, limit=5):
    clean_items = _synthesis_bullets(items, limit=limit, filter_low_signal=True)
    if not clean_items:
        return []
    focus = _human_join(clean_items, limit=4, filter_low_signal=True)
    if not focus:
        return []
    lines = [
        f"The useful interpretation is how the relevant system elements, including {focus}, support the document's stated purpose and practical use.",
        "Repeated specification fields should be read as constraints or implementation details, not as the main message.",
        "The strongest answers should explain what the system does, how its parts fit together, and why the document exists.",
    ]
    return lines[:limit]


def _component_context_text(items):
    clean_items = _synthesis_bullets(items, limit=6, filter_low_signal=True)
    if not clean_items:
        return ""
    focus = _human_join(clean_items, limit=4, filter_low_signal=True)
    return (
        f"The available context points to {focus} as relevant system elements. "
        "They should be interpreted through their role in the architecture, workflow, and engineering purpose rather than as a standalone module list."
    )


def normalize_synthesis_text(text):
    """Remove metadata labels so synthesized answers sound natural."""
    clean = document_intelligence_clean_line(text)
    clean = re.sub(r"(?i)\bdocument-level summary:\s*", "", clean)
    clean = re.sub(r"(?i)\bdominant topics:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\bkey topics:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\btopics:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\bentities:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\btechnical domains:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\bfile type:\s*[^.]{0,120}\.?", "", clean)
    clean = re.sub(r"(?i)\bthe document's strongest semantic signals are:\s*[^.]{0,260}\.?", "", clean)
    clean = re.sub(r"(?i)\bit includes \d+ structured table\(s\) and \d+ diagram reference\(s\)\.?", "", clean)
    clean = re.sub(r"(?i)\bstructured tables indexed:\s*\d+\.\s*", "", clean)
    clean = re.sub(r"(?i)\bdiagram references indexed:\s*\d+\.\s*", "", clean)
    clean = re.sub(r"(?i)\btables indexed:\s*\d+\.\s*", "", clean)
    clean = re.sub(r"(?i)\bdiagram references:\s*\d+\.\s*", "", clean)
    return document_intelligence_clean_line(clean)


def synthesize_document_response(query, document_intent, brains, docs=None, sources_text=""):
    """Produce non-repetitive, intent-specific document synthesis without metadata dumping."""
    inputs = collect_synthesis_inputs(brains, docs=docs)
    files = inputs["files"]
    evidence = inputs["evidence"]
    if not files:
        return (
            "Answer:\nI could not build enough document-level context from the selected files. "
            "Try re-uploading the document or ask for a specific page, section, table, or exact term.\n\n"
            f"Sources:\n{sources_text or '- No sources found'}"
        )

    all_summaries = _synthesis_bullets([f["summary"] for f in files], limit=3)
    all_technical = _synthesis_bullets([f["technical_summary"] for f in files], limit=3)
    all_topics = _synthesis_bullets([topic for f in files for topic in f["topics"]], limit=8, filter_low_signal=True)
    all_concepts = _synthesis_bullets([concept for f in files for concept in f["concepts"]], limit=8, filter_low_signal=True)
    all_domains = _synthesis_bullets([domain for f in files for domain in f["domains"]], limit=5)
    all_components = _synthesis_bullets([component for f in files for component in f["components"]], limit=10, filter_low_signal=True)
    relationships = [
        rel.get("text", "")
        for f in files
        for rel in f["relationships"]
        if rel.get("text")
    ]
    section_lines = []
    for f in files:
        for section in f["sections"][:4]:
            summary = normalize_synthesis_text(section.get("summary", ""))
            if summary:
                section_lines.append(f"- {summary[:260]}")
    evidence_lines = [f"- {item['source']}: {item['snippet']}" for item in evidence[:4]]
    sources = sources_text or "\n".join(f"- {f['file_name']}" for f in files)
    focus_items = all_components or all_concepts or all_topics
    focus_sentence = _natural_focus_sentence(focus_items)
    insight_lines = _natural_insight_lines(focus_items, limit=5)
    core_purpose = " ".join(all_summaries)[:1400] or focus_sentence
    technical_readout = " ".join(all_technical or all_summaries)[:1400] or focus_sentence

    if document_intent == "summarization_request":
        response_parts = [
            "Answer:",
            "**Executive Summary**",
            core_purpose,
        ]
        component_context = _component_context_text(all_components)
        if component_context:
            response_parts.append("**System Context**\n" + component_context)
        if insight_lines:
            response_parts.append("**Key Insights**\n" + "\n".join(f"- {item}" for item in insight_lines[:4]))
        response_parts.append(
            "**Key Takeaways**\n"
            + "\n".join(f"- {item}" for item in (
                insight_lines[:3] or ["Use this as a quick orientation before asking for details on a specific area."]
            ))
        )
    elif document_intent == "overview_request":
        response_parts = [
            "Answer:",
            "**What This Document Is About**",
            core_purpose[:1100],
            "**Quick Orientation**",
            "\n".join(f"- {item}" for item in (insight_lines[:4] or [focus_sentence])),
        ]
        if section_lines:
            response_parts.append("**Useful Context**\n" + "\n".join(section_lines[:4]))
    elif document_intent in {"technical_overview", "themes_request"}:
        heading = "**Technical Overview**" if document_intent == "technical_overview" else "**Main Themes**"
        response_parts = [
            "Answer:",
            heading,
            technical_readout,
        ]
        component_context = _component_context_text(all_components)
        if component_context:
            response_parts.append("**System / Architecture Context**\n" + component_context)
        if relationships:
            response_parts.append("**Relationships / Flow**\n" + "\n".join(f"- {item}" for item in relationships[:5]))
        if insight_lines:
            response_parts.append("**What The Pattern Means**\n" + "\n".join(f"- {item}" for item in insight_lines[:5]))
    elif document_intent == "analysis_request":
        response_parts = [
            "Answer:",
            "**Executive Readout**",
            core_purpose[:1200],
            "**Deep Analysis**",
            focus_sentence,
        ]
        component_context = _component_context_text(all_components)
        if component_context:
            response_parts.append("**Architecture / Structure**\n" + component_context)
        if relationships:
            response_parts.append("**Conceptual Relationships**\n" + "\n".join(f"- {item}" for item in relationships[:5]))
        response_parts.append(
            "**Strengths / Useful Insights**\n"
            + "\n".join(f"- {item}" for item in (insight_lines[:4] or [focus_sentence]))
        )
        response_parts.append(
            "**Limitations / Uncertainty**\n"
            "- Details that are only visible inside images or diagrams may be incomplete.\n"
            "- Exact technical values are not inferred unless present in the document context.\n"
            "- Repeated specification fields are treated as supporting detail, not automatically as the document's purpose."
        )
        response_parts.append(
            "**Actionable Takeaways**\n"
            + "\n".join(f"- {item}" for item in (
                insight_lines[:3] or ["Ask for a component, workflow, table, or limitation to go deeper."]
            ))
        )
    elif document_intent == "comparison_request":
        response_parts = [
            "Answer:",
            "**Comparison Readout**",
            "| Area | Evidence / Interpretation |",
            "|---|---|",
        ]
        for f in files:
            response_parts.append(f"| {f['file_name']} | {(f['summary'] or f['technical_summary'])[:260]} |")
        if all_concepts:
            response_parts.append("**Interpretation**\n" + "\n".join(f"- {item}" for item in _natural_insight_lines(all_concepts[:6], limit=6)))
    else:
        response_parts = [
            "Answer:",
            " ".join(all_summaries or all_technical)[:1200] or focus_sentence,
        ]
        if evidence_lines:
            response_parts.append("**What Supports This**\n" + "\n".join(evidence_lines[:2]))

    response_parts.append(f"Sources:\n{sources}")
    return "\n\n".join(part for part in response_parts if str(part or "").strip())


def build_best_effort_response(query, document_intent, brains, sources_text):
    """Graceful fallback for broad questions when retrieval/LLM coverage is thin."""
    return synthesize_document_response(
        query=query,
        document_intent=document_intent,
        brains=brains,
        docs=[],
        sources_text=sources_text,
    )


def response_looks_like_metadata_dump(response):
    """Detect low-quality generated answers that repeat metadata instead of synthesizing."""
    text = str(response or "").lower()
    if not text.strip():
        return True
    dump_markers = [
        "dominant topics:",
        "structured tables indexed:",
        "diagram references indexed:",
        "key topics:",
        "entities:",
        "metadata:",
        "document-level summary:",
        "semantic signals",
        "retrieval_layer",
        "chunk_index",
        "selected_chunk_index",
        "bm25",
        "faiss",
        "keyword list",
    ]
    hard_markers = [
        "semantic signals",
        "retrieval_layer",
        "chunk_index",
        "selected_chunk_index",
        "bm25",
        "faiss",
        "ocr artifacts",
        "metadata object",
    ]
    if any(marker in text for marker in hard_markers):
        return True
    marker_hits = sum(1 for marker in dump_markers if marker in text)
    if marker_hits >= 2:
        return True
    if text.count("not specified in the provided context") >= 5:
        return True
    if len(text) < 120 and any(marker in text for marker in dump_markers):
        return True
    return False

MASTER_SYSTEM_PROMPT = """You are an Enterprise Document Intelligence Agent.

Your role is to analyze, understand, and extract information from uploaded documents with high accuracy and professional structure.

You write like an expert analyst: natural, clear, insightful, and grounded in the selected document.
Never expose OCR artifacts, extracted heading dumps, keyword lists, semantic diagnostics, metadata objects,
retrieval mechanics, chunking details, or internal pipeline information.

---------------------------------------------------------------------

## 1. CORE RESPONSIBILITY

For every user query:

1. Identify the intent
2. Select the correct processing mode
3. Use the appropriate strategy
4. Produce a structured, accurate response

---------------------------------------------------------------------

## 2. INTENT CLASSIFICATION (MANDATORY)

Classify the user request into ONE of these:

- FULL_DOCUMENT_ANALYSIS -> "analyze"
- SHORT_SUMMARY -> "summary"
- OVERVIEW -> "overview"
- QUESTION_ANSWERING -> specific questions
- TABLE_EXTRACTION -> tables/data
- COMPONENT_EXTRACTION -> modules/components
- DIAGRAM_ANALYSIS -> diagrams/flows
- COMPARISON -> comparing items/files
- SEARCH -> find/highlight/count

If unclear -> default to QUESTION_ANSWERING

---------------------------------------------------------------------

## 3. PROCESSING MODES

### A. FULL DOCUMENT ANALYSIS

Use when user says:
"analyze", "full analysis"

Strategy:
- Use the available document understanding across the selected content
- Combine all available context

Output:

### Overview
### Purpose
### Core Concept
### Key Features / Capabilities
### Architecture / Structure
### Components / Modules
### Workflow / Usage
### Use Cases / Applications
### Important Notes / Constraints
### Key Takeaways

---------------------------------------------------------------------

### B. SHORT SUMMARY

Output:
- What the document is about
- Purpose
- 3-5 key insights
- 2-3 key takeaways

---------------------------------------------------------------------

### C. QUESTION ANSWERING

Strategy:
- Use only relevant context
- Be precise

Rules:
- Do NOT hallucinate
- If missing info:
  "Not specified in the provided context"

---------------------------------------------------------------------

### D. TABLE EXTRACTION

Output:
- Clean structured table
- Preserve rows and columns

---------------------------------------------------------------------

### E. COMPONENT EXTRACTION

Output:
For each component:
- Name
- Purpose
- Key features
- Interfaces (if available)

---------------------------------------------------------------------

### F. DIAGRAM ANALYSIS

Output:
- Components in diagram
- Relationships
- Flow (signals/data)
- Purpose

Optional:
- ASCII diagram

---------------------------------------------------------------------

### G. COMPARISON

Output table:

| Criteria | Item 1 | Item 2 | Difference |

---------------------------------------------------------------------

### H. SEARCH / FIND

- Highlight or extract relevant matches
- Be precise

---------------------------------------------------------------------

## 4. CONTEXT HANDLING

Documents may contain noise:
- headers/footers
- table of contents
- page numbers
- OCR fragments

You MUST ignore these.

Focus on:
- meaningful content
- technical explanations
- features
- structure
- usage

---------------------------------------------------------------------

## 5. PARTIAL CONTEXT RULE

If context is incomplete:

- Answer what is available
- For missing parts:
  "Not specified in the provided context"

DO NOT reject entire answer unless NOTHING useful exists.

---------------------------------------------------------------------

## 6. STRICT RULES

- Do NOT hallucinate
- Do NOT invent specs or values
- Do NOT repeat raw text
- Do NOT output noise
- Do NOT over-generalize
- Do NOT mention internal retrieval, chunking, embeddings, OCR diagnostics, semantic signals, or metadata objects

---------------------------------------------------------------------

## 7. RESPONSE STYLE

- Professional
- Structured headings
- Bullet points
- Clear and concise
- No unnecessary repetition

---------------------------------------------------------------------

## 8. MEMORY AWARENESS

If prior context or summaries are available:
- Use them to improve answers
- Maintain consistency

---------------------------------------------------------------------

## 9. CONFIDENCE

At the end of every response, include:

Confidence: High / Medium / Low

Based on:
- completeness of context
- clarity of information

---------------------------------------------------------------------

## 10. FINAL RULE

Always produce the BEST POSSIBLE answer using available document context, even if incomplete.

Never default to refusal unless absolutely necessary."""

CHUNK_LEVEL_SUMMARY_PROMPT = """You are an expert technical document analyzer.

You are given a PARTIAL chunk of a large document. This is NOT the full document.

-----------------------------

## YOUR GOAL

Extract meaningful information from this chunk and produce a clean structured summary.

-----------------------------

## IMPORTANT CONTEXT

The input may contain:
- incomplete sentences
- OCR noise
- headers/footers
- page numbers
- table fragments
- repeated titles

You MUST ignore these.

-----------------------------

## FOCUS ONLY ON

- core ideas
- features or capabilities
- components/modules
- workflows/processes
- technical details (only if clearly stated)
- constraints or important notes

-----------------------------

## OUTPUT FORMAT

### Section Summary
1-2 sentences describing what this section is about

### Key Points
- bullet points of important information

### Technical Details
- only if clearly present

### Notes
- warnings, constraints, or observations

-----------------------------

## STRICT RULES

- Do NOT assume missing context
- Do NOT hallucinate
- Do NOT repeat text verbatim
- Do NOT preserve extracted headings, keyword lists, or internal extraction artifacts as the summary
- If content is mostly noise, return:

"This section contains minimal useful information."

CONTEXT:
{context}
"""

FINAL_DOCUMENT_ANALYSIS_PROMPT = """You are a senior enterprise document intelligence system.

You are given multiple summarized sections of a document. Together they represent the full document.

-----------------------------

## YOUR TASK

Combine these into a complete, structured, professional document analysis.
Write naturally, like an expert analyst explaining what matters and why.

-----------------------------

## OUTPUT STRUCTURE (MANDATORY)

### Overview
What the document is about

### Purpose
Why this document exists

### Core Concept
Main idea or system explained

### Key Features / Capabilities
- bullet list

### Architecture / Structure
System design or organization (if available)

### Components / Modules
Important parts and roles

### Workflow / Usage
Step-by-step or logical flow

### Use Cases / Applications
Where and how this is used

### Important Notes / Constraints
Limitations, warnings, or important details

### Key Takeaways
Concise insights

-----------------------------

## RULES

- Use ONLY provided summaries
- Merge overlapping information
- Avoid repetition
- Do NOT hallucinate
- Do NOT expose extracted headings, keyword lists, metadata, or internal processing details
- If missing info, write:
  "Not specified in the provided context"

-----------------------------

## STYLE

- Professional
- Clear headings
- Bullet points
- Clean structure

SECTION SUMMARIES:
{context}
"""

FAST_SUMMARY_PROMPT = """You are a document summarization assistant.

Provide a concise executive summary of the document.

-----------------------------

## OUTPUT

- What the document is about
- Purpose
- 3-5 key insights
- 2-3 key takeaways

-----------------------------

## RULES

- Focus only on meaningful content
- Ignore noise (headers, TOC, metadata)
- Do NOT hallucinate
- Do NOT concatenate headings or output keyword lists
- Keep it concise

CONTENT:
{context}
"""

RAG_QA_PROMPT = MASTER_SYSTEM_PROMPT + """

You are an expert assistant answering questions from document context.

-----------------------------

## YOUR TASK

Answer the user's question using ONLY the provided context.

-----------------------------

## RULES

- Be precise and direct
- Use only relevant information
- Do NOT hallucinate
- If answer is missing, say:
  "Not specified in the provided context"
- Include Sources using the readable source labels supplied in the context
- Do not mention retrieval, chunks, embeddings, keyword search, semantic search, or pipeline details

-----------------------------

## STYLE

- Clear
- Professional
- Structured if needed

Conversation memory for these selected documents:
{memory}

CONTEXT:
{context}

QUESTION:
{question}
"""

# ==============================
# NEW ANALYSIS BUTTON PROMPTS
# ==============================
ANALYSIS_PROMPT = """You are an Enterprise Document Intelligence Engine for technical documentation.

CORE CAPABILITIES:
- Analyze PDF, DOCX, PPTX, XLSX, CSV, TXT, HTML, images, manuals, specifications, reports
- Filter OCR noise, metadata, TOC, headers/footers automatically
- Classify intent: FULL_DOCUMENT_ANALYSIS | SHORT_SUMMARY | OVERVIEW | FEATURES_ONLY | SPECIFIC_COMPONENT_DETAILS | PIN_DIAGRAMS | WORKFLOW | USE_CASES | COMPARISON | TABLE_EXTRACTION | DIAGRAM_EXPLANATION | REPORT | TROUBLESHOOTING | REQUIREMENTS_EXTRACTION
- Generate structured, professional responses
- Extract tables/images as downloads
- Support pin diagrams, connectors, signal tables
- Handle VN devices, CAPL scripts, automotive specs

CONTEXT FILTERING (MANDATORY):
IGNORE unless directly relevant:
- PDF metadata (author, title, dates)
- Copyright/imprint/warranty/trademark
- Table of contents/page numbers
- Headers/footers
- Repeated section titles
- Raw OCR fragments
- Lines like "Main Features 13"

FOCUS ON:
- Explanatory paragraphs
- Purpose/use case
- Architecture/components
- Features/capabilities
- Workflow/steps
- Tables/diagrams
- Connectors/pinouts
- Safety/constraints

RESPONSE RULES BY BUTTON/INTENT:

🔍 ANALYZE (FULL_DOCUMENT_ANALYSIS):
1. Overview | 2. Purpose | 3. Core Concept | 4. Architecture | 5. Key Features | 6. Capabilities | 7. Components | 8. Workflow | 9. Use Cases | 10. Notes | 11. Takeaways

📋 SUMMARY (SHORT_SUMMARY):
Short summary | What it is | Purpose | Key points | Takeaways

👁️ OVERVIEW:
What it is | Who for | Usage | Main concept | Areas covered

⭐ FEATURES:
Table: Feature | What it does | Why matters | Component

SPECIFIC COMPONENT:
Only about requested item - Overview | Purpose | Features | Details | Interfaces | Usage | Notes

PIN/CONNECTOR:
Pin table | Diagram | Notes | "Not specified" if missing

NEVER:
- Copy raw text
- Show page numbers/metadata
- Repeat info
- Invent specs/pins
- Use TOC headings as content

QUALITY CHECK:
If context = metadata/TOC only → "Insufficient meaningful content"

FINAL OUTPUT:
Professional | Structured | Markdown tables | CSV-ready | Downloadable

The user request is:
"{USER_QUERY}"

The uploaded document content may contain OCR noise, metadata, page headers, footers, copyright text, table of contents, and repeated section titles. You must ignore those unless they are directly useful."""

SUMMARY_PROMPT = """Summarize this document clearly and professionally.

Ignore metadata, table of contents, headers, footers, copyright text, and OCR noise.

Give only:
- Short summary
- What the document is about
- Main purpose
- Most important points
- Key takeaways

Keep it concise.
Do not include architecture, long module lists, raw extracted text, or page-wise content unless necessary."""

OVERVIEW_PROMPT = """Give a high-level overview of this document.

Ignore metadata, table of contents, headers, footers, copyright text, and OCR noise.

Explain:
- What this document/product/system is
- Who it is for
- What it is used for
- Main concept
- Main areas covered

Keep it simple, clean, and professional.
Do not list raw headings or page numbers."""

FEATURES_PROMPT = """Extract the real features and capabilities described in this document.

Ignore metadata, table of contents, headers, footers, copyright text, and OCR noise.

Do not list headings such as "Main Features 13".
Instead, identify actual functional features from the explanatory content.

Output:
- Feature name
- What it does
- Why it matters
- Related component/module, if applicable

Use a clean table if possible.
Do not invent missing details."""

# Reader-facing analysis prompts. These override the legacy button prompts above
# so chat answers sound like expert synthesis instead of extraction diagnostics.
ANALYSIS_PROMPT = """You are an expert document analyst.

{USER_QUERY}

Produce a natural, human-like, insightful analysis grounded only in the document content.

Do not expose OCR artifacts, extracted headings, keyword lists, semantic diagnostics, metadata objects,
retrieval mechanics, chunking details, or internal pipeline information.

For Analyze, go beyond summary:
- Explain what the document is really about.
- Identify the purpose and core concept.
- Connect architecture, components, workflow, capabilities, constraints, and implications when supported.
- Discuss strengths, risks, design patterns, and practical takeaways.
- Clearly mark anything missing as "Not specified in the provided context."

Never dump raw extracted text unless the user explicitly asks for exact text. Never concatenate headings into
the answer. Write like ChatGPT, Claude, Gemini, or an advanced AI document assistant."""

SUMMARY_PROMPT = """Summarize this document as a concise executive summary.

Ignore metadata, table of contents, headers, footers, copyright text, and extraction noise.

Explain naturally:
- What the document is about
- The main purpose
- The most important insights
- The key takeaways

Do not output keyword lists, raw extracted text, page-wise content, or concatenated headings."""

OVERVIEW_PROMPT = """Give a high-level overview of this document for quick understanding.

Ignore metadata, table of contents, headers, footers, copyright text, and extraction noise.

Explain naturally:
- What it is
- Who or what it is for
- What it is used for
- The main concept
- The major areas it covers

Keep it simple, polished, and human-readable. Do not list raw headings or page numbers."""

FEATURES_PROMPT = """Extract the real features and capabilities described in this document.

Ignore metadata, table of contents, headers, footers, copyright text, and extraction noise.

Do not list headings such as "Main Features 13" and do not output a keyword list. Identify actual
functional features from explanatory content and explain what each feature does and why it matters.

Use a clean table when useful. Do not invent missing details."""

CREATOR_USERNAME = "Vignesh"
CREATOR_PASSWORD = "Rider@100"

# Query params are updated by app.py before helpers read them.
query_params = {}

# ==============================
# BACKEND / SHARED FUNCTION BLOCKS
# Sections below come from the original monolith: preview persistence, document
# extraction, workspace memory, chatbot helpers, dashboard parsers, compare
# helpers, CAPL analysis engine, and shared UI-adjacent helpers.
# ==============================
def get_file_hash(file_bytes):
    """Generate SHA256 hash of file contents for change detection"""
    return hashlib.sha256(file_bytes).hexdigest()


def file_has_changed(file_name, file_bytes):
    """Check if file has been modified since last processing"""
    new_hash = get_file_hash(file_bytes)
    cache_key = f"{file_name}_hash"
    old_hash = FILE_HASH_CACHE.get(cache_key)
    FILE_HASH_CACHE[cache_key] = new_hash
    return old_hash != new_hash


def load_preview_data():
    """Load preview data from file"""
    global PREVIEW_TOKENS, PREVIEW_STORE
    if os.path.exists(PREVIEW_DATA_FILE):
        try:
            with open(PREVIEW_DATA_FILE, "rb") as f:
                data = pickle.load(f)
            if not isinstance(data, dict):
                raise ValueError("preview data is not a dictionary")
            PREVIEW_TOKENS = data.get("tokens", {}) if isinstance(data.get("tokens", {}), dict) else {}
            PREVIEW_STORE = data.get("store", {}) if isinstance(data.get("store", {}), dict) else {}
        except Exception as e:
            backup_path = f"{PREVIEW_DATA_FILE}.corrupt.{datetime.now().strftime('%Y%m%d%H%M%S')}"
            try:
                os.replace(PREVIEW_DATA_FILE, backup_path)
                st.warning("Preview cache was corrupted and has been reset. Please open the document preview again from the sidebar.")
            except Exception:
                st.warning(f"Could not load preview data: {e}")
            PREVIEW_TOKENS = {}
            PREVIEW_STORE = {}


def save_preview_data():
    """Save preview data to file"""
    temp_file = None
    try:
        data = {
            "tokens": PREVIEW_TOKENS,
            "store": PREVIEW_STORE
        }
        preview_dir = os.path.dirname(PREVIEW_DATA_FILE) or "."
        os.makedirs(preview_dir, exist_ok=True)
        temp_file = f"{PREVIEW_DATA_FILE}.{os.getpid()}.{uuid.uuid4().hex}.tmp"
        with open(temp_file, "wb") as f:
            pickle.dump(data, f)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_file, PREVIEW_DATA_FILE)
    except Exception as e:
        # Streamlit reruns can overlap during upload/sidebar rendering. If an
        # atomic temp replace fails, keep the in-memory preview store working and
        # fall back to a direct write before showing a warning.
        try:
            with open(PREVIEW_DATA_FILE, "wb") as f:
                pickle.dump({
                    "tokens": PREVIEW_TOKENS,
                    "store": PREVIEW_STORE
                }, f)
        except Exception as fallback_error:
            st.warning(f"Could not save preview data: {fallback_error}")
    finally:
        if temp_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except Exception:
                pass


def cleanup_expired_preview_tokens():
    """Remove preview tokens older than 1 hour to prevent memory accumulation."""
    now = datetime.now()
    expired_tokens = []
    for token, data in PREVIEW_TOKENS.items():
        if now - data['timestamp'] > timedelta(hours=1):
            expired_tokens.append(token)
    
    for token in expired_tokens:
        del PREVIEW_TOKENS[token]
        if token in PREVIEW_STORE:
            del PREVIEW_STORE[token]
    
    if expired_tokens:
        save_preview_data()


def render_html_frame(html_content, height="content", width="stretch"):
    """Render inline HTML with Streamlit's supported components API."""
    if height == "content":
        height = 240
    if isinstance(height, int) and height < 1:
        height = 1
    component_width = None if width in (None, "stretch") else width
    components.html(str(html_content), width=component_width, height=height, scrolling=True)


# ==============================
# PAGINATION SCROLL ANCHOR HELPERS
# Streamlit-safe scroll reset for paginated viewers across all tabs.
# Stores the requested anchor in session_state before rerun, then scrolls after
# the next render using a tiny same-page script.
# ==============================
def request_scroll_to_anchor(anchor_id):
    """Ask the next Streamlit rerun to smoothly scroll to an anchor."""
    st.session_state.pending_scroll_anchor = str(anchor_id or "page-viewer-top")


def render_scroll_anchor(anchor_id):
    """Render a scroll target and consume pending scroll state if it matches."""
    anchor_id = str(anchor_id or "page-viewer-top")
    pending_anchor = st.session_state.get("pending_scroll_anchor")
    st.markdown(
        f"<div id='{html.escape(anchor_id)}' style='position:relative; top:-12px; height:1px;'></div>",
        unsafe_allow_html=True,
    )
    if pending_anchor == anchor_id:
        render_html_frame(
            f"""
            <script>
            const scrollTargetId = {json.dumps(anchor_id)};
            const scrollToAnchor = () => {{
                const root = window.parent ? window.parent.document : document;
                const target = root.getElementById(scrollTargetId);
                if (target) {{
                    target.scrollIntoView({{ behavior: "smooth", block: "start", inline: "nearest" }});
                }}
            }};
            requestAnimationFrame(scrollToAnchor);
            setTimeout(scrollToAnchor, 80);
            </script>
            """,
            height=1,
        )
        st.session_state.pending_scroll_anchor = None


def set_paginated_index(state_key, value, minimum, maximum, scroll_anchor_id):
    """Update a page index safely and request a top-of-viewer scroll."""
    bounded_value = max(minimum, min(int(value), maximum))
    if int(st.session_state.get(state_key, minimum)) != bounded_value:
        st.session_state[state_key] = bounded_value
        request_scroll_to_anchor(scroll_anchor_id)
    else:
        st.session_state[state_key] = bounded_value


@st.cache_data(show_spinner=False)
def get_needle_minimalist_logo():
    frames = []

    silver_grey = "#A0A0A0"
    star_light = "#DCDCDC"
    star_shadow = "#B8B8B8"
    canvas_size = 220
    center = canvas_size // 2
    radius = 86

    for angle_deg in range(360, 0, -15):
        image = Image.new("RGBA", (canvas_size, canvas_size), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image)
        raw_scale = math.cos(math.radians(angle_deg))
        flip_scale = raw_scale if abs(raw_scale) > 0.08 else (0.08 if raw_scale >= 0 else -0.08)

        ellipse_box = [
            center - int(radius * abs(flip_scale)),
            center - radius,
            center + int(radius * abs(flip_scale)),
            center + radius,
        ]
        draw.ellipse(ellipse_box, outline=silver_grey, width=4)

        for base_angle in [90, 210, 330]:
            angle = math.radians(base_angle)
            tip = (
                center + int(radius * 0.88 * math.cos(angle) * flip_scale),
                center - int(radius * 0.88 * math.sin(angle)),
            )
            side_l = (
                center + int(radius * 0.13 * math.cos(angle + 2.15) * flip_scale),
                center - int(radius * 0.13 * math.sin(angle + 2.15)),
            )
            side_r = (
                center + int(radius * 0.13 * math.cos(angle - 2.15) * flip_scale),
                center - int(radius * 0.13 * math.sin(angle - 2.15)),
            )
            c_l, c_r = (star_light, star_shadow) if flip_scale > 0 else (star_shadow, star_light)
            draw.polygon([(center, center), tip, side_l], fill=c_l)
            draw.polygon([(center, center), tip, side_r], fill=c_r)

        frames.append(image)

    gif_buf = BytesIO()
    if frames:
        frames[0].save(
            gif_buf,
            format='GIF',
            save_all=True,
            append_images=frames[1:],
            duration=80,
            loop=0,
            disposal=2
        )

    return base64.b64encode(gif_buf.getvalue()).decode('utf-8')


def render_status_strip():
    if not st.session_state.get("is_authenticated"):
        return

    # Get current elapsed time
    if 'start_time' not in st.session_state or st.session_state.start_time is None:
        st.session_state.start_time = time.time()
    
    elapsed = int(time.time() - st.session_state.start_time)
    hours, rem = divmod(elapsed, 3600)
    mins, secs = divmod(rem, 60)
    timer_str = f"{hours:02d}:{mins:02d}:{secs:02d}"

    username = st.session_state.get("logged_in_username") or "Vignesh"
    role = st.session_state.get("user_role") or "User"
    available_files = len(st.session_state.get("selected_files", []))

    # Create a placeholder for the live timer
    timer_placeholder = st.empty()
    
    # JavaScript for live timer
    live_timer_js = f"""
    <script>
        // Get the start time from Python
        var startTime = {st.session_state.start_time * 1000}; // Convert to milliseconds
        
        function updateTimer() {{
            var now = new Date().getTime();
            var elapsed = Math.floor((now - startTime) / 1000);
            
            var hours = Math.floor(elapsed / 3600);
            var minutes = Math.floor((elapsed % 3600) / 60);
            var seconds = elapsed % 60;
            
            var timerStr = 
                hours.toString().padStart(2, '0') + ':' +
                minutes.toString().padStart(2, '0') + ':' +
                seconds.toString().padStart(2, '0');
            
            // Update the timer display
            var timerElement = document.getElementById('live-timer');
            if (timerElement) {{
                timerElement.textContent = timerStr;
            }}
        }}
        
        // Update immediately and then every second
        updateTimer();
        setInterval(updateTimer, 1000);
    </script>
    """

    status_html = f"""
    <style>
        body {{
            margin: 0;
            background: transparent;
            font-family: 'Segoe UI', Tahoma, sans-serif;
        }}
        .dashboard-grid {{
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 8px;
            padding: 0 2px 4px;
        }}
        .metric-card {{
            position: relative;
            overflow: hidden;
            min-height: 58px;
            border-radius: 12px;
            padding: 8px 10px;
            border: 1px solid rgba(255, 255, 255, 0.5) !important;
            box-shadow: 0 6px 14px rgba(15, 23, 42, 0.12);
            transform: translateY(8px) scale(0.98);
            opacity: 0;
            animation: riseIn 0.65s ease-out forwards, cardFloat 4.5s ease-in-out infinite;
        }}
        .metric-card:nth-child(1) {{ animation-delay: 0.05s, 0.9s; }}
        .metric-card:nth-child(2) {{ animation-delay: 0.15s, 1.05s; }}
        .metric-card:nth-child(3) {{ animation-delay: 0.25s, 1.2s; }}
        .metric-card:nth-child(4) {{ animation-delay: 0.35s, 1.35s; }}
        .metric-card::before {{
            content: "";
            position: absolute;
            width: 100px;
            height: 100px;
            top: -40px;
            right: -20px;
            border-radius: 50%;
            background: rgba(255, 255, 255, 0.22);
            filter: blur(4px);
            animation: bubbleDrift 8s ease-in-out infinite;
        }}
        .metric-card::after {{
            content: "";
            position: absolute;
            width: 60px;
            height: 60px;
            bottom: -20px;
            left: -8px;
            border-radius: 50%;
            background: rgba(255, 255, 255, 0.16);
        }}
        .card-label, .card-value {{
            position: relative;
            z-index: 1;
            display: block;
        }}
        .card-label {{
            font-size: 0.58rem !important;
            font-weight: 700 !important;
            letter-spacing: 0.05em;
            text-transform: uppercase;
            color: rgba(15, 23, 42, 0.68) !important;
            margin-bottom: 4px;
        }}
        .card-value {{
            font-size: 1rem !important;
            font-weight: 800 !important;
            line-height: 1.2;
            word-break: break-word;
        }}
        #live-timer {{
            letter-spacing: 0.06em;
            animation: timerGlow 1.8s ease-in-out infinite;
        }}
        @keyframes riseIn {{
            from {{ opacity: 0; transform: translateY(8px) scale(0.98); }}
            to {{ opacity: 1; transform: translateY(0) scale(1); }}
        }}
        @keyframes cardFloat {{
            0%, 100% {{ transform: translateY(0); }}
            50% {{ transform: translateY(-3px); }}
        }}
        @keyframes bubbleDrift {{
            0%, 100% {{ transform: translate(0, 0); }}
            50% {{ transform: translate(-8px, 10px); }}
        }}
        @keyframes timerGlow {{
            0%, 100% {{ text-shadow: 0 0 0 rgba(46, 125, 50, 0); }}
            50% {{ text-shadow: 0 0 12px rgba(46, 125, 50, 0.25); }}
        }}
        @media (max-width: 900px) {{
            .dashboard-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }}
        }}
        @media (max-width: 560px) {{
            .dashboard-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
                gap: 8px;
                padding: 4px 2px 8px;
            }}
            .metric-card {{
                min-height: 74px;
                border-radius: 12px;
                padding: 9px 8px;
                box-shadow: 0 6px 14px rgba(15, 23, 42, 0.12);
            }}
            .card-label {{
                font-size: 0.54rem !important;
                letter-spacing: 0.04em;
                margin-bottom: 4px;
                white-space: normal;
            }}
            .card-value {{
                font-size: clamp(0.82rem, 4vw, 1rem) !important;
                line-height: 1.15;
                overflow-wrap: anywhere;
            }}
            #live-timer {{
                font-size: clamp(0.78rem, 3.8vw, 0.95rem) !important;
                letter-spacing: 0.02em;
            }}
        }}
    </style>
    <script>
        // Dynamically override button colors on page load and mutations
        function applyButtonColors() {{
            const buttons = document.querySelectorAll('button, [role="button"]');
            buttons.forEach(btn => {{
                btn.style.backgroundColor = '#e8f6ff !important';
                btn.style.color = '#1e293b !important';
                btn.style.borderColor = '#c0dff0 !important';
                btn.style.border = '2px solid #c0dff0 !important';
            }});
        }}
        
        // Apply on page load
        document.addEventListener('DOMContentLoaded', applyButtonColors);
        
        // Apply on mutation (when new buttons are added)
        const observer = new MutationObserver(applyButtonColors);
        observer.observe(document.body, {{ childList: true, subtree: true }});
        
        // Apply immediately
        applyButtonColors();
    </script>
    {live_timer_js}
    <div class="dashboard-grid">
        <div class="metric-card" style="background: linear-gradient(135deg, #e8eaf6 0%, #c5cae9 100%); color: #3c4f7e; border: 1px solid #e8eaf6;">
            <span class="card-label" style="color: #666;">👤 User</span>
            <span class="card-value" style="color: #3c4f7e;">{html.escape(username)}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #fce4ec 0%, #f8bbd9 100%); color: #7b1fa2; border: 1px solid #fce4ec;">
            <span class="card-label" style="color: #666;">🔑 Role</span>
            <span class="card-value" style="color: #7b1fa2;">{html.escape(str(role).title())}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #e3f2fd 0%, #bbdefb 100%); color: #1565c0; border: 1px solid #e3f2fd;">
            <span class="card-label" style="color: #666;">📁 Available Files</span>
            <span class="card-value" style="color: #1565c0;">{available_files}</span>
        </div>
        <div class="metric-card" style="background: linear-gradient(135deg, #e8f5e8 0%, #c8e6c9 100%); color: #2e7d32; border: 1px solid #e8f5e8;">
            <span class="card-label" style="color: #666;">⏱️ Usage Time</span>
            <span class="card-value" id="live-timer" style="color: #2e7d32; font-family: 'Courier New', monospace;">{timer_str}</span>
        </div>
    </div>
    """

    render_html_frame(status_html, height=118)


def _help_state_key(tab_name):
    return f"show_help_popup_{tab_name}"


def _help_query_param_key(tab_name):
    return f"help_popup_{tab_name}"


def ensure_help_popup_state(tab_name):
    key = _help_state_key(tab_name)
    query_key = _help_query_param_key(tab_name)
    if key not in st.session_state:
        st.session_state[key] = False
    if query_key in query_params and query_params[query_key]:
        query_value = query_params[query_key]
        if isinstance(query_value, list):
            query_value = query_value[0] if query_value else ""
        st.session_state[key] = str(query_value).strip().lower() in {"1", "true", "yes", "open"}
    return key


def _set_query_params(params):
    try:
        if hasattr(st, "query_params"):
            st.query_params.clear()
            for param_key, param_value in params.items():
                if isinstance(param_value, list):
                    st.query_params[param_key] = [str(v) for v in param_value]
                else:
                    st.query_params[param_key] = str(param_value)
        elif hasattr(st, "experimental_set_query_params"):
            st.experimental_set_query_params(**params)
        elif hasattr(st, "set_query_params"):
            st.set_query_params(**params)
    except Exception:
        pass


def set_help_popup_state(tab_name, is_open):
    state_key = ensure_help_popup_state(tab_name)
    query_key = _help_query_param_key(tab_name)
    st.session_state[state_key] = is_open

    updated_params = {}
    try:
        for param_key in query_params.keys():
            param_value = query_params[param_key]
            if isinstance(param_value, list):
                updated_params[param_key] = list(param_value)
            else:
                updated_params[param_key] = param_value
    except Exception:
        updated_params = dict(query_params) if isinstance(query_params, dict) else {}

    if is_open:
        updated_params[query_key] = "1"
    else:
        updated_params.pop(query_key, None)

    _set_query_params(updated_params)


def render_mobile_workspace_controls():
    if not st.session_state.get("is_authenticated"):
        return

    show_sidebar = st.session_state.get("mobile_sidebar_visible", False)
    if show_sidebar:
        mobile_mode_css = """
        <style>
        @media (max-width: 767px) {
            section.main,
            [data-testid="stMain"],
            div[data-testid="stMain"] {
                display: none !important;
            }
            [data-testid="stSidebar"],
            .stSidebar {
                display: block !important;
                width: 100% !important;
                min-width: 0 !important;
                max-width: 100% !important;
                position: relative !important;
                transform: none !important;
                visibility: visible !important;
                opacity: 1 !important;
            }
            [data-testid="stSidebar"] > div {
                width: 100% !important;
                max-width: 100% !important;
            }
        }
        </style>
        """
    else:
        mobile_mode_css = """
        <style>
        @media (max-width: 767px) {
            [data-testid="stSidebar"],
            .stSidebar {
                display: none !important;
                visibility: hidden !important;
                width: 0 !important;
                min-width: 0 !important;
                max-width: 0 !important;
                transform: translateX(-100%) !important;
            }
            section.main,
            [data-testid="stMain"],
            div[data-testid="stMain"],
            div[data-testid="stAppViewContainer"] {
                display: block !important;
                width: 100% !important;
                max-width: 100% !important;
                margin-left: 0 !important;
                padding-left: 0 !important;
            }
        }
        </style>
        """

    st.markdown(
        mobile_mode_css
        + """
        <style>
        .st-key-mobile_show_files_btn,
        .st-key-mobile_open_workspace_btn {
            display: none;
        }
        @media (max-width: 767px) {
            .st-key-mobile_show_files_btn,
            .st-key-mobile_open_workspace_btn {
                display: block !important;
                margin-bottom: 0.75rem !important;
            }
            .st-key-mobile_show_files_btn button,
            .st-key-mobile_open_workspace_btn button {
                width: 100% !important;
                min-height: 46px !important;
                border-radius: 12px !important;
                border: 2px solid #93c5fd !important;
                background: #eff6ff !important;
                color: #1e3a8a !important;
                font-weight: 800 !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def lazy_load_file_section(file_name, section_id, loader_func):
    """Lazily load file sections on demand to reduce initial load time"""
    cache_key = f"{file_name}_{section_id}"
    cached_result = FILE_TEXT_CACHE.get(cache_key, ttl_seconds=7200)
    if cached_result is not None:
        return cached_result
    
    result = loader_func()
    FILE_TEXT_CACHE.set(cache_key, result)
    return result


def optimize_tab_rendering():
    """Optimize rendering by deferring non-active tab loading"""
    active_tab = st.session_state.get("active_main_tab", "💬 Chat")
    return active_tab


def update_uploaded_file_status(file_name, status):
    """Update the sidebar status badge for an uploaded file."""
    for file_info in st.session_state.get("uploaded_files", []):
        if file_info.get("name") == file_name:
            file_info["status"] = status
            break


def ensure_file_processed(file_name):
    """Process file with caching to avoid redundant extraction"""
    file_info = get_uploaded_file_entry(file_name)
    if not file_info:
        return
    update_uploaded_file_status(file_name, "processing")
    file_name_lower = file_name.lower()
    file_bytes = file_info["bytes"]
    new_hash = get_file_hash(file_bytes)
    hash_cache_key = f"{file_name}_hash"
    has_changed = FILE_HASH_CACHE.get(hash_cache_key) != new_hash
    FILE_HASH_CACHE[hash_cache_key] = new_hash
    
    # Check cache first
    text_cache_key = f"{TEXT_EXTRACTION_SCHEMA_VERSION}:{file_name}"
    cached_text = FILE_TEXT_CACHE.get(text_cache_key)
    if cached_text is not None and not has_changed:
        st.session_state.file_texts[file_name] = cached_text
        if detect_file_type(file_name) == "excel":
            cached_excel = EXCEL_DATA_CACHE.get(file_name)
            if cached_excel is not None:
                st.session_state.excel_data_by_file[file_name] = cached_excel
        update_uploaded_file_status(file_name, "ready")
        return
    
    # Process file if not in cache
    if file_name not in st.session_state.file_texts or has_changed:
        extracted_text = extract_text(file_name, file_bytes)
        st.session_state.file_texts[file_name] = extracted_text
        FILE_TEXT_CACHE.set(text_cache_key, extracted_text)

    if detect_file_type(file_name) == "excel" and (file_name not in st.session_state.excel_data_by_file or has_changed):
        excel_data = extract_excel_data(file_name, file_bytes)
        st.session_state.excel_data_by_file[file_name] = excel_data
        EXCEL_DATA_CACHE.set(file_name, excel_data)

    update_uploaded_file_status(file_name, "ready")


def detect_file_type(filename):
    """Detect the normalized document type from a filename extension."""
    ext = os.path.splitext(str(filename or ""))[1].lower()
    mapping = {
        ".pdf": "pdf",
        ".docx": "word", ".doc": "word", ".odt": "word", ".rtf": "word",
        ".pptx": "ppt", ".ppt": "ppt",
        ".xlsx": "excel", ".xls": "excel",
        ".csv": "csv",
        ".txt": "text", ".md": "text", ".log": "text",
        ".html": "html", ".htm": "html",
        ".pages": "pages",
        ".capl": "code", ".can": "code",
        ".png": "image", ".jpg": "image", ".jpeg": "image", ".gif": "image", ".bmp": "image", ".webp": "image",
    }
    return mapping.get(ext, "unknown")


def _read_uploaded_bytes(file):
    """Return uploaded-file bytes without relying on the current file pointer."""
    if file is None:
        return b""
    if isinstance(file, bytes):
        return file
    if isinstance(file, bytearray):
        return bytes(file)
    if hasattr(file, "getvalue"):
        try:
            return file.getvalue()
        except Exception:
            pass
    try:
        if hasattr(file, "seek"):
            file.seek(0)
        data = file.read()
        return data.encode("utf-8", errors="ignore") if isinstance(data, str) else (data or b"")
    except Exception:
        return b""


def _to_bytes_io(file):
    return BytesIO(_read_uploaded_bytes(file))


def clean_text(text):
    """Normalize noisy extracted text while preserving paragraph boundaries for chunking."""
    text = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"(?im)^\s*Page\s+\d+\s*$", "", text)
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n", text)
    text = re.sub(r" *\n *", "\n", text)
    return text.strip()


def combine_extracted_content(text_parts):
    """Flatten typed extraction records into searchable text."""
    combined = []
    for content_type, content in text_parts or []:
        if not content:
            continue
        if content_type == "TEXT":
            combined.append(str(content))
        elif content_type == "TABLE":
            combined.append(f"TABLE:\n{content}")
        elif content_type == "IMAGE":
            combined.append(f"[IMAGE: {content}]")
        elif content_type == "EMBEDDED_IMAGE":
            combined.append(f"[EMBEDDED_IMAGE: {content}]")
        elif content_type == "METADATA":
            combined.append(str(content))
        elif content_type == "ERROR":
            combined.append(f"ERROR: {content}")
        elif content_type == "UNSUPPORTED":
            combined.append(str(content))
    return clean_text("\n".join(combined))


def extract_pdf(file):
    """Extract PDF text, preferring PyMuPDF when available and falling back to pdfplumber."""
    file_bytes = _read_uploaded_bytes(file)
    if not file_bytes:
        return ""

    fitz_module = None
    try:
        fitz_module = importlib.import_module("fitz")
    except Exception:
        fitz_module = None

    if fitz_module is not None:
        try:
            pdf_doc = fitz_module.open(stream=file_bytes, filetype="pdf")
            pages = [page.get_text() for page in pdf_doc]
            pdf_doc.close()
            return clean_text("\n".join(pages))
        except Exception:
            pass

    return combine_extracted_content(extract_pdf_content(BytesIO(file_bytes)))


def extract_word(file, filename=None):
    """Extract Word-like documents: DOCX directly, ODT/RTF via XML/plain-text helpers, DOC best-effort."""
    file_bytes = _read_uploaded_bytes(file)
    ext = os.path.splitext(str(filename or ""))[1].lower()
    if ext == ".docx" or not ext:
        try:
            document = docx.Document(BytesIO(file_bytes))
            paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            table_blocks = []
            for table_index, table in enumerate(document.tables, start=1):
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                if any(row.strip() for row in rows):
                    table_blocks.append(f"Table {table_index}:\n" + "\n".join(rows))
            return clean_text("\n".join(paragraphs + table_blocks))
        except Exception:
            if not ext:
                return ""
    if ext == ".odt":
        return combine_extracted_content(extract_odt_content(BytesIO(file_bytes)))
    if ext == ".rtf":
        return combine_extracted_content(extract_rtf_content(BytesIO(file_bytes)))
    if ext == ".doc":
        return combine_extracted_content(extract_legacy_office_content(BytesIO(file_bytes), "Legacy Word document"))
    return ""


def extract_ppt(file, filename=None):
    """Extract PowerPoint slide text and simple table content."""
    file_bytes = _read_uploaded_bytes(file)
    ext = os.path.splitext(str(filename or ""))[1].lower()
    if ext == ".ppt":
        return combine_extracted_content(extract_legacy_office_content(BytesIO(file_bytes), "Legacy PowerPoint document"))
    try:
        presentation = Presentation(BytesIO(file_bytes))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.append(shape.text)
                if hasattr(shape, "table"):
                    rows = []
                    for row in shape.table.rows:
                        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                    if rows:
                        content.append("Table:\n" + "\n".join(rows))
            slides.append(f"Slide {index}:\n" + "\n".join(content))
        return clean_text("\n\n".join(slides))
    except Exception:
        return combine_extracted_content(extract_pptx_content(BytesIO(file_bytes)))


def extract_excel(file, filename=None):
    """Extract spreadsheet data from every sheet into text."""
    file_bytes = _read_uploaded_bytes(file)
    try:
        sheets = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
        text = []
        for name, df in sheets.items():
            text.append(f"Sheet: {name}")
            text.append(df.fillna("").to_string(index=False))
        return clean_text("\n\n".join(text))
    except Exception:
        ext = os.path.splitext(str(filename or ""))[1].lower()
        if ext == ".xls":
            return combine_extracted_content(extract_legacy_office_content(BytesIO(file_bytes), "Legacy Excel workbook"))
        return combine_extracted_content(extract_xlsx_content(BytesIO(file_bytes)))


def extract_csv(file):
    """Extract CSV data into plain table text."""
    try:
        df = pd.read_csv(_to_bytes_io(file))
        return clean_text(df.fillna("").to_string(index=False))
    except Exception:
        return combine_extracted_content(extract_csv_content(_to_bytes_io(file)))


def extract_html(file):
    """Extract visible text from HTML."""
    try:
        soup = BeautifulSoup(_read_uploaded_bytes(file), "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return clean_text(soup.get_text("\n", strip=True))
    except Exception:
        return combine_extracted_content(extract_html_content(_to_bytes_io(file)))


def extract_image_ocr(file):
    """Extract text from an image with OCR when pytesseract/Tesseract are available."""
    try:
        pytesseract_module = importlib.import_module("pytesseract")
        image = Image.open(_to_bytes_io(file))
        return clean_text(pytesseract_module.image_to_string(image))
    except Exception as exc:
        return f"OCR unavailable for this image: {str(exc)[:160]}"


def extract_code(file):
    """Extract source/code-like files without aggressively normalizing syntax."""
    return _read_uploaded_bytes(file).decode("utf-8", errors="ignore").strip()


def extract_pages(file):
    """Apple Pages is proprietary; ask the user for a reliable export format."""
    return "Please convert the .pages file to PDF or DOCX for reliable chat analysis."


def extract_text(file_or_name, filename_or_bytes=None):
    """Unified loader for upload -> detect type -> format-specific extraction.

    Supports both the new `(file, filename)` shape and the existing app calls
    that pass `(file_name, file_bytes)`.
    """
    if isinstance(file_or_name, (str, os.PathLike)) and isinstance(filename_or_bytes, (bytes, bytearray)):
        filename = str(file_or_name)
        file = BytesIO(bytes(filename_or_bytes))
    else:
        file = file_or_name
        filename = str(filename_or_bytes or getattr(file_or_name, "name", "") or "")

    file_type = detect_file_type(filename)
    try:
        if file_type == "pdf":
            return extract_pdf(file)
        if file_type == "word":
            return extract_word(file, filename)
        if file_type == "ppt":
            return extract_ppt(file, filename)
        if file_type == "excel":
            return extract_excel(file, filename)
        if file_type == "csv":
            return extract_csv(file)
        if file_type == "html":
            return extract_html(file)
        if file_type == "text":
            return clean_text(_read_uploaded_bytes(file).decode("utf-8", errors="ignore"))
        if file_type == "image":
            return extract_image_ocr(file)
        if file_type == "code":
            return extract_code(file)
        if file_type == "pages":
            return extract_pages(file)
        return ""
    except Exception as exc:
        return f"ERROR: Error extracting content: {str(exc)}"


def smart_chunk(text, chunk_size=1200):
    """Paragraph-aware chunking before embedding."""
    paragraphs = [paragraph.strip() for paragraph in str(text or "").split("\n") if paragraph.strip()]
    chunks = []
    current = ""

    for paragraph in paragraphs:
        if not current:
            current = paragraph + "\n"
        elif len(current) + len(paragraph) < chunk_size:
            current += paragraph + "\n"
        else:
            chunks.append(current.strip())
            current = paragraph + "\n"

    if current.strip():
        chunks.append(current.strip())
    return chunks


def create_documents(chunks, filename, file_type):
    """Create LangChain documents with normalized and legacy-compatible metadata."""
    return [
        Document(
            page_content=chunk,
            metadata={
                "source": filename,
                "type": file_type,
                "file_name": filename,
                "file_type": file_type,
                "page_number": "1",
                "page_or_sheet": "Document",
                "section": "Document",
                "chunk_index": index,
            },
        )
        for index, chunk in enumerate(chunks, start=1)
    ]


def process_file(file, filename):
    """Run the upload -> detect -> extract -> normalize -> chunk -> document pipeline."""
    file_type = detect_file_type(filename)
    raw_text = extract_text(file, filename)
    normalized_text = clean_text(raw_text)
    if not normalized_text:
        return []
    return create_documents(smart_chunk(normalized_text), filename, file_type)


def extract_pdf_content(bio):
    """Extract searchable PDF text quickly.

    Full table/image extraction is intentionally deferred to preview/download
    helpers because scanning every page of a large manual is slow.
    """
    content = []
    try:
        with pdfplumber.open(bio) as pdf:
            # Add metadata
            if pdf.metadata:
                metadata = []
                for key, value in pdf.metadata.items():
                    if value:
                        metadata.append(f"{key}: {value}")
                if metadata:
                    content.append(("METADATA", "PDF Metadata:\n" + "\n".join(metadata)))
            
            content.append(("METADATA", f"Total Pages: {len(pdf.pages)}"))
            
            # Extract text from each page. Avoid table/image scans here so Chat
            # can load large manuals without blocking on expensive page parsing.
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    content.append(("TEXT", f"Page {i+1} Text:\n{page_text}"))
    
    except Exception as e:
        content.append(("ERROR", f"PDF extraction failed: {str(e)}"))
    
    return content


def table_to_png_bytes(table_data, title=None):
    """Render table rows as a PNG image and return the bytes."""
    try:
        font = ImageFont.load_default()
    except Exception:
        font = None

    padding = 10
    row_height = 22
    col_padding = 18

    # Normalize table data
    normalized_table = [[str(cell) for cell in row] for row in table_data]
    col_widths = []
    for col_idx in range(len(normalized_table[0])):
        col_width = max(len(row[col_idx]) for row in normalized_table) * 7 + col_padding
        col_widths.append(col_width)

    width = sum(col_widths) + padding * 2
    height = row_height * len(normalized_table) + padding * 2
    if title:
        height += row_height

    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    y = padding

    if title:
        draw.text((padding, y), title, fill="black", font=font)
        y += row_height

    for row in normalized_table:
        x = padding
        for col_idx, cell in enumerate(row):
            draw.text((x, y), cell, fill="black", font=font)
            x += col_widths[col_idx]
        y += row_height

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def image_bytes_to_png_bytes(image_bytes):
    """Convert an uploaded image to PNG bytes."""
    with Image.open(BytesIO(image_bytes)) as image:
        png_buffer = BytesIO()
        image.save(png_buffer, format="PNG")
        return png_buffer.getvalue()


def dataframe_to_table_rows(df):
    """Convert a dataframe to table rows suitable for PNG rendering."""
    safe_df = df.fillna("")
    rows = [list(map(str, safe_df.columns.tolist()))]
    rows.extend([list(map(str, row)) for row in safe_df.values.tolist()])
    return rows


def crop_pdf_region_to_png(page, bbox, resolution=150):
    """Crop a rectangular region from a PDF page and return it as PNG bytes."""
    cropped_page = page.crop(bbox)
    cropped_image = cropped_page.to_image(resolution=resolution)
    output = BytesIO()
    cropped_image.original.save(output, format="PNG")
    return output.getvalue()


def extract_docx_content(bio):
    """Extract text, tables, and images from DOCX."""
    content = []
    try:
        doc = docx.Document(bio)
        
        # Extract metadata
        core_props = doc.core_properties
        metadata = []
        if core_props.title:
            metadata.append(f"Title: {core_props.title}")
        if core_props.author:
            metadata.append(f"Author: {core_props.author}")
        if core_props.created:
            metadata.append(f"Created: {core_props.created}")
        if metadata:
            content.append(("METADATA", "Document Metadata:\n" + "\n".join(metadata)))
        
        # Extract embedded images for preview and downloads
        image_count = 0
        if 'extracted_images' not in st.session_state:
            st.session_state.extracted_images = {}
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    image_ext = image_part.content_type.split('/')[-1]
                    if image_ext == 'jpeg':
                        image_ext = 'jpg'
                    image_key = f"docx_image_{image_count}"
                    st.session_state.extracted_images[image_key] = {
                        'bytes': image_bytes,
                        'ext': image_ext,
                        'filename': f"image_{image_count}.{image_ext}"
                    }
                    content.append(("EMBEDDED_IMAGE", f"Embedded Image {image_count + 1}: {image_key}"))
                    image_count += 1
                except Exception as e:
                    content.append(("ERROR", f"Could not extract image {image_count + 1}: {e}"))
        if image_count > 0:
            content.append(("METADATA", f"Total Images: {image_count}"))

        # Walk the document body in order and preserve headings, tables, and text.
        current_section_title = None
        current_section_lines = []
        table_count = 0

        def flush_current_section():
            nonlocal current_section_title, current_section_lines
            if current_section_title or current_section_lines:
                if current_section_title:
                    section_text = "\n\n".join(current_section_lines).strip()
                    if section_text:
                        content.append(("TEXT", f"Heading: {current_section_title}\n{section_text}"))
                    else:
                        content.append(("TEXT", f"Heading: {current_section_title}"))
                else:
                    section_text = "\n\n".join(current_section_lines).strip()
                    if section_text:
                        content.append(("TEXT", section_text))
                current_section_title = None
                current_section_lines = []

        def is_docx_heading(para):
            text = para.text.strip()
            if not text:
                return False
            try:
                style_name = (para.style.name or "").lower()
            except Exception:
                style_name = ""
            if "heading" in style_name or style_name.startswith("title") or style_name.startswith("subtitle"):
                return True
            if re.match(r'^\d+(?:\.\d+)*\s+.+', text) and len(text) <= 120:
                return True
            return False

        for element in doc.element.body:
            if element.tag.endswith('}p'):
                paragraph = Paragraph(element, doc)
                paragraph_text = paragraph.text.strip()
                if not paragraph_text:
                    continue
                if is_docx_heading(paragraph):
                    flush_current_section()
                    current_section_title = paragraph_text
                else:
                    current_section_lines.append(paragraph_text)
            elif element.tag.endswith('}tbl'):
                flush_current_section()
                table_count += 1
                table = Table(element, doc)
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    content.append(("TABLE", f"Table {table_count}:\n{table_text}"))

        flush_current_section()

        if table_count > 0:
            content.append(("METADATA", f"Total Tables: {table_count}"))

        if not any(item[0] in ("TEXT", "TABLE", "EMBEDDED_IMAGE", "IMAGE", "METADATA") for item in content):
            paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            if paragraphs:
                content.append(("TEXT", "\n\n".join(paragraphs)))
    except Exception as e:
        content.append(("ERROR", f"DOCX extraction failed: {str(e)}"))
    return content


def extract_pptx_content(bio):
    """Extract text, tables, and images from PPTX."""
    content = []
    try:
        prs = Presentation(bio)
        
        content.append(("METADATA", f"Total Slides: {len(prs.slides)}"))
        
        # Extract and display images
        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_ext = shape.image.content_type.split('/')[-1]
                        if image_ext == 'jpeg':
                            image_ext = 'jpg'
                        
                        # Create a unique key for the image
                        image_key = f"pptx_image_{image_count}"
                        
                        # Store image data for display
                        if 'extracted_images' not in st.session_state:
                            st.session_state.extracted_images = {}
                        st.session_state.extracted_images[image_key] = {
                            'bytes': image_bytes,
                            'ext': image_ext,
                            'filename': f"slide_image_{image_count}.{image_ext}"
                        }
                        
                        content.append(("EMBEDDED_IMAGE", f"Slide Image {image_count + 1}: {image_key}"))
                        image_count += 1
                    except Exception as e:
                        content.append(("ERROR", f"Could not extract slide image {image_count + 1}: {e}"))
        
        if image_count > 0:
            content.append(("METADATA", f"Total Images: {image_count}"))
        
        # Extract text and tables from each slide
        for i, slide in enumerate(prs.slides):
            slide_content = []
            
            # Extract text shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            
            # Extract tables
            for shape in slide.shapes:
                if hasattr(shape, 'table'):
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    
                    if table_data:
                        table_text = "\n".join([" | ".join(row) for row in table_data])
                        slide_content.append(f"Table:\n{table_text}")
            
            if slide_content:
                content.append(("TEXT", f"Slide {i+1}:\n" + "\n\n".join(slide_content)))
    
    except Exception as e:
        content.append(("ERROR", f"PPTX extraction failed: {str(e)}"))
    
    return content


def extract_xlsx_content(bio):
    """Extract data from all sheets in XLSX."""
    content = []
    try:
        wb = openpyxl.load_workbook(bio, data_only=True)
        
        content.append(("METADATA", f"Workbook contains {len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)}"))
        
        # Extract data from each sheet
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_data = []
            
            # Get all rows with data
            for row in sheet.iter_rows(values_only=True):
                if any(cell for cell in row):  # Skip empty rows
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    sheet_data.append(row_data)
            
            if sheet_data:
                table_text = "\n".join([" | ".join(row) for row in sheet_data])
                content.append(("TABLE", f"Sheet '{sheet_name}':\n{table_text}"))
    
    except Exception as e:
        content.append(("ERROR", f"XLSX extraction failed: {str(e)}"))
    
    return content


def extract_csv_content(bio):
    """Extract CSV rows as table text for Chat, Compare, and Preview."""
    content = []
    try:
        df = pd.read_csv(bio)
        content.append(("METADATA", f"CSV rows: {len(df)} columns: {len(df.columns)}"))
        if not df.empty:
            preview_df = df.fillna("").head(500)
            table_text = "\n".join(
                " | ".join(map(str, row))
                for row in [preview_df.columns.tolist()] + preview_df.values.tolist()
            )
            content.append(("TABLE", f"CSV Data:\n{table_text}"))
    except Exception as e:
        content.append(("ERROR", f"CSV extraction failed: {str(e)}"))
    return content


def extract_html_content(bio):
    """Extract text and metadata from HTML."""
    content = []
    try:
        html_content = bio.read()
        soup = BeautifulSoup(html_content, "html.parser")
        
        # Extract title
        title = soup.title.string if soup.title else "No title"
        content.append(("METADATA", f"Title: {title}"))
        
        # Extract meta tags
        meta_info = []
        for meta in soup.find_all('meta'):
            if meta.get('name') and meta.get('content'):
                meta_info.append(f"{meta['name']}: {meta['content']}")
        if meta_info:
            content.append(("METADATA", "Meta Tags:\n" + "\n".join(meta_info)))
        
        # Count images
        images = soup.find_all('img')
        if images:
            content.append(("IMAGE", f"{len(images)} images found in HTML"))
        
        # Extract text content
        text = soup.get_text(separator="\n")
        if text.strip():
            content.append(("TEXT", text))
    
    except Exception as e:
        content.append(("ERROR", f"HTML extraction failed: {str(e)}"))
    
    return content


def xml_text_content(xml_bytes):
    """Extract readable text from XML-based office documents."""
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        soup = BeautifulSoup(xml_bytes, "xml")
        return soup.get_text("\n", strip=True)

    text_items = []
    for element in root.iter():
        tag_name = element.tag.split("}", 1)[-1].lower()
        if tag_name in {"p", "h", "span", "line-break", "tab"} and element.text:
            text_items.append(element.text.strip())
        if element.tail:
            text_items.append(element.tail.strip())
    return "\n".join(item for item in text_items if item)


def extract_odt_content(bio):
    """Extract text and simple tables from OpenDocument Text files."""
    content = []
    try:
        with zipfile.ZipFile(bio) as odt_zip:
            if "meta.xml" in odt_zip.namelist():
                meta_text = xml_text_content(odt_zip.read("meta.xml"))
                if meta_text.strip():
                    content.append(("METADATA", "ODT Metadata:\n" + meta_text[:2000]))

            if "content.xml" not in odt_zip.namelist():
                content.append(("ERROR", "ODT content.xml was not found."))
                return content

            content_xml = odt_zip.read("content.xml")
            soup = BeautifulSoup(content_xml, "xml")
            text_blocks = []
            for node in soup.find_all(["text:h", "text:p"]):
                text_value = node.get_text(" ", strip=True)
                if text_value:
                    text_blocks.append(text_value)
            if not text_blocks:
                fallback_text = xml_text_content(content_xml)
                if fallback_text.strip():
                    text_blocks.append(fallback_text)
            if text_blocks:
                content.append(("TEXT", "\n".join(text_blocks)))

            for table_index, table in enumerate(soup.find_all("table:table"), start=1):
                rows = []
                for row in table.find_all("table:table-row"):
                    cells = [cell.get_text(" ", strip=True) for cell in row.find_all("table:table-cell")]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    content.append(("TABLE", f"Table {table_index}:\n" + "\n".join(rows)))
    except Exception as e:
        content.append(("ERROR", f"ODT extraction failed: {str(e)}"))
    return content


def strip_rtf_to_text(rtf_text):
    """Best-effort RTF to plain text conversion without external dependencies."""
    text = rtf_text
    text = re.sub(r"\\'[0-9a-fA-F]{2}", lambda m: bytes.fromhex(m.group(0)[2:]).decode("latin-1", errors="ignore"), text)
    text = re.sub(r"\\(par|line)\b", "\n", text)
    text = re.sub(r"\\tab\b", "\t", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("\\{", "{").replace("\\}", "}").replace("\\\\", "\\")
    text = re.sub(r"[{}]", "", text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return html.unescape(text).strip()


def extract_rtf_content(bio):
    """Extract text from Rich Text Format documents."""
    content = []
    try:
        raw = bio.read()
        rtf_text = raw.decode("utf-8", errors="ignore")
        if not rtf_text.strip():
            rtf_text = raw.decode("latin-1", errors="ignore")
        plain_text = strip_rtf_to_text(rtf_text)
        if plain_text:
            content.append(("TEXT", plain_text))
        else:
            content.append(("ERROR", "No readable text was found in the RTF file."))
    except Exception as e:
        content.append(("ERROR", f"RTF extraction failed: {str(e)}"))
    return content


def extract_legacy_office_content(bio, label):
    """Best-effort text recovery for legacy binary Office files."""
    content = []
    try:
        raw = bio.read()
        decoded = raw.decode("utf-16le", errors="ignore") + "\n" + raw.decode("latin-1", errors="ignore")
        strings = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s.,;:!?()/_+\-]{3,}", decoded)
        cleaned = []
        seen = set()
        for value in strings:
            value = re.sub(r"\s+", " ", value).strip()
            if len(value) < 4 or len(value) > 240:
                continue
            if value.lower() in seen:
                continue
            seen.add(value.lower())
            cleaned.append(value)
            if len(cleaned) >= 1000:
                break

        if cleaned:
            content.append(("METADATA", f"{label}: recovered readable text using best-effort binary extraction."))
            content.append(("TEXT", "\n".join(cleaned)))
        else:
            content.append(("ERROR", f"{label} text could not be recovered. Save/export as DOCX, PPTX, XLSX, PDF, RTF, or TXT for full analysis."))
    except Exception as e:
        content.append(("ERROR", f"{label} extraction failed: {str(e)}"))
    return content


def extract_pages_content(bio):
    """Extract readable text from Apple Pages files when XML/text previews exist."""
    content = []
    try:
        with zipfile.ZipFile(bio) as pages_zip:
            names = pages_zip.namelist()
            readable_parts = [
                name for name in names
                if name.lower().endswith((".xml", ".txt", ".html", ".xhtml"))
                and not name.lower().startswith(("metadata/", "quicklook/thumbnail"))
            ]

            extracted_blocks = []
            for name in readable_parts[:20]:
                try:
                    part_bytes = pages_zip.read(name)
                    if name.lower().endswith((".html", ".xhtml")):
                        text_value = BeautifulSoup(part_bytes, "html.parser").get_text("\n", strip=True)
                    elif name.lower().endswith(".xml"):
                        text_value = xml_text_content(part_bytes)
                    else:
                        text_value = part_bytes.decode("utf-8", errors="ignore")
                    text_value = re.sub(r"\s+", " ", text_value).strip()
                    if text_value:
                        extracted_blocks.append(f"{name}\n{text_value}")
                except Exception:
                    pass

            if extracted_blocks:
                content.append(("TEXT", "\n\n".join(extracted_blocks)))
            else:
                content.append(("ERROR", "No readable text preview was found in this Pages file. Export it as DOCX/PDF for full analysis."))
    except zipfile.BadZipFile:
        content.append(("ERROR", "Pages extraction failed because this file is not a readable Pages ZIP package."))
    except Exception as e:
        content.append(("ERROR", f"Pages extraction failed: {str(e)}"))
    return content


def get_uploaded_file_entry(file_name):
    for file_info in st.session_state.uploaded_files:
        if file_info["name"] == file_name:
            return file_info
    return None


def create_preview_link(file_name, highlight_term=None, page_num=None):
    file_entry = get_uploaded_file_entry(file_name)
    if not file_entry:
        return None

    token = None
    for existing_token, token_data in list(PREVIEW_TOKENS.items()):
        if token_data.get("file_name") == file_name and existing_token in PREVIEW_STORE:
            token = existing_token
            token_data["timestamp"] = datetime.now()
            PREVIEW_STORE[existing_token] = file_entry
            break

    if token is None:
        token = str(uuid.uuid4())
        PREVIEW_TOKENS[token] = {'file_name': file_name, 'timestamp': datetime.now()}
        PREVIEW_STORE[token] = file_entry
        save_preview_data()

    params = [f"preview_token={token}"]
    if highlight_term:
        params.append(f"highlight={urllib.parse.quote_plus(highlight_term)}")
    if page_num is not None:
        params.append(f"page={urllib.parse.quote_plus(str(page_num))}")
    return "?" + "&".join(params)


def create_heading_anchor(text):
    anchor_text = str(text or "").strip().lower()
    anchor_text = re.sub(r'[^a-z0-9]+', '-', anchor_text)
    anchor_text = re.sub(r'-{2,}', '-', anchor_text).strip('-')
    if not anchor_text:
        anchor_text = 'preview'
    return f"heading-{anchor_text}"


def highlight_for_preview(text, highlight_term=None):
    if not highlight_term:
        return html.escape(text)
    escaped = html.escape(text)
    pattern = re.compile(re.escape(highlight_term), re.IGNORECASE)
    highlighted = pattern.sub(
        lambda m: f"<mark style='background:#fff3a3; padding:0 2px;'>{html.escape(m.group(0))}</mark>",
        escaped
    )
    return highlighted


def render_text_block(text, highlight_term=None, anchor_id=None):
    if not highlight_term:
        return None

    escaped = html.escape(text)
    pattern = re.compile(re.escape(highlight_term), re.IGNORECASE)
    first_match = True

    def replace_match(match):
        nonlocal first_match
        content = html.escape(match.group(0))
        if first_match and anchor_id:
            first_match = False
            return f"<span id='{anchor_id}'></span><mark style='background:#fff3a3; padding:0 2px;'>{content}</mark>"
        first_match = False
        return f"<mark style='background:#fff3a3; padding:0 2px;'>{content}</mark>"

    highlighted = pattern.sub(replace_match, escaped)
    return f"<pre style='white-space: pre-wrap; word-break: break-word; background:#f4f7fb; padding:12px; border-radius:8px; font-family: inherit;'>{highlighted}</pre>"


def render_document_preview(file_name, file_entry=None, highlight_term=None, highlight_page=None):
    """Render document preview with caching and error handling"""
    st.markdown(f"**Preview: {file_name}**")
    file_name_lower = file_name.lower()
    
    # Ensure file entry exists
    if file_entry is None:
        if not file_name_lower.endswith(".pdf"):
            ensure_file_processed(file_name)
        file_entry = get_uploaded_file_entry(file_name)
    else:
        # Keep PDF rendering lazy. Full PDF text extraction happens only when an
        # analysis panel explicitly asks for it.
        if file_name not in st.session_state.file_texts and not file_name_lower.endswith(".pdf"):
            extracted = extract_text(file_name, file_entry["bytes"])
            st.session_state.file_texts[file_name] = extracted
            FILE_TEXT_CACHE.set(file_name, extracted)
        if file_name.lower().endswith(".xlsx") and file_name not in st.session_state.excel_data_by_file:
            excel_data = extract_excel_data(file_name, file_entry["bytes"])
            st.session_state.excel_data_by_file[file_name] = excel_data
            EXCEL_DATA_CACHE.set(file_name, excel_data)
    
    if not file_entry:
        st.error("❌ File preview unavailable - file could not be loaded.")
        return

    image_download_items = []
    table_download_items = []

    # Special handling for PDF files: render actual page images for true preview
    if file_name_lower.endswith(".pdf"):
        try:
            pdf_bio = BytesIO(file_entry["bytes"])
            with pdfplumber.open(pdf_bio) as pdf:
                total_pages = len(pdf.pages)
                if total_pages == 0:
                    st.warning("⚠️ PDF has no readable pages")
                    return
                    
                st.markdown(f"**PDF Pages: {total_pages}**")
                
                preview_key_base = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
                batch_key = f"pdf_preview_batch_{preview_key_base}"
                scroll_anchor_id = f"pdf-viewer-top-{preview_key_base}"
                batch_size = 5
                max_batch = max(0, (total_pages - 1) // batch_size)
                default_batch = 0
                if highlight_page and 1 <= highlight_page <= total_pages:
                    default_batch = (highlight_page - 1) // batch_size
                elif highlight_page is not None:
                    st.warning(f"⚠️ Page {highlight_page} not found. Showing the first page batch.")

                if batch_key not in st.session_state:
                    st.session_state[batch_key] = default_batch
                st.session_state[batch_key] = max(0, min(int(st.session_state[batch_key]), max_batch))

                batch_start = st.session_state[batch_key] * batch_size
                batch_end = min(total_pages, batch_start + batch_size)
                pages_to_show = range(batch_start, batch_end)

                render_scroll_anchor(scroll_anchor_id)
                render_tables = st.checkbox(
                    "Detect tables on previewed pages",
                    value=False,
                    key=f"pdf_preview_tables_{preview_key_base}",
                    help="Table detection is slower, so it only runs when enabled.",
                )
                st.caption(f"Showing pages {batch_start + 1} to {batch_end} of {total_pages}. Five pages are rendered per batch for faster previews.")

                # Render pages
                highlight_found = False
                for i in pages_to_show:
                    page = pdf.pages[i]
                    
                    try:
                        # Extract page text for highlighting
                        page_text = page.extract_text() or ""
                        page_anchor_id = None
                        
                        if highlight_term and highlight_term.lower() in page_text.lower():
                            page_anchor_id = create_heading_anchor(highlight_term)
                            highlight_found = True
                            if page_anchor_id:
                                st.markdown(f"<div id='{page_anchor_id}'></div>", unsafe_allow_html=True)
                        
                        # Render page image with caching
                        page_cache_key = f"{file_name}_page_{i}_image_{PDF_PREVIEW_RESOLUTION}"
                        cached_image = FILE_TEXT_CACHE.get(page_cache_key)
                        
                        if cached_image is None:
                            page_image = page.to_image(resolution=PDF_PREVIEW_RESOLUTION)
                            image_bytes_io = BytesIO()
                            page_image.original.save(image_bytes_io, format="PNG")
                            image_bytes = image_bytes_io.getvalue()
                            FILE_TEXT_CACHE.set(page_cache_key, image_bytes)
                        else:
                            image_bytes = cached_image
                        
                        st.image(image_bytes, caption=f"Page {i+1}", use_container_width=True)
                        
                        # Show highlighted text if match found
                        if page_anchor_id and highlight_term:
                            st.markdown("### Highlighted Text", unsafe_allow_html=True)
                            st.markdown(render_text_block(page_text, highlight_term, anchor_id=None), unsafe_allow_html=True)
                        
                        # Extract tables only when requested; this is slow on large manuals.
                        tables = page.extract_tables() if render_tables else []
                        if tables:
                            for j, table in enumerate(tables):
                                if table and any(any(cell for cell in row) for row in table):
                                    table_cache_key = f"{file_name}_page_{i}_table_{j}"
                                    cached_table = FILE_TEXT_CACHE.get(table_cache_key)
                                    
                                    if cached_table is None:
                                        table_png = table_to_png_bytes(table, title=f"Page {i+1} Table {j+1}")
                                        FILE_TEXT_CACHE.set(table_cache_key, table_png)
                                    else:
                                        table_png = cached_table
                                    
                                    st.image(table_png, caption=f"Page {i+1} Table {j+1}", use_container_width=True)
                                    table_download_items.append({
                                        "label": f"📥 Download Table {j+1} as PNG",
                                        "data": table_png,
                                        "file_name": f"{os.path.splitext(file_name)[0]}_page_{i+1}_table_{j+1}.png",
                                        "mime": "image/png",
                                        "key": f"download_pdf_table_{file_name}_{i}_{j}"
                                    })
                        
                        image_download_items.append({
                            "label": f"📥 Download Page {i+1} as PNG",
                            "data": image_bytes,
                            "file_name": f"{os.path.splitext(file_name)[0]}_page_{i+1}.png",
                            "mime": "image/png",
                            "key": f"download_pdf_page_{file_name}_{i}"
                        })
                    
                    except Exception as page_err:
                        st.warning(f"⚠️ Could not render page {i+1} as image: {str(page_err)[:100]}")
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            st.markdown(f"#### Page {i+1} Text")
                            st.code(page_text[:1000], language="text")

                st.markdown("---")
                nav_cols = st.columns([1, 1], vertical_alignment="center")
                with nav_cols[0]:
                    if st.button(
                        "⬅ Previous",
                        key=f"pdf_prev_batch_{preview_key_base}",
                        use_container_width=True,
                        disabled=st.session_state[batch_key] <= 0,
                    ):
                        set_paginated_index(batch_key, int(st.session_state[batch_key]) - 1, 0, max_batch, scroll_anchor_id)
                        st.rerun()
                with nav_cols[1]:
                    if st.button(
                        "Next ➡",
                        key=f"pdf_next_batch_{preview_key_base}",
                        use_container_width=True,
                        disabled=batch_end >= total_pages,
                    ):
                        set_paginated_index(batch_key, int(st.session_state[batch_key]) + 1, 0, max_batch, scroll_anchor_id)
                        st.rerun()
                
                # Download sections
                if image_download_items:
                    with st.expander("🖼️ Image Downloads", expanded=False):
                        for item in image_download_items[:10]:  # Limit to first 10
                            st.download_button(
                                label=item["label"],
                                data=item["data"],
                                file_name=item["file_name"],
                                mime=item["mime"],
                                key=item["key"]
                            )
                
                if table_download_items:
                    with st.expander("📊 Table Downloads", expanded=False):
                        for item in table_download_items[:10]:
                            st.download_button(
                                label=item["label"],
                                data=item["data"],
                                file_name=item["file_name"],
                                mime=item["mime"],
                                key=item["key"]
                            )
            
            return
        
        except Exception as pdf_err:
            st.error(f"❌ PDF rendering error: {str(pdf_err)[:200]}")
            st.info("Falling back to text-based document preview...")

    # Special handling for images
    if file_name_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
        try:
            st.image(file_entry["bytes"], caption=file_name, use_container_width=True)
            png_bytes = image_bytes_to_png_bytes(file_entry["bytes"])
            ext = file_name_lower.split('.')[-1]
            mime_type = f"image/{ext}"
            if ext == "jpg":
                mime_type = "image/jpeg"
            st.download_button(
                "📥 Download Image",
                png_bytes,
                file_name=file_name,
                mime=mime_type
            )
            return
        except Exception as img_err:
            st.error(f"❌ Could not display image: {str(img_err)[:100]}")

    # Fallback: text-based preview
    try:
        file_text = st.session_state.file_texts.get(file_name, "")
        if not file_text.strip():
            file_text = extract_text(file_name, file_entry["bytes"])
        
        if file_text.strip():
            preview_length = 2000
            preview_text = file_text[:preview_length]
            if len(file_text) > preview_length:
                preview_text += f"\n\n... ({len(file_text) - preview_length} more characters)"
            
            if highlight_term:
                st.markdown(render_text_block(preview_text, highlight_term), unsafe_allow_html=True)
            else:
                st.code(preview_text, language="text")
        else:
            st.info("📄 No readable content found in document.")
    
    except Exception as fallback_err:
        st.error(f"❌ Preview error: {str(fallback_err)[:100]}")
    # Special handling for Excel files (show as table)
    if file_name_lower.endswith(".xlsx"):
        with st.spinner("Loading preview..."):
            data = st.session_state.excel_data_by_file.get(file_name, [])
            if data:
                st.markdown("### Excel Preview")
                preview_df = pd.DataFrame(data).head(20)
                st.dataframe(preview_df, use_container_width=True, hide_index=True)
                st.download_button(
                    label="Download table as CSV",
                    data=preview_df.to_csv(index=False),
                    file_name=f"{os.path.splitext(file_name)[0]}_preview.csv",
                    mime="text/csv",
                    key=f"download_excel_csv_{file_name}"
                )
                try:
                    table_png = table_to_png_bytes(
                        dataframe_to_table_rows(preview_df),
                        title=f"{file_name} Preview"
                    )
                    st.download_button(
                        label="Download table as PNG",
                        data=table_png,
                        file_name=f"{os.path.splitext(file_name)[0]}_preview.png",
                        mime="image/png",
                        key=f"download_excel_png_{file_name}"
                    )
                except Exception:
                    pass
            else:
                st.info("No preview data available for this spreadsheet.")
        return

    # For all other files, show comprehensive extracted content
    with st.spinner("Loading preview..."):
        # Ensure text extraction has run
        if file_name not in st.session_state.file_texts:
            st.session_state.file_texts[file_name] = extract_text(file_name, file_entry["bytes"])
        
        full_content = st.session_state.file_texts.get(file_name, "")
        
        if not full_content.strip():
            st.warning("No content could be extracted from this file. This might indicate an issue with the file format or content extraction.")
            # Try to show basic file info
            file_size = len(file_entry["bytes"])
            st.info(f"File size: {file_size} bytes")
            st.info(f"File type: {file_name.split('.')[-1].upper()}")
            return

        # Parse the content into sections
        sections = parse_extracted_content(full_content)
        
        # For PDF files, if parsing doesn't produce good content, show raw text
        if file_name_lower.endswith('.pdf') and sections:
            # Check if we have meaningful text content
            has_meaningful_text = any(
                section_type == "TEXT" and section_content.strip()
                for section_type, _, section_content in sections
            )
            if not has_meaningful_text:
                # Fall back to showing raw content
                st.markdown("### Document Content")
                anchor_id = create_heading_anchor(highlight_term) if highlight_term else None
                with st.expander("Show extracted text", expanded=True):
                    if highlight_term:
                        st.markdown(render_text_block(full_content.strip(), highlight_term, anchor_id=anchor_id), unsafe_allow_html=True)
                    else:
                        st.code(full_content.strip(), language="text")
                return
        
        if not sections:
            st.warning("Content was extracted but could not be parsed into displayable sections.")
            anchor_id = create_heading_anchor(highlight_term) if highlight_term else None
            if highlight_term:
                st.markdown(render_text_block(full_content[:1000] + ("..." if len(full_content) > 1000 else ""), highlight_term, anchor_id=anchor_id), unsafe_allow_html=True)
            else:
                st.code(full_content[:1000] + ("..." if len(full_content) > 1000 else ""), language="text")
            return
        
        # Display each section
        for section_type, section_title, section_content in sections:
            if section_type == "METADATA":
                with st.expander(f"📋 {section_title}", expanded=False):
                    if section_content.strip():
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
                    else:
                        st.info("No metadata available")
            elif section_type == "TEXT":
                section_anchor_id = create_heading_anchor(section_title)
                st.markdown(f"<div id='{section_anchor_id}'></div>", unsafe_allow_html=True)
                st.markdown(f"### {section_title}")
                if section_content.strip():
                    if len(section_content) > 2000:
                        with st.expander("Show text content", expanded=False):
                            if highlight_term:
                                st.markdown(render_text_block(section_content, highlight_term, anchor_id=None), unsafe_allow_html=True)
                            else:
                                st.code(section_content, language="text")
                    else:
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term, anchor_id=None), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
                else:
                    st.info("No text content available for this section.")
            elif section_type == "TABLE":
                with st.expander(f"📊 {section_title}", expanded=True):
                    # Try to parse table and display as dataframe
                    try:
                        lines = section_content.strip().split('\n')
                        if lines:
                            # Parse table data
                            table_data = []
                            for line in lines:
                                if ' | ' in line:
                                    row = [cell.strip() for cell in line.split(' | ')]
                                    table_data.append(row)
                            
                            if table_data:
                                df = pd.DataFrame(table_data[1:] if len(table_data) > 1 else table_data, 
                                                columns=table_data[0] if len(table_data) > 1 else None)
                                st.dataframe(df, use_container_width=True, hide_index=True)
                                
                                # Add download button for table as CSV
                                csv_data = df.to_csv(index=False)
                                st.download_button(
                                    label="Download as CSV",
                                    data=csv_data,
                                    file_name=f"{section_title.replace(' ', '_')}.csv",
                                    mime="text/csv",
                                    key=f"download_table_{file_name}_{section_title}"
                                )
                                try:
                                    table_png = table_to_png_bytes(table_data, title=section_title)
                                    table_download_items.append({
                                        "label": f"Download {section_title} as PNG",
                                        "data": table_png,
                                        "file_name": f"{section_title.replace(' ', '_')}.png",
                                        "mime": "image/png",
                                        "key": f"download_table_png_{file_name}_{section_title}"
                                    })
                                except Exception:
                                    pass
                            else:
                                if highlight_term:
                                    st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                                else:
                                    st.code(section_content, language="text")
                    except Exception:
                        if highlight_term:
                            st.markdown(render_text_block(section_content, highlight_term), unsafe_allow_html=True)
                        else:
                            st.code(section_content, language="text")
            elif section_type == "EMBEDDED_IMAGE":
                # Display embedded image
                image_key = section_content.split(": ")[-1] if ": " in section_content else section_content
                if 'extracted_images' in st.session_state and image_key in st.session_state.extracted_images:
                    image_data = st.session_state.extracted_images[image_key]
                    try:
                        st.image(image_data['bytes'], caption=image_data['filename'], use_container_width=True)
                        # Add download button
                        mime_type = f"image/{image_data['ext']}"
                        if image_data['ext'] == "jpg":
                            mime_type = "image/jpeg"
                        st.download_button(
                            label="Download Image",
                            data=image_data['bytes'],
                            file_name=image_data['filename'],
                            mime=mime_type,
                            key=f"download_embedded_{image_key}"
                        )
                        try:
                            png_bytes = image_bytes_to_png_bytes(image_data['bytes'])
                            st.download_button(
                                label="Download as PNG",
                                data=png_bytes,
                                file_name=f"{os.path.splitext(image_data['filename'])[0]}.png",
                                mime="image/png",
                                key=f"download_embedded_png_{image_key}"
                            )
                        except Exception:
                            pass
                    except Exception as e:
                        st.error(f"Could not display image: {e}")
                else:
                    st.info(f"🖼️ {section_content}")
            elif section_type == "IMAGE":
                st.info(f"🖼️ {section_content}")
            elif section_type == "ERROR":
                st.error(f"❌ {section_content}")
            elif section_type == "UNSUPPORTED":
                st.warning(f"⚠️ {section_content}")

        if image_download_items:
            with st.expander("🖼️ Image Downloads", expanded=False):
                for item in image_download_items:
                    st.download_button(
                        label=item["label"],
                        data=item["data"],
                        file_name=item["file_name"],
                        mime=item["mime"],
                        key=item["key"]
                    )
        if table_download_items:
            with st.expander("📊 Table Downloads", expanded=False):
                for item in table_download_items:
                    st.download_button(
                        label=item["label"],
                        data=item["data"],
                        file_name=item["file_name"],
                        mime=item["mime"],
                        key=item["key"]
                    )


def parse_extracted_content(content):
    """Parse the extracted content into sections for display."""
    sections = []
    lines = content.split('\n')
    current_section = None
    current_content = []

    def flush_section():
        nonlocal current_section, current_content
        if not current_section:
            return

        section_type, section_title, section_value = current_section
        if section_type in ("TEXT", "TABLE"):
            sections.append((section_type, section_title, '\n'.join(current_content).strip()))
        else:
            final_value = section_value
            if current_content:
                final_value = '\n'.join(current_content).strip()
            sections.append((section_type, section_title, final_value))

        current_section = None
        current_content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check for section markers
        if line.startswith('TABLE:'):
            # Save previous section
            flush_section()
            
            # Start new table section
            current_section = ("TABLE", "Table Content", "")
            current_content = []
            
        elif line.startswith('[IMAGE:'):
            # Save previous section
            flush_section()
            
            # Add image info
            sections.append(("IMAGE", "Images", line))
            current_section = None
            current_content = []
            
        elif line.startswith('[EMBEDDED_IMAGE:'):
            # Save previous section
            flush_section()
            
            # Add embedded image info
            image_info = line.replace('[EMBEDDED_IMAGE: ', '').replace(']', '')
            sections.append(("EMBEDDED_IMAGE", "Embedded Image", image_info))
            current_section = None
            current_content = []
            
        elif line.startswith('PDF Metadata:') or line.startswith('Document Metadata:') or line.startswith('Meta Tags:') or line.startswith('Title:') or 'Pages:' in line or 'Slides:' in line or 'sheets:' in line:
            # Save previous section
            flush_section()
            
            # Start metadata section
            if 'Metadata:' in line or 'Tags:' in line:
                section_title = line.split(':')[0] + " Information"
            else:
                section_title = "Document Information"
            
            current_section = ("METADATA", section_title, line)
            current_content = [line]
            
        elif line.startswith('Heading:'):
            # Save previous section
            flush_section()
            
            heading_title = line.replace('Heading:', '', 1).strip()
            current_section = ("TEXT", heading_title or "Heading", "")
            current_content = []

        elif line.startswith('Page ') and 'Text:' in line:
            # Save previous section
            flush_section()
            
            # Start new text section
            current_section = ("TEXT", f"Page {line.split()[1]} Content", "")
            current_content = []
            
        elif line.startswith('Slide ') and ':' in line:
            # Save previous section
            flush_section()
            
            # Start new slide section
            slide_num = line.split(':')[0]
            current_section = ("TEXT", f"{slide_num} Content", "")
            current_content = []
            
        elif line.startswith('Sheet ') and ':' in line:
            # Save previous section
            flush_section()
            
            # Start new sheet section
            sheet_name = line.split(':')[0].replace("'", "")
            current_section = ("TABLE", f"{sheet_name} Data", "")
            current_content = []
            
        elif current_section:
            # Add to current section
            if current_section[0] == "METADATA":
                current_content.append(line)
                current_section = (current_section[0], current_section[1], '\n'.join(current_content))
            else:
                current_content.append(line)
        else:
            # Start default text section
            if not current_section:
                current_section = ("TEXT", "Document Content", "")
                current_content = []
            current_content.append(line)
    
    # Save final section
    flush_section()
    
    # If no sections were created but we have content, create a default text section
    if not sections and content.strip():
        sections.append(("TEXT", "Document Content", content.strip()))
    
    return sections


@st.cache_data(show_spinner=False)
def build_summary_download_assets(file_name, file_bytes):
    file_name_lower = file_name.lower()
    image_items = []
    table_items = []

    try:
        if file_name_lower.endswith(".pdf"):
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for page_index, page in enumerate(pdf.pages[:PDF_ASSET_SCAN_PAGE_LIMIT], start=1):
                    for image_index, image_info in enumerate(page.images or [], start=1):
                        try:
                            bbox = (
                                image_info.get("x0", 0),
                                image_info.get("top", 0),
                                image_info.get("x1", 0),
                                image_info.get("bottom", 0)
                            )
                            if bbox[0] < bbox[2] and bbox[1] < bbox[3]:
                                image_items.append({
                                    "label": f"{file_name} - Page {page_index} Image {image_index}",
                                    "data": crop_pdf_region_to_png(page, bbox),
                                    "file_name": f"{os.path.splitext(file_name)[0]}_page_{page_index}_image_{image_index}.png",
                                    "mime": "image/png"
                                })
                        except Exception:
                            pass

                    for table_index, table in enumerate(page.extract_tables() or [], start=1):
                        if table and any(any(cell for cell in row) for row in table):
                            try:
                                table_items.append({
                                    "label": f"{file_name} - Page {page_index} Table {table_index}",
                                    "data": table_to_png_bytes(table, title=f"Page {page_index} Table {table_index}"),
                                    "file_name": f"{os.path.splitext(file_name)[0]}_page_{page_index}_table_{table_index}.png",
                                    "mime": "image/png"
                                })
                            except Exception:
                                pass

        elif file_name_lower.endswith(".docx"):
            doc = docx.Document(BytesIO(file_bytes))
            image_index = 1
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_bytes = rel.target_part.blob
                        image_items.append({
                            "label": f"{file_name} - Image {image_index}",
                            "data": image_bytes_to_png_bytes(image_bytes),
                            "file_name": f"{os.path.splitext(file_name)[0]}_image_{image_index}.png",
                            "mime": "image/png"
                        })
                        image_index += 1
                    except Exception:
                        pass

            for table_index, table in enumerate(doc.tables, start=1):
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if table_data:
                    try:
                        table_items.append({
                            "label": f"{file_name} - Table {table_index}",
                            "data": table_to_png_bytes(table_data, title=f"Table {table_index}"),
                            "file_name": f"{os.path.splitext(file_name)[0]}_table_{table_index}.png",
                            "mime": "image/png"
                        })
                    except Exception:
                        pass

        elif file_name_lower.endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            image_index = 1
            table_index = 1
            for slide_index, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_items.append({
                                "label": f"{file_name} - Slide {slide_index} Image {image_index}",
                                "data": image_bytes_to_png_bytes(shape.image.blob),
                                "file_name": f"{os.path.splitext(file_name)[0]}_slide_{slide_index}_image_{image_index}.png",
                                "mime": "image/png"
                            })
                            image_index += 1
                        except Exception:
                            pass
                    if hasattr(shape, "table"):
                        try:
                            table_data = [
                                [cell.text.strip() for cell in row.cells]
                                for row in shape.table.rows
                            ]
                            table_items.append({
                                "label": f"{file_name} - Slide {slide_index} Table {table_index}",
                                "data": table_to_png_bytes(table_data, title=f"Slide {slide_index} Table {table_index}"),
                                "file_name": f"{os.path.splitext(file_name)[0]}_slide_{slide_index}_table_{table_index}.png",
                                "mime": "image/png"
                            })
                            table_index += 1
                        except Exception:
                            pass

        elif file_name_lower.endswith(".xlsx"):
            data = extract_excel_data(file_name, file_bytes)
            if data:
                preview_df = pd.DataFrame(data).head(20)
                table_items.append({
                    "label": f"{file_name} - Preview Table",
                    "data": table_to_png_bytes(dataframe_to_table_rows(preview_df), title=f"{file_name} Preview"),
                    "file_name": f"{os.path.splitext(file_name)[0]}_preview.png",
                    "mime": "image/png"
                })

        elif file_name_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
            image_items.append({
                "label": f"{file_name} - Image",
                "data": image_bytes_to_png_bytes(file_bytes),
                "file_name": f"{os.path.splitext(file_name)[0]}.png",
                "mime": "image/png"
            })
    except Exception:
        return {"images": [], "tables": []}

    return {"images": image_items, "tables": table_items}


def render_extracted_assets_preview(file_name, file_entry):
    st.markdown(f"**Preview: {file_name}**")
    assets = build_summary_download_assets(file_name, file_entry["bytes"])
    image_items = assets.get("images", [])
    table_items = assets.get("tables", [])

    if not image_items and not table_items:
        st.info("No extractable images or tables were found in this file.")
        return

    if image_items:
        st.markdown("### Extracted Images")
        for index, item in enumerate(image_items):
            st.image(item["data"], caption=item["label"], use_container_width=True)
            st.download_button(
                label=f"Download {item['label']} as PNG",
                data=item["data"],
                file_name=item["file_name"],
                mime=item["mime"],
                key=f"preview_image_download_{index}_{item['file_name']}"
            )

    if table_items:
        st.markdown("### Extracted Tables")
        for index, item in enumerate(table_items):
            st.image(item["data"], caption=item["label"], use_container_width=True)
            st.download_button(
                label=f"Download {item['label']} as PNG",
                data=item["data"],
                file_name=item["file_name"],
                mime=item["mime"],
                key=f"preview_table_download_{index}_{item['file_name']}"
            )


def format_file_size(byte_count):
    """Return a readable file size label."""
    try:
        size = float(byte_count or 0)
        for unit in ["B", "KB", "MB", "GB"]:
            if size < 1024 or unit == "GB":
                return f"{size:.1f} {unit}" if unit != "B" else f"{int(size)} {unit}"
            size /= 1024
    except Exception:
        return "Unknown"


@st.cache_data(show_spinner=False)
def get_preview_metadata(file_name, file_bytes, extracted_text):
    """Collect lightweight metadata without rendering the whole document."""
    metadata = {
        "File name": file_name,
        "File type": os.path.splitext(file_name)[1].lower().lstrip(".").upper() or "Unknown",
        "File size": format_file_size(len(file_bytes or b"")),
        "Extracted text": f"{len(str(extracted_text or '')):,} characters",
        "Pages / Slides / Sheets": "Not available",
        "Tables": "0",
        "Images": "0",
    }

    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".pdf"):
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                metadata["Pages / Slides / Sheets"] = f"{len(pdf.pages)} pages"
        elif file_name_lower.endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            metadata["Pages / Slides / Sheets"] = f"{len(prs.slides)} slides"
        elif file_name_lower.endswith(".xlsx"):
            wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            metadata["Pages / Slides / Sheets"] = f"{len(wb.sheetnames)} sheets"
        elif file_name_lower.endswith(".docx"):
            doc = docx.Document(BytesIO(file_bytes))
            metadata["Tables"] = str(len(doc.tables))

        page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, extracted_text)
        if page_count and metadata["Pages / Slides / Sheets"] == "Not available":
            metadata["Pages / Slides / Sheets"] = str(page_count)
        metadata["Tables"] = str(max(int(metadata.get("Tables", "0") or 0), table_count))
        metadata["Images"] = str(image_count)
    except Exception:
        pass

    return metadata


def render_preview_metadata_cards(metadata):
    """Render compact metadata cards for the document viewer."""
    try:
        labels = list(metadata.items())
        cols = st.columns(3)
        for index, (label, value) in enumerate(labels):
            with cols[index % 3]:
                st.metric(label, value)
    except Exception as e:
        st.warning(f"Could not render preview metadata: {e}")


def chunk_preview_text(text, chunk_size=1200, overlap=180):
    """Split extracted text into chunks for search and Q&A."""
    try:
        clean_text = re.sub(r"\s+", " ", str(text or "")).strip()
        if not clean_text:
            return []
        chunks = []
        start = 0
        while start < len(clean_text):
            end = min(len(clean_text), start + chunk_size)
            chunk = clean_text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(clean_text):
                break
            start = max(0, end - overlap)
        return chunks
    except Exception:
        return []


@st.cache_data(show_spinner=False)
def keyword_search_preview_chunks(text, query, limit=8):
    """Fast fallback search for preview Q&A and Search tabs."""
    try:
        query_terms = [
            term.lower()
            for term in re.findall(r"[A-Za-z0-9_+\-/]{2,}", str(query or ""))
        ]
        if not query_terms:
            return []

        scored = []
        for chunk in chunk_preview_text(text):
            chunk_lower = chunk.lower()
            score = sum(chunk_lower.count(term) for term in query_terms)
            if score:
                scored.append((score, chunk))
        scored.sort(key=lambda item: item[0], reverse=True)
        return [chunk for _, chunk in scored[:limit]]
    except Exception:
        return []


def semantic_search_preview_chunks(file_name, text, query, limit=5):
    """Retrieve relevant chunks with FAISS when available, then fall back to keyword search."""
    try:
        if not str(query or "").strip() or not str(text or "").strip():
            return []
        vector_key = f"preview_vector::{file_name}"
        if vector_key not in st.session_state.vector_stores:
            st.session_state.vector_stores[vector_key] = create_vector_store(str(text)[:MAX_VECTOR_TEXT_CHARS])
        docs = st.session_state.vector_stores[vector_key].similarity_search(query, k=limit)
        return [doc.page_content for doc in docs if getattr(doc, "page_content", "")]
    except Exception:
        return keyword_search_preview_chunks(text, query, limit=limit)


def build_preview_answer(file_name, text, question):
    """Create an extractive answer from retrieved chunks."""
    try:
        chunks = semantic_search_preview_chunks(file_name, text, question, limit=5)
        if not chunks:
            return "No relevant information was found in the document text."
        answer_parts = []
        for chunk in chunks[:3]:
            sentences = re.split(r"(?<=[.!?])\s+", chunk)
            useful = [normalize_extracted_line(sentence) for sentence in sentences if len(sentence.strip()) > 30]
            answer_parts.extend(useful[:2])
        answer = "\n".join(f"- {part}" for part in answer_parts[:6])
        return answer or "\n\n".join(chunks[:2])
    except Exception as e:
        return f"Could not answer the question: {e}"


def extract_preview_tables(file_name, file_bytes, extracted_text):
    """Extract interactive tables from spreadsheets or pipe-delimited extracted text."""
    tables = []
    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".csv"):
            tables.append(("CSV Data", pd.read_csv(BytesIO(file_bytes))))
        elif file_name_lower.endswith(".xlsx"):
            workbook = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            for sheet_name, df in workbook.items():
                tables.append((sheet_name, df))
        elif file_name_lower.endswith(".xls"):
            try:
                workbook = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
                for sheet_name, df in workbook.items():
                    tables.append((sheet_name, df))
            except Exception:
                pass
        else:
            current_rows = []
            table_index = 1
            for line in str(extracted_text or "").splitlines():
                if " | " in line:
                    current_rows.append([cell.strip() for cell in line.split(" | ")])
                elif current_rows:
                    width = max(len(row) for row in current_rows)
                    normalized = [row + [""] * (width - len(row)) for row in current_rows]
                    tables.append((f"Extracted Table {table_index}", pd.DataFrame(normalized)))
                    table_index += 1
                    current_rows = []
            if current_rows:
                width = max(len(row) for row in current_rows)
                normalized = [row + [""] * (width - len(row)) for row in current_rows]
                tables.append((f"Extracted Table {table_index}", pd.DataFrame(normalized)))
    except Exception:
        pass
    return tables


def dataframe_to_xlsx_bytes(df):
    """Convert a dataframe to an XLSX download."""
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Data")
    return output.getvalue()


def ocr_image_best_effort(file_bytes):
    """Run optional OCR for image files when pytesseract is installed."""
    try:
        pytesseract = importlib.import_module("pytesseract")
        with Image.open(BytesIO(file_bytes)) as image:
            return pytesseract.image_to_string(image).strip()
    except Exception as e:
        return f"OCR is unavailable or failed: {e}"


def build_preview_summary_markdown(file_name, file_bytes, extracted_text):
    """Build a concise markdown summary for downloads."""
    try:
        plain_lines = [
            normalize_extracted_line(line)
            for line in str(extracted_text or "").splitlines()
            if 20 <= len(line.strip()) <= 220
        ]
        words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{2,}", str(extracted_text or ""))
        keyword_counts = Counter(
            word.lower()
            for word in words
            if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS
        )
        keywords = [word.title() for word, _ in keyword_counts.most_common(8)]

        sections = [
            f"# Document Intelligence Summary: {file_name}",
            "## Overview",
            f"- File type: {os.path.splitext(file_name)[1].lower() or 'unknown'}",
            f"- Extracted text length: {len(str(extracted_text or '')):,} characters",
            f"- Main themes: {', '.join(keywords) if keywords else 'Not enough readable text detected'}",
            "## Key Insights",
        ]
        key_lines = plain_lines[:5] or ["No reliable readable text was extracted."]
        sections.extend(f"- {line}" for line in key_lines)
        sections.extend([
            "## Practical Use",
            "- Use the Viewer tab for visual inspection.",
            "- Use Search and Q&A for targeted analysis over relevant passages.",
            "- Use Tables, Images, and Downloads to export reusable assets.",
        ])
        return "\n".join(sections)
    except Exception as e:
        return f"# Summary unavailable\n\n{e}"


def render_paginated_text_document_viewer(file_name, extracted_text, page_size=9000):
    """Render long text-like documents with explicit pagination or full-load mode."""
    try:
        text = str(extracted_text or "").strip()
        if not text:
            st.info("No readable text was extracted for this document.")
            return

        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
        scroll_anchor_id = f"text-viewer-top-{base_key}"
        full_key = f"text_full_view_{base_key}"
        render_scroll_anchor(scroll_anchor_id)
        st.checkbox("Load full document", key=full_key, help="Shows all extracted text. For very large files this may be slower.")

        if st.session_state.get(full_key):
            st.info("Full document mode is active. No content is truncated.")
            st.text_area("Full extracted document", value=text, height=720, key=f"text_full_area_{base_key}")
            return

        total_sections = max(1, (len(text) + page_size - 1) // page_size)
        st.info("Preview mode: showing one section at a time. Use Load full document to view everything at once.")

        nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
        page_key = f"text_section_{base_key}"
        if page_key not in st.session_state:
            st.session_state[page_key] = 1
        with nav_cols[0]:
            if st.button("Previous section", key=f"text_prev_{base_key}", use_container_width=True):
                set_paginated_index(page_key, int(st.session_state[page_key]) - 1, 1, total_sections, scroll_anchor_id)
                st.rerun()
        with nav_cols[1]:
            if st.button("Next section", key=f"text_next_{base_key}", use_container_width=True):
                set_paginated_index(page_key, int(st.session_state[page_key]) + 1, 1, total_sections, scroll_anchor_id)
                st.rerun()
        with nav_cols[2]:
            section_number = st.number_input(
                "Jump to section",
                min_value=1,
                max_value=total_sections,
                value=int(st.session_state[page_key]),
                key=f"text_jump_{base_key}",
            )
            if int(section_number) != int(st.session_state[page_key]):
                set_paginated_index(page_key, int(section_number), 1, total_sections, scroll_anchor_id)
                st.rerun()

        start = (int(st.session_state[page_key]) - 1) * page_size
        end = min(len(text), start + page_size)
        st.caption(f"Showing section {st.session_state[page_key]} of {total_sections}. Characters {start + 1:,}-{end:,} of {len(text):,}.")
        st.text_area("Document section", value=text[start:end], height=620, key=f"text_section_area_{base_key}_{st.session_state[page_key]}")
    except Exception as e:
        st.error(f"Could not render text document viewer: {e}")


def render_spreadsheet_document_viewer(file_name, file_bytes):
    """Render CSV/XLS/XLSX with sheet selection and table pagination."""
    try:
        file_name_lower = file_name.lower()
        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
        scroll_anchor_id = f"sheet-viewer-top-{base_key}"
        render_scroll_anchor(scroll_anchor_id)

        if file_name_lower.endswith(".csv"):
            sheets = {"CSV Data": pd.read_csv(BytesIO(file_bytes))}
        else:
            try:
                sheets = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            except Exception as e:
                st.warning(f"Spreadsheet preview is best-effort for this format: {e}")
                rows = extract_excel_data(file_name, file_bytes)
                sheets = {"Recovered Data": pd.DataFrame(rows)}

        if not sheets:
            st.info("No sheets or rows were found.")
            return

        sheet_name = st.selectbox("Sheet", list(sheets.keys()), key=f"sheet_select_{base_key}")
        df = sheets[sheet_name]
        total_rows = len(df)
        rows_per_page = st.selectbox("Rows per page", [25, 50, 100, 250, 500, 1000], index=1, key=f"rows_page_{base_key}")
        total_pages = max(1, (total_rows + rows_per_page - 1) // rows_per_page)

        page_key = f"sheet_page_{base_key}_{sheet_name}"
        if page_key not in st.session_state:
            st.session_state[page_key] = 1

        nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
        with nav_cols[0]:
            if st.button("Previous rows", key=f"sheet_prev_{base_key}", use_container_width=True):
                set_paginated_index(page_key, int(st.session_state[page_key]) - 1, 1, total_pages, scroll_anchor_id)
                st.rerun()
        with nav_cols[1]:
            if st.button("Next rows", key=f"sheet_next_{base_key}", use_container_width=True):
                set_paginated_index(page_key, int(st.session_state[page_key]) + 1, 1, total_pages, scroll_anchor_id)
                st.rerun()
        with nav_cols[2]:
            jump_page = st.number_input(
                "Jump to row page",
                min_value=1,
                max_value=total_pages,
                value=int(st.session_state[page_key]),
                key=f"sheet_jump_{base_key}",
            )
            if int(jump_page) != int(st.session_state[page_key]):
                set_paginated_index(page_key, int(jump_page), 1, total_pages, scroll_anchor_id)
                st.rerun()

        start = (int(st.session_state[page_key]) - 1) * rows_per_page
        end = min(total_rows, start + rows_per_page)
        st.caption(f"Sheet '{sheet_name}': showing rows {start + 1:,}-{end:,} of {total_rows:,}. All sheets are selectable above.")
        st.dataframe(df.iloc[start:end], use_container_width=True, hide_index=True)
    except Exception as e:
        st.error(f"Could not render spreadsheet viewer: {e}")


def render_presentation_document_viewer(file_name, file_bytes, extracted_text):
    """Render PPTX slides on demand, with best-effort text pagination for PPT."""
    try:
        base_key = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:12]
        scroll_anchor_id = f"slide-viewer-top-{base_key}"
        render_scroll_anchor(scroll_anchor_id)
        if file_name.lower().endswith(".pptx"):
            prs = Presentation(BytesIO(file_bytes))
            total_slides = len(prs.slides)
            if total_slides == 0:
                st.info("No slides were found.")
                return

            slide_key = f"slide_number_{base_key}"
            if slide_key not in st.session_state:
                st.session_state[slide_key] = 1

            nav_cols = st.columns([1, 1, 2], vertical_alignment="center")
            with nav_cols[0]:
                if st.button("Previous slide", key=f"slide_prev_{base_key}", use_container_width=True):
                    set_paginated_index(slide_key, int(st.session_state[slide_key]) - 1, 1, total_slides, scroll_anchor_id)
                    st.rerun()
            with nav_cols[1]:
                if st.button("Next slide", key=f"slide_next_{base_key}", use_container_width=True):
                    set_paginated_index(slide_key, int(st.session_state[slide_key]) + 1, 1, total_slides, scroll_anchor_id)
                    st.rerun()
            with nav_cols[2]:
                slide_number = st.number_input(
                    "Jump to slide",
                    min_value=1,
                    max_value=total_slides,
                    value=int(st.session_state[slide_key]),
                    key=f"slide_jump_{base_key}",
                )
                if int(slide_number) != int(st.session_state[slide_key]):
                    set_paginated_index(slide_key, int(slide_number), 1, total_slides, scroll_anchor_id)
                    st.rerun()

            slide = prs.slides[int(st.session_state[slide_key]) - 1]
            st.caption(f"Showing slide {st.session_state[slide_key]} of {total_slides}. Slides are rendered on demand.")
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
                if hasattr(shape, "table"):
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text.strip() for cell in row.cells])
                    if table_rows:
                        st.dataframe(pd.DataFrame(table_rows), use_container_width=True, hide_index=True)
            if slide_lines:
                st.markdown("### Slide Content")
                st.markdown("\n\n".join(html.escape(line) for line in slide_lines), unsafe_allow_html=True)
            else:
                st.info("This slide has no readable text. Native PPTX image rendering is not available without an external converter.")
        else:
            st.warning("Legacy PPT preview uses best-effort extracted text. Convert to PPTX for slide-level rendering.")
            render_paginated_text_document_viewer(file_name, extracted_text, page_size=7000)
    except Exception as e:
        st.error(f"Could not render presentation viewer: {e}")


def render_universal_document_viewer(file_name, file_entry, extracted_text, highlight_term=None, highlight_page=None):
    """Route the Viewer tab to a full-navigation renderer by file type."""
    try:
        file_bytes = file_entry["bytes"]
        file_name_lower = file_name.lower()
        if file_name_lower.endswith(".pdf"):
            render_document_preview(file_name, file_entry=file_entry, highlight_term=highlight_term, highlight_page=highlight_page)
        elif file_name_lower.endswith((".pptx", ".ppt")):
            render_presentation_document_viewer(file_name, file_bytes, extracted_text)
        elif file_name_lower.endswith((".xlsx", ".xls", ".csv")):
            render_spreadsheet_document_viewer(file_name, file_bytes)
        elif file_name_lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp")):
            st.image(file_bytes, caption=file_name, use_container_width=True)
        else:
            render_paginated_text_document_viewer(file_name, extracted_text)
    except Exception as e:
        st.error(f"Could not render universal document viewer: {e}")


def get_preview_extracted_text(file_name, file_bytes):
    """Extract full text only when an analysis panel actually needs it."""
    try:
        if file_name not in st.session_state.file_texts:
            with st.spinner("Extracting full document text for this analysis panel..."):
                extracted = extract_text(file_name, file_bytes)
                st.session_state.file_texts[file_name] = extracted
                FILE_TEXT_CACHE.set(file_name, extracted)
        return st.session_state.file_texts.get(file_name, "")
    except Exception as e:
        st.warning(f"Could not extract full document text: {e}")
        return ""


def render_professional_document_preview(file_name, file_entry=None, highlight_term=None, highlight_page=None):
    """Render a professional multi-tab document intelligence preview."""
    global PDF_PREVIEW_RESOLUTION

    try:
        if file_entry is None:
            file_entry = get_uploaded_file_entry(file_name)
        if not file_entry:
            st.error("File preview unavailable - file could not be loaded.")
            return

        file_bytes = file_entry["bytes"]
        file_name_lower = file_name.lower()
        reliable_formats = (".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".odt", ".rtf", ".txt", ".md", ".html", ".htm", ".png", ".jpg", ".jpeg", ".webp")
        best_effort_formats = (".doc", ".ppt", ".xls", ".pages")

        extracted_text = st.session_state.file_texts.get(file_name, "")

        metadata = get_preview_metadata(file_name, file_bytes, extracted_text)

        st.markdown(
            """
            <style>
                .preview-shell {
                    padding: 12px 0 2px;
                }
                .preview-note {
                    border: 1px solid #dbeafe;
                    background: #f8fbff;
                    border-radius: 10px;
                    padding: 12px 14px;
                    color: #173152;
                    margin: 8px 0 16px;
                }
            </style>
            """,
            unsafe_allow_html=True,
        )

        st.sidebar.markdown("### Preview Controls")
        zoom_percent = st.sidebar.slider("PDF render zoom", 80, 180, 100, 10)
        PDF_PREVIEW_RESOLUTION = max(72, int(zoom_percent))
        quick_search = st.sidebar.text_input("What ", value=highlight_term or "")
        st.sidebar.caption("Large files are rendered in pages, chunks, and cached extraction layers.")

        st.markdown(f"## {html.escape(file_name)}")
        render_preview_metadata_cards(metadata)

        if file_name_lower.endswith(best_effort_formats):
            st.warning(
                "This is a binary or proprietary format. Preview and extraction use best-effort recovery. "
                "For reliable formatting and analysis, convert to DOCX, PPTX, XLSX, PDF, ODT, RTF, TXT, or HTML."
            )
        elif not file_name_lower.endswith(reliable_formats):
            st.warning("This file type is not fully supported. The app will attempt best-effort extraction.")

        active_preview_panel = st.radio(
            "Preview panel",
            ["Viewer", "Summary", "Search", "Q&A", "Tables", "Images", "Downloads"],
            horizontal=True,
            key=f"preview_panel_{hashlib.md5(file_name.encode('utf-8')).hexdigest()[:12]}",
        )

        if active_preview_panel == "Viewer":
            st.markdown("<div class='preview-note'>Viewer mode supports full-document navigation. PDFs, slides, sheets, and long text documents are rendered on demand instead of silently truncating content.</div>", unsafe_allow_html=True)
            viewer_text = extracted_text
            if file_name_lower.endswith((".doc", ".docx", ".odt", ".rtf", ".txt", ".md", ".html", ".htm", ".pages", ".ppt")):
                viewer_text = get_preview_extracted_text(file_name, file_bytes)
            render_universal_document_viewer(
                file_name,
                file_entry,
                viewer_text,
                highlight_term=quick_search or highlight_term,
                highlight_page=highlight_page,
            )

        elif active_preview_panel == "Summary":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            summary_md = build_preview_summary_markdown(file_name, file_bytes, extracted_text)
            st.markdown(summary_md)
            st.download_button(
                "Download summary as Markdown",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.md",
                mime="text/markdown",
                key=f"preview_summary_md_{file_name}",
            )
            st.download_button(
                "Download summary as TXT",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.txt",
                mime="text/plain",
                key=f"preview_summary_txt_{file_name}",
            )

        elif active_preview_panel == "Search":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            search_query = st.text_input("Search extracted text", value=quick_search, key=f"preview_search_{file_name}")
            if search_query:
                chunks = keyword_search_preview_chunks(extracted_text, search_query, limit=20)
                st.caption(f"{len(chunks)} relevant passages found.")
                for index, chunk in enumerate(chunks, start=1):
                    with st.expander(f"Match {index}", expanded=index == 1):
                        st.write(chunk)
            else:
                st.info("Enter a search term to search extracted text chunks.")

        elif active_preview_panel == "Q&A":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            question = st.text_input("Ask a question about this document", key=f"preview_qa_{file_name}")
            if question:
                answer = build_preview_answer(file_name, extracted_text, question)
                st.markdown("### Answer")
                st.markdown(answer)
                st.download_button(
                    "Download Q&A result",
                    data=f"Question: {question}\n\nAnswer:\n{answer}".encode("utf-8"),
                    file_name=f"{os.path.splitext(file_name)[0]}_qa.txt",
                    mime="text/plain",
                    key=f"preview_qa_download_{file_name}",
                )
            else:
                st.info("Ask a focused question. The app retrieves only relevant chunks before answering.")

        elif active_preview_panel == "Tables":
            if not file_name_lower.endswith((".csv", ".xlsx", ".xls")):
                extracted_text = get_preview_extracted_text(file_name, file_bytes)
            tables = extract_preview_tables(file_name, file_bytes, extracted_text)
            if not tables:
                st.info("No tables were detected. For scanned PDFs, enable page-level table detection in the Viewer tab.")
            for index, (table_name, df) in enumerate(tables, start=1):
                with st.expander(table_name, expanded=index == 1):
                    st.dataframe(df, use_container_width=True, hide_index=True)
                    st.download_button(
                        "Download CSV",
                        data=df.to_csv(index=False).encode("utf-8"),
                        file_name=f"{os.path.splitext(file_name)[0]}_{index}.csv",
                        mime="text/csv",
                        key=f"preview_table_csv_{file_name}_{index}",
                    )
                    st.download_button(
                        "Download XLSX",
                        data=dataframe_to_xlsx_bytes(df),
                        file_name=f"{os.path.splitext(file_name)[0]}_{index}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"preview_table_xlsx_{file_name}_{index}",
                    )

        elif active_preview_panel == "Images":
            if file_name_lower.endswith((".png", ".jpg", ".jpeg", ".webp")):
                st.image(file_bytes, caption=file_name, use_container_width=True)
                if st.checkbox("Run OCR if available", key=f"preview_ocr_{file_name}"):
                    st.text_area("OCR text", value=ocr_image_best_effort(file_bytes), height=220)
            render_extracted_assets_preview(file_name, file_entry)

        elif active_preview_panel == "Downloads":
            extracted_text = get_preview_extracted_text(file_name, file_bytes)
            summary_md = build_preview_summary_markdown(file_name, file_bytes, extracted_text)
            st.download_button(
                "Summary - Markdown",
                data=summary_md.encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_summary.md",
                mime="text/markdown",
                key=f"preview_download_summary_md_{file_name}",
            )
            st.download_button(
                "Extracted text - TXT",
                data=str(extracted_text or "").encode("utf-8"),
                file_name=f"{os.path.splitext(file_name)[0]}_extracted_text.txt",
                mime="text/plain",
                key=f"preview_download_text_{file_name}",
            )
            report_bytes = None
            try:
                report_doc = docx.Document()
                report_doc.add_heading(f"Document Intelligence Report: {file_name}", level=1)
                report_doc.add_paragraph(summary_md)
                report_output = BytesIO()
                report_doc.save(report_output)
                report_bytes = report_output.getvalue()
            except Exception:
                report_bytes = None
            if report_bytes:
                st.download_button(
                    "Generated report - DOCX",
                    data=report_bytes,
                    file_name=f"{os.path.splitext(file_name)[0]}_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key=f"preview_download_report_{file_name}",
                )
            with st.expander("Recommended install commands"):
                st.code(
                    "pip install streamlit pandas openpyxl python-docx python-pptx pdfplumber beautifulsoup4 pillow plotly faiss-cpu langchain-community langchain-text-splitters sentence-transformers pytesseract",
                    language="bash",
                )
    except Exception as e:
        st.error(f"Error rendering professional document preview: {e}")


@st.cache_data(show_spinner=False)
def extract_excel_data(file_name, file_bytes):
    data = []
    bio = BytesIO(file_bytes)
    file_name_lower = file_name.lower()
    try:
        if file_name_lower.endswith(".xlsx"):
            wb = openpyxl.load_workbook(bio, data_only=True)
            for sheet in wb:
                headers = None
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i == 0:
                        headers = list(row)
                    else:
                        if row and any(cell is not None for cell in row):
                            data.append(dict(zip(headers, row)))
    except Exception:
        data = []
    return data


@st.cache_data(show_spinner=False)
def create_vector_store(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    chunks = splitter.split_text(text[:MAX_VECTOR_TEXT_CHARS])
    emb = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(chunks, emb)


@st.cache_resource(show_spinner=False)
def get_chatpdf_embeddings():
    """Shared embedding model for metadata-preserving ChatPDF retrieval."""
    return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")


def get_active_user_id():
    """Stable local user id used to isolate document vectors and memory."""
    return str(st.session_state.get("logged_in_username") or "local_user").strip() or "local_user"


def get_document_id(file_name, file_bytes):
    """Stable document id from file name + content hash."""
    digest = hashlib.sha1(file_bytes or b"").hexdigest()[:16]
    clean_name = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(file_name or "document")).strip("_")
    return f"{clean_name}_{digest}"


def get_chatpdf_file_type(file_name):
    """Normalize supported document extensions for ChatPDF metadata."""
    file_type = detect_file_type(file_name)
    return "document" if file_type == "unknown" else file_type


def build_chatpdf_citation_label(metadata):
    """Format citations in ChatPDF style for each supported file type."""
    meta = metadata or {}
    file_name = str(meta.get("file_name") or "document")
    file_type = str(meta.get("file_type") or get_chatpdf_file_type(file_name)).lower()
    locator = str(meta.get("page_or_sheet") or meta.get("page_number") or "").strip()
    section = re.sub(r"\s+", " ", str(meta.get("section") or "")).strip()

    if file_type == "pdf" and locator:
        return f"{file_name} (PDF Page {locator})"
    if file_type in {"ppt", "pptx"} and locator:
        return f"{file_name} (Slide {locator})"
    if file_type in {"excel", "xlsx"} and locator:
        return f"{file_name} (Sheet: {locator})"

    generic_sections = {
        "",
        "document",
        "text content",
        "extracted content",
        "html",
        "odt content",
        "rtf content",
    }
    if section.lower() not in generic_sections:
        return f"{file_name} (Section: {section[:90]})"
    return file_name


def get_chatpdf_collection_id(user_id, file_names):
    selection_parts = []
    for name in sorted(str(file_name) for file_name in file_names):
        file_entry = get_uploaded_file_entry(name)
        file_hash = get_file_hash(file_entry.get("bytes", b"")) if file_entry else "missing"
        selection_parts.append(f"{name}:{file_hash}")
    selection = "|".join(selection_parts)
    return hashlib.sha1(f"{CHATPDF_SCHEMA_VERSION}|{user_id}|{selection}".encode("utf-8")).hexdigest()[:24]


def get_chatpdf_vector_dir(user_id, file_names):
    collection_id = get_chatpdf_collection_id(user_id, file_names)
    safe_user = re.sub(r"[^A-Za-z0-9_.-]+", "_", user_id)
    return os.path.join(APP_DIR, "chatpdf_vectorstores", safe_user, collection_id)


def get_chatpdf_memory_key(user_id, file_names):
    collection_id = get_chatpdf_collection_id(user_id, file_names)
    return f"{user_id}:{collection_id}"


def init_chatpdf_memory():
    if "document_chat_memory" not in st.session_state or not isinstance(st.session_state.document_chat_memory, dict):
        st.session_state.document_chat_memory = {}


def get_chatpdf_memory(user_id, file_names):
    init_chatpdf_memory()
    key = get_chatpdf_memory_key(user_id, file_names)
    return st.session_state.document_chat_memory.setdefault(key, [])


def append_chatpdf_memory(user_id, file_names, question, answer):
    memory = get_chatpdf_memory(user_id, file_names)
    memory.append({"question": str(question), "answer": str(answer)})
    st.session_state.document_chat_memory[get_chatpdf_memory_key(user_id, file_names)] = memory[-10:]


def init_file_brain_registry():
    """Session-local global memory registry for uploaded file brains."""
    if "file_brains" not in st.session_state or not isinstance(st.session_state.file_brains, dict):
        st.session_state.file_brains = {}
    if "global_memory_registry" not in st.session_state or not isinstance(st.session_state.global_memory_registry, dict):
        st.session_state.global_memory_registry = {"files": {}}
    st.session_state.global_memory_registry.setdefault("files", {})
    return st.session_state.global_memory_registry


def normalize_brain_cell(value):
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return ""
    if isinstance(value, (datetime,)):
        return value.isoformat()
    return str(value)


def summarize_table_for_index(table):
    headers = " | ".join(str(header) for header in table.get("headers", [])[:20])
    sample_rows = []
    for row in table.get("rows", [])[:8]:
        sample_rows.append(" | ".join(str(cell) for cell in row[:20]))
    table_text = "\n".join([headers] + sample_rows).strip()
    locator = table.get("sheet") or table.get("page") or table.get("table") or "Table"
    return clean_text(f"{locator}\n{table_text}")


def extract_pdf_pages(file):
    """Extract PDF pages one at a time while preserving basic block structure."""
    file_bytes = _read_uploaded_bytes(file)
    pages = []
    diagrams = []
    if not file_bytes:
        return pages, diagrams

    try:
        fitz_module = importlib.import_module("fitz")
    except Exception:
        fitz_module = None

    if fitz_module is not None:
        pdf_doc = None
        try:
            pdf_doc = fitz_module.open(stream=file_bytes, filetype="pdf")
            for index, page in enumerate(pdf_doc, start=1):
                text = clean_text(page.get_text("text"))
                blocks = page.get_text("blocks") or []
                page_record = {"page": index, "text": text, "blocks": blocks}
                pages.append(page_record)
                if any(keyword in text.lower() for keyword in ["figure", "diagram", "block", "flow", "layout"]):
                    diagrams.append({
                        "page": index,
                        "text": text[:500],
                        "kind": "pdf-text-reference",
                    })
            return pages, diagrams
        except Exception:
            pages = []
            diagrams = []
        finally:
            try:
                if pdf_doc is not None:
                    pdf_doc.close()
            except Exception:
                pass

    try:
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = clean_text(page.extract_text() or "")
                pages.append({"page": index, "text": text, "blocks": []})
                if any(keyword in text.lower() for keyword in ["figure", "diagram", "block", "flow", "layout"]):
                    diagrams.append({
                        "page": index,
                        "text": text[:500],
                        "kind": "pdf-text-reference",
                    })
    except Exception as exc:
        pages.append({"page": 1, "text": f"PDF extraction failed: {str(exc)[:160]}", "blocks": []})
    return pages, diagrams


def extract_word_pages_and_tables(file, filename=None):
    file_bytes = _read_uploaded_bytes(file)
    pages = []
    tables = []
    ext = os.path.splitext(str(filename or ""))[1].lower()
    if ext and ext != ".docx":
        text = extract_word(BytesIO(file_bytes), filename)
        return [{"page": 1, "text": text, "blocks": []}], tables

    try:
        document = docx.Document(BytesIO(file_bytes))
        section = "Document"
        buffer = []
        page_number = 1
        for paragraph in document.paragraphs:
            paragraph_text = paragraph.text.strip()
            if not paragraph_text:
                continue
            style_name = str(getattr(paragraph.style, "name", "") or "")
            if style_name.lower().startswith("heading") and buffer:
                pages.append({"page": page_number, "section": section, "text": clean_text("\n".join(buffer)), "blocks": []})
                buffer = []
                page_number += 1
                section = paragraph_text[:120]
            buffer.append(paragraph_text)
            if len("\n".join(buffer)) >= FILE_BRAIN_PAGE_CONTEXT_CHARS:
                pages.append({"page": page_number, "section": section, "text": clean_text("\n".join(buffer)), "blocks": []})
                buffer = []
                page_number += 1
        if buffer:
            pages.append({"page": page_number, "section": section, "text": clean_text("\n".join(buffer)), "blocks": []})

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                rows.append([clean_text(cell.text) for cell in row.cells])
            if rows:
                headers = rows[0]
                body = rows[1:FILE_BRAIN_MAX_TABLE_ROWS + 1]
                tables.append({
                    "table": f"Table {table_index}",
                    "headers": headers,
                    "rows": body,
                    "row_count": max(len(rows) - 1, 0),
                    "truncated": len(rows) - 1 > FILE_BRAIN_MAX_TABLE_ROWS,
                })
        return pages or [{"page": 1, "text": extract_word(BytesIO(file_bytes), filename), "blocks": []}], tables
    except Exception:
        return [{"page": 1, "text": extract_word(BytesIO(file_bytes), filename), "blocks": []}], tables


def extract_ppt_pages(file, filename=None):
    file_bytes = _read_uploaded_bytes(file)
    ext = os.path.splitext(str(filename or ""))[1].lower()
    if ext == ".ppt":
        return [{"page": 1, "text": extract_ppt(BytesIO(file_bytes), filename), "blocks": []}], []

    pages = []
    tables = []
    try:
        presentation = Presentation(BytesIO(file_bytes))
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_lines.append(shape.text)
                if hasattr(shape, "table"):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([clean_text(cell.text) for cell in row.cells])
                    if rows:
                        tables.append({
                            "sheet": f"Slide {slide_index}",
                            "table": f"Slide {slide_index} table",
                            "headers": rows[0],
                            "rows": rows[1:FILE_BRAIN_MAX_TABLE_ROWS + 1],
                            "row_count": max(len(rows) - 1, 0),
                            "truncated": len(rows) - 1 > FILE_BRAIN_MAX_TABLE_ROWS,
                        })
            pages.append({"page": slide_index, "text": clean_text("\n".join(slide_lines)), "blocks": []})
    except Exception:
        pages.append({"page": 1, "text": extract_ppt(BytesIO(file_bytes), filename), "blocks": []})
    return pages, tables


def extract_tables_excel(file, filename=None):
    """Extract spreadsheet tables as structured rows instead of flattened text."""
    file_bytes = _read_uploaded_bytes(file)
    tables = []
    try:
        sheets = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
        for name, df in sheets.items():
            normalized_df = df.fillna("")
            rows = [
                [normalize_brain_cell(cell) for cell in row]
                for row in normalized_df.values.tolist()[:FILE_BRAIN_MAX_TABLE_ROWS]
            ]
            tables.append({
                "sheet": str(name),
                "headers": [normalize_brain_cell(header) for header in list(normalized_df.columns)],
                "rows": rows,
                "row_count": len(normalized_df.index),
                "truncated": len(normalized_df.index) > FILE_BRAIN_MAX_TABLE_ROWS,
            })
    except Exception:
        if str(filename or "").lower().endswith(".xls"):
            text = extract_excel(BytesIO(file_bytes), filename)
            tables.append({"sheet": "Workbook", "headers": [], "rows": [[text]], "row_count": 1, "truncated": False})
    return tables


def extract_csv_tables(file):
    file_bytes = _read_uploaded_bytes(file)
    try:
        df = pd.read_csv(BytesIO(file_bytes)).fillna("")
        return [{
            "sheet": "CSV",
            "headers": [normalize_brain_cell(header) for header in list(df.columns)],
            "rows": [
                [normalize_brain_cell(cell) for cell in row]
                for row in df.values.tolist()[:FILE_BRAIN_MAX_TABLE_ROWS]
            ],
            "row_count": len(df.index),
            "truncated": len(df.index) > FILE_BRAIN_MAX_TABLE_ROWS,
        }]
    except Exception:
        text = extract_csv(BytesIO(file_bytes))
        return [{"sheet": "CSV", "headers": [], "rows": [[text]], "row_count": 1, "truncated": False}]


def extract_generic_pages(file, filename=None):
    text = extract_text(BytesIO(_read_uploaded_bytes(file)), filename)
    chunks = smart_chunk(text, chunk_size=FILE_BRAIN_PAGE_CONTEXT_CHARS)
    if not chunks:
        chunks = [text]
    return [{"page": index, "text": clean_text(chunk), "blocks": []} for index, chunk in enumerate(chunks, start=1)]


def extract_file_structure(file_name, file_bytes):
    file_type = detect_file_type(file_name)
    bio = BytesIO(file_bytes or b"")
    if file_type == "pdf":
        pages, diagrams = extract_pdf_pages(bio)
        return pages, [], diagrams
    if file_type == "word":
        pages, tables = extract_word_pages_and_tables(bio, file_name)
        return pages, tables, []
    if file_type == "ppt":
        pages, tables = extract_ppt_pages(bio, file_name)
        diagrams = [
            {"page": p.get("page"), "text": p.get("text", "")[:500], "kind": "slide-visual-reference"}
            for p in pages
            if any(keyword in str(p.get("text", "")).lower() for keyword in ["figure", "diagram", "flow", "block", "architecture"])
        ]
        return pages, tables, diagrams
    if file_type == "excel":
        tables = extract_tables_excel(bio, file_name)
        pages = [
            {"page": table.get("sheet") or index, "section": table.get("sheet") or "Sheet", "text": summarize_table_for_index(table), "blocks": []}
            for index, table in enumerate(tables, start=1)
        ]
        return pages, tables, []
    if file_type == "csv":
        tables = extract_csv_tables(bio)
        pages = [{"page": 1, "section": "CSV", "text": summarize_table_for_index(tables[0]) if tables else "", "blocks": []}]
        return pages, tables, []
    pages = extract_generic_pages(bio, file_name)
    diagrams = [
        {"page": p.get("page"), "text": p.get("text", "")[:500], "kind": "text-reference"}
        for p in pages
        if any(keyword in str(p.get("text", "")).lower() for keyword in ["figure", "diagram", "flow", "block", "architecture"])
    ]
    return pages, [], diagrams


def extract_lightweight_entities(text):
    entities = set()
    for token in re.findall(r"\b[A-Za-z][A-Za-z0-9_-]{2,}\b", str(text or "")):
        if token.isupper() or re.search(r"[A-Z]{2,}|\d", token):
            entities.add(token[:80])
    return entities


def extract_lightweight_facts(text, page, limit=8):
    fact_markers = [
        "supports", "provides", "requires", "shall", "must", "can", "enables",
        "contains", "includes", "used for", "configured", "connects", "consists",
    ]
    facts = []
    sentences = re.split(r"(?<=[.!?])\s+|\n+", str(text or ""))
    for sentence in sentences:
        clean_sentence = clean_text(sentence)
        if len(clean_sentence) < 35:
            continue
        if any(marker in clean_sentence.lower() for marker in fact_markers):
            facts.append({"page": page, "fact": clean_sentence[:350]})
        if len(facts) >= limit:
            break
    return facts


def build_file_brain(file_name, file_bytes):
    """Build a mini knowledge system for one uploaded file."""
    pages, tables, diagrams = extract_file_structure(file_name, file_bytes)
    brain = {
        "schema": FILE_BRAIN_SCHEMA_VERSION,
        "file_name": file_name,
        "file_type": detect_file_type(file_name),
        "file_hash": get_file_hash(file_bytes or b""),
        "pages": pages,
        "page_index": [],
        "facts": [],
        "entities": set(),
        "tables": tables,
        "diagrams": diagrams,
    }

    for page_record in pages:
        page_text = clean_text(page_record.get("text", ""))
        page_id = page_record.get("page")
        token_counts = Counter(tokenize_chatpdf_text(page_text))
        keywords = [word for word, _ in token_counts.most_common(30)]
        brain["page_index"].append({
            "page": page_id,
            "section": page_record.get("section", f"Page {page_id}"),
            "preview": page_text[:FILE_BRAIN_PREVIEW_CHARS],
            "keywords": keywords,
            "char_count": len(page_text),
        })
        brain["entities"].update(extract_lightweight_entities(page_text))
        brain["facts"].extend(extract_lightweight_facts(page_text, page_id))

    for table_index, table in enumerate(tables, start=1):
        table_text = summarize_table_for_index(table)
        brain["entities"].update(extract_lightweight_entities(table_text))
        if table_text:
            brain["facts"].append({
                "page": table.get("sheet") or table.get("page") or f"Table {table_index}",
                "fact": f"Structured table available: {table_text[:300]}",
            })

    brain["entities"] = sorted(brain["entities"])[:500]
    brain["semantic_metadata"] = build_semantic_metadata(
        file_name=file_name,
        file_type=brain["file_type"],
        pages=pages,
        tables=tables,
        diagrams=diagrams,
        entities=brain["entities"],
    )
    return brain


def register_file(file_id, brain):
    registry = init_file_brain_registry()
    registry["files"][file_id] = brain
    return brain


def file_brain_has_current_semantics(brain):
    semantic = brain.get("semantic_metadata", {}) if isinstance(brain, dict) else {}
    if not isinstance(semantic, dict):
        return False
    required_keys = {
        "executive_summary",
        "technical_summary",
        "key_concepts",
        "architecture_components",
        "semantic_relationships",
        "section_summaries",
    }
    return required_keys.issubset(set(semantic.keys()))


def ensure_file_brain(file_name):
    init_file_brain_registry()
    file_entry = get_uploaded_file_entry(file_name)
    if not file_entry:
        return None
    file_bytes = file_entry.get("bytes", b"")
    file_hash = get_file_hash(file_bytes)
    cached = st.session_state.file_brains.get(file_name)
    if (
        isinstance(cached, dict)
        and cached.get("hash") == file_hash
        and isinstance(cached.get("brain"), dict)
        and file_brain_has_current_semantics(cached["brain"])
    ):
        register_file(file_name, cached["brain"])
        return cached["brain"]

    brain = build_file_brain(file_name, file_bytes)
    st.session_state.file_brains[file_name] = {"hash": file_hash, "brain": brain}
    register_file(file_name, brain)
    return brain


def get_file_brains(file_names):
    brains = {}
    for file_name in file_names or []:
        brain = ensure_file_brain(file_name)
        if brain:
            brains[file_name] = brain
    return brains


def retrieve_pages(query, brain, top_k=8):
    """BM25 + keyword page retrieval over the page index, not full file text."""
    query_tokens = set(tokenize_chatpdf_text(query))
    if not query_tokens or not brain:
        return []

    bm25_scores = []
    bm25_okapi_class = get_bm25_okapi_class()
    page_index = brain.get("page_index", [])
    if bm25_okapi_class is not None and page_index:
        try:
            corpus_tokens = [
                tokenize_chatpdf_text(
                    " ".join([
                        str(page_record.get("section", "")),
                        str(page_record.get("preview", "")),
                        " ".join(str(keyword) for keyword in page_record.get("keywords", [])),
                    ])
                )
                for page_record in page_index
            ]
            bm25 = bm25_okapi_class(corpus_tokens)
            bm25_scores = list(bm25.get_scores(list(query_tokens)))
        except Exception:
            bm25_scores = []

    scored = []
    for index, page_record in enumerate(page_index):
        preview = str(page_record.get("preview", "")).lower()
        section = str(page_record.get("section", "")).lower()
        keywords = set(page_record.get("keywords", []))
        score = len(query_tokens.intersection(keywords)) * 3
        score += sum(1 for token in query_tokens if token in preview)
        score += sum(1 for token in query_tokens if token in section)
        if bm25_scores:
            score += float(bm25_scores[index])
        if str(query or "").lower().strip() and str(query or "").lower().strip() in preview:
            score += 5
        if score:
            scored.append((score, page_record))
    scored.sort(key=lambda item: item[0], reverse=True)
    return [page for _, page in scored[:top_k]]


def search_tables(query, tables, top_k=5):
    query_tokens = set(tokenize_chatpdf_text(query))
    if not query_tokens:
        return []
    results = []
    for table in tables or []:
        headers_text = " ".join(str(header) for header in table.get("headers", []))
        rows_text = " ".join(" ".join(str(cell) for cell in row) for row in table.get("rows", [])[:120])
        table_text = f"{table.get('sheet', '')} {headers_text} {rows_text}".lower()
        score = sum(1 for token in query_tokens if token in table_text)
        if any(term in str(query).lower() for term in ["table", "sheet", "row", "column", "csv", "excel"]):
            score += 1 if score else 0
        if score:
            results.append((score, table))
    results.sort(key=lambda item: item[0], reverse=True)
    return [table for _, table in results[:top_k]]


def search_diagrams(query, diagrams, top_k=5):
    query_tokens = set(tokenize_chatpdf_text(query))
    diagram_terms = {"diagram", "figure", "image", "flow", "block", "architecture", "layout", "pin", "connector"}
    wants_diagram = bool(query_tokens.intersection(diagram_terms))
    results = []
    for diagram in diagrams or []:
        diagram_text = str(diagram.get("text", "")).lower()
        score = sum(1 for token in query_tokens if token in diagram_text)
        if wants_diagram:
            score += 1
        if score:
            results.append((score, diagram))
    results.sort(key=lambda item: item[0], reverse=True)
    return [diagram for _, diagram in results[:top_k]]


def cross_file_search(query, file_names=None, top_k=5):
    registry = init_file_brain_registry()
    query_tokens = set(tokenize_chatpdf_text(query))
    if not query_tokens:
        return []
    allowed = set(file_names or registry.get("files", {}).keys())
    results = []
    for file_id, brain in registry.get("files", {}).items():
        if file_id not in allowed:
            continue
        score = 0
        for fact in brain.get("facts", []):
            fact_text = str(fact.get("fact", "")).lower()
            if any(token in fact_text for token in query_tokens):
                score += 2
        for entity in brain.get("entities", []):
            if str(entity).lower() in query_tokens:
                score += 3
        if score:
            results.append((score, file_id))
    results.sort(key=lambda item: item[0], reverse=True)
    return results[:top_k]


def get_page_record_by_id(brain, page_id):
    for page_record in brain.get("pages", []):
        if str(page_record.get("page")) == str(page_id):
            return page_record
    return None


def table_to_context_text(table, max_rows=30):
    headers = [str(header) for header in table.get("headers", [])]
    lines = []
    if headers:
        lines.append(" | ".join(headers))
    for row in table.get("rows", [])[:max_rows]:
        lines.append(" | ".join(str(cell) for cell in row))
    if table.get("truncated"):
        lines.append(f"... table truncated after {max_rows} displayed rows; total rows indexed: {table.get('row_count')}")
    return clean_text("\n".join(lines))


def chunk_selected_documents(docs, chunk_chars=FILE_BRAIN_SELECTED_CHUNK_CHARS):
    """Chunk only the retrieved candidates before optional lazy embedding."""
    selected_chunks = []
    for doc in docs or []:
        content = clean_text(getattr(doc, "page_content", ""))
        if not content:
            continue
        if len(content) <= chunk_chars:
            selected_chunks.append(doc)
            continue

        chunks = smart_chunk(content, chunk_size=chunk_chars)
        for chunk_index, chunk in enumerate(chunks, start=1):
            metadata = dict(getattr(doc, "metadata", {}) or {})
            metadata["selected_chunk_index"] = chunk_index
            selected_chunks.append(Document(page_content=chunk, metadata=metadata))
    return selected_chunks


def selected_content_cache_key(question, docs):
    signature_parts = [str(question or "")]
    for doc in docs or []:
        meta = getattr(doc, "metadata", {}) or {}
        signature_parts.append(
            "|".join([
                str(meta.get("file_name", "")),
                str(meta.get("page_or_sheet") or meta.get("page_number") or ""),
                str(meta.get("section", "")),
                hashlib.sha1(str(getattr(doc, "page_content", "")).encode("utf-8", errors="ignore")).hexdigest()[:12],
            ])
        )
    return hashlib.sha1("\n".join(signature_parts).encode("utf-8", errors="ignore")).hexdigest()[:24]


def lazy_embed_selected_documents(question, docs, top_k=8):
    """Embed only already-retrieved page/table/diagram candidates for final ordering."""
    docs = chunk_selected_documents([doc for doc in (docs or []) if getattr(doc, "page_content", "")])
    if len(docs) <= 1:
        return docs[:top_k]

    cache = st.session_state.setdefault("selected_chunk_vector_cache", {})
    cache_key = selected_content_cache_key(question, docs)
    vector_store = cache.get(cache_key)
    if vector_store is None:
        try:
            vector_store = FAISS.from_documents(docs, get_chatpdf_embeddings())
            cache[cache_key] = vector_store
            # Keep the lazy selected-content cache small across reruns.
            if len(cache) > 12:
                for old_key in list(cache.keys())[:-12]:
                    cache.pop(old_key, None)
            st.session_state.selected_chunk_vector_cache = cache
        except Exception:
            return rerank_chatpdf_documents(question, docs, top_k=top_k)

    try:
        embedded_docs = vector_store.similarity_search(str(question or ""), k=min(top_k, len(docs)))
        if embedded_docs:
            return embedded_docs
    except Exception:
        pass
    return rerank_chatpdf_documents(question, docs, top_k=top_k)


def retrieve_file_brain_documents(question, file_names, user_id=None, top_k=8):
    """Retrieve only selected page/table/diagram content, then convert to Documents."""
    del user_id  # User isolation is already handled by session state and uploaded file selection.
    brains = get_file_brains(file_names)
    if not brains:
        return []

    selected_documents = []
    per_file_page_k = max(2, min(top_k, math.ceil(top_k / max(len(brains), 1)) + 2))
    for file_name, brain in brains.items():
        file_type = brain.get("file_type") or get_chatpdf_file_type(file_name)
        for page_index_record in retrieve_pages(question, brain, top_k=per_file_page_k):
            page_record = get_page_record_by_id(brain, page_index_record.get("page"))
            if not page_record:
                continue
            page_text = clean_text(page_record.get("text", ""))[:FILE_BRAIN_PAGE_CONTEXT_CHARS]
            if not page_text:
                continue
            selected_documents.append(Document(
                page_content=page_text,
                metadata={
                    "file_name": file_name,
                    "file_type": file_type,
                    "page_number": str(page_record.get("page")),
                    "page_or_sheet": str(page_record.get("page")),
                    "document_id": get_document_id(file_name, get_uploaded_file_entry(file_name).get("bytes", b"")),
                    "section": page_index_record.get("section") or f"Page {page_record.get('page')}",
                    "retrieval_layer": "file_brain_page",
                },
            ))

        for table in search_tables(question, brain.get("tables", []), top_k=3):
            table_text = table_to_context_text(table)
            if not table_text:
                continue
            locator = table.get("sheet") or table.get("table") or "Table"
            selected_documents.append(Document(
                page_content=table_text,
                metadata={
                    "file_name": file_name,
                    "file_type": "excel" if file_type == "excel" else file_type,
                    "page_number": str(locator),
                    "page_or_sheet": str(locator),
                    "document_id": get_document_id(file_name, get_uploaded_file_entry(file_name).get("bytes", b"")),
                    "section": f"Structured table: {locator}",
                    "retrieval_layer": "file_brain_table",
                },
            ))

        for diagram in search_diagrams(question, brain.get("diagrams", []), top_k=2):
            diagram_text = clean_text(diagram.get("text", ""))[:1200]
            if not diagram_text:
                continue
            locator = diagram.get("page") or "Diagram"
            selected_documents.append(Document(
                page_content=diagram_text,
                metadata={
                    "file_name": file_name,
                    "file_type": file_type,
                    "page_number": str(locator),
                    "page_or_sheet": str(locator),
                    "document_id": get_document_id(file_name, get_uploaded_file_entry(file_name).get("bytes", b"")),
                    "section": f"Diagram metadata: {diagram.get('kind', 'visual reference')}",
                    "retrieval_layer": "file_brain_diagram",
                },
            ))

    if not selected_documents:
        for _, file_name in cross_file_search(question, file_names=file_names, top_k=top_k):
            brain = brains.get(file_name)
            if not brain:
                continue
            for fact in brain.get("facts", [])[:2]:
                selected_documents.append(Document(
                    page_content=str(fact.get("fact", "")),
                    metadata={
                        "file_name": file_name,
                        "file_type": brain.get("file_type") or get_chatpdf_file_type(file_name),
                        "page_number": str(fact.get("page") or "Memory"),
                        "page_or_sheet": str(fact.get("page") or "Memory"),
                        "document_id": get_document_id(file_name, get_uploaded_file_entry(file_name).get("bytes", b"")),
                        "section": "File brain fact",
                        "retrieval_layer": "file_brain_fact",
                    },
                ))

    selected_documents = rerank_chatpdf_documents(question, selected_documents, top_k=max(top_k * 2, top_k))
    return lazy_embed_selected_documents(question, selected_documents, top_k=top_k)


def retrieve_document_understanding_documents(question, file_names, user_id=None, top_k=10):
    """Retrieve document-level understanding notes for broad requests."""
    del user_id
    brains = get_file_brains(file_names)
    documents = []
    for file_name, brain in brains.items():
        file_entry = get_uploaded_file_entry(file_name)
        file_bytes = file_entry.get("bytes", b"") if file_entry else b""
        document_id = get_document_id(file_name, file_bytes)
        semantic = brain.get("semantic_metadata", {}) or {}
        file_type = brain.get("file_type") or get_chatpdf_file_type(file_name)

        summary = semantic.get("document_summary", "")
        if summary:
            documents.append(Document(
                page_content=summary,
                metadata={
                    "file_name": file_name,
                    "file_type": file_type,
                    "page_number": "Document",
                    "page_or_sheet": "Document",
                    "document_id": document_id,
                    "section": "Document summary",
                    "retrieval_layer": "semantic_document_summary",
                },
            ))

        understanding_notes = []
        concepts = [
            str(item) for item in semantic.get("key_concepts", [])[:8]
            if str(item).strip() and not _is_low_signal_theme(item)
        ]
        components = [
            str(item) for item in semantic.get("architecture_components", [])[:8]
            if str(item).strip() and not _is_low_signal_theme(item)
        ]
        if concepts:
            understanding_notes.append(
                "Important ideas to connect in the answer: " + _human_join(concepts, limit=5) + "."
            )
        if components:
            understanding_notes.append(
                "Important elements that may shape the explanation: " + _human_join(components, limit=5) + "."
            )
        if brain.get("tables"):
            understanding_notes.append("The document includes structured data that may support table-oriented questions.")
        if brain.get("diagrams"):
            understanding_notes.append("The document includes visual references that may support diagram-oriented questions.")
        if understanding_notes:
            documents.append(Document(
                page_content=clean_text(" ".join(understanding_notes)),
                metadata={
                    "file_name": file_name,
                    "file_type": file_type,
                    "page_number": "Document",
                    "page_or_sheet": "Document",
                    "document_id": document_id,
                    "section": "Document understanding",
                    "retrieval_layer": "document_understanding",
                },
            ))

        for section_summary in semantic.get("section_summaries", [])[:max(top_k, 8)]:
            summary_text = section_summary.get("summary", "")
            if not summary_text:
                continue
            documents.append(Document(
                page_content=summary_text,
                metadata={
                    "file_name": file_name,
                    "file_type": file_type,
                    "page_number": str(section_summary.get("page") or "Section"),
                    "page_or_sheet": str(section_summary.get("page") or "Section"),
                    "document_id": document_id,
                    "section": section_summary.get("section") or "Section summary",
                    "retrieval_layer": "semantic_section_summary",
                },
            ))

    return rerank_chatpdf_documents(question, documents, top_k=top_k)


def extract_chatpdf_structured_pages(file_name, file_bytes):
    """Extract page/slide/sheet-aware text records with traceable metadata."""
    records = []
    bio = BytesIO(file_bytes)
    file_lower = str(file_name).lower()
    file_type = get_chatpdf_file_type(file_name)
    document_id = get_document_id(file_name, file_bytes)

    def add_record(text, page_number="1", section="Document", page_or_sheet=None):
        clean_text = re.sub(r"\s+", " ", str(text or "")).strip()
        if not clean_text:
            return
        locator = str(page_or_sheet if page_or_sheet is not None else page_number)
        records.append({
            "text": clean_text,
            "metadata": {
                "file_name": file_name,
                "file_type": file_type,
                "page_number": str(page_number),
                "page_or_sheet": locator,
                "document_id": document_id,
                "section": section or "Document",
            },
        })

    try:
        if file_lower.endswith(".pdf"):
            with pdfplumber.open(bio) as pdf:
                for page_index, page in enumerate(pdf.pages, start=1):
                    add_record(page.extract_text() or "", page_index, f"Page {page_index}", page_index)

        elif file_lower.endswith(".docx"):
            document = docx.Document(bio)
            section = "Document"
            paragraph_buffer = []
            page_number = 1
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                if paragraph.style and str(paragraph.style.name).lower().startswith("heading"):
                    if paragraph_buffer:
                        add_record("\n".join(paragraph_buffer), page_number, section, section)
                        paragraph_buffer = []
                    section = text[:120]
                paragraph_buffer.append(text)
                if len(" ".join(paragraph_buffer)) > 3500:
                    add_record("\n".join(paragraph_buffer), page_number, section, section)
                    paragraph_buffer = []
                    page_number += 1
            if paragraph_buffer:
                add_record("\n".join(paragraph_buffer), page_number, section, section)
            for table_index, table in enumerate(document.tables, start=1):
                rows = [" | ".join(str(cell.text).strip() for cell in row.cells) for row in table.rows]
                add_record("\n".join(rows), f"T{table_index}", f"Table {table_index}", f"Table {table_index}")

        elif file_lower.endswith(".pptx"):
            presentation = Presentation(bio)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_lines.append(shape.text)
                add_record("\n".join(slide_lines), slide_index, f"Slide {slide_index}", slide_index)

        elif file_lower.endswith(".xlsx"):
            workbook = openpyxl.load_workbook(bio, data_only=True, read_only=True)
            for sheet in workbook.worksheets:
                rows = []
                for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        rows.append(" | ".join(values))
                    if len(rows) >= 120:
                        add_record("\n".join(rows), sheet.title, f"Sheet {sheet.title}", sheet.title)
                        rows = []
                if rows:
                    add_record("\n".join(rows), sheet.title, f"Sheet {sheet.title}", sheet.title)

        elif file_lower.endswith((".html", ".htm")):
            soup = BeautifulSoup(bio.read().decode("utf-8", errors="ignore"), "html.parser")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            section = soup.title.get_text(" ", strip=True) if soup.title else "HTML"
            add_record(soup.get_text("\n", strip=True), 1, section, section)

        elif file_lower.endswith(".odt"):
            with zipfile.ZipFile(bio) as odt_zip:
                xml_text = odt_zip.read("content.xml").decode("utf-8", errors="ignore")
            root = ET.fromstring(xml_text)
            text_nodes = [node.text for node in root.iter() if node.text and node.text.strip()]
            add_record("\n".join(text_nodes), 1, "ODT content", "ODT content")

        elif file_lower.endswith(".rtf"):
            raw = bio.read().decode("utf-8", errors="ignore")
            plain = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
            plain = re.sub(r"\\[a-zA-Z]+\d* ?", " ", plain)
            plain = re.sub(r"[{}]", " ", plain)
            add_record(plain, 1, "RTF content", "RTF content")

        elif file_lower.endswith((".txt", ".md", ".log", ".can")):
            add_record(bio.read().decode("utf-8", errors="ignore"), 1, "Text content", "Text content")

        else:
            fallback_text = extract_text(file_name, file_bytes)
            for match in re.finditer(r"Page\s+(\d+)\s+Text:\s*(.*?)(?=Page\s+\d+\s+Text:|\Z)", fallback_text, re.DOTALL | re.IGNORECASE):
                add_record(match.group(2), match.group(1), f"Page {match.group(1)}", match.group(1))
            if not records:
                add_record(fallback_text, 1, "Extracted content", "Extracted content")

    except Exception as exc:
        fallback_text = st.session_state.get("file_texts", {}).get(file_name) or extract_text(file_name, file_bytes)
        add_record(fallback_text, 1, f"Fallback extraction: {str(exc)[:80]}", "Extracted content")

    return records


def build_chatpdf_documents(file_names):
    """Create metadata-preserving LangChain documents for selected files."""
    documents = []
    for file_name in file_names:
        file_entry = get_uploaded_file_entry(file_name)
        if not file_entry:
            continue
        file_bytes = file_entry.get("bytes", b"")
        document_id = get_document_id(file_name, file_bytes)
        for document in process_file(BytesIO(file_bytes), file_name):
            metadata = dict(getattr(document, "metadata", {}) or {})
            metadata["document_id"] = document_id
            document.metadata = metadata
            if len(document.page_content) > MAX_VECTOR_TEXT_CHARS:
                document.page_content = document.page_content[:MAX_VECTOR_TEXT_CHARS]
            documents.append(document)
    return documents


def get_chatpdf_documents_for_selection(file_names, user_id=None):
    """Session-cached structured chunks for hybrid retrieval."""
    user_id = user_id or get_active_user_id()
    cache = st.session_state.setdefault("chatpdf_documents", {})
    cache_key = get_chatpdf_collection_id(user_id, file_names)
    if cache_key not in cache:
        cache[cache_key] = build_chatpdf_documents(file_names)
        st.session_state.chatpdf_documents = cache
    return cache.get(cache_key, [])


def tokenize_chatpdf_text(text):
    """Tokenize text for BM25 and fallback reranking."""
    tokens = re.findall(r"[A-Za-z0-9_]+", str(text or "").lower())
    return [token for token in tokens if len(token) > 1 and token not in SUMMARY_STOPWORDS]


def get_bm25_okapi_class():
    """Load rank_bm25 dynamically so missing optional dependency does not break imports."""
    try:
        return getattr(importlib.import_module("rank_bm25"), "BM25Okapi", None)
    except Exception:
        return None


def get_chatpdf_bm25_index(file_names, user_id=None):
    """Build a lightweight sparse keyword index for selected documents."""
    bm25_okapi_class = get_bm25_okapi_class()
    if bm25_okapi_class is None:
        return None, []
    user_id = user_id or get_active_user_id()
    cache = st.session_state.setdefault("chatpdf_bm25_indexes", {})
    cache_key = get_chatpdf_collection_id(user_id, file_names)
    if cache_key in cache:
        entry = cache[cache_key]
        return entry.get("bm25"), entry.get("documents", [])

    documents = get_chatpdf_documents_for_selection(file_names, user_id=user_id)
    corpus_tokens = [tokenize_chatpdf_text(doc.page_content) for doc in documents]
    usable_pairs = [(doc, tokens) for doc, tokens in zip(documents, corpus_tokens) if tokens]
    if not usable_pairs:
        return None, documents

    indexed_documents = [doc for doc, _ in usable_pairs]
    indexed_tokens = [tokens for _, tokens in usable_pairs]
    bm25 = bm25_okapi_class(indexed_tokens)
    cache[cache_key] = {"bm25": bm25, "documents": indexed_documents}
    st.session_state.chatpdf_bm25_indexes = cache
    return bm25, indexed_documents


def chatpdf_document_key(doc):
    meta = getattr(doc, "metadata", {}) or {}
    content_hash = hashlib.sha1(str(getattr(doc, "page_content", "")).encode("utf-8", errors="ignore")).hexdigest()[:12]
    return (
        meta.get("document_id", ""),
        meta.get("file_name", ""),
        str(meta.get("page_or_sheet") or meta.get("page_number") or ""),
        str(meta.get("section", "")),
        str(meta.get("chunk_index", "")),
        content_hash,
    )


def sparse_chatpdf_search(question, file_names, user_id=None, top_k=12):
    """BM25 keyword search over the same chunks used by dense retrieval."""
    bm25, documents = get_chatpdf_bm25_index(file_names, user_id=user_id)
    if bm25 is None or not documents:
        return []
    query_tokens = tokenize_chatpdf_text(question)
    if not query_tokens:
        return []
    scores = bm25.get_scores(query_tokens)
    ranked_indices = sorted(range(len(scores)), key=lambda index: scores[index], reverse=True)
    return [documents[index] for index in ranked_indices[:top_k] if scores[index] > 0]


def merge_chatpdf_results(*result_groups):
    """Merge dense and sparse candidates while preserving first-seen rank order."""
    merged = []
    seen = set()
    for group in result_groups:
        for doc in group or []:
            key = chatpdf_document_key(doc)
            if key in seen:
                continue
            seen.add(key)
            merged.append(doc)
    return merged


@st.cache_resource(show_spinner=False)
def load_chatpdf_reranker():
    """Optional cross-encoder reranker; falls back silently if unavailable."""
    try:
        from sentence_transformers import CrossEncoder
        model_name = os.environ.get("CHATPDF_RERANKER_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")
        return CrossEncoder(model_name)
    except Exception:
        return None


def lexical_rerank_score(question, doc):
    query_tokens = set(tokenize_chatpdf_text(question))
    doc_tokens = tokenize_chatpdf_text(getattr(doc, "page_content", ""))
    if not query_tokens or not doc_tokens:
        return 0
    doc_token_set = set(doc_tokens)
    overlap = len(query_tokens.intersection(doc_token_set))
    phrase_bonus = 0
    question_l = str(question or "").lower().strip()
    content_l = str(getattr(doc, "page_content", "") or "").lower()
    if question_l and question_l in content_l:
        phrase_bonus = 4
    return (overlap * 3) + phrase_bonus + min(len(doc_tokens), 500) / 1000


def rerank_chatpdf_documents(question, docs, top_k=8):
    """Rerank merged candidates with a cross-encoder when possible."""
    if not docs:
        return []
    reranker = load_chatpdf_reranker()
    if reranker is not None:
        try:
            pairs = [[str(question), str(getattr(doc, "page_content", ""))[:4000]] for doc in docs]
            scores = reranker.predict(pairs)
            ranked = [doc for _, doc in sorted(zip(scores, docs), key=lambda item: item[0], reverse=True)]
            return ranked[:top_k]
        except Exception:
            pass
    return sorted(docs, key=lambda doc: lexical_rerank_score(question, doc), reverse=True)[:top_k]


def rewrite_chatpdf_query(llm, question):
    """Optionally rewrite the user question into a compact retrieval query."""
    original_question = re.sub(r"\s+", " ", str(question or "")).strip()
    if not original_question:
        return original_question
    lower_question = original_question.lower()
    if re.fullmatch(r"(analy[sz]e|analysis|summary|summari[sz]e|overview)", lower_question):
        return (
            "introduction overview purpose main features capabilities architecture components "
            "workflow applications use cases technical details key takeaways"
        )
    if llm is None:
        return original_question

    rewrite_prompt = f"""Rewrite this question to improve document search.

Rules:
- Preserve the user's intent.
- Keep important names, numbers, acronyms, and technical terms.
- Return only one concise search query.
- Do not answer the question.

Question: {original_question}

Search query:"""
    try:
        rewritten = str(llm.invoke(rewrite_prompt)).strip()
    except Exception:
        return original_question

    rewritten = re.sub(r"(?is)^.*?search query:\s*", "", rewritten).strip()
    rewritten = re.sub(r"\s+", " ", rewritten).strip(" `\"'")
    if not rewritten:
        return original_question
    if len(rewritten) > 260:
        rewritten = rewritten[:260].rsplit(" ", 1)[0].strip() or original_question
    if rewritten.lower() in {"none", "n/a", "not available"}:
        return original_question
    return rewritten


def hybrid_chatpdf_retrieve(question, file_names, user_id=None, dense_k=12, sparse_k=12, final_k=8, search_query=None):
    """Dense FAISS + sparse BM25 retrieval followed by reranking."""
    user_id = user_id or get_active_user_id()
    vector_store = get_chatpdf_vector_store(file_names, user_id=user_id)
    search_queries = [str(question or "").strip()]
    if search_query and str(search_query).strip() not in search_queries:
        search_queries.append(str(search_query).strip())

    result_groups = []
    for query_text in [query for query in search_queries if query]:
        dense_results = []
        if vector_store is not None:
            try:
                dense_results = vector_store.similarity_search(query_text, k=dense_k)
            except Exception:
                dense_results = []
        sparse_results = sparse_chatpdf_search(query_text, file_names, user_id=user_id, top_k=sparse_k)
        result_groups.extend([dense_results, sparse_results])

    candidates = merge_chatpdf_results(*result_groups)
    return rerank_chatpdf_documents(question, candidates, top_k=final_k)


def get_chatpdf_vector_store(file_names, user_id=None):
    """Build or load persistent FAISS collection for user/document selection."""
    user_id = user_id or get_active_user_id()
    ensure_files_processed(file_names)
    collection_key = f"chatpdf::{user_id}::{get_chatpdf_collection_id(user_id, file_names)}"
    cached_vs = VECTOR_STORE_CACHE.get(collection_key)
    if cached_vs is not None:
        return cached_vs
    if collection_key in st.session_state.get("vector_stores", {}):
        return st.session_state.vector_stores[collection_key]

    embeddings = get_chatpdf_embeddings()
    vector_dir = get_chatpdf_vector_dir(user_id, file_names)
    if os.path.exists(vector_dir):
        try:
            vs = FAISS.load_local(vector_dir, embeddings, allow_dangerous_deserialization=True)
            st.session_state.vector_stores[collection_key] = vs
            VECTOR_STORE_CACHE.set(collection_key, vs)
            return vs
        except Exception:
            pass

    documents = get_chatpdf_documents_for_selection(file_names, user_id=user_id)
    if not documents:
        return None
    vs = FAISS.from_documents(documents, embeddings)
    os.makedirs(vector_dir, exist_ok=True)
    try:
        vs.save_local(vector_dir)
    except Exception:
        pass
    st.session_state.vector_stores[collection_key] = vs
    VECTOR_STORE_CACHE.set(collection_key, vs)
    return vs


def format_chatpdf_context(docs):
    blocks = []
    for index, doc in enumerate(docs, start=1):
        meta = getattr(doc, "metadata", {}) or {}
        source = build_chatpdf_citation_label(meta)
        blocks.append(
            f"[Reference {index}: {source}]\n"
            f"{getattr(doc, 'page_content', str(doc))}"
        )
    return "\n\n---\n\n".join(blocks)


def format_chatpdf_sources(docs, include_snippets=False):
    sources = []
    seen = set()
    for doc in docs:
        meta = getattr(doc, "metadata", {}) or {}
        line_label = build_chatpdf_citation_label(meta)
        key = (
            meta.get("file_name", "document"),
            meta.get("file_type", ""),
            str(meta.get("page_or_sheet") or meta.get("page_number") or ""),
            meta.get("section", ""),
        )
        if key in seen:
            continue
        seen.add(key)
        line = f"- {line_label}"
        if include_snippets:
            snippet = re.sub(r"\s+", " ", getattr(doc, "page_content", "")).strip()[:220]
            if snippet:
                line += f"\n  Snippet: {snippet}"
        sources.append(line)
    return "\n".join(sources) if sources else "- No sources found"


def build_chatpdf_prompt(question, docs, memory):
    memory_text = "\n".join(
        f"User: {item.get('question', '')}\nAssistant: {item.get('answer', '')}"
        for item in memory[-4:]
    )
    context = format_chatpdf_context(docs)
    return RAG_QA_PROMPT.format(
        memory=memory_text or "No previous conversation.",
        context=context,
        question=question,
    )


SMART_QUERY_INTENT_TERMS = {
    "FULL_DOCUMENT_ANALYSIS": (
        "introduction overview purpose main features capabilities architecture components "
        "workflow applications use cases technical details constraints key takeaways"
    ),
    "SHORT_SUMMARY": "overview purpose key points main findings summary takeaways",
    "OVERVIEW": "overview purpose usage audience main concept areas covered",
    "FEATURES_ONLY": "features capabilities functions benefits modules components",
    "SPECIFIC_COMPONENT_DETAILS": "component module interface connector configuration usage limitations details",
    "PIN_DIAGRAMS_CONNECTORS_TABLES": "pin connector pinout signal mapping table channel diagram figure",
    "WORKFLOW_OR_PROCESS": "workflow process steps procedure input output operation usage flow",
    "USE_CASES_APPLICATIONS": "use cases applications users scenarios benefits practical usage",
    "COMPARISON": "compare difference similarities criteria usage capabilities limitations",
    "TABLE_EXTRACTION": "table sheet rows columns csv data structured values",
    "IMAGE_OR_DIAGRAM_EXPLANATION": "diagram figure image flow block architecture visual layout",
    "DOWNLOADABLE_REPORT": "overview purpose findings features workflow tables recommendations report",
    "TROUBLESHOOTING_OR_LIMITATIONS": "troubleshooting limitations constraints issue problem error cause recommendation",
    "REQUIREMENTS_OR_SPECIFICATION_EXTRACTION": "requirements specifications shall must value condition applies notes",
}


def build_retrieval_query_for_intent(question, intent):
    """Blend user terms with intent terms so retrieval stays focused but robust."""
    question_text = re.sub(r"\s+", " ", str(question or "")).strip()
    intent_terms = SMART_QUERY_INTENT_TERMS.get(str(intent or ""), "")
    if not intent_terms:
        return question_text
    if not question_text:
        return intent_terms
    return f"{question_text} {intent_terms}"


def build_smart_document_prompt(question, intent, docs, memory, cross_file_hits=None):
    document_intent = classify_query_intent(question, previous_messages=memory)
    memory_text = "\n".join(
        f"User: {item.get('question', '')}\nAssistant: {item.get('answer', '')}"
        for item in (memory or [])[-4:]
    ) or "No previous conversation."
    context = format_chatpdf_context(docs)
    cross_file_text = "\n".join(
        f"- Related prior context from {file_id}"
        for score, file_id in (cross_file_hits or [])
    ) or "No additional cross-file memory hits."
    return build_document_intelligence_prompt(
        query=question,
        document_intent=document_intent,
        context=context,
        memory=memory_text,
        cross_file_hints=cross_file_text,
    )


def smart_file_brain_query(question, file_names, user_id=None, intent=None, top_k=8):
    """End-to-end query engine: intent-aware retrieval, selected chunking, lazy embedding, citations."""
    user_id = user_id or get_active_user_id()
    file_names = list(dict.fromkeys(file_names or []))
    if not file_names:
        return "Answer:\nPlease select one or more documents before asking a question.\n\nSources:\n- No sources found", []

    document_intent = classify_query_intent(question, previous_messages=get_chatpdf_memory(user_id, file_names))
    intent = intent or to_technical_intent(document_intent)
    if intent in {"SHORT_SUMMARY"} and document_intent == "factual_query":
        intent = classify_technical_document_request(question)
    if document_intent == "analysis_request":
        top_k = max(top_k, 12)
    elif requires_document_scope(document_intent):
        top_k = max(top_k, 9)
    retrieval_query = build_retrieval_query_for_intent(question, intent)
    docs = []
    if requires_document_scope(document_intent) or intent in {
        "FULL_DOCUMENT_ANALYSIS",
        "SHORT_SUMMARY",
        "OVERVIEW",
        "DOWNLOADABLE_REPORT",
    }:
        docs.extend(retrieve_document_understanding_documents(retrieval_query, file_names, user_id=user_id, top_k=top_k))

    docs = merge_chatpdf_results(
        docs,
        retrieve_file_brain_documents(retrieval_query, file_names, user_id=user_id, top_k=top_k),
    )

    if str(retrieval_query).strip().lower() != str(question or "").strip().lower():
        direct_docs = retrieve_file_brain_documents(question, file_names, user_id=user_id, top_k=max(3, top_k // 2))
        docs = merge_chatpdf_results(docs, direct_docs)
        docs = rerank_chatpdf_documents(question, docs, top_k=max(top_k * 2, top_k))
        docs = lazy_embed_selected_documents(question, docs, top_k=top_k)

    if not docs:
        docs = sparse_chatpdf_search(retrieval_query or question, file_names, user_id=user_id, top_k=top_k)
        docs = rerank_chatpdf_documents(question, docs, top_k=top_k)

    if not docs:
        brains = get_file_brains(file_names)
        return build_best_effort_response(
            question,
            document_intent,
            brains,
            "\n".join(f"- {file_name}" for file_name in file_names),
        ), []

    memory = get_chatpdf_memory(user_id, file_names)
    cross_hits = cross_file_search(question, file_names=file_names, top_k=5)
    llm = load_llm()
    if llm is None:
        if requires_document_scope(document_intent):
            answer = synthesize_document_response(
                query=question,
                document_intent=document_intent,
                brains=get_file_brains(file_names),
                docs=docs,
                sources_text=format_chatpdf_sources(docs),
            )
        else:
            answer = build_extractive_chatpdf_answer(question, docs)
    else:
        prompt = build_smart_document_prompt(question, intent, docs, memory, cross_hits)
        try:
            answer = str(llm.invoke(prompt)).strip()
        except Exception:
            answer = build_extractive_chatpdf_answer(question, docs)

        if not answer:
            answer = synthesize_document_response(
                query=question,
                document_intent=document_intent,
                brains=get_file_brains(file_names),
                docs=docs,
                sources_text=format_chatpdf_sources(docs),
            )
        if requires_document_scope(document_intent) and response_looks_like_metadata_dump(answer):
            answer = synthesize_document_response(
                query=question,
                document_intent=document_intent,
                brains=get_file_brains(file_names),
                docs=docs,
                sources_text=format_chatpdf_sources(docs),
            )
        if "Sources:" not in answer:
            answer = answer.rstrip() + "\n\nSources:\n" + format_chatpdf_sources(docs)
        elif not any((getattr(doc, "metadata", {}) or {}).get("file_name", "") in answer for doc in docs):
            answer = answer.rstrip() + "\n" + format_chatpdf_sources(docs)

    append_chatpdf_memory(user_id, file_names, question, answer)
    return strip_llm_suggestions_from_response(answer), docs


def build_extractive_chatpdf_answer(question, docs):
    """Fallback answer when no LLM is available; remains grounded without dumping raw context blocks."""
    if not docs:
        return (
            "Answer:\nI could not find strong matching evidence in the selected documents. "
            "Try a broader wording or ask for a summary/overview so the document-level memory can be used.\n\n"
            "Sources:\n- No sources found"
        )
    bullets = []
    for doc in docs[:4]:
        text = getattr(doc, "page_content", "")
        for sentence in document_intelligence_meaningful_sentences(text, limit=3):
            sentence = normalize_synthesis_text(sentence)
            if sentence and sentence not in bullets:
                bullets.append(sentence[:360])
            if len(bullets) >= 5:
                break
        if len(bullets) >= 5:
            break
    if not bullets:
        return (
            "Answer:\nI found candidate context, but it did not contain enough readable detail to answer confidently.\n\n"
            "Sources:\n- No sources found"
        )
    answer_intro = "Based on the available document context, here is the most relevant synthesis:"
    return "Answer:\n" + answer_intro + "\n\n" + "\n".join(f"- {item}" for item in bullets) + "\n\nSources:\n" + format_chatpdf_sources(docs)


def answer_chatpdf_question(question, file_names, user_id=None, top_k=8):
    """File-brain first answer with citations and per-document memory.

    The dense FAISS path is retained as a fallback. Normal queries retrieve
    page/table/diagram candidates from the file brain and avoid embedding the
    whole uploaded selection.
    """
    return smart_file_brain_query(
        question,
        file_names,
        user_id=user_id,
        intent=classify_technical_document_request(question),
        top_k=top_k,
    )


@st.cache_resource(show_spinner=False)
def load_llm():
    try:
        from transformers import pipeline
    except Exception:
        st.warning(
            "LLM is unavailable because transformers could not be imported (torchvision unmet).\nInstall torch & torchvision if you want AI features.")
        return None

    candidate_tasks = ["text2text-generation", "text-generation", "image-text-to-text", "table-question-answering"]
    task_errors = []

    for task in candidate_tasks:
        try:
            pipe = pipeline(task, model="google/flan-t5-small", max_new_tokens=384, return_full_text=False)
            st.session_state.llm_task = task
            return HuggingFacePipeline(pipeline=pipe)
        except Exception as e:
            task_errors.append((task, str(e)))

    # Only show a single warning if no task could be initialized
    if task_errors:
        st.warning("LLM initialization failed for all candidate tasks; AI features are unavailable.")
        st.session_state.llm_task = None

    return None


def ensure_files_processed(file_names):
    for file_name in file_names:
        ensure_file_processed(file_name)
        ensure_file_brain(file_name)


def process_selected_files():
    """Fully extract all currently selected files so every tab uses the same data."""
    ensure_files_processed(st.session_state.selected_files)


def init_workspace_db():
    """Initialize persistent workspace storage for memory and logs."""
    os.makedirs(APP_DIR, exist_ok=True)
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workspace_meta (
            meta_key TEXT PRIMARY KEY,
            meta_value TEXT
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workspace_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            timestamp TEXT,
            log_type TEXT,
            message TEXT,
            details TEXT
        )
        """
    )
    conn.commit()
    conn.close()


def default_workspace_memory():
    return {
        "chat": [],
        "agent_runs": [],
        "indexed_files": [],
        "memory_events": [],
        "summary": {},
        "metadata": {},
    }


def normalize_workspace_memory(memory):
    """Keep older saved memory compatible with the autonomous workspace schema."""
    normalized = default_workspace_memory()
    if isinstance(memory, dict):
        for key, value in memory.items():
            normalized[key] = value
    for list_key in ["chat", "agent_runs", "indexed_files", "memory_events"]:
        if not isinstance(normalized.get(list_key), list):
            normalized[list_key] = []
    for dict_key in ["summary", "metadata"]:
        if not isinstance(normalized.get(dict_key), dict):
            normalized[dict_key] = {}
    return normalized


def load_workspace_memory():
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "SELECT meta_value FROM workspace_meta WHERE meta_key = ?",
        (WORKSPACE_MEMORY_KEY,)
    )
    row = cursor.fetchone()
    conn.close()
    if row:
        try:
            return normalize_workspace_memory(json.loads(row[0]))
        except Exception:
            pass
    return default_workspace_memory()


def save_workspace_memory():
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "INSERT OR REPLACE INTO workspace_meta (meta_key, meta_value) VALUES (?, ?)",
        (WORKSPACE_MEMORY_KEY, json.dumps(st.session_state.workspace_memory, default=str))
    )
    conn.commit()
    conn.close()


def save_memory_log(log_type, message, details=None):
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    details_json = json.dumps(details, default=str) if details is not None else None
    cursor.execute(
        "INSERT INTO workspace_logs (timestamp, log_type, message, details) VALUES (?, ?, ?, ?)",
        (datetime.now().isoformat(), log_type, message, details_json)
    )
    conn.commit()
    conn.close()


def get_memory_logs(limit=50):
    init_workspace_db()
    conn = sqlite3.connect(WORKSPACE_DB_FILE, check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute(
        "SELECT timestamp, log_type, message, details FROM workspace_logs ORDER BY id DESC LIMIT ?",
        (limit,)
    )
    rows = cursor.fetchall()
    conn.close()
    results = []
    for timestamp, log_type, message, details in rows:
        try:
            details = json.loads(details) if details else None
        except Exception:
            details = details
        results.append({
            "timestamp": timestamp,
            "type": log_type,
            "message": message,
            "details": details,
        })
    return results


def record_workspace_memory_event(event_type, title, content, source=None):
    """Append a compact memory event that every module can retrieve later."""
    event = {
        "type": event_type,
        "title": title,
        "content": str(content or "")[:4000],
        "source": source or "workspace",
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    st.session_state.workspace_memory["memory_events"].append(event)
    st.session_state.workspace_memory["memory_events"] = st.session_state.workspace_memory["memory_events"][-200:]
    return event


def append_chat_to_workspace_memory(user_input, assistant_response, file_names):
    """Store chat as durable workspace memory instead of isolated chat history."""
    chat_entry = {
        "user": user_input,
        "assistant": assistant_response,
        "files": list(file_names or []),
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    st.session_state.workspace_memory["chat"].append(chat_entry)
    st.session_state.workspace_memory["chat"] = st.session_state.workspace_memory["chat"][-200:]
    record_workspace_memory_event(
        "chat",
        "Conversation memory",
        f"User: {user_input}\nAssistant: {assistant_response}",
        source=", ".join(file_names or []) or "chat",
    )
    return chat_entry


def build_unified_memory_text(file_names=None, include_chat=True, include_agents=True, max_chars=MAX_VECTOR_TEXT_CHARS):
    """Compose documents, conversations, and agent runs into one AI brain text."""
    st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)
    sections = []

    candidate_files = file_names if file_names is not None else [f["name"] for f in st.session_state.get("uploaded_files", [])]
    candidate_files = [file_name for file_name in candidate_files if file_name in st.session_state.file_texts]
    for file_name in candidate_files:
        text = str(st.session_state.file_texts.get(file_name, "")).strip()
        if text:
            sections.append(f"[DOCUMENT: {file_name}]\n{text[:60000]}")

    if include_chat:
        for entry in st.session_state.workspace_memory.get("chat", [])[-80:]:
            sections.append(
                "[CHAT MEMORY]\n"
                f"User: {entry.get('user', '')}\n"
                f"Assistant: {entry.get('assistant', '')}\n"
                f"Files: {', '.join(entry.get('files', []) or [])}"
            )

    for event in st.session_state.workspace_memory.get("memory_events", [])[-100:]:
        sections.append(
            "[MEMORY EVENT]\n"
            f"Type: {event.get('type', '')}\n"
            f"Title: {event.get('title', '')}\n"
            f"Source: {event.get('source', '')}\n"
            f"{event.get('content', '')}"
        )

    if include_agents:
        for run in st.session_state.workspace_memory.get("agent_runs", [])[-40:]:
            sections.append(
                "[CAPL AGENT RUN]\n"
                f"Goal: {run.get('goal', '')}\n"
                f"Plan: {', '.join(run.get('plan', []) or [])}\n"
                f"Final: {run.get('final_response', '')[:2500]}"
            )

    memory_text = "\n\n".join(sections).strip()
    return memory_text[:max_chars]


def get_unified_workspace_vector_store(file_names=None):
    memory_text = build_unified_memory_text(file_names=file_names)
    if not memory_text.strip():
        return None

    digest = hashlib.md5(memory_text.encode("utf-8", errors="ignore")).hexdigest()
    selection_key = f"unified_memory::{digest}"
    cached_vs = VECTOR_STORE_CACHE.get(selection_key)
    if cached_vs is not None:
        st.session_state.vector_stores[selection_key] = cached_vs
        return cached_vs

    try:
        vs = create_vector_store(memory_text)
        st.session_state.vector_stores[selection_key] = vs
        VECTOR_STORE_CACHE.set(selection_key, vs)
        return vs
    except Exception:
        return None


def get_workspace_vector_store(file_names=None):
    return get_unified_workspace_vector_store(file_names)


def search_workspace_memory(query, limit=4):
    vector_store = get_workspace_vector_store()
    if vector_store is None:
        return []
    try:
        docs = vector_store.similarity_search(query, k=limit)
        return [getattr(doc, "page_content", str(doc)) for doc in docs]
    except Exception:
        return []


def extract_risk_signals(text, limit=12):
    risk_terms = [
        "risk", "failure", "hazard", "issue", "problem", "danger", "alert",
        "warning", "fault", "breach", "leak", "vulnerability", "impact",
        "delay", "downtime", "non-compliance",
    ]
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    results = []
    for line in lines:
        lower_line = line.lower()
        if any(term in lower_line for term in risk_terms):
            results.append(line)
        if len(results) >= limit:
            break
    return results or ["No explicit risk signals were found in the selected documents."]


def extract_entities(text, limit=20):
    raw_text = str(text or "")
    candidates = re.findall(r"\b[A-Z][A-Za-z0-9]{2,}(?: [A-Z][A-Za-z0-9]{2,})*\b", raw_text)
    counts = Counter(candidates)
    return [entity for entity, _ in counts.most_common(limit)]


def extract_key_themes(text, limit=8):
    words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{3,}", str(text or ""))
    counts = Counter(
        word.lower()
        for word in words
        if word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    return [word.title() for word, _ in counts.most_common(limit)]


def build_workspace_intelligence_summary(file_names=None):
    """Produce the live AI Insight Engine view from shared memory."""
    if file_names is None:
        file_names = st.session_state.get("selected_files", []) or [
            f["name"] for f in st.session_state.get("uploaded_files", [])
        ]
    ensure_files_processed(file_names)
    memory_text = build_unified_memory_text(file_names=file_names, max_chars=180000)
    chat_count = len(st.session_state.workspace_memory.get("chat", []))
    agent_count = len(st.session_state.workspace_memory.get("agent_runs", []))
    indexed_files = sorted(set(st.session_state.workspace_memory.get("indexed_files", []) + list(file_names or [])))

    themes = extract_key_themes(memory_text, limit=10)
    entities = extract_entities(memory_text, limit=16)
    risks = extract_risk_signals(memory_text, limit=8)
    recent_logs = get_memory_logs(limit=6)

    insights = []
    if themes:
        insights.append(f"Dominant knowledge themes: {', '.join(themes[:5])}.")
    if chat_count:
        insights.append(f"Conversation memory is active with {chat_count} stored exchange(s).")
    if indexed_files:
        insights.append(f"{len(indexed_files)} document(s) are connected to the shared AI memory.")
    if agent_count:
        insights.append(f"CAPL agents have completed {agent_count} autonomous run(s).")
    if not insights:
        insights.append("Upload documents or start a conversation to grow the workspace memory.")

    return {
        "themes": themes,
        "insights": insights,
        "entities": entities,
        "risks": risks,
        "state": {
            "indexed_files": len(indexed_files),
            "chat_entries": chat_count,
            "agent_runs": agent_count,
            "memory_events": len(st.session_state.workspace_memory.get("memory_events", [])),
            "memory_chars": len(memory_text),
        },
        "logs": recent_logs,
    }


def render_workspace_intelligence_panel(file_names=None):
    summary = build_workspace_intelligence_summary(file_names=file_names)
    state = summary["state"]

    st.markdown("### AI Insight Engine")
    metric_cols = st.columns(4)
    metric_cols[0].metric("Indexed Files", state["indexed_files"])
    metric_cols[1].metric("Chat Memory", state["chat_entries"])
    metric_cols[2].metric("CAPL Runs", state["agent_runs"])
    metric_cols[3].metric("Memory Events", state["memory_events"])

    cols = st.columns(2)
    with cols[0]:
        st.markdown("#### Key Themes")
        st.markdown(", ".join(summary["themes"]) if summary["themes"] else "No themes detected yet.")
        st.markdown("#### Entities")
        st.markdown(", ".join(summary["entities"][:12]) if summary["entities"] else "No entities detected yet.")
    with cols[1]:
        st.markdown("#### Insights")
        for item in summary["insights"]:
            st.markdown(f"- {html.escape(item)}")
        st.markdown("#### Risks / Signals")
        for item in summary["risks"][:6]:
            st.markdown(f"- {html.escape(str(item)[:240])}")

    if summary["logs"]:
        with st.expander("Live Memory Loop", expanded=False):
            for log in summary["logs"]:
                st.markdown(f"**{html.escape(log.get('type', 'log'))}** - {html.escape(log.get('message', ''))}")


def plan_autonomous_task(goal):
    lower_goal = str(goal or "").lower()
    tasks = ["retrieve_relevant_memory"]
    if any(term in lower_goal for term in ["compare", "difference", "diff", "semantic diff"]):
        tasks.append("compare_documents")
    if any(term in lower_goal for term in ["risk", "issue", "signal", "warning", "hazard"]):
        tasks.append("extract_risks")
    if any(term in lower_goal for term in ["entity", "entities", "extract entities", "parts", "components"]):
        tasks.append("extract_entities")
    if any(term in lower_goal for term in ["summarize", "summary", "overview", "insight", "analyze"]):
        tasks.append("summarize_findings")
    if any(term in lower_goal for term in ["analyze", "review", "inspect"]) and "summarize_findings" not in tasks:
        tasks.append("analyze_documents")
    if "summarize_findings" not in tasks and "analyze_documents" not in tasks:
        tasks.append("summarize_findings")
    return list(dict.fromkeys(tasks))


def planning_agent(goal):
    steps = plan_autonomous_task(goal)
    return {
        "agent": "Planning Agent",
        "role": "Brain / Orchestrator",
        "strategy": "Break the goal into memory retrieval, tool execution, reasoning, and coordination steps.",
        "steps": steps,
    }


def retrieve_autonomous_context(goal, file_names):
    if not goal:
        return []
    vector_store = get_workspace_vector_store(file_names)
    if vector_store is None:
        return []
    try:
        docs = vector_store.similarity_search(goal, k=4)
        return [getattr(doc, "page_content", str(doc)) for doc in docs]
    except Exception:
        return []


def retrieval_agent(goal, file_names):
    context = retrieve_autonomous_context(goal, file_names)
    return {
        "agent": "Retrieval Agent",
        "role": "Memory Brain",
        "context": context,
        "summary": f"Retrieved {len(context)} relevant memory fragment(s) from shared FAISS memory.",
    }


def execute_autonomous_tool(task, file_names, context):
    file_names = file_names or [f["name"] for f in st.session_state.uploaded_files]
    if task == "retrieve_relevant_memory":
        return {"context": context}
    if task == "analyze_documents":
        results = []
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            file_entry = get_uploaded_file_entry(file_name)
            if file_text.strip() and file_entry is not None:
                results.append(build_detailed_document_summary(file_name, file_entry["bytes"], file_text))
        return {"analysis": results or ["No document content available for analysis."]}
    if task == "summarize_findings":
        context_text = "\n\n".join(context[:3]) if context else "No relevant memory context found."
        return {"summary": context_text}
    if task == "extract_risks":
        results = {}
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            if file_text.strip():
                results[file_name] = extract_risk_signals(file_text)
        return {"risks": results or {"message": "No risk-related content found."}}
    if task == "extract_entities":
        results = {}
        for file_name in file_names:
            file_text = st.session_state.file_texts.get(file_name, "")
            if file_text.strip():
                results[file_name] = extract_entities(file_text)
        return {"entities": results or {"message": "No entities could be extracted."}}
    if task == "compare_documents":
        if len(file_names) < 2:
            return {"compare": "At least two files are required for document comparison."}
        file_texts = {f: st.session_state.file_texts.get(f, "") for f in file_names}
        return {"compare": highlight_multi_file_differences(file_texts)}
    return {"result": f"No tool implementation exists for '{task}'."}


def execution_agent(steps, file_names, context):
    outputs = {}
    for task in steps:
        if task == "retrieve_relevant_memory":
            outputs[task] = {"context": context}
        else:
            outputs[task] = execute_autonomous_tool(task, file_names, context)
    return {
        "agent": "Execution Agent",
        "role": "Tool Runner",
        "outputs": outputs,
        "summary": f"Executed {len(outputs)} autonomous tool step(s).",
    }


def reason_over_results(agent_outputs):
    summary_parts = []
    for step, output in agent_outputs.items():
        if isinstance(output, dict):
            if "summary" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** {output['summary']}")
            elif "context" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Retrieved {len(output['context'])} memory fragments.")
            elif "risks" in output:
                risk_lines = sum(len(v) for v in output['risks'].values() if isinstance(v, list))
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Extracted {risk_lines} risk lines.")
            elif "entities" in output:
                entity_count = sum(len(v) for v in output['entities'].values() if isinstance(v, list))
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Extracted {entity_count} entities.")
            elif "compare" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Comparison results are available in the agent output." )
            elif "analysis" in output:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Document analysis completed for {len(output['analysis'])} files.")
            else:
                summary_parts.append(f"**{step.replace('_', ' ').title()}:** Output produced.")
        else:
            summary_parts.append(f"**{step.replace('_', ' ').title()}:** {str(output)[:200]}")
    return "\n\n".join(summary_parts)


def reasoning_agent(goal, file_names, context, tool_outputs):
    memory_text = build_unified_memory_text(file_names=file_names, max_chars=80000)
    combined_reasoning_text = "\n\n".join(context or []) + "\n\n" + memory_text[:40000]
    themes = extract_key_themes(combined_reasoning_text, limit=8)
    entities = extract_entities(combined_reasoning_text, limit=12)
    risks = extract_risk_signals(combined_reasoning_text, limit=8)
    interpretation = [
        f"Goal interpreted as: {goal}",
        f"Key themes: {', '.join(themes) if themes else 'None detected'}",
        f"Important entities: {', '.join(entities[:8]) if entities else 'None detected'}",
        f"Risk/signal count: {len(risks)}",
        reason_over_results(tool_outputs),
    ]
    return {
        "agent": "Reasoning Agent",
        "role": "Analyst",
        "themes": themes,
        "entities": entities,
        "risks": risks,
        "interpretation": "\n\n".join(interpretation),
    }


def coordinate_agents(goal, steps, outputs):
    sections = [f"### Autonomous CAPL Agent Result\n**Goal:** {html.escape(goal)}\n"]
    sections.append(reason_over_results(outputs))
    for step in steps:
        output = outputs.get(step)
        if isinstance(output, dict) and step == "compare_documents":
            sections.append("### Comparison Output")
            sections.append(output.get("compare", "No comparison result."))
        elif isinstance(output, dict) and step == "analyze_documents":
            analysis = output.get("analysis", [])
            sections.append("### Analysis Output")
            sections.extend(analysis[:2] if isinstance(analysis, list) else [str(analysis)])
    return "\n\n".join(sections)


def coordination_agent(goal, planning, retrieval, execution, reasoning):
    tool_outputs = execution.get("outputs", {})
    sections = [
        "### Autonomous CAPL Agent Result",
        f"**Goal:** {html.escape(goal)}",
        "#### Agent Trace",
        f"- **Planning Agent:** {html.escape(planning.get('strategy', ''))}",
        f"- **Retrieval Agent:** {html.escape(retrieval.get('summary', ''))}",
        f"- **Execution Agent:** {html.escape(execution.get('summary', ''))}",
        "- **Reasoning Agent:** Interpreted tool outputs against shared memory.",
        "- **Coordination Agent:** Merged agent outputs into this final response.",
        "#### Execution Plan",
    ]
    sections.extend(f"- {html.escape(step.replace('_', ' ').title())}" for step in planning.get("steps", []))
    sections.extend(["#### Reasoned Findings", reasoning.get("interpretation", "")])

    if reasoning.get("risks"):
        sections.append("#### Risks / Signals")
        sections.extend(f"- {html.escape(str(item)[:240])}" for item in reasoning.get("risks", [])[:8])

    for step in planning.get("steps", []):
        output = tool_outputs.get(step)
        if isinstance(output, dict) and step == "compare_documents":
            sections.append("#### Comparison Output")
            sections.append(output.get("compare", "No comparison result."))
        elif isinstance(output, dict) and step == "analyze_documents":
            analysis = output.get("analysis", [])
            sections.append("#### Analysis Output")
            sections.extend(analysis[:2] if isinstance(analysis, list) else [str(analysis)])

    return "\n\n".join(sections)


def run_capl_agent(goal, file_names):
    if not goal or not str(goal).strip():
        return "Provide a goal for the autonomous CAPL agents."
    target_files = file_names or [f["name"] for f in st.session_state.uploaded_files]
    target_files = [f for f in target_files if get_uploaded_file_entry(f)]
    if not target_files:
        return "No processed files are available. Upload files and allow the system to extract them first."

    ensure_files_processed(target_files)
    planning = planning_agent(goal)
    plan = planning["steps"]
    retrieval = retrieval_agent(goal, target_files)
    execution = execution_agent(plan, target_files, retrieval.get("context", []))
    reasoning = reasoning_agent(goal, target_files, retrieval.get("context", []), execution.get("outputs", {}))
    outputs = {
        "planning": planning,
        "retrieval": retrieval,
        "execution": execution,
        "reasoning": reasoning,
    }
    final_response = coordination_agent(goal, planning, retrieval, execution, reasoning)

    run_entry = {
        "goal": goal,
        "files": target_files,
        "plan": plan,
        "outputs": outputs,
        "final_response": final_response,
        "timestamp": datetime.now().isoformat(),
    }
    st.session_state.agent_run_history.append(run_entry)
    st.session_state.workspace_memory["agent_runs"].append(run_entry)
    record_workspace_memory_event("capl_agent", f"Autonomous goal: {goal}", final_response, source="CAPL")
    st.session_state.workspace_memory["indexed_files"] = sorted(set(st.session_state.workspace_memory.get("indexed_files", []) + target_files))
    save_workspace_memory()
    save_memory_log("capl_agent", f"Ran autonomous CAPL goal: {goal}", {"files": target_files, "plan": plan})

    return final_response


def update_workspace_memory_selection(file_names):
    file_names = sorted(set(file_names or []))
    st.session_state.workspace_memory["indexed_files"] = file_names
    save_workspace_memory()


def ensure_workspace_memory_loaded():
    if not st.session_state.workspace_memory_loaded:
        st.session_state.workspace_memory = normalize_workspace_memory(load_workspace_memory())
        st.session_state.workspace_memory_loaded = True
    else:
        st.session_state.workspace_memory = normalize_workspace_memory(st.session_state.workspace_memory)


def get_selection_signature(file_names):
    digest = hashlib.md5()
    for file_name in sorted(file_names):
        digest.update(file_name.encode("utf-8"))
        digest.update(st.session_state.file_texts.get(file_name, "").encode("utf-8"))
    return f"combined::{digest.hexdigest()}"


@st.cache_data(show_spinner=False)
def get_document_asset_counts(file_name, file_bytes, extracted_text):
    file_name_lower = file_name.lower()
    page_count = 0
    table_count = 0
    image_count = 0

    if file_name_lower.endswith(".pdf"):
        page_match = re.search(r"Total Pages:\s*(\d+)", extracted_text)
        page_count = int(page_match.group(1)) if page_match else len(re.findall(r"Page \d+ Text:", extracted_text))
        table_count = len(re.findall(r"Page \d+ Table \d+:", extracted_text))
        image_count = len(re.findall(r"\[IMAGE:", extracted_text))
    elif file_name_lower.endswith(".pptx"):
        slide_match = re.search(r"Total Slides:\s*(\d+)", extracted_text)
        page_count = int(slide_match.group(1)) if slide_match else 0
        table_count = len(re.findall(r"\bTable:\n", extracted_text))
        image_count = len(re.findall(r"\[EMBEDDED_IMAGE:", extracted_text))
    elif file_name_lower.endswith(".docx"):
        table_count = len(re.findall(r"Table \d+:", extracted_text))
        image_count = len(re.findall(r"\[EMBEDDED_IMAGE:", extracted_text))
    elif file_name_lower.endswith(".xlsx"):
        sheet_match = re.search(r"Workbook contains (\d+) sheets", extracted_text)
        page_count = int(sheet_match.group(1)) if sheet_match else 0
        table_count = len(re.findall(r"Sheet '.*?':", extracted_text))
    elif file_name_lower.endswith((".html", ".htm")):
        image_match = re.search(r"(\d+) images found in HTML", extracted_text)
        image_count = int(image_match.group(1)) if image_match else 0

    return page_count, image_count, table_count


def empty_chat_summary_downloads():
    return {"images": [], "tables": [], "csv": [], "diagrams": []}


def generate_analysis_response(chat_files, analysis_type):
    """Generate structured document analysis using specialized prompts."""
    if not chat_files:
        return "No files selected for analysis."
    
    # Get combined text
    selected_file_texts = {f: st.session_state.file_texts.get(f, "") for f in chat_files}
    combined_text = "\n\n".join(selected_file_texts.values())
    if not combined_text.strip():
        return "No readable text found in selected files."
    
    # Select prompt
    prompts = {
        "analyze": ANALYSIS_PROMPT,
        "summary": SUMMARY_PROMPT,
        "overview": OVERVIEW_PROMPT,
        "features": FEATURES_PROMPT,
    }
    prompt_template = prompts.get(analysis_type)
    if not prompt_template:
        return f"Unknown analysis type: {analysis_type}"
    
    llm = load_llm()
    if llm is None:
        return "AI analysis unavailable (LLM not loaded)."
    
    try:
        # Build chain: context + prompt + empty query (document analysis)
        system_prompt = prompt_template.format(USER_QUERY="Provide a complete analysis/summary/overview/features of this document.")
        full_prompt = f"""SYSTEM: {system_prompt}

DOCUMENT CONTENT:
{combined_text[:MAX_VECTOR_TEXT_CHARS]}"""
        
        response = llm.invoke(full_prompt)
        response = str(response).strip()
        
        # Reset analysis state
        st.session_state.chat_analysis_type = None
        
        return response or "No response generated."
    
    except Exception as e:
        st.error(f"Analysis failed: {e}")
        return f"Analysis error: {str(e)}"


def render_chat_summary_downloads():
    downloads = st.session_state.get("chat_summary_downloads", empty_chat_summary_downloads())
    image_items = downloads.get("images", [])
    table_items = downloads.get("tables", [])
    csv_items = downloads.get("csv", [])
    diagram_items = downloads.get("diagrams", [])

    if not image_items and not table_items and not csv_items and not diagram_items:
        return

    st.markdown("### Summary Downloads")

    if image_items:
        with st.expander("🖼️ Image PNG Downloads", expanded=False):
            for index, item in enumerate(image_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_image_{index}_{item['file_name']}"
                )

    if table_items:
        with st.expander("📊 Table PNG Downloads", expanded=False):
            for index, item in enumerate(table_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_table_{index}_{item['file_name']}"
                )

    if csv_items:
        with st.expander("Pin Table CSV Downloads", expanded=False):
            for index, item in enumerate(csv_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_csv_{index}_{item['file_name']}"
                )

    if diagram_items:
        with st.expander("ASCII Diagram Downloads", expanded=False):
            for index, item in enumerate(diagram_items):
                st.download_button(
                    label=item["label"],
                    data=item["data"],
                    file_name=item["file_name"],
                    mime=item["mime"],
                    key=f"chat_summary_diagram_{index}_{item['file_name']}"
                )


# ==============================
# DOCUMENT-AWARE CHAT REASONING ENGINE
# Shared chat helpers that classify intent before generating responses.
# This prevents extraction prompts from producing generic follow-up buttons.
# ==============================
def classify_document_chat_intent(user_query):
    """Classify the user's document-chat intent with lightweight keyword rules."""
    query = str(user_query or "").strip().lower()
    if not query:
        return "UNKNOWN"

    guidance_terms = ["what can i ask", "what should i do", "suggest", "next step", "guide me", "help me"]
    comparison_terms = ["compare", "difference", "differences", "diff", "versus", "vs ", "between"]
    summary_terms = ["summary", "summarize", "summarise", "overview", "brief", "recap"]
    analysis_terms = ["why", "how", "explain", "analyze", "analyse", "insight", "reason", "impact", "meaning"]
    extraction_terms = [
        "list", "list out", "find", "show", "get", "extract", "give me", "display",
        "what are", "which are", "where is", "all the", "all ", "names of",
    ]

    if any(term in query for term in guidance_terms):
        return "GUIDANCE"
    if any(term in query for term in comparison_terms):
        return "COMPARISON"
    if any(term in query for term in summary_terms):
        return "SUMMARY"
    if any(term in query for term in analysis_terms):
        return "ANALYSIS"
    if any(term in query for term in extraction_terms):
        return "EXTRACTION"
    if extract_bare_item_name(user_query):
        return "EXTRACTION"
    return "UNKNOWN"


def detect_document_chat_profile(file_names, document_context):
    """Detect a broad document profile for prompt and extraction routing."""
    lower_names = " ".join(file_names or []).lower()
    lower_context = str(document_context or "").lower()
    combined = f"{lower_names}\n{lower_context[:50000]}"

    if ".can" in lower_names or "capl" in combined or "on message" in combined:
        return "CAPL"
    if "vn" in combined and re.search(r"\bvn\s*[- ]?\d{3,5}[a-z]?\b", combined, re.IGNORECASE):
        return "VN_DEVICE"
    if any(marker in combined for marker in ["table:", "sheet '", "csv rows:", "|"]):
        return "TABLE"
    if ".pdf" in lower_names or "pdf metadata" in combined or "page 1 text" in combined:
        return "PDF"
    return "MIXED"


def normalize_technical_identifier(value):
    """Normalize identifiers such as 'VN 1630A' to 'VN1630A'."""
    raw_value = str(value or "").strip(" .,:;()[]{}")
    raw_value = re.sub(r"\s+", " ", raw_value)
    raw_value = re.sub(r"\b(VN)\s*[- ]?\s*(\d{3,5}[A-Za-z]?)\b", lambda m: f"{m.group(1).upper()}{m.group(2).upper()}", raw_value, flags=re.IGNORECASE)
    raw_value = re.sub(r"\b([A-Z]{1,6})\s+(\d{2,}[A-Z0-9]*)\b", lambda m: f"{m.group(1).upper()}{m.group(2).upper()}", raw_value, flags=re.IGNORECASE)
    raw_value = re.sub(r"\b(D-SUB)\s*(\d+)\b", lambda m: f"{m.group(1).upper()}{m.group(2)}", raw_value, flags=re.IGNORECASE)
    return raw_value.strip()


def extract_vn_devices_from_text(text):
    """Extract unique VN device identifiers from technical manuals and PDFs."""
    devices = []
    seen = set()
    for match in re.finditer(r"\bVN\s*[- ]?\s*(\d{3,5}[A-Za-z]?)\b", str(text or ""), re.IGNORECASE):
        device = f"VN{match.group(1).upper()}"
        if device.lower() not in seen:
            seen.add(device.lower())
            devices.append(device)
    return devices


def derive_extraction_topic(user_query):
    """Best-effort extraction topic for generic list/show/get prompts."""
    query = str(user_query or "").strip()
    patterns = [
        r"\b(?:list(?: out)?|show|get|extract|give me|display)\s+(?:all\s+|the\s+|all the\s+)?(.+)$",
        r"\b(?:what are|which are)\s+(?:all\s+|the\s+|all the\s+)?(.+)$",
        r"\b(?:find|search|locate)\s+(?:all\s+|the\s+|all the\s+)?(.+)$",
    ]
    for pattern in patterns:
        match = re.search(pattern, query, re.IGNORECASE)
        if match:
            topic = match.group(1).strip(" ?.:-")
            topic = re.sub(r"\b(?:in|from|inside|within)\s+(?:this|the|selected)?\s*(?:document|file|files|pdf)?\b.*$", "", topic, flags=re.IGNORECASE).strip()
            return topic
    return ""


def extract_lines_for_topic(text, topic, limit=25):
    """Return relevant lines for a generic extraction topic."""
    topic = str(topic or "").strip()
    if not topic:
        return []

    topic_terms = [
        term.lower()
        for term in re.findall(r"[A-Za-z0-9_+\-/]{2,}", topic)
        if term.lower() not in {"all", "the", "a", "an", "of", "device", "devices", "list"}
    ]
    if not topic_terms:
        topic_terms = [topic.lower()]

    results = []
    seen = set()
    for raw_line in str(text or "").splitlines():
        line = normalize_extracted_line(raw_line)
        if len(line) < 3 or len(line) > 260:
            continue
        lower_line = line.lower()
        if all(term in lower_line for term in topic_terms) or any(term in lower_line for term in topic_terms):
            key = lower_line
            if key not in seen:
                seen.add(key)
                results.append(line)
        if len(results) >= limit:
            break
    return results


def build_extraction_response_for_query(user_query, file_texts):
    """Build a direct extraction answer with no suggestions."""
    query = str(user_query or "")
    query_lower = query.lower()
    file_texts = file_texts or {}

    if "vn" in query_lower and any(term in query_lower for term in ["device", "devices", "interface", "module", "modules"]):
        rows = []
        all_devices = []
        for file_name, text in file_texts.items():
            devices = extract_vn_devices_from_text(text)
            all_devices.extend(devices)
            if devices:
                rows.append(f"**{html.escape(file_name)}**\n" + "\n".join(f"- {html.escape(device)}" for device in devices))
        unique_devices = list(dict.fromkeys(all_devices))
        if unique_devices:
            return "**VN devices found:**\n\n" + "\n\n".join(rows)
        return "No VN device identifiers were found in the selected document text."

    topic = derive_extraction_topic(query)
    if not topic:
        return "What exact information should I extract from the selected document?"

    response_blocks = []
    for file_name, text in file_texts.items():
        lines = extract_lines_for_topic(text, topic)
        if lines:
            response_blocks.append(
                f"**{html.escape(file_name)}**\n"
                + "\n".join(f"- {html.escape(line)}" for line in lines)
            )

    if response_blocks:
        return "\n\n---\n\n".join(response_blocks)
    return f"No direct matches were found for **{html.escape(topic)}** in the selected document text."


def strip_llm_suggestions_from_response(response):
    """Remove model-produced Suggestions blocks from final chat output."""
    text = str(response or "").strip()
    if not text:
        return text
    text = re.split(r"\n\s*-{3,}\s*\n\s*Suggestions\s*:", text, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    text = re.split(r"\n\s*Suggestions\s*:", text, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    return text



def should_show_chat_suggestions(intent, user_query):
    """Only show suggestion buttons when the user explicitly asks for guidance."""
    query = str(user_query or "").lower()
    if intent == "GUIDANCE":
        return True
    if any(term in query for term in ["suggest", "next step", "what can i ask", "guide me"]):
        return True
    return False


# ==============================
# ANALYSIS BUTTON HANDLER
# ==============================
def generate_analysis_response(chat_files, analysis_type):
    """Generate structured document analysis using specialized prompts."""
    if not chat_files:
        return "No files selected for analysis."
    
    # Get combined text
    selected_file_texts = {f: st.session_state.file_texts.get(f, "") for f in chat_files}
    combined_text = "\n\n".join(selected_file_texts.values())
    if not combined_text.strip():
        return "No readable text found in selected files."
    
    # Select prompt
    prompts = {
        "analyze": ANALYSIS_PROMPT,
        "summary": SUMMARY_PROMPT,
        "overview": OVERVIEW_PROMPT,
        "features": FEATURES_PROMPT,
    }
    prompt_template = prompts.get(analysis_type)
    if not prompt_template:
        return f"Unknown analysis type: {analysis_type}"
    
    llm = load_llm()
    if llm is None:
        return "AI analysis unavailable (LLM not loaded)."
    
    try:
        # Build chain: context + prompt + empty query (document analysis)
        system_prompt = prompt_template.format(USER_QUERY="Provide a complete analysis/summary/overview/features of this document.")
        full_prompt = f"""SYSTEM: {system_prompt}

DOCUMENT CONTENT:
{combined_text[:MAX_VECTOR_TEXT_CHARS]}"""
        
        response = llm.invoke(full_prompt)
        response = str(response).strip()
        
        # Reset analysis state
        st.session_state.chat_analysis_type = None
        
        return response or "No response generated."
    
    except Exception as e:
        st.error(f"Analysis failed: {e}")
        return f"Analysis error: {str(e)}"



# ==============================
# PREMIUM TECHNICAL DOCUMENT RESPONSE ROUTER
# Classifies user requests into documentation-grade response types and builds
# concise, structured answers without raw page-wise extraction.
# ==============================
def extract_specific_component_name(user_query):
    """Detect a requested component/module/item name from a technical query."""
    quoted = extract_quoted_item_name(user_query)
    if quoted:
        return normalize_technical_identifier(quoted)

    bare_item = extract_bare_item_name(user_query)
    if bare_item:
        return normalize_technical_identifier(bare_item)

    component_patterns = [
        r"\b(?:component|module|device|item|part|interface|connector)\s+([A-Za-z][A-Za-z0-9_+\-/]{2,40})\b",
        r"\b(?:about|details(?:\s+about)?|information(?:\s+about)?|explain)\s+([A-Za-z][A-Za-z0-9_+\-/]{2,40})\b",
    ]
    for pattern in component_patterns:
        match = re.search(pattern, str(user_query or ""), re.IGNORECASE)
        if match:
            candidate = normalize_technical_identifier(match.group(1))
            if candidate.lower() not in {"component", "module", "device", "item", "part", "document"}:
                return candidate

    identifiers = re.findall(r"\b[A-Z]{1,6}[-_ ]?\d{2,}[A-Z0-9_+\-/]*\b", str(user_query or ""))
    ignored = {"PDF", "DOCX", "PPTX", "XLSX", "HTML", "CSV"}
    for identifier in identifiers:
        candidate = normalize_technical_identifier(identifier)
        if candidate.upper() not in ignored:
            return candidate
    return ""


def extract_multiple_component_names(user_query):
    """Detect multiple named items for comparison-style prompts."""
    text = str(user_query or "")
    quoted_items = [normalize_technical_identifier(item) for match in re.findall(r"'(.*?)'|\"(.*?)\"", text) for item in match if item]
    identifiers = [normalize_technical_identifier(item) for item in re.findall(r"\b[A-Z]{1,6}[-_ ]?\d{2,}[A-Z0-9_+\-/]*\b", text)]
    ignored = {"PDF", "DOCX", "PPTX", "XLSX", "HTML", "CSV"}
    items = [item for item in quoted_items + identifiers if item and item.upper() not in ignored]
    return list(dict.fromkeys(items))


def classify_technical_document_request(user_query):
    """Classify user queries into the exact enterprise document intent categories requested by the user."""
    query = str(user_query or "").strip().lower()
    if not query:
        return "SHORT_SUMMARY"  # Default to SHORT_SUMMARY if unclear
    compact_query = re.sub(r"[^a-z0-9]+", " ", query).strip()

    if compact_query in {"analyze", "analyse", "analysis", "full analysis", "detailed analysis"}:
        return "FULL_DOCUMENT_ANALYSIS"
    if compact_query in {"summary", "summarize", "summarise", "short summary", "brief summary"}:
        return "SHORT_SUMMARY"
    if compact_query == "overview":
        return "OVERVIEW"
    if any(term in query for term in ["key insights", "deep analysis", "document analysis", "analyze the document", "analyse the document"]):
        return "FULL_DOCUMENT_ANALYSIS"
    if any(term in query for term in ["main themes", "key themes", "themes", "key topics", "main topics"]):
        return "OVERVIEW"
    if any(term in query for term in ["architecture", "technical overview", "system design", "technical explanation"]):
        return "OVERVIEW"

    # Priority: Specific Component > Comparison > Full Analysis > Features > Workflow > Use Cases > Table Extraction > Image/Diagram > Report > Troubleshooting > Requirements > Overview > Short Summary

    # Check for specific component first
    if extract_specific_component_name(user_query):
        return "SPECIFIC_COMPONENT_DETAILS"

    # Check for comparison
    multiple_items = extract_multiple_component_names(user_query)
    if any(term in query for term in ["compare", "difference", "differences", " vs ", " versus "]):
        return "COMPARISON"
    if len(multiple_items) >= 2 and any(term in query for term in ["between", "which", "better", "different"]):
        return "COMPARISON"

    # Check for troubleshooting or limitations
    if any(term in query for term in ["troubleshooting", "troubleshoot", "debug", "fix", "issue", "problem", "limitation", "limitations", "constraint", "constraints", "error", "failure"]):
        return "TROUBLESHOOTING_OR_LIMITATIONS"

    # Check for requirements or specifications
    if any(term in query for term in ["requirements", "requirement", "specifications", "specification", "specs", "reqs"]):
        return "REQUIREMENTS_OR_SPECIFICATION_EXTRACTION"

    # Check for report
    if any(term in query for term in ["downloadable report", "export report", "generate report", "report download", "create report", "report"]):
        return "DOWNLOADABLE_REPORT"

    # Check for table extraction
    if any(term in query for term in ["table extract", "extract table", "extract data", "table data", "table rows", "csv extract", "spreadsheet", "table only"]):
        return "TABLE_EXTRACTION"

    # Check for image or diagram explanation
    if any(term in query for term in ["image", "diagram", "visual", "figure", "schematic", "illustration", "drawing", "visual extraction"]):
        return "IMAGE_OR_DIAGRAM_EXPLANATION"

    # Check for pin diagrams connectors tables
    if any(term in query for term in ["pin", "diagram", "connector", "mapping", "pinout", "visual structure", "technical table", "structured data"]):
        return "PIN_DIAGRAMS_CONNECTORS_TABLES"

    # Check for workflow or process
    if any(term in query for term in ["workflow", "process", "process flow", "how does", "how it works", "steps", "procedure"]):
        return "WORKFLOW_OR_PROCESS"

    # Check for use cases or applications
    if any(term in query for term in ["use case", "use cases", "application", "applications", "real usage", "practical use"]):
        return "USE_CASES_APPLICATIONS"

    # Check for features only
    if any(term in query for term in ["feature", "features", "capability", "capabilities", "functional behavior"]):
        return "FEATURES_ONLY"

    # Check for full document analysis
    if any(term in query for term in ["full analysis", "complete document", "analyze document", "analyse document", "full document", "explain document", "detailed analysis"]):
        return "FULL_DOCUMENT_ANALYSIS"

    # Check for overview
    if any(term in query for term in ["overview"]):
        return "OVERVIEW"

    # Check for short summary
    if any(term in query for term in ["short summary", "brief summary", "concise summary", "main points", "key points", "3 key takeaways", "summary", "summarize", "summarise"]):
        return "SHORT_SUMMARY"

    # Default to SHORT_SUMMARY if unclear
    return "SHORT_SUMMARY"


def evaluate_context_quality(text):
    """Evaluate if the provided context contains enough meaningful content to answer accurately."""
    if not text or not str(text).strip():
        return False, "No text content provided."

    lines = [line.strip() for line in str(text).splitlines() if line.strip()]
    if not lines:
        return False, "No readable lines found."

    # Count low-quality indicators
    low_quality_count = 0
    total_lines = len(lines)

    low_quality_patterns = [
        r"^\s*(page\s+\d+|slide\s+\d+|table\s+\d+|figure\s+\d+|image\s+\d+)\s*$",
        r"^\s*(copyright|imprint|trademark|warranty|legal|disclaimer)\s*$",
        r"^\s*(table\s+of\s+contents|index|contents)\s*$",
        r"^\s*(header|footer|repeated|section\s+\d+)\s*$",
        r"^\s*(metadata|author|title|creation|date|version)\s*$",
        r"^\s*[-=]{3,}\s*$",
        r"^\s*\d+\.\s*$",
    ]

    for line in lines:
        lower_line = line.lower()
        if len(line) < 5 or any(re.search(pattern, lower_line, re.IGNORECASE) for pattern in low_quality_patterns):
            low_quality_count += 1

    # If more than 70% is low-quality, reject
    if low_quality_count / total_lines > 0.7:
        return False, "The provided context contains mostly metadata, headers, footers, table of contents, or repeated sections without explanatory content."

    # Check for meaningful content
    meaningful_lines = [line for line in lines if len(line) > 10 and not any(re.search(pattern, line.lower()) for pattern in low_quality_patterns)]
    if len(meaningful_lines) < 3:
        return False, "Insufficient meaningful explanatory content found."

    return True, "Context appears to contain useful content."


def join_response_blocks(blocks):
    """Join non-empty response blocks with a clean divider."""
    return "\n\n---\n\n".join(block for block in blocks if str(block or "").strip())


def init_document_summary_cache():
    """Session cache for expensive per-document summary/analysis responses."""
    if "doc_cache" not in st.session_state or not isinstance(st.session_state.doc_cache, dict):
        st.session_state.doc_cache = {}
    if "document_summary" in st.session_state and isinstance(st.session_state.document_summary, dict):
        for file_name, cached_value in st.session_state.document_summary.items():
            st.session_state.doc_cache.setdefault(file_name, cached_value)
    st.session_state.document_summary = st.session_state.doc_cache


def get_document_summary_cache_entry(file_name, text, mode):
    init_document_summary_cache()
    file_cache = st.session_state.doc_cache.get(file_name, {})
    if not isinstance(file_cache, dict):
        return None
    text_hash = hashlib.sha1(str(text or "").encode("utf-8", errors="ignore")).hexdigest()
    entry = file_cache.get(mode)
    if isinstance(entry, dict) and entry.get("text_hash") == text_hash:
        return entry.get("response")
    return None


def set_document_summary_cache_entry(file_name, text, mode, response):
    init_document_summary_cache()
    text_hash = hashlib.sha1(str(text or "").encode("utf-8", errors="ignore")).hexdigest()
    file_cache = st.session_state.doc_cache.setdefault(file_name, {})
    if not isinstance(file_cache, dict):
        file_cache = {}
        st.session_state.doc_cache[file_name] = file_cache
    file_cache[mode] = {
        "text_hash": text_hash,
        "response": response,
        "updated_at": datetime.now().isoformat(timespec="seconds"),
    }
    st.session_state.document_summary = st.session_state.doc_cache


def split_text(text, chunk_size=1000, overlap=150):
    """Split document text into word-based chunks for map-reduce summarization."""
    words = re.findall(r"\S+", str(text or ""))
    if not words:
        return []
    chunk_size = max(100, int(chunk_size or 1000))
    overlap = max(0, min(int(overlap or 0), chunk_size - 1))
    chunks = []
    start = 0
    while start < len(words):
        end = min(len(words), start + chunk_size)
        chunks.append(" ".join(words[start:end]))
        if end >= len(words):
            break
        start = end - overlap
    return chunks


def build_summary_source_text(text, max_lines=900):
    """Prepare meaningful content for summarization while removing obvious extraction noise."""
    meaningful_lines = get_meaningful_document_lines(
        text,
        min_len=12,
        max_len=320,
        limit=max_lines,
    )
    if meaningful_lines:
        return "\n".join(meaningful_lines)
    return str(text or "")[:MAX_VECTOR_TEXT_CHARS]


def invoke_document_summary_prompt(llm, prompt_template, context):
    if llm is None:
        return ""
    try:
        return str(llm.invoke(prompt_template.replace("{context}", str(context or "")))).strip()
    except Exception:
        return ""


def process_summary_chunk(llm, chunk):
    """Map-step worker for one chunk."""
    summary = invoke_document_summary_prompt(llm, CHUNK_LEVEL_SUMMARY_PROMPT, chunk)
    if summary and "minimal useful information" not in summary.lower():
        return summary
    return ""


def parallel_map_summary_chunks(llm, chunks, max_workers=6):
    """Parallel chunk processing for large documents, with sequential fallback."""
    if llm is None or not chunks:
        return []
    workers = max(1, min(int(max_workers or 1), len(chunks)))
    try:
        with ThreadPoolExecutor(max_workers=workers) as executor:
            results = list(executor.map(lambda chunk: process_summary_chunk(llm, chunk), chunks))
    except Exception:
        results = [process_summary_chunk(llm, chunk) for chunk in chunks]
    return [summary for summary in results if summary]


def generated_summary_is_useful(response):
    text = str(response or "").strip()
    lower = text.lower()
    if len(text) < 80:
        return False
    if "this information is not available in the uploaded documents" in lower:
        return False
    if lower.count("not specified in the provided context") >= 6:
        return False
    return True


def build_map_reduce_document_analysis(file_name, text):
    """Analyze a full document through chunk summaries, final reduction, and session caching."""
    cached = get_document_summary_cache_entry(file_name, text, "analysis")
    if cached:
        return cached

    source_text = build_summary_source_text(text)
    if not source_text.strip():
        response = f"**{html.escape(file_name)}**\n\nNo readable content found in this document."
        set_document_summary_cache_entry(file_name, text, "analysis", response)
        return response

    llm = load_llm()
    final_analysis = ""
    if llm is not None:
        chunks = split_text(source_text, chunk_size=1000, overlap=150)[:48]
        chunk_summaries = parallel_map_summary_chunks(llm, chunks, max_workers=6)
        if chunk_summaries:
            final_input = "\n\n".join(chunk_summaries)
            candidate = invoke_document_summary_prompt(llm, FINAL_DOCUMENT_ANALYSIS_PROMPT, final_input)
            if generated_summary_is_useful(candidate):
                final_analysis = f"### Full Analysis: {file_name}\n\n{candidate}"

    if not final_analysis:
        file_entry = get_uploaded_file_entry(file_name)
        file_bytes = file_entry["bytes"] if file_entry else b""
        final_analysis = build_product_documentation_analysis(file_name, file_bytes, text)

    set_document_summary_cache_entry(file_name, text, "analysis", final_analysis)
    return final_analysis


def build_fast_document_summary(file_name, text):
    """Build a quick cached summary using the fast prompt, with deterministic fallback."""
    cached = get_document_summary_cache_entry(file_name, text, "summary")
    if cached:
        return cached

    source_text = get_document_summary_cache_entry(file_name, text, "analysis")
    if not source_text:
        source_text = build_summary_source_text(text, max_lines=220)
    if not source_text.strip():
        response = f"No readable content found in {html.escape(file_name)}."
        set_document_summary_cache_entry(file_name, text, "summary", response)
        return response

    llm = load_llm()
    candidate = invoke_document_summary_prompt(llm, FAST_SUMMARY_PROMPT, source_text[:50000]) if llm is not None else ""
    if generated_summary_is_useful(candidate):
        response = f"### Summary: {file_name}\n\n{candidate}"
    else:
        file_entry = get_uploaded_file_entry(file_name)
        file_bytes = file_entry["bytes"] if file_entry else b""
        response = build_short_document_summary(file_name, file_bytes, text)

    set_document_summary_cache_entry(file_name, text, "summary", response)
    return response


def classify_intent(query):
    """Classify a chat request into the production document-agent intent set."""
    q = re.sub(r"\s+", " ", str(query or "").lower()).strip()
    compact = re.sub(r"[^a-z0-9]+", " ", q).strip()
    if compact in {"analyze", "analyse", "analysis", "full analysis"}:
        return "FULL_DOCUMENT_ANALYSIS"
    if compact in {"summary", "summarize", "summarise", "short summary"}:
        return "SHORT_SUMMARY"
    if compact == "overview" or re.search(r"\boverview\b", q):
        return "OVERVIEW"
    if any(term in q for term in ["find", "search", "highlight", "count", "occurrence", "locate"]):
        return "SEARCH"
    if any(term in q for term in ["compare", "difference", "versus", " vs "]):
        return "COMPARISON"
    if any(term in q for term in ["table", "spreadsheet", "sheet", "row", "column", "csv", "data"]):
        return "TABLE_EXTRACTION"
    if any(term in q for term in ["diagram", "flowchart", "flow", "figure", "image", "visual", "pin diagram"]):
        return "DIAGRAM_ANALYSIS"
    if any(term in q for term in ["component", "module", "interface", "connector", "device", "part"]):
        return "COMPONENT_EXTRACTION"
    return "QUESTION_ANSWERING"


def handle_query(user_input, combined_text, file_name):
    """Master router for a single document where no vector file list is available."""
    intent = classify_intent(user_input)
    if intent == "FULL_DOCUMENT_ANALYSIS":
        return build_map_reduce_document_analysis(file_name, combined_text)
    if intent == "SHORT_SUMMARY":
        return build_fast_document_summary(file_name, combined_text)
    if intent == "OVERVIEW":
        return build_overview_response({file_name: combined_text})
    if intent == "TABLE_EXTRACTION":
        return build_table_extraction_response({file_name: combined_text})
    if intent == "DIAGRAM_ANALYSIS":
        return build_image_or_diagram_extraction_response({file_name: combined_text}, user_input)
    if intent == "COMPONENT_EXTRACTION":
        return build_component_extraction_response({file_name: combined_text}, user_input)
    if intent == "COMPARISON":
        return build_component_comparison_response({file_name: combined_text}, user_input)
    if intent == "SEARCH":
        return build_extraction_response_for_query(user_input, {file_name: combined_text})
    return None


def handle_document_chat_query(user_input, file_texts, file_names=None, user_id=None):
    """Agent router backed by the file-brain query engine."""
    file_texts = file_texts or {}
    target_files = file_names or list(file_texts.keys())
    intent = classify_technical_document_request(user_input)
    return smart_file_brain_query(user_input, target_files, user_id=user_id, intent=intent, top_k=8)


def infer_response_confidence(response, file_texts=None, citation_docs=None):
    """Estimate answer confidence from context availability and answer quality."""
    response_text = str(response or "").lower()
    context_chars = sum(len(str(text or "")) for text in (file_texts or {}).values())
    citation_count = len(citation_docs or [])
    if not context_chars or "no readable content" in response_text:
        return "Low"
    if "not specified in the provided context" in response_text:
        return "Medium" if context_chars > 2000 else "Low"
    if "this information is not available" in response_text or "insufficient meaningful" in response_text:
        return "Low"
    if citation_count >= 2 or context_chars > 8000:
        return "High"
    return "Medium"


def append_confidence_to_response(response, file_texts=None, citation_docs=None):
    """Append the required confidence footer if it is not already present."""
    response = str(response or "").rstrip()
    if re.search(r"(?im)^confidence:\s*(high|medium|low)\s*$", response):
        return response
    confidence = infer_response_confidence(response, file_texts=file_texts, citation_docs=citation_docs)
    return response + f"\n\nConfidence: {confidence}"


def build_component_extraction_response(file_texts, user_query=""):
    """Extract named components/modules/devices and summarize their available roles."""
    requested_item = extract_specific_component_name(user_query) or extract_quoted_item_name(user_query) or extract_bare_item_name(user_query)
    if requested_item:
        return build_specific_component_response(file_texts, user_query)

    blocks = []
    for file_name, text in (file_texts or {}).items():
        meaningful = get_meaningful_document_lines(text, min_len=12, max_len=280, limit=260)
        if not meaningful:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable component or module information found.")
            continue

        candidates = []
        seen = set()
        for line in meaningful:
            matches = re.findall(
                r"\b(?:[A-Z]{2,}[A-Za-z0-9_+\-/]{1,}|[A-Z][A-Za-z0-9_+\-/]+(?:\s+[A-Z][A-Za-z0-9_+\-/]+){0,2})\b",
                line,
            )
            if re.search(r"\b(component|module|device|interface|connector|unit|adapter|channel|driver|application|tool)\b", line, re.IGNORECASE):
                words = re.findall(r"\b[A-Z][A-Za-z0-9_+\-/]{2,}\b", line)
                matches.extend(words[:2])
            for candidate in matches:
                name = normalize_extracted_line(candidate).strip(" :-")
                lower_name = name.lower()
                if (
                    len(name) < 3
                    or len(name) > 48
                    or lower_name in seen
                    or lower_name in SUMMARY_STOPWORDS
                    or lower_name in {"pdf", "table", "figure", "page", "chapter", "section"}
                ):
                    continue
                seen.add(lower_name)
                candidates.append(name)
                break
            if len(candidates) >= 12:
                break

        if not candidates:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo explicit components or modules were identified in the provided context.")
            continue

        rows = []
        for component in candidates[:10]:
            context_lines = collect_item_context_lines(text, component, window=4, limit=30)
            if not context_lines:
                context_lines = [line for line in meaningful if component.lower() in line.lower()][:8]
            purpose = select_relevant_lines(context_lines, ["purpose", "used", "provides", "supports", "enables", "allows"], limit=2)
            features = select_relevant_lines(context_lines, ["feature", "capability", "function", "supports", "configuration"], limit=2)
            interfaces = select_relevant_lines(context_lines, ["interface", "connector", "port", "channel", "pin", "network"], limit=2)
            rows.append([
                component,
                "; ".join(purpose[:2]) if purpose else "Not specified in the provided context",
                "; ".join(features[:2]) if features else "Not specified in the provided context",
                "; ".join(interfaces[:2]) if interfaces else "Not specified in the provided context",
            ])

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Components / Modules: {html.escape(file_name)}</h3>"
            f"{html_table(['Name', 'Purpose', 'Key features', 'Interfaces'], rows)}"
            f"</div>"
        )
    return join_response_blocks(blocks)


def build_overview_response(file_texts):
    """Provide high-level overview: What it is, Who it is for, What it is used for, Main concept, Main areas covered."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        meaningful = get_meaningful_document_lines(text, min_len=18, max_len=260, limit=180)
        if len(meaningful) < 3:
            blocks.append(f"**{html.escape(file_name)}**\n\nInsufficient meaningful explanatory content found.")
            continue

        what_it_is = ""
        for line in meaningful[:30]:
            if any(term in line.lower() for term in ["is a", "provides", "enables", "supports", "system", "device", "module", "interface"]):
                what_it_is = line
                break
        if not what_it_is:
            what_it_is = "This document provides technical reference and operational guidance."

        who_for = ""
        for line in meaningful:
            if any(term in line.lower() for term in ["for users", "target", "intended for", "engineers", "developers", "operators"]):
                who_for = line
                break
        if not who_for:
            who_for = "Technical professionals and system integrators."

        used_for = ""
        for line in meaningful:
            if any(term in line.lower() for term in ["used to", "used for", "application", "purpose", "function"]):
                used_for = line
                break
        if not used_for:
            used_for = "Configuration, operation, and troubleshooting of technical systems."

        main_concept = ""
        for line in meaningful:
            if any(term in line.lower() for term in ["concept", "overview", "introduction", "main idea"]):
                main_concept = line
                break
        if not main_concept:
            main_concept = "Technical specification and operational reference."

        main_areas = []
        seen = set()
        for line in meaningful:
            lower_line = line.lower()
            if any(term in lower_line for term in ["features", "capabilities", "components", "interfaces", "configuration", "usage", "troubleshooting"]) and lower_line not in seen:
                main_areas.append(line)
                seen.add(lower_line)
            if len(main_areas) >= 5:
                break
        if not main_areas:
            main_areas = ["Technical specifications", "Operational guidance", "Configuration details"]

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Overview: {html.escape(file_name)}</h3>"
            f"<p><b>What it is:</b> {html.escape(what_it_is)}</p>"
            f"<p><b>Who it is for:</b> {html.escape(who_for)}</p>"
            f"<p><b>What it is used for:</b> {html.escape(used_for)}</p>"
            f"<p><b>Main concept:</b> {html.escape(main_concept)}</p>"
            f"<p><b>Main areas covered:</b> {', '.join(html.escape(area) for area in main_areas)}</p>"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_features_only_response(file_texts):
    """Extract actual functional features and capabilities in a table format."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        quality_ok, quality_msg = evaluate_context_quality(text)
        if not quality_ok:
            blocks.append(f"**{html.escape(file_name)}**\n\n{quality_msg}")
            continue

        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 10 <= len(line) <= 250]

        features = []
        for line in meaningful:
            lower_line = line.lower()
            if any(term in lower_line for term in ["feature", "function", "capability", "enable", "allows", "provide", "interface", "communication", "diagnostic", "support"]):
                # Parse feature, what it does, why it matters, related component
                feature_name = line.split(":")[0].strip() if ":" in line else line[:50].strip()
                description = line.split(":", 1)[1].strip() if ":" in line else "Provides specific functionality."
                why_matters = "Enhances system capabilities."  # Placeholder, as not always specified
                related = "Not specified in the provided context."
                features.append({
                    "feature": feature_name,
                    "what_it_does": description,
                    "why_it_matters": why_matters,
                    "related_component": related,
                })
            if len(features) >= 10:
                break

        if not features:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo explicit features were identified in the provided context.")
            continue

        table_rows = "".join(
            f"<tr><td>{html.escape(f['feature'])}</td><td>{html.escape(f['what_it_does'])}</td><td>{html.escape(f['why_it_matters'])}</td><td>{html.escape(f['related_component'])}</td></tr>"
            for f in features
        )

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Features: {html.escape(file_name)}</h3>"
            "<table style='border-collapse:collapse; width:100%; margin:8px 0;'>"
            "<thead><tr><th>Feature</th><th>What it does</th><th>Why it matters</th><th>Related component/module</th></tr></thead>"
            f"<tbody>{table_rows}</tbody>"
            "</table>"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_workflow_or_process_response(file_texts):
    """Provide process overview, step-by-step workflow, inputs, outputs, tools/components involved, practical notes."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        quality_ok, quality_msg = evaluate_context_quality(text)
        if not quality_ok:
            blocks.append(f"**{html.escape(file_name)}**\n\n{quality_msg}")
            continue

        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 10 <= len(line) <= 250]

        process_overview = ""
        for line in meaningful[:10]:
            if any(term in line.lower() for term in ["process", "workflow", "overview", "steps", "procedure"]):
                process_overview = line
                break
        if not process_overview:
            process_overview = "Describes operational procedures and workflows."

        steps = []
        for line in meaningful:
            if re.match(r"^\s*\d+\.", line) or any(term in line.lower() for term in ["step", "first", "then", "next", "finally"]):
                steps.append(line)
            if len(steps) >= 8:
                break

        inputs = []
        outputs = []
        tools = []
        for line in meaningful:
            lower_line = line.lower()
            if "input" in lower_line:
                inputs.append(line)
            if "output" in lower_line:
                outputs.append(line)
            if any(term in lower_line for term in ["tool", "component", "module", "interface", "software", "hardware"]):
                tools.append(line)
            if len(inputs) >= 5 and len(outputs) >= 5 and len(tools) >= 5:
                break

        practical_notes = []
        for line in meaningful:
            if any(term in line.lower() for term in ["note", "important", "caution", "tip", "warning"]):
                practical_notes.append(line)
            if len(practical_notes) >= 5:
                break

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Workflow / Process: {html.escape(file_name)}</h3>"
            f"<p><b>Process overview:</b> {html.escape(process_overview)}</p>"
            f"{f'<p><b>Step-by-step workflow:</b></p><ul>{''.join(f'<li>{html.escape(step)}</li>' for step in steps[:10])}</ul>' if steps else '<p><b>Step-by-step workflow:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Inputs:</b></p><ul>{''.join(f'<li>{html.escape(inp)}</li>' for inp in inputs[:5])}</ul>' if inputs else '<p><b>Inputs:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Outputs:</b></p><ul>{''.join(f'<li>{html.escape(out)}</li>' for out in outputs[:5])}</ul>' if outputs else '<p><b>Outputs:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Tools/components involved:</b></p><ul>{''.join(f'<li>{html.escape(tool)}</li>' for tool in tools[:5])}</ul>' if tools else '<p><b>Tools/components involved:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Practical notes:</b></p><ul>{''.join(f'<li>{html.escape(note)}</li>' for note in practical_notes[:5])}</ul>' if practical_notes else '<p><b>Practical notes:</b> Not specified in the provided context.</p>'}"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_use_cases_applications_response(file_texts):
    """Provide primary use cases, real-world applications, target users, benefits, example scenarios."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        quality_ok, quality_msg = evaluate_context_quality(text)
        if not quality_ok:
            blocks.append(f"**{html.escape(file_name)}**\n\n{quality_msg}")
            continue

        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 10 <= len(line) <= 250]

        use_cases = []
        applications = []
        target_users = []
        benefits = []
        scenarios = []

        for line in meaningful:
            lower_line = line.lower()
            if "use case" in lower_line or "used to" in lower_line or "application" in lower_line:
                use_cases.append(line)
            if "real-world" in lower_line or "practical" in lower_line or "industry" in lower_line:
                applications.append(line)
            if "user" in lower_line or "engineer" in lower_line or "developer" in lower_line or "operator" in lower_line:
                target_users.append(line)
            if "benefit" in lower_line or "advantage" in lower_line or "improve" in lower_line:
                benefits.append(line)
            if "example" in lower_line or "scenario" in lower_line or "for instance" in lower_line:
                scenarios.append(line)
            if len(use_cases) >= 5 and len(applications) >= 5 and len(target_users) >= 5 and len(benefits) >= 5 and len(scenarios) >= 5:
                break

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Use Cases / Applications: {html.escape(file_name)}</h3>"
            f"{f'<p><b>Primary use cases:</b></p><ul>{''.join(f'<li>{html.escape(uc)}</li>' for uc in use_cases[:5])}</ul>' if use_cases else '<p><b>Primary use cases:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Real-world applications:</b></p><ul>{''.join(f'<li>{html.escape(app)}</li>' for app in applications[:5])}</ul>' if applications else '<p><b>Real-world applications:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Target users:</b></p><ul>{''.join(f'<li>{html.escape(user)}</li>' for user in target_users[:5])}</ul>' if target_users else '<p><b>Target users:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Benefits:</b></p><ul>{''.join(f'<li>{html.escape(ben)}</li>' for ben in benefits[:5])}</ul>' if benefits else '<p><b>Benefits:</b> Not specified in the provided context.</p>'}"
            f"{f'<p><b>Example scenarios:</b></p><ul>{''.join(f'<li>{html.escape(scen)}</li>' for scen in scenarios[:5])}</ul>' if scenarios else '<p><b>Example scenarios:</b> Not specified in the provided context.</p>'}"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_troubleshooting_or_limitations_response(file_texts):
    """Provide problems/limitations, causes, constraints, actions, unspecified items."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        quality_ok, quality_msg = evaluate_context_quality(text)
        if not quality_ok:
            blocks.append(f"**{html.escape(file_name)}**\n\n{quality_msg}")
            continue

        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 10 <= len(line) <= 250]

        problems = []
        causes = []
        constraints = []
        actions = []
        unspecified = []

        for line in meaningful:
            lower_line = line.lower()
            if any(term in lower_line for term in ["problem", "issue", "error", "failure", "limitation", "constraint", "troubleshooting"]):
                problems.append(line)
            if "cause" in lower_line or "due to" in lower_line or "because" in lower_line:
                causes.append(line)
            if "constraint" in lower_line or "limitation" in lower_line or "cannot" in lower_line or "not supported" in lower_line:
                constraints.append(line)
            if "action" in lower_line or "fix" in lower_line or "resolve" in lower_line or "recommend" in lower_line:
                actions.append(line)
            if "not specified" in lower_line or "unknown" in lower_line or "not available" in lower_line:
                unspecified.append(line)
            if len(problems) >= 5 and len(causes) >= 5 and len(constraints) >= 5 and len(actions) >= 5 and len(unspecified) >= 5:
                break

        table_rows = ""
        for i, prob in enumerate(problems[:10]):
            cause = causes[i] if i < len(causes) else "Not specified in the provided context."
            constraint = constraints[i] if i < len(constraints) else "Not specified in the provided context."
            action = actions[i] if i < len(actions) else "Not specified in the provided context."
            table_rows += f"<tr><td>{html.escape(prob)}</td><td>{html.escape(cause)}</td><td>{html.escape(constraint)}</td><td>{html.escape(action)}</td></tr>"

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Troubleshooting / Limitations: {html.escape(file_name)}</h3>"
            "<table style='border-collapse:collapse; width:100%; margin:8px 0;'>"
            "<thead><tr><th>Problem / Limitation</th><th>Likely Cause</th><th>Relevant Constraints</th><th>Recommended Action</th></tr></thead>"
            f"<tbody>{table_rows}</tbody>"
            "</table>"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_requirements_or_specification_extraction_response(file_texts):
    """Extract requirements/specifications in a structured table."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        quality_ok, quality_msg = evaluate_context_quality(text)
        if not quality_ok:
            blocks.append(f"**{html.escape(file_name)}**\n\n{quality_msg}")
            continue

        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 10 <= len(line) <= 250]

        specs = []
        id_counter = 1
        for line in meaningful:
            lower_line = line.lower()
            if any(term in lower_line for term in ["requirement", "spec", "must", "shall", "should", "specification", "parameter", "value", "condition"]):
                req = line
                category = "Functional"  # Default
                if "performance" in lower_line:
                    category = "Performance"
                elif "interface" in lower_line or "connector" in lower_line:
                    category = "Interface"
                elif "environmental" in lower_line:
                    category = "Environmental"
                applies_to = "System"  # Default
                value = "Not specified"
                notes = "Extracted from document"
                specs.append({
                    "id": f"REQ-{id_counter}",
                    "requirement": req,
                    "category": category,
                    "applies_to": applies_to,
                    "value": value,
                    "notes": notes,
                })
                id_counter += 1
            if len(specs) >= 15:
                break

        if not specs:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo explicit requirements or specifications were identified in the provided context.")
            continue

        table_rows = "".join(
            f"<tr><td>{html.escape(s['id'])}</td><td>{html.escape(s['requirement'])}</td><td>{html.escape(s['category'])}</td><td>{html.escape(s['applies_to'])}</td><td>{html.escape(s['value'])}</td><td>{html.escape(s['notes'])}</td></tr>"
            for s in specs
        )

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Requirements / Specifications: {html.escape(file_name)}</h3>"
            "<table style='border-collapse:collapse; width:100%; margin:8px 0;'>"
            "<thead><tr><th>ID</th><th>Requirement / Specification</th><th>Category</th><th>Applies to</th><th>Value / Condition</th><th>Notes</th></tr></thead>"
            f"<tbody>{table_rows}</tbody>"
            "</table>"
            "</div>"
        )
    return join_response_blocks(blocks)


def build_full_document_summary_response(file_texts):
    """Build cached map-reduce document analysis for selected files."""
    blocks = []
    for file_name, file_text in (file_texts or {}).items():
        if file_text and str(file_text).strip():
            blocks.append(build_map_reduce_document_analysis(file_name, file_text))
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable content found in this document.")
    return join_response_blocks(blocks)


def build_short_document_summary(file_name, file_bytes, text):
    raw_text = str(text or "")
    lines = get_meaningful_document_lines(raw_text, min_len=18, max_len=260, limit=180)
    if not lines:
        return f"No readable content found in {html.escape(file_name)}."

    title = get_document_display_title(file_name, raw_text)

    main_purpose = ""
    for line in lines:
        ll = line.lower()
        if any(term in ll for term in ["purpose", "provides", "supports", "enables", "used for", "application", "allows", "designed", "test system", "assembled"]):
            main_purpose = line
            break
    if not main_purpose:
        main_purpose = lines[0]

    key_points = []
    seen = set()
    for line in lines:
        ll = line.lower()
        if any(term in ll for term in ["provides", "supports", "enables", "used for", "feature", "capability", "function", "application", "module", "interface", "measurement", "simulation", "integration", "configuration", "test"]):
            cleaned = line.strip()
            if cleaned.lower() not in seen:
                key_points.append(cleaned)
                seen.add(cleaned.lower())
        if len(key_points) >= 5:
            break
    if not key_points:
        key_points = lines[:5]

    key_takeaways = []
    if key_points:
        key_takeaways.append("The document describes the selected system from a practical technical/product-information perspective.")
        key_takeaways.append("The most useful content is in the meaningful feature, application, module, and usage sections rather than the cover pages or table of contents.")
        if any("CANoe" in point or "ECU" in point for point in key_points):
            key_takeaways.append("The document is relevant to ECU test setups, vehicle-network testing, and CANoe-integrated workflows.")
    key_takeaways = key_takeaways[:3] or key_points[:3]
    key_points_html = "".join(f"<li>{html.escape(point)}</li>" for point in key_points[:5])
    takeaways_html = "".join(f"<li>{html.escape(point)}</li>" for point in key_takeaways)

    return (
        f"<div style='margin-bottom:18px; line-height:1.5;'>"
        f"<p><b>What it is:</b> {html.escape(title)}</p>"
        f"<p><b>Purpose:</b> {html.escape(main_purpose)}</p>"
        f"<p><b>Key insights:</b></p><ul>{key_points_html}</ul>"
        f"<p><b>Key takeaways:</b></p><ul>{takeaways_html}</ul>"
        f"</div>"
    )


def build_short_summary_response(file_texts):
    blocks = []
    for file_name, file_text in (file_texts or {}).items():
        if file_text and str(file_text).strip():
            blocks.append(build_fast_document_summary(file_name, file_text))
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable content found in this document.")
    return join_response_blocks(blocks)


def build_table_extraction_response(file_texts):
    blocks = []
    for file_name, text in (file_texts or {}).items():
        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        table_lines = select_relevant_lines(lines, ["table", "row", "column", "csv", "sheet", "spreadsheet", "cells", "header", "entry"], limit=10)
        if table_lines:
            rows = "".join(f"<li>{html.escape(line)}</li>" for line in table_lines)
            blocks.append(
                f"<div style='margin-bottom:18px; line-height:1.5;'>"
                f"<h3 style='margin:0 0 10px 0; color:#173152;'>Table Extraction: {html.escape(file_name)}</h3>"
                f"<p>Extracted table or tabular structure lines from the document text.</p>"
                f"<ul>{rows}</ul>"
                f"</div>"
            )
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo table-like data was found in the extracted document text.")
    return join_response_blocks(blocks)


def build_image_or_diagram_extraction_response(file_texts, user_query):
    blocks = []
    for file_name, text in (file_texts or {}).items():
        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        image_lines = select_relevant_lines(lines, ["figure", "image", "diagram", "illustration", "schematic", "drawing", "visual"], limit=12)
        pin_rows = extract_pin_rows(lines)
        ascii_diagram = build_ascii_pin_diagram(pin_rows, os.path.splitext(file_name)[0]) if pin_rows else None
        if image_lines or ascii_diagram:
            image_block = "".join(f"<li>{html.escape(line)}</li>" for line in image_lines) if image_lines else "<li>No direct image references were found.</li>"
            diagram_block = f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>" if ascii_diagram else ""
            blocks.append(
                f"<div style='margin-bottom:18px; line-height:1.5;'>"
                f"<h3 style='margin:0 0 10px 0; color:#173152;'>Image / Diagram Extraction: {html.escape(file_name)}</h3>"
                f"<h4 style='margin:16px 0 6px 0; color:#173152;'>Image / Figure References</h4><ul>{image_block}</ul>"
                f"{diagram_block}"
                f"</div>"
            )
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo image or diagram references were found in the readable document content.")
    return join_response_blocks(blocks)


def build_strict_extraction_response(file_texts, user_query):
    """Return direct verbatim extraction output with no explanation or added structure."""
    blocks = []
    query_lower = str(user_query or "").lower()
    for file_name, text in (file_texts or {}).items():
        raw_lines = [line.rstrip() for line in str(text or "").splitlines()]
        if not raw_lines:
            blocks.append("Not available in the document")
            continue

        if any(term in query_lower for term in ["table", "csv", "spreadsheet", "tabular", "rows", "columns", "column"]):
            selected = [line for line in raw_lines if any(keyword in line.lower() for keyword in ["table", "row", "column", "csv", "spreadsheet", "cell", "header", "|", ","]) ]
        elif any(term in query_lower for term in ["image", "diagram", "visual", "figure", "schematic", "illustration", "drawing", "pin", "connector"]):
            selected = [line for line in raw_lines if any(keyword in line.lower() for keyword in ["figure", "image", "diagram", "schematic", "illustration", "drawing", "pin", "connector", "socket", "port", "cable"])]
        else:
            selected = raw_lines

        if not selected:
            blocks.append("Not available in the document")
        else:
            blocks.append("\n".join(selected).strip())

    return "\n\n".join(blocks)


def build_downloadable_report_response(file_texts):
    blocks = []
    for file_name, file_text in (file_texts or {}).items():
        if file_text and str(file_text).strip():
            blocks.append(
                f"<div style='margin-bottom:18px; line-height:1.5;'>"
                f"<h3 style='margin:0 0 10px 0; color:#173152;'>Downloadable Report: {html.escape(file_name)}</h3>"
                f"<p>This response is prepared for export-style delivery. Use the document preview Downloads tab to generate a DOCX or Markdown report from the readable content.</p>"
                f"</div>"
            )
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable content found to build a downloadable report.")
    return join_response_blocks(blocks)


def build_specific_component_response(file_texts, user_query):
    """Answer only for the requested component, ignoring unrelated document content."""
    component_name = extract_specific_component_name(user_query)
    if not component_name:
        return "Which specific component, module, device, or item should I focus on?"

    blocks = []
    for file_name, file_text in (file_texts or {}).items():
        if file_text and str(file_text).strip():
            blocks.append(build_item_information_response(file_name, file_text, component_name))
        else:
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable content found in this document.")
    return join_response_blocks(blocks)


def build_document_visual_response(file_name, text, item_name=None):
    """Build document-wide pin/connector/table output when no exact item is supplied."""
    context_lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
    if item_name:
        context_lines = collect_item_context_lines(text, item_name, window=8, limit=160) or context_lines

    pin_rows = extract_pin_rows(context_lines)
    connector_lines = select_relevant_lines(context_lines, ["connector", "port", "d-sub", "usb", "channel", "plug", "socket", "interface"], limit=14)
    table_lines = select_relevant_lines(context_lines, ["table", "pin", "signal", "configuration", "mapping", "assignment"], limit=14)
    visual_lines = select_relevant_lines(context_lines, ["figure", "image", "diagram", "layout", "visual", "pin assignment"], limit=10)
    display_name = item_name or os.path.splitext(file_name)[0]
    ascii_diagram = build_ascii_pin_diagram(pin_rows, display_name)
    pin_table_rows = [[row["pin"], row["signal"], row["description"], row.get("notes", "")] for row in pin_rows]

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Diagrams / Pin Details: {html.escape(display_name)}</h3>",
        f"<p><b>Source:</b> {html.escape(file_name)}</p>",
        "<h4 style='margin:16px 0 6px 0; color:#173152;'>Pin Table</h4>",
        html_table(["Pin Number", "Signal Name", "Description", "Notes"], pin_table_rows) if pin_rows else "<p>No explicit pin rows were found in the provided document context.</p>",
        "<h4 style='margin:16px 0 6px 0; color:#173152;'>Diagram</h4>",
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>",
        html_section("Connector Mapping", connector_lines),
        html_section("Tables / Structured References", table_lines),
        html_section("Visual Structure", visual_lines),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(section for section in sections if section) + "</div>"


def build_diagram_pin_details_response(file_texts, user_query):
    """Build pin/diagram/table focused output and CSV/diagram downloads."""
    item_name = extract_specific_component_name(user_query) or extract_quoted_item_name(user_query)
    blocks = []
    csv_downloads = []
    diagram_downloads = []
    for file_name, file_text in (file_texts or {}).items():
        if not str(file_text or "").strip():
            blocks.append(f"**{html.escape(file_name)}**\n\nNo readable content found in this document.")
            continue
        if item_name:
            blocks.append(build_item_visual_response(file_name, file_text, item_name))
            visual_assets = build_item_visual_assets(file_name, file_text, item_name)
        else:
            blocks.append(build_document_visual_response(file_name, file_text))
            pin_rows = extract_pin_rows([normalize_extracted_line(line) for line in str(file_text).splitlines() if line.strip()])
            safe_name = re.sub(r"[^A-Za-z0-9_-]+", "_", os.path.splitext(file_name)[0]).strip("_") or "document"
            visual_assets = {
                "csv": [{
                    "label": f"{file_name} - pin table CSV",
                    "data": build_pin_csv(pin_rows).encode("utf-8"),
                    "file_name": f"{safe_name}_pin_table.csv",
                    "mime": "text/csv",
                }] if pin_rows else [],
                "diagrams": [{
                    "label": f"{file_name} - ASCII diagram",
                    "data": build_ascii_pin_diagram(pin_rows, safe_name).encode("utf-8"),
                    "file_name": f"{safe_name}_diagram.txt",
                    "mime": "text/plain",
                }] if pin_rows else [],
            }
        csv_downloads.extend(visual_assets.get("csv", []))
        diagram_downloads.extend(visual_assets.get("diagrams", []))
    return join_response_blocks(blocks), csv_downloads, diagram_downloads


def build_features_workflow_response(file_texts):
    """Build a functional response: features, capabilities, workflow, inputs/outputs, applications, benefits."""
    blocks = []
    for file_name, text in (file_texts or {}).items():
        lower_text = str(text or "").lower()
        lines = [normalize_extracted_line(line) for line in str(text or "").splitlines() if line.strip()]
        meaningful = [line for line in lines if 12 <= len(line) <= 220]

        def collect(terms, fallback, limit=7):
            selected = []
            seen = set()
            for line in meaningful:
                lower_line = line.lower()
                if any(term in lower_line for term in terms) and lower_line not in seen:
                    selected.append(line)
                    seen.add(lower_line)
                if len(selected) >= limit:
                    break
            return selected or fallback

        features = collect(
            ["feature", "function", "capability", "enable", "allows", "provide", "interface", "communication", "diagnostic"],
            ["Provides functional reference information for understanding and using the documented system."],
        )
        capabilities = collect(
            ["support", "capability", "function", "enable", "allows", "provide", "interface", "communication", "diagnostic", "configuration"],
            ["Supports technical reference and operational guidance."],
        )
        workflow = collect(
            ["install", "configure", "connect", "select", "execute", "run", "start", "use", "download", "export", "review"],
            ["Identify the relevant function, configure required inputs, execute the workflow, then review or export results."],
        )
        inputs_outputs = collect(
            ["input", "output", "signal", "data", "channel", "port", "interface", "protocol", "format"],
            ["Accepts configuration inputs and produces structured outputs or reports."],
        )
        applications = collect(
            ["application", "used for", "used to", "use case", "measurement", "testing", "diagnostic", "monitoring", "analysis", "report"],
            ["Technical reference, configuration planning, validation, troubleshooting, and documentation support."],
        )
        benefits = collect(
            ["benefit", "advantage", "improve", "enhance", "optimize", "efficient", "reliable", "accurate", "fast", "easy"],
            ["Provides reliable technical reference and operational efficiency."],
        )

        blocks.append(
            f"<div style='margin-bottom:18px; line-height:1.5;'>"
            f"<h3 style='margin:0 0 10px 0; color:#173152;'>Functional Analysis: {html.escape(file_name)}</h3>"
            f"{html_section('Features', features[:7])}"
            f"{html_section('Capabilities', capabilities[:7])}"
            f"{html_section('Workflow', workflow[:7])}"
            f"{html_section('Inputs/Outputs', inputs_outputs[:7])}"
            f"{html_section('Applications', applications[:7])}"
            f"{html_section('Benefits', benefits[:7])}"
            f"</div>"
        )
    return join_response_blocks(blocks)


def build_component_comparison_response(file_texts, user_query):
    """Compare named components/items inside selected documents without repeating shared content."""
    items = extract_multiple_component_names(user_query)
    if len(items) < 2:
        return "Which two or more components/items should I compare?"

    sections = [
        "<div style='margin-bottom:18px; line-height:1.5;'>",
        "<h3 style='margin:0 0 10px 0; color:#173152;'>Comparison</h3>",
        "<table style='border-collapse:collapse; width:100%; margin:8px 0;'>",
        "<thead><tr><th>Item</th><th>Purpose / Context</th><th>Technical Signals</th><th>Interfaces / Notes</th></tr></thead><tbody>",
    ]

    for item in items:
        item_lines = []
        for file_text in (file_texts or {}).values():
            item_lines.extend(collect_item_context_lines(file_text, item, window=5, limit=60))
        item_lines = list(dict.fromkeys(item_lines))
        purpose = select_relevant_lines(item_lines, ["used", "purpose", "application", "support", "provide", "allows"], limit=3) or item_lines[:3]
        technical = select_relevant_lines(item_lines, ["channel", "protocol", "mbit", "kbit", "volt", "can", "lin", "flexray", "ethernet", "diagnostic"], limit=4)
        interfaces = select_relevant_lines(item_lines, ["connector", "interface", "port", "pin", "d-sub", "usb", "configuration"], limit=4)
        sections.append(
            "<tr>"
            f"<td><b>{html.escape(item)}</b></td>"
            f"<td>{html.escape('; '.join(purpose[:3]) if purpose else 'No focused purpose context found.')}</td>"
            f"<td>{html.escape('; '.join(technical[:4]) if technical else 'No focused technical signal found.')}</td>"
            f"<td>{html.escape('; '.join(interfaces[:4]) if interfaces else 'No focused interface note found.')}</td>"
            "</tr>"
        )

    sections.extend(["</tbody></table>"])

    # Add additional sections
    all_item_lines = []
    for item in items:
        for file_text in (file_texts or {}).values():
            all_item_lines.extend(collect_item_context_lines(file_text, item, window=5, limit=60))
    all_item_lines = list(dict.fromkeys(all_item_lines))

    similarities = select_relevant_lines(all_item_lines, ["same", "similar", "common", "shared", "both", "equivalent"], limit=5)
    differences = select_relevant_lines(all_item_lines, ["different", "differs", "unique", "specific", "only", "versus", "vs"], limit=5)
    key_insights = select_relevant_lines(all_item_lines, ["important", "key", "note", "critical", "main", "primary"], limit=5)
    best_use = select_relevant_lines(all_item_lines, ["best for", "recommended", "ideal", "suitable", "use case", "application"], limit=5)

    if similarities:
        sections.append(html_section("Similarities", similarities))
    if differences:
        sections.append(html_section("Differences", differences))
    if key_insights:
        sections.append(html_section("Key Insights", key_insights))
    if best_use:
        sections.append(html_section("Best-Use Scenarios", best_use))

    sections.append("</div>")
    return "".join(sections)


def build_adaptive_document_analysis(file_name, file_bytes, text):
    raw_text = str(text or "")
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    words = re.findall(r"\w+", raw_text)
    title_match = re.search(r"Title:\s*(.+)", raw_text)
    title = title_match.group(1).strip() if title_match and title_match.group(1).strip() else file_name

    keyword_counts = Counter(
        word.lower()
        for word in words
        if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    keywords = [word.title() for word, _ in keyword_counts.most_common(8)]
    keyword_text = ", ".join(keywords) if keywords else "Not available"
    page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, raw_text)

    ignored_prefixes = (
        "pdf metadata:", "document metadata:", "meta tags:", "total pages:", "total slides:",
        "workbook contains", "error:", "[image:", "[embedded_image:", "table:"
    )
    metadata_prefixes = (
        "producer:", "creationdate:", "moddate:", "author:", "creator:", "title:",
        "subject:", "keywords:", "trapped:", "pdfversion:"
    )

    def prettify_extracted_text(value):
        value = str(value or "").strip()
        if not value:
            return value
        value = re.sub(r"([a-z])([A-Z])", r"\1 \2", value)
        value = re.sub(r"([A-Za-z])(\d)", r"\1 \2", value)
        value = re.sub(r"(\d)([A-Za-z])", r"\1 \2", value)
        value = re.sub(r"\s+", " ", value)
        value = value.replace("e. g.", "e.g.").replace("i. e.", "i.e.")
        return value.strip()

    keywords = [prettify_extracted_text(keyword) for keyword in keywords]
    keywords = [keyword for keyword in keywords if not _is_low_signal_theme(keyword)]
    keyword_text = ", ".join(keywords) if keywords else "Not available"

    def clean_content_lines(max_items=12):
        cleaned = []
        seen = set()
        for line in lines:
            line_lower = line.lower()
            if line_lower.startswith(ignored_prefixes):
                continue
            if line_lower.startswith(metadata_prefixes):
                continue
            if len(line) < 8 or len(line) > 240:
                continue
            if re.fullmatch(r"[\W_]+", line):
                continue
            line = prettify_extracted_text(line)
            line_lower = line.lower()
            if line in seen:
                continue
            seen.add(line)
            cleaned.append(line)
            if len(cleaned) >= max_items:
                break
        return cleaned

    key_lines = clean_content_lines(12)
    headings = extract_document_headings(raw_text)
    toc_entries = extract_toc_with_page_numbers(raw_text)
    lower_text = raw_text.lower()
    file_name_lower = file_name.lower()
    type_scores = {
        "technical": sum(1 for term in [
            "architecture", "system", "module", "component", "workflow", "api", "interface",
            "configuration", "requirement", "software", "hardware", "capl", "diagnostic", "test"
        ] if term in lower_text),
        "business": sum(1 for term in [
            "strategy", "market", "customer", "revenue", "business", "goal", "objective",
            "stakeholder", "risk", "cost", "benefit", "performance", "operation"
        ] if term in lower_text),
        "research": sum(1 for term in [
            "abstract", "methodology", "experiment", "hypothesis", "dataset", "findings",
            "results", "conclusion", "references", "study", "analysis"
        ] if term in lower_text),
    }
    if file_name_lower.endswith((".can", ".capl")):
        type_scores["technical"] += 4
    if file_name_lower.endswith((".xlsx", ".html", ".htm")):
        type_scores["business"] += 1
    document_type = max(type_scores, key=type_scores.get) if max(type_scores.values() or [0]) > 0 else "general"

    def pick_lines(patterns, limit=5):
        selected = []
        seen = set()
        for line in key_lines + lines:
            line = prettify_extracted_text(line)
            if len(line) < 8 or len(line) > 260:
                continue
            line_lower = line.lower()
            if line_lower.startswith(ignored_prefixes) or line_lower.startswith(metadata_prefixes):
                continue
            if any(pattern in line_lower for pattern in patterns) and line not in seen:
                selected.append(line)
                seen.add(line)
            if len(selected) >= limit:
                break
        return selected

    feature_lines = pick_lines(["feature", "component", "module", "function", "capability", "system", "interface", "configuration"])
    workflow_lines = pick_lines(["step", "process", "workflow", "flow", "first", "then", "after", "before", "execute", "upload", "select"])
    use_case_lines = pick_lines(["use case", "application", "used for", "used to", "can be used", "supports", "helps", "enables"])
    important_note_lines = pick_lines(["warning", "caution", "note", "limit", "constraint", "assumption", "must", "shall", "required", "error"])

    context_bits = []
    if page_count:
        context_bits.append(f"{page_count} pages/sections")
    if image_count:
        context_bits.append(f"{image_count} images")
    if table_count:
        context_bits.append(f"{table_count} tables")
    context_text = ", ".join(context_bits) if context_bits else f"{len(lines)} content lines"

    def bullet_list(items, fallback):
        usable_items = [item for item in items if item] or fallback
        return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in usable_items) + "</ul>"

    def section(title_text, body_html):
        if not body_html:
            return ""
        return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title_text)}</h4>{body_html}"

    structure_items = []
    if toc_entries:
        structure_items = [
            f"{num + ' ' if num else ''}{prettify_extracted_text(heading)}" + (f" - page {page_num}" if page_num else "")
            for num, heading, page_num in toc_entries[:6]
        ]
    elif headings:
        structure_items = [f"{num + ' ' if num else ''}{prettify_extracted_text(heading)}" for num, heading in headings[:6]]
    else:
        structure_items = [
            "The readable content is best understood by grouping related ideas rather than following the original layout.",
            f"Available context includes {context_text}.",
        ]

    purpose_by_type = {
        "technical": "to describe a system, process, implementation, test, or technical capability",
        "business": "to communicate objectives, operational context, metrics, or decision-oriented information",
        "research": "to explain a problem, method, evidence, findings, and conclusions",
        "general": "to present information in a readable and referenceable form",
    }
    summary_focus = ", ".join(keywords[:4]) if keywords else title
    key_point_items = key_lines[:5] or ["Not enough explanatory content was available to identify detailed points."]
    insight_items = [
        f"The document is best read as a {document_type} reference centered on {summary_focus}.",
        f"It should be read as a {document_type} document.",
    ]
    if image_count or table_count:
        insight_items.append("Visual or tabular material may contain supporting details that complement the readable text.")
    if important_note_lines:
        insight_items.append("Several lines contain requirements, constraints, warnings, or operational notes.")

    simplified_items = [
        f"In simple terms, this document is about {summary_focus}.",
        "It collects the main information a reader needs to understand the topic, context, and next actions.",
    ]
    takeaway_items = []
    if keywords:
        takeaway_items.append("The strongest takeaway is to connect the recurring ideas into a practical understanding of the document, not treat them as isolated terms.")
    takeaway_items.extend(key_point_items[:4])
    takeaway_items = takeaway_items[:5]

    summary_html = (
        "<div>"
        f"<div><b>What the document is about:</b> {html.escape(title)}</div>"
        f"<div><b>Main purpose:</b> {html.escape(purpose_by_type[document_type])}.</div>"
        f"<div><b>Key context:</b> {html.escape(context_text)}.</div>"
        f"<div><b>Reader focus:</b> Use the meaningful points below to understand the document's purpose, structure, and practical value.</div>"
        "</div>"
    )

    optional_sections = ""
    if feature_lines or document_type == "technical":
        optional_sections += section("Features / Concepts / Components", bullet_list(
            feature_lines,
            ["The document contains related concepts, but no explicit feature list was detected in the readable content."]
        ))
    if workflow_lines or document_type == "technical":
        optional_sections += section("Workflow / Process", bullet_list(
            workflow_lines,
            ["No clear step-by-step workflow was identified in the readable content."]
        ))
    if use_case_lines or document_type in ("technical", "business"):
        optional_sections += section("Use Cases / Applications", bullet_list(
            use_case_lines,
            ["Use this document as a reference for understanding the topic, validating details, or planning related work."]
        ))
    if important_note_lines:
        optional_sections += section("Important Notes", bullet_list(important_note_lines, []))

    analysis_sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Document Analysis: {html.escape(file_name)}</h3>",
        section("Summary", summary_html),
        section("Key Points", bullet_list(key_point_items, [])),
        section("Structure Breakdown", bullet_list(structure_items, [])),
        section("Key Insights / Core Insights", bullet_list(insight_items, [])),
        optional_sections,
        section("Simplified Explanation", bullet_list(simplified_items, [])),
        section("Key Takeaways", bullet_list(takeaway_items, [])),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(analysis_sections) + "</div>"


def build_product_documentation_analysis(file_name, file_bytes, text):
    raw_text = str(text or "")
    lines = [normalize_extracted_line(line) for line in raw_text.splitlines() if line.strip()]
    lower_text = raw_text.lower()

    ignored_prefixes = (
        "pdf metadata:", "document metadata:", "odt metadata:", "meta tags:", "total pages:",
        "total slides:", "workbook contains", "csv rows:", "error:", "[image:", "[embedded_image:",
        "page ", "slide ", "sheet ", "table:"
    )
    metadata_prefixes = (
        "producer:", "creationdate:", "moddate:", "author:", "creator:", "title:",
        "subject:", "keywords:", "trapped:", "pdfversion:"
    )

    def clean_sentence(value):
        value = normalize_extracted_line(value)
        value = re.sub(r"^(?:page|slide|sheet)\s+\d+\s*(?:text|content)?\s*:?", "", value, flags=re.IGNORECASE).strip()
        value = re.sub(r"\s+", " ", value)
        return value.strip(" -:")

    def meaningful_lines(max_items=180):
        selected = []
        seen = set()
        for line in lines:
            cleaned = clean_sentence(line)
            if not cleaned:
                continue
            lowered = cleaned.lower()
            if lowered.startswith(ignored_prefixes) or lowered.startswith(metadata_prefixes) or is_document_noise_line(cleaned):
                continue
            if len(cleaned) < 18 or len(cleaned) > 220:
                continue
            if re.fullmatch(r"[\W_]+", cleaned):
                continue
            if re.fullmatch(r"\d+(?:\.\d+)*\s+.+", cleaned) and len(cleaned.split()) <= 7:
                continue
            key = lowered
            if key in seen:
                continue
            seen.add(key)
            selected.append(cleaned)
            if len(selected) >= max_items:
                break
        return selected

    clean_lines = meaningful_lines()
    clean_text_for_keywords = " ".join(clean_lines) or raw_text
    words = re.findall(r"[A-Za-z][A-Za-z0-9_+\-/]{2,}", clean_text_for_keywords)
    keyword_counts = Counter(
        word.lower()
        for word in words
        if len(word) > 3 and word.lower() not in SUMMARY_STOPWORDS and not word.isdigit()
    )
    keywords = [normalize_extracted_line(word).title() for word, _ in keyword_counts.most_common(10)]
    keywords = [keyword for keyword in keywords if not _is_low_signal_theme(keyword)]

    type_scores = {
        "technical system": sum(1 for term in [
            "interface", "module", "device", "hardware", "software", "configuration", "channel",
            "connector", "signal", "protocol", "diagnostic", "architecture", "firmware", "driver"
        ] if term in lower_text),
        "process or workflow": sum(1 for term in [
            "process", "workflow", "procedure", "step", "approval", "operation", "execute", "setup",
            "install", "configure", "use", "report"
        ] if term in lower_text),
        "data or report": sum(1 for term in [
            "metric", "statistics", "result", "dashboard", "table", "test case", "passed", "failed",
            "executed", "summary", "analysis"
        ] if term in lower_text),
        "business document": sum(1 for term in [
            "customer", "market", "objective", "stakeholder", "cost", "risk", "benefit", "strategy",
            "requirement", "decision"
        ] if term in lower_text),
    }
    document_kind = max(type_scores, key=type_scores.get) if max(type_scores.values() or [0]) > 0 else "reference document"

    page_count, image_count, table_count = get_document_asset_counts(file_name, file_bytes, raw_text)
    title_match = re.search(r"Title:\s*(.+)", raw_text)
    detected_title = clean_sentence(title_match.group(1)) if title_match else ""
    display_name = detected_title if detected_title and len(detected_title) < 120 else os.path.splitext(file_name)[0]
    topic_terms = keywords[:5] or [display_name]
    topic_phrase = ", ".join(topic_terms[:4])

    def collect_by_terms(terms, limit=6):
        selected = []
        seen = set()
        for line in clean_lines:
            lowered = line.lower()
            if any(term in lowered for term in terms):
                simplified = synthesize_line(line)
                if simplified and simplified.lower() not in seen:
                    selected.append(simplified)
                    seen.add(simplified.lower())
            if len(selected) >= limit:
                break
        return selected

    def synthesize_line(line):
        line = clean_sentence(line)
        if not line:
            return ""
        line = re.sub(r"\b(?:note|warning|caution)\s*[:\-]\s*", "", line, flags=re.IGNORECASE)
        line = re.sub(r"\s*\.+\s*\d+\s*$", "", line)
        if len(line) > 150:
            line = line[:147].rsplit(" ", 1)[0] + "..."
        return line[0].upper() + line[1:] if line else line

    capability_terms = [
        "support", "supports", "feature", "function", "capability", "enable", "allows", "provide",
        "communication", "measurement", "analysis", "diagnostic", "configuration", "export"
    ]
    architecture_terms = [
        "component", "module", "interface", "channel", "connector", "port", "device", "unit",
        "software", "hardware", "driver", "network", "table", "sheet", "slide"
    ]
    workflow_terms = [
        "install", "configure", "connect", "select", "upload", "execute", "start", "use",
        "create", "open", "set", "download", "export", "analyze"
    ]
    use_case_terms = [
        "application", "used for", "used to", "use case", "measurement", "testing", "diagnostic",
        "report", "automation", "monitoring", "analysis", "configuration"
    ]

    capabilities = collect_by_terms(capability_terms, 7)
    architecture_evidence = collect_by_terms(architecture_terms, 7)
    workflow_evidence = collect_by_terms(workflow_terms, 6)
    use_case_evidence = collect_by_terms(use_case_terms, 6)

    def has_any(*terms):
        return any(term in lower_text for term in terms)

    generated_capabilities = []
    if has_any("can ", "can-fd", "can fd", "lin", "flexray", "ethernet", "protocol", "interface", "communication"):
        generated_capabilities.append("Supports communication-oriented work through documented protocols, interfaces, or channels.")
    if has_any("configuration", "configure", "setup", "install", "driver", "software"):
        generated_capabilities.append("Provides configuration and setup guidance so the system can be prepared for practical use.")
    if has_any("measurement", "test", "diagnostic", "analysis", "monitor", "report"):
        generated_capabilities.append("Supports analysis, measurement, diagnostics, reporting, or validation activities.")
    if has_any("connector", "pin", "port", "socket", "plug", "channel"):
        generated_capabilities.append("Documents physical or logical connectivity details needed for integration.")
    if has_any("table", "sheet", "csv", "dashboard", "statistics", "result"):
        generated_capabilities.append("Contains structured data or results that can be reviewed, summarized, or exported.")
    if has_any("image", "figure", "diagram", "visual", "illustration"):
        generated_capabilities.append("Includes visual or diagram-like information that can support engineering reference work.")
    if generated_capabilities:
        capabilities = generated_capabilities

    generated_architecture = []
    if has_any("hardware", "device", "unit", "module", "component"):
        generated_architecture.append("Hardware or device layer: the physical units, modules, or components described by the source.")
    if has_any("software", "driver", "application", "tool", "configuration"):
        generated_architecture.append("Software and configuration layer: tools, drivers, settings, and setup behavior around the system.")
    if has_any("interface", "protocol", "channel", "network", "communication"):
        generated_architecture.append("Communication layer: interfaces, protocols, and channels that connect the system to other tools or networks.")
    if has_any("connector", "pin", "port", "socket", "plug"):
        generated_architecture.append("Connectivity layer: ports, connectors, pin assignments, or wiring-related details.")
    if has_any("table", "sheet", "report", "result", "metadata"):
        generated_architecture.append("Information layer: tables, results, and reference data used to interpret the document.")
    if generated_architecture:
        architecture_evidence = generated_architecture

    generated_workflow = []
    if has_any("upload", "select", "open", "choose"):
        generated_workflow.append("Select the relevant file, section, component, or dataset.")
    if has_any("install", "setup", "driver", "connect"):
        generated_workflow.append("Prepare the environment by installing, connecting, or setting up the required parts.")
    if has_any("configure", "configuration", "setting", "parameter"):
        generated_workflow.append("Configure the required options, channels, interfaces, or parameters.")
    if has_any("execute", "run", "start", "measurement", "test", "analysis"):
        generated_workflow.append("Run the intended operation such as measurement, testing, communication, analysis, or review.")
    if has_any("result", "report", "export", "download", "table"):
        generated_workflow.append("Review outputs, results, tables, or reports and export anything needed for reference.")
    if generated_workflow:
        workflow_evidence = generated_workflow

    generated_use_cases = []
    if has_any("measurement", "canalyzer", "canoe", "diagnostic", "test"):
        generated_use_cases.append("Vehicle/network measurement, diagnostics, testing, and validation workflows.")
    if has_any("configuration", "install", "setup", "driver"):
        generated_use_cases.append("Setup and configuration reference for engineers or technicians.")
    if has_any("interface", "connector", "pin", "channel", "protocol"):
        generated_use_cases.append("Integration reference for ports, channels, protocols, connectors, or pin mappings.")
    if has_any("report", "dashboard", "statistics", "result", "table"):
        generated_use_cases.append("Report review, structured data analysis, and documentation support.")
    if has_any("warning", "caution", "note", "safety", "required"):
        generated_use_cases.append("Operational guidance where constraints, warnings, or required practices matter.")
    if generated_use_cases:
        use_case_evidence = generated_use_cases

    components = []
    component_candidates = []
    for pattern in [
        r"\b[A-Z]{2,}[A-Za-z0-9_+\-/]*\b",
        r"\b[A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z0-9]+){0,2}\b",
    ]:
        component_candidates.extend(re.findall(pattern, raw_text))
    component_counts = Counter(
        normalize_extracted_line(candidate).strip()
        for candidate in component_candidates
        if 3 <= len(normalize_extracted_line(candidate).strip()) <= 45
        and normalize_extracted_line(candidate).lower() not in SUMMARY_STOPWORDS
    )
    for candidate, _ in component_counts.most_common(8):
        lowered = candidate.lower()
        if lowered in {"pdf", "metadata", "page", "text", "table", "figure"}:
            continue
        components.append(f"{candidate}: appears to be a major referenced part, concept, interface, or artifact in the document.")
        if len(components) >= 5:
            break

    if not components and architecture_evidence:
        components = architecture_evidence[:5]

    assets = []
    if page_count:
        assets.append(f"about {page_count} pages or sections")
    if table_count:
        assets.append(f"{table_count} table-like data areas")
    if image_count:
        assets.append(f"{image_count} visual assets")
    asset_phrase = ", ".join(assets) if assets else "the available extracted content"

    overview_items = [
        f"This is a {document_kind} centered on {topic_phrase}.",
        f"It serves as a practical reference for understanding the subject, its purpose, and how the relevant pieces fit together.",
    ]
    if assets:
        overview_items.append(f"The source contains {asset_phrase}, but the summary below reorganizes the content by meaning rather than document order.")

    core_concept_items = [
        f"In simple terms, the document explains how {topic_terms[0] if topic_terms else 'the subject'} is used, configured, or understood in context.",
        "The important ideas are grouped into purpose, structure, capabilities, usage flow, and practical value so a reader can act on them quickly.",
    ]
    if capabilities:
        core_concept_items.append("The central behavior is reflected in the capabilities and usage evidence described below.")

    if not architecture_evidence:
        architecture_evidence = [
            "The document content is best understood as a set of related concepts, interfaces, configuration details, and operational notes.",
            "Related elements are grouped logically instead of following the original document layout."
        ]

    if not capabilities:
        capabilities = [
            "Provides reference information needed to understand and apply the documented subject.",
            "Combines functional context with technical details where the source provides them."
        ]

    if not workflow_evidence:
        workflow_evidence = [
            "Identify the relevant subject or component.",
            "Review its purpose, interfaces, configuration needs, and constraints.",
            "Apply the information in implementation, testing, documentation, or troubleshooting work."
        ]

    if not use_case_evidence:
        use_case_evidence = [
            "Engineering reference and onboarding.",
            "Configuration or implementation planning.",
            "Troubleshooting, validation, and documentation support."
        ]

    takeaway_items = []
    if keywords:
        takeaway_items.append("The most useful reading is to connect the recurring technical ideas into purpose, structure, workflow, and practical value.")
    takeaway_items.append(f"The document is most useful as a {document_kind} rather than as a narrative document.")
    if capabilities:
        takeaway_items.append("The key value is translating scattered technical or functional details into usable reference knowledge.")
    if architecture_evidence:
        takeaway_items.append("Understanding the relationships between components, interfaces, and usage flow is more important than memorizing the original section order.")
    takeaway_items = takeaway_items[:5]

    def bullet_list(items):
        clean_items = [item for item in items if item]
        return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in clean_items) + "</ul>"

    def section(title_text, items):
        if not items:
            return ""
        return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title_text)}</h4>{bullet_list(items)}"

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Full Analysis: {html.escape(file_name)}</h3>",
        section("Overview", overview_items),
        section("Core Concept", core_concept_items),
        section("Structure / Architecture", architecture_evidence[:7]),
        section("Key Elements", capabilities[:7] + components[:6]),
        section("Workflow / Logic", workflow_evidence[:6]),
        section("Applications / Use Cases", use_case_evidence[:6]),
        section("Key Takeaways", takeaway_items[:5]),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(part for part in sections if part) + "</div>"


def build_detailed_document_summary(file_name, file_bytes, text):
    return build_product_documentation_analysis(file_name, file_bytes, text)


def extract_quoted_item_name(user_input):
    match = re.search(r"'(.*?)'|\"(.*?)\"", str(user_input or ""))
    if match:
        return (match.group(1) or match.group(2) or "").strip()

    patterns = [
        r"\b(?:item|about|for|related to)\s+([A-Za-z0-9][A-Za-z0-9 _./+\-]{1,80})",
        r"\b(?:pin(?:s)?|diagram|connector|visual)\s+([A-Za-z0-9][A-Za-z0-9 _./+\-]{1,80})",
    ]
    for pattern in patterns:
        match = re.search(pattern, str(user_input or ""), re.IGNORECASE)
        if match:
            item = re.split(r"\b(?:from|in|with|please|and|details?|info|information)\b", match.group(1), 1, flags=re.IGNORECASE)[0]
            return item.strip(" :-")
    return ""


def extract_bare_item_name(user_input):
    """Detect terse part-number queries such as VN1671 or VN 1671."""
    text = str(user_input or "").strip()
    match = re.search(r"\b(VN)\s*[- ]?\s*(\d{4}[A-Za-z]?)\b", text, re.IGNORECASE)
    if match:
        return f"{match.group(1).upper()}{match.group(2).upper()}"
    return ""


def normalize_extracted_line(line):
    line = str(line or "").strip()
    line = re.sub(r"([a-z])([A-Z])", r"\1 \2", line)
    line = re.sub(r"([A-Za-z])(\d)", r"\1 \2", line)
    line = re.sub(r"(\d)([A-Za-z])", r"\1 \2", line)
    line = re.sub(r"\s+", " ", line)
    return line.strip()


def is_document_noise_line(line):
    """Filter metadata, TOC entries, headers/footers, and extraction noise."""
    cleaned = normalize_extracted_line(line)
    lower = cleaned.lower().strip()
    if not lower:
        return True
    metadata_prefixes = (
        "pdf metadata:", "document metadata:", "creation date:", "creator:", "mod date:",
        "producer:", "subject:", "title:", "total pages:", "pdfversion:", "author:",
        "keywords:", "trapped:", "page ", "slide ", "sheet "
    )
    if lower.startswith(metadata_prefixes):
        return True
    if lower in {"table of contents", "contents", "index", "product information"}:
        return True
    if re.fullmatch(r"v\s*\d+(?:\.\d+)*\s*/\s*\d{4}", lower):
        return True
    if re.fullmatch(r"\d{1,3}", lower):
        return True
    if re.search(r"\.{5,}", cleaned):
        return True
    if re.fullmatch(r"\d+(?:\.\d+)*\s+.+\s+\d{1,3}", cleaned) and len(cleaned.split()) <= 12:
        return True
    if re.fullmatch(r"[\W_]+", cleaned):
        return True
    return False


def get_document_display_title(file_name, text):
    """Prefer document title metadata when it looks useful."""
    raw_text = str(text or "")
    title_match = re.search(r"^Title:\s*(.+)$", raw_text, re.IGNORECASE | re.MULTILINE)
    if title_match:
        title = normalize_extracted_line(title_match.group(1)).strip(" :-")
        if title and not is_document_noise_line(f"Not metadata {title}") and len(title) <= 120:
            return title
    return os.path.splitext(str(file_name or "document"))[0]


def get_meaningful_document_lines(text, min_len=18, max_len=260, limit=220):
    """Return clean, unique content lines suitable for summaries and analysis."""
    selected = []
    seen = set()
    for raw_line in str(text or "").splitlines():
        line = normalize_extracted_line(raw_line)
        line = re.sub(r"^>\s*", "", line).strip()
        line = re.sub(r"^(?:page|slide|sheet)\s+\d+\s*(?:text|content)?\s*:?", "", line, flags=re.IGNORECASE).strip()
        if is_document_noise_line(line):
            continue
        if len(line) < min_len or len(line) > max_len:
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        selected.append(line)
        if len(selected) >= limit:
            break
    return selected


def collect_item_context_lines(text, item_name, window=4, limit=80):
    item_name = str(item_name or "").strip()
    if not item_name:
        return []

    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    item_tokens = [token.lower() for token in re.findall(r"[A-Za-z0-9]+", item_name) if len(token) > 1]
    if not item_tokens:
        return []

    selected = []
    seen = set()
    for index, line in enumerate(lines):
        line_lower = line.lower()
        compact_line = re.sub(r"\s+", "", line_lower)
        compact_item = re.sub(r"\s+", "", item_name.lower())
        has_match = compact_item in compact_line or all(token in line_lower for token in item_tokens)
        if not has_match:
            continue

        start = max(0, index - window)
        end = min(len(lines), index + window + 1)
        for context_line in lines[start:end]:
            pretty_line = normalize_extracted_line(context_line)
            if len(pretty_line) < 3 or len(pretty_line) > 300:
                continue
            key = pretty_line.lower()
            if key in seen:
                continue
            seen.add(key)
            selected.append(pretty_line)
            if len(selected) >= limit:
                return selected
    return selected


def select_relevant_lines(context_lines, patterns, limit=8):
    selected = []
    seen = set()
    for line in context_lines:
        line_lower = line.lower()
        if any(pattern in line_lower for pattern in patterns) and line not in seen:
            selected.append(line)
            seen.add(line)
        if len(selected) >= limit:
            break
    return selected


def html_bullet_list(items):
    if not items:
        return ""
    return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in items) + "</ul>"


def html_section(title, items):
    if not items:
        return ""
    return f"<h4 style='margin:16px 0 6px 0; color:#173152;'>{html.escape(title)}</h4>{html_bullet_list(items)}"


def build_item_information_response(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=5, limit=100)
    if not context_lines:
        return f"<div><h3>Item Information: {html.escape(item_name)}</h3><p>No relevant information for this item was found in {html.escape(file_name)}.</p></div>"

    overview = context_lines[:5]
    purpose = select_relevant_lines(context_lines, ["purpose", "used for", "provides", "supports", "enables", "allows", "designed"])
    features = select_relevant_lines(context_lines, ["feature", "support", "capability", "function", "operation", "application"])
    technical = select_relevant_lines(context_lines, ["mbit", "kbit", "volt", "channel", "standard", "protocol", "can", "lin", "flexray", "interface", "specification", "iso", "structure", "component", "module", "piggy", "channel", "internal", "family", "device"])
    interfaces = select_relevant_lines(context_lines, ["connector", "port", "pin", "d-sub", "usb", "channel", "plug", "socket", "interface", "relationship", "connects to"])
    usage = select_relevant_lines(context_lines, ["configure", "configuration", "install", "insert", "setup", "use", "driver", "software", "hardware", "role"])
    notes = select_relevant_lines(context_lines, ["special", "unique", "only", "limitation", "difference", "optional", "available", "not supported", "note", "warning", "caution", "must", "shall", "important", "avoid", "required"])
    takeaways = select_relevant_lines(context_lines, ["key", "important", "main", "critical", "takeaway"], limit=3)

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Component: {html.escape(item_name)}</h3>",
        f"<p><b>Source:</b> {html.escape(file_name)}</p>",
        html_section("Overview", overview),
        html_section("Purpose", purpose),
        html_section("Key Features", features),
        html_section("Technical / Contextual Details", technical),
        html_section("Interfaces / Relationships (if applicable)", interfaces),
        html_section("Usage / Role", usage),
        html_section("Notes", notes),
        html_section("Key Takeaways", takeaways),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(section for section in sections if section) + "</div>"


def extract_pin_rows(context_lines):
    rows = []
    seen = set()
    pin_patterns = [
        r"\bpin\s*(\d+)\b\s*[:\-]?\s*([A-Za-z0-9_+/.\- ]{0,40})\s*(.*)",
        r"^\s*(\d{1,2})\s+([A-Za-z][A-Za-z0-9_+/.\-]*)\s*(.*)",
    ]
    for line in context_lines:
        line_lower = line.lower()
        if not any(term in line_lower for term in ["pin", "signal", "d-sub", "connector", "ground", "shield", "can", "lin", "vbat"]):
            continue
        for pattern in pin_patterns:
            match = re.search(pattern, line, re.IGNORECASE)
            if not match:
                continue
            pin_no = match.group(1).strip()
            signal = (match.group(2) or "").strip(" :-") or "Not specified"
            description = (match.group(3) or "").strip(" :-") or line
            key = (pin_no, signal.lower(), description.lower())
            if key in seen:
                break
            seen.add(key)
            rows.append({
                "pin": pin_no,
                "signal": signal,
                "description": description,
                "notes": ""
            })
            break
    return rows[:40]


def build_pin_csv(pin_rows):
    lines = ["Pin Number,Signal Name,Description,Notes"]
    for row in pin_rows:
        values = [row["pin"], row["signal"], row["description"], row.get("notes", "")]
        escaped_values = ['"' + str(value).replace('"', '""') + '"' for value in values]
        lines.append(",".join(escaped_values))
    return "\n".join(lines)


def build_ascii_pin_diagram(pin_rows, item_name):
    if not pin_rows:
        return f"+------------------------------+\n| {item_name[:28]:<28} |\n| Pin diagram not available    |\n+------------------------------+"
    left = pin_rows[::2]
    right = pin_rows[1::2]
    width = 34
    lines = [f"+{'-' * width}+", f"| {item_name[:width-4]:<{width-4}} |", f"+{'-' * width}+"]
    max_len = max(len(left), len(right))
    for index in range(max_len):
        left_text = ""
        right_text = ""
        if index < len(left):
            left_text = f"{left[index]['pin']}:{left[index]['signal']}"[:15]
        if index < len(right):
            right_text = f"{right[index]['pin']}:{right[index]['signal']}"[:15]
        lines.append(f"| {left_text:<15}  {right_text:>15} |")
    lines.append(f"+{'-' * width}+")
    return "\n".join(lines)


def build_item_visual_assets(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=8, limit=140)
    pin_rows = extract_pin_rows(context_lines)
    if not pin_rows:
        return {"csv": [], "diagrams": []}

    safe_item_name = re.sub(r"[^A-Za-z0-9_-]+", "_", str(item_name)).strip("_") or "item"
    file_base = re.sub(r"[^A-Za-z0-9_-]+", "_", os.path.splitext(file_name)[0]).strip("_") or "document"
    csv_text = build_pin_csv(pin_rows)
    ascii_diagram = build_ascii_pin_diagram(pin_rows, item_name)

    return {
        "csv": [{
            "label": f"{file_name} - {item_name} pin table CSV",
            "data": csv_text.encode("utf-8"),
            "file_name": f"{file_base}_{safe_item_name}_pin_table.csv",
            "mime": "text/csv",
        }],
        "diagrams": [{
            "label": f"{file_name} - {item_name} ASCII diagram",
            "data": ascii_diagram.encode("utf-8"),
            "file_name": f"{file_base}_{safe_item_name}_diagram.txt",
            "mime": "text/plain",
        }],
    }


def html_table(headers, rows):
    if not rows:
        return ""
    head_html = "".join(f"<th>{html.escape(header)}</th>" for header in headers)
    body_html = ""
    for row in rows:
        body_html += "<tr>" + "".join(f"<td>{html.escape(str(cell))}</td>" for cell in row) + "</tr>"
    return f"<table style='border-collapse:collapse; width:100%; margin:8px 0;'><thead><tr>{head_html}</tr></thead><tbody>{body_html}</tbody></table>"


def build_item_visual_response(file_name, text, item_name):
    context_lines = collect_item_context_lines(text, item_name, window=8, limit=140)
    if not context_lines:
        return f"<div><h3>Visual / Pin Reference: {html.escape(item_name)}</h3><p>No relevant visual or structural information for this item was found in {html.escape(file_name)}.</p></div>"

    pin_rows = extract_pin_rows(context_lines)
    connector_lines = select_relevant_lines(context_lines, ["connector", "port", "d-sub", "usb", "channel", "plug", "socket", "interface"], limit=12)
    image_lines = select_relevant_lines(context_lines, ["figure", "image", "diagram", "pin assignment", "illustration"], limit=10)
    table_lines = select_relevant_lines(context_lines, ["table", "specification", "signal", "configuration", "pin"], limit=12)
    csv_text = build_pin_csv(pin_rows) if pin_rows else "Pin Number,Signal Name,Description,Notes\n"
    ascii_diagram = build_ascii_pin_diagram(pin_rows, item_name)
    pin_table_rows = [[row["pin"], row["signal"], row["description"], row.get("notes", "")] for row in pin_rows]

    sections = [
        f"<h3 style='margin:0 0 10px 0; color:#173152;'>Visual / Pin Reference: {html.escape(item_name)}</h3>",
        f"<p><b>Source:</b> {html.escape(file_name)}</p>",
        html_section("Pin Diagrams", ["Recreated below from available pin/signal information." if pin_rows else "No explicit pin diagram was found in the provided document context."]),
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>",
        f"<h4 style='margin:16px 0 6px 0; color:#173152;'>Pin Configuration Table</h4>",
        html_table(["Pin Number", "Signal Name", "Description", "Notes"], pin_table_rows) if pin_rows else "<p>No pin table data was found.</p>",
        html_section("Connector Details", connector_lines),
        html_section("Images & Visuals", image_lines),
        html_section("Technical Tables", table_lines),
        f"<h4 style='margin:16px 0 6px 0; color:#173152;'>Downloadable Outputs</h4>",
        "<p><b>a) Pin table as CSV</b></p>",
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(csv_text)}</pre>",
        "<p><b>b) Diagram as ASCII / structured format</b></p>",
        f"<pre style='white-space:pre-wrap; background:#f4f7fb; padding:12px; border-radius:8px;'>{html.escape(ascii_diagram)}</pre>",
        "<p><b>c) Image references or recreated diagrams</b></p>",
        html_bullet_list(image_lines or ["No direct image reference was found in the provided document context; use the recreated ASCII diagram above when pin rows are available."]),
    ]
    return "<div style='margin-bottom:18px; line-height:1.5;'>" + "".join(section for section in sections if section) + "</div>"


def extract_page_text(text, page_number=1):
    text = str(text)
    pattern = rf"Page {page_number}\s+Text:\s*(.*?)(?=Page \d+\s+Text:|\Z)"
    match = re.search(pattern, text, re.S | re.IGNORECASE)
    if match:
        return match.group(1).strip()

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines[:80])


def find_heading_page_number(text, heading):
    text = str(text)
    lines = [line for line in text.splitlines()]
    heading_pattern = re.escape(str(heading).strip())
    for index, line in enumerate(lines):
        if re.search(rf"\b{heading_pattern}\b", line, re.IGNORECASE):
            for j in range(index, -1, -1):
                page_match = re.search(r'Page\s+(\d+)\s+Text:', lines[j], re.IGNORECASE)
                if page_match:
                    return int(page_match.group(1))
    return None


def resolve_heading_page_number(text, heading, toc_entries=None):
    if not heading:
        return None
    heading_text = str(heading).strip()
    if toc_entries is None:
        toc_entries = extract_toc_with_page_numbers(text)
    for num, title, page_num in toc_entries:
        if title.strip().lower() == heading_text.lower():
            return page_num
        if heading_text.lower() in title.strip().lower() or title.strip().lower() in heading_text.lower():
            return page_num
    return find_heading_page_number(text, heading_text)


def extract_document_headings(text):
    """Extract numbered headings and explicit DOCX headings from extracted text."""
    headings = []
    text = str(text)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    
    for line in lines:
        # Skip lines that are too long
        if len(line) > 120:
            continue
        
        # Skip metadata, page markers, and special content
        if (line.isupper() or line.endswith(":") or 
            "Page" in line or "PDF Metadata" in line or 
            "Total Pages" in line or "TABLE:" in line):
            continue

        # Match explicit heading markers from DOCX extraction
        if line.startswith("Heading:"):
            heading_text = line.replace("Heading:", "", 1).strip()
            if 3 <= len(heading_text) <= 120:
                headings.append(("", heading_text))
            continue
        
        # Match numbered headings at start: "1 Overview", "1.1 Introduction", etc.
        match = re.match(r'^(\d+(?:\.\d+)*)\s+([A-Za-z\s][^.]*?)(?:\s*\.+\s*\d+)?\s*$', line)
        if match:
            num = match.group(1)
            title = match.group(2).strip()
            
            # Clean up any trailing dots or page numbers
            title = re.sub(r'\s*\.+\s*\d*\s*$', '', title).strip()
            
            if 3 <= len(title) <= 120:
                headings.append((num, title))
    
    # Remove duplicates while preserving order
    seen = set()
    deduped = []
    for num, title in headings:
        key = f"{num}:{title}"
        if key not in seen:
            seen.add(key)
            deduped.append((num, title))
    
    return deduped


def extract_toc_with_page_numbers(text):
    """Extract table of contents entries with page numbers from document."""
    toc_entries = []
    text = str(text)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    
    # First, try explicit TOC patterns on full text
    for regex in [
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s+\.{2,}\s*(\d+)\s*$',
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s{3,}(\d+)\s*$',
        r'(?m)^\s*(\d+(?:\.\d+)*)\s+(.+?)\s+(\d+)\s*$'
    ]:
        for match in re.finditer(regex, text):
            num = match.group(1)
            title = match.group(2).strip()
            page_num = match.group(3)
            if 3 <= len(title) <= 120 and len(re.findall(r'\d+', title)) <= 2:
                toc_entries.append((num, title, page_num))
        if toc_entries:
            return toc_entries

    # Fallback: build TOC from detected headings and page markers
    headings = extract_document_headings(text)
    if headings:
        for num, title in headings:
            page_num = None
            search_pattern = re.escape(title)
            for i, line in enumerate(lines):
                if title in line or re.search(search_pattern, line, re.IGNORECASE):
                    for j in range(i, max(0, i - 20), -1):
                        page_match = re.search(r'Page\s+(\d+)\s+Text:', lines[j])
                        if page_match:
                            page_num = page_match.group(1)
                            break
                    if page_num:
                        break
            toc_entries.append((num, title, page_num))
    return toc_entries


def build_file_overview(file_name, text):
    text = str(text)
    toc_entries = extract_toc_with_page_numbers(text)
    all_headings = extract_document_headings(text)

    overview_parts = [f"📄 **{file_name}**"]
    
    # Table of Contents section
    overview_parts.append("### Table of Contents")
    if toc_entries:
        overview_parts.append("| Contents | Page No |")
        overview_parts.append("|----------|---------|")
        for num, title, page_num in toc_entries:
            content_str = f"{num} {title}" if num else title
            display_text = f"{content_str} (Page {page_num})" if page_num else content_str
            preview_link = create_preview_link(file_name, highlight_term=title, page_num=page_num)
            anchor_id = create_heading_anchor(title)
            if preview_link:
                page_display = page_num if page_num else "-"
                overview_parts.append(f"| <a href='{preview_link}#{anchor_id}' target='_blank'>{html.escape(display_text)}</a> | {page_display} |")
            else:
                page_display = page_num if page_num else "-"
                overview_parts.append(f"| {html.escape(display_text)} | {page_display} |")
    else:
        overview_parts.append("- No table of contents found with page numbers.")

    # Document Headings section
    overview_parts.append("### Document Headings")
    if all_headings:
        for num, title in all_headings:
            content_str = f"{num} {title}" if num else title
            anchor_id = create_heading_anchor(title)
            page_num = resolve_heading_page_number(text, title, toc_entries)
            preview_link = create_preview_link(file_name, highlight_term=title, page_num=page_num)
            if preview_link:
                overview_parts.append(f"- <a href='{preview_link}#{anchor_id}' target='_blank'>{html.escape(content_str)}</a>")
            else:
                overview_parts.append(f"- {content_str}")
    else:
        overview_parts.append("- No document headings were detected.")

    return "\n".join(overview_parts)


@st.cache_data(show_spinner=False)
def build_highlighted_search_results(file_name, text, query):
    if not query:
        return ""

    pattern = re.compile(re.escape(query), re.IGNORECASE)
    matches = []

    for line_no, raw_line in enumerate(str(text).splitlines(), 1):
        if pattern.search(raw_line):
            escaped_line = html.escape(raw_line)
            highlighted_line = pattern.sub(
                lambda match: f"<mark style='background:#fff3a3; padding:0 2px;'>{html.escape(match.group(0))}</mark>",
                escaped_line
            )
            matches.append(
                f"<div style='margin:0 0 8px 0;'><b>Line {line_no}</b>: {highlighted_line}</div>"
            )

    if not matches:
        return f"<div><b>{html.escape(file_name)}</b><br>No matches found for <code>{html.escape(query)}</code>.</div>"

    return (
        f"<div style='margin-bottom:14px;'>"
        f"<h4 style='margin:0 0 8px 0; color:#a8d8f0;'>{html.escape(file_name)} ({len(matches)} matches)</h4>"
        f"{''.join(matches)}"
        f"</div>"
    )


@st.cache_data(show_spinner=False)
def extract_login_name_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    text = soup.get_text(" ", strip=True)
    match = re.search(r'login name[:\s]+(.+?)(version|$)', text, re.IGNORECASE)
    if match:
        name = match.group(1).strip()
        parts = name.split()
        return " ".join(parts[:1])
    return "Not found"


@st.cache_data(show_spinner=False)
def extract_statistics_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    stats = {
        "Executed": 0,
        "Passed": 0,
        "Failed": 0,
        "Inconclusive": 0,
        "Error": 0
    }

    text = soup.get_text(" ", strip=True).lower()
    patterns = {
        "Executed": r'executed test cases[:\s]+(\d+)',
        "Passed": r'passed[:\s]+(\d+)',
        "Failed": r'failed[:\s]+(\d+)',
        "Inconclusive": r'inconclusive[:\s]+(\d+)',
        "Error": r'error[:\s]+(\d+)'
    }

    for key, pattern in patterns.items():
        match = re.search(pattern, text)
        if match:
            stats[key] = int(match.group(1))

    return stats


@st.cache_data(show_spinner=False)
def extract_test_results_grouped_from_html(file_bytes):
    soup = BeautifulSoup(BytesIO(file_bytes), "html.parser")
    results = {}

    group_tables = soup.find_all('table', class_='GroupHeadingTable')

    for group_table in group_tables:
        try:
            rows = group_table.find_all('tr')
            if len(rows) >= 2:
                first_row = rows[0]
                heading = first_row.find('big', class_='Heading3')

                if heading:
                    heading_text = heading.get_text(strip=True)
                    fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', heading_text, re.IGNORECASE)

                    if fixture_match:
                        fixture_name = fixture_match.group(1).strip()
                        second_row = rows[1]
                        overview_table = second_row.find('table', class_='OverviewResultTable')

                        if overview_table:
                            count_cell = overview_table.find('td')
                            if count_cell:
                                try:
                                    count = int(count_cell.get_text(strip=True))

                                    if fixture_name not in results:
                                        results[fixture_name] = {
                                            "name": fixture_name,
                                            "test_cases": [],
                                            "pass": count,
                                            "fail": 0,
                                            "error": 0,
                                            "not executed": 0,
                                            "inconclusive": 0,
                                            "total": count,
                                            "count_cell_class": count_cell.get('class', [''])[0]
                                        }
                                except ValueError:
                                    pass
        except Exception:
            pass

    full_text = soup.get_text("\n", strip=True)
    lines = [l.strip() for l in full_text.split("\n") if l.strip()]

    current_fixture = None

    for i, line in enumerate(lines):
        line_lower = line.lower()

        if "test fixture:" in line_lower:
            fixture_match = re.search(r'Test Fixture:\s*(.+?)(?:\s|$)', line, re.IGNORECASE)
            if fixture_match:
                current_fixture = fixture_match.group(1).strip()
                if current_fixture not in results:
                    results[current_fixture] = {
                        "name": current_fixture,
                        "test_cases": [],
                        "pass": 0,
                        "fail": 0,
                        "error": 0,
                        "not executed": 0,
                        "inconclusive": 0,
                        "total": 0
                    }

        elif re.match(r'^\d+\.\d+', line) and current_fixture:
            verdict_match = re.search(r':\s*(Passed|Failed|Pass|Fail|Error|Not Executed|Inconclusive)\s*$', line,
                                      re.IGNORECASE)

            if verdict_match:
                verdict_word = verdict_match.group(1).lower()

                if "pass" in verdict_word:
                    verdict_type = "Pass"
                    results[current_fixture]["pass"] += 1
                elif "fail" in verdict_word:
                    verdict_type = "Failed"
                    results[current_fixture]["fail"] += 1
                elif "error" in verdict_word:
                    verdict_type = "Error"
                    results[current_fixture]["error"] += 1
                elif "not executed" in verdict_word:
                    verdict_type = "Not Executed"
                    results[current_fixture]["not executed"] += 1
                elif "inconclusive" in verdict_word:
                    verdict_type = "Inconclusive"
                    results[current_fixture]["inconclusive"] += 1
                else:
                    continue

                timestamp = None
                test_step = "Step"
                failure_step_id = ""

                def score_timestamp(candidate):
                    if not candidate:
                        return -1
                    parts = candidate.split('.')
                    if len(parts) != 2 or not parts[0].isdigit() or not parts[1].isdigit():
                        return -1
                    leading_num = int(parts[0])
                    decimal_places = len(parts[1])
                    decimal_bonus = 10000 if decimal_places >= 3 else (100 if decimal_places == 2 else 0)
                    return decimal_bonus + leading_num

                def find_best_timestamp(text):
                    matches = re.findall(r'\b(\d+\.\d+)\b', text)
                    if not matches:
                        return None
                    return max(matches, key=score_timestamp)

                def find_first_relevant_timestamp(text):
                    for m in re.findall(r'\b(\d+\.\d+)\b', text):
                        if len(m.split('.')[1]) >= 3:
                            return m
                    for m in re.findall(r'\b(\d+\.\d+)\b', text):
                        if len(m.split('.')[1]) >= 2:
                            return m
                    return None

                def consider_timestamp(candidate):
                    nonlocal timestamp
                    if not candidate:
                        return
                    if not timestamp:
                        timestamp = candidate
                        return
                    if len(timestamp.split('.')[1]) >= 3:
                        return
                    if len(candidate.split('.')[1]) > len(timestamp.split('.')[1]):
                        timestamp = candidate
                        return
                    if score_timestamp(candidate) > score_timestamp(timestamp):
                        timestamp = candidate

                same_line_step = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', line,
                                           re.IGNORECASE)
                if same_line_step:
                    failure_step_id = same_line_step.group(1)
                    action_text = same_line_step.group(2).strip()
                    test_step = action_text
                    consider_timestamp(find_first_relevant_timestamp(line) or find_best_timestamp(line))

                for k in range(i + 1, min(i + 150, len(lines))):
                    next_line = lines[k]

                    if re.match(r'^\d+\.\d+(?:\s|$)', next_line) and k > i + 5:
                        break

                    consider_timestamp(find_first_relevant_timestamp(next_line) or find_best_timestamp(next_line))

                    if verdict_type in ["Failed", "Error"] and not failure_step_id:
                        next_line_lower = next_line.lower()

                        step_match = re.search(r'(\d+(?:\.\d+)+)\.\s+([^:]+):\s*(failed|fail|error)', next_line,
                                               re.IGNORECASE)
                        if step_match:
                            failure_step_id = step_match.group(1)
                            action_text = step_match.group(2).strip()
                            test_step = action_text
                            consider_timestamp(find_best_timestamp(next_line))
                        else:
                            if any(keyword in next_line_lower for keyword in
                                   ["condition", "value", "expected", "actual", "mismatch", "not found",
                                    "exception", "error", "failed to", "failed"]):
                                if not re.match(r'^\d+\.\d+', next_line):
                                    step_num_match = re.match(r'^(\d+(?:\.\d+)*)', next_line.strip())
                                    if step_num_match:
                                        failure_step_id = step_num_match.group(1)
                                        test_step = next_line[:80]

                    if verdict_type == "Pass":
                        next_line_lower = next_line.lower()

                        if "execute" in next_line_lower:
                            match = re.search(r'execute\s+(\w+)', next_line_lower)
                            if match:
                                test_step = match.group(1).capitalize()
                        elif "question" in next_line_lower and "text" in next_line_lower:
                            test_step = "Question/Text"
                        elif "await" in next_line_lower or "wait" in next_line_lower:
                            test_step = "Await Value Match"
                        elif "resume" in next_line_lower:
                            test_step = "Resume"
                        elif "set" in next_line_lower:
                            test_step = "Set"
                        elif "tester" in next_line_lower and "confirmed" in next_line_lower:
                            test_step = "Tester Confirmation"

                if timestamp:
                    results[current_fixture]["test_cases"].append({
                        "timestamp": timestamp,
                        "verdict": verdict_type,
                        "details": test_step
                    })

    for fixture_name in results:
        parsed_count = len(results[fixture_name]["test_cases"])
        initial_count = results[fixture_name].get("total", 0)
        results[fixture_name]["total"] = max(parsed_count, initial_count)

    return results


def get_column_counts(data, column):
    counts = defaultdict(int)
    for row in data:
        val = row.get(column)
        if val is not None:
            counts[val] += 1
    return dict(counts)


def plot_pie_chart(counts, title):
    fig = px.pie(
        names=list(counts.keys()),
        values=list(counts.values()),
        title=title,
        hole=0.3,
    )
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(margin=dict(t=50, b=20, l=20, r=20))
    return fig


def plot_bar_chart(counts, title, horizontal=False):
    labels = list(counts.keys())
    values = list(counts.values())
    if horizontal:
        fig = px.bar(x=values, y=labels, orientation="h", title=title)
    else:
        fig = px.bar(x=labels, y=values, title=title)
    fig.update_layout(margin=dict(t=50, b=80, l=40, r=20))
    return fig


@st.cache_data(show_spinner=False)
def highlight_multi_file_differences_cached(file_items, comparison_mode="Exact inline word diff", reference_file=None):
    if len(file_items) < 2:
        return "Select at least two files to compare."

    files = [fname for fname, _ in file_items]
    if reference_file is None or reference_file not in files:
        reference_file = files[0]

    css = """
    <style>
        body { font-family: Arial; margin: 20px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid black; padding: 4px; vertical-align: top; white-space: pre-wrap; }
        th { background-color: #f0f0f0; }
        td.line-number { background-color: #f0f0f0; font-weight: bold; text-align: center; }
        .match { background-color: #ccffcc; }
        .mismatch { background-color: #ffcccc; }
        .scrollable { overflow:auto; max-height:800px; }
        p.legend span { display:inline-block; width:20px; height:20px; margin-right:5px; vertical-align:middle; }
    </style>
    """
    html_parts = [
        "<html><head>", css, "</head><body><div class='scrollable'>",
        "<p class='legend'><b>Legend:</b> <span class='match'></span> Matched word, <span class='mismatch'></span> Different or missing word</p>",
        "<table><tr><th>Line #</th>",
        "".join(f"<th>{html.escape(fname)}</th>" for fname in files),
        "</tr>",
    ]

    file_lines = {fname: text.splitlines() for fname, text in file_items}
    max_lines = max(len(lines) for lines in file_lines.values())

    for i in range(max_lines):
        html_parts.append(f"<tr><td class='line-number'>{i + 1}</td>")

        line_word_lists = {}
        ordered_words = []
        word_presence = defaultdict(int)

        for fname in files:
            raw_line = file_lines[fname][i] if i < len(file_lines[fname]) else ""
            words = raw_line.split()
            line_word_lists[fname] = words
            for word in words:
                if word not in ordered_words:
                    ordered_words.append(word)
            for word in set(words):
                word_presence[word] += 1

        reference_words = line_word_lists.get(reference_file, [])

        for fname in files:
            words = line_word_lists[fname]
            if comparison_mode == "Word presence summary":
                highlighted = []
                word_set = set(words)
                for word in ordered_words:
                    escaped_word = html.escape(word)
                    if word in word_set and word_presence[word] == len(files):
                        highlighted.append(f"<span class='match'>{escaped_word}</span>")
                    else:
                        highlighted.append(f"<span class='mismatch'>{escaped_word}</span>")
                cell_html = ' '.join(highlighted) if highlighted else '&nbsp;'
            else:
                highlighted = []
                matcher = SequenceMatcher(None, reference_words, words)
                for tag, i1, i2, j1, j2 in matcher.get_opcodes():
                    if tag == "equal":
                        highlighted.extend(f"<span class='match'>{html.escape(w)}</span>" for w in words[j1:j2])
                    else:
                        highlighted.extend(f"<span class='mismatch'>{html.escape(w)}</span>" for w in words[j1:j2])
                cell_html = ' '.join(highlighted) if highlighted else '&nbsp;'
            html_parts.append(f"<td>{cell_html}</td>")

        html_parts.append("</tr>")

    html_parts.append("</table></div></body></html>")
    return "".join(html_parts)


def highlight_side_by_side_differences_cached(file_items, reference_file=None):
    files = [fname for fname, _ in file_items]
    if len(files) < 2:
        return "Select at least two files to compare."
    if reference_file is None or reference_file not in files:
        reference_file = files[0]

    file_lines = {fname: text.splitlines() for fname, text in file_items}
    max_lines = max(len(lines) for lines in file_lines.values())

    css = """
    <style>
        body { font-family: Arial; margin: 20px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid black; padding: 4px; vertical-align: top; white-space: pre-wrap; }
        th { background-color: #f0f0f0; }
        td.line-number { background-color: #f0f0f0; font-weight: bold; text-align: center; }
        .line-match { background-color: #ccffcc; display: block; width: 100%; }
        .line-mismatch { background-color: #ffcccc; display: block; width: 100%; }
        .scrollable { overflow:auto; max-height:800px; }
        p.legend span { display:inline-block; width:20px; height:20px; margin-right:5px; vertical-align:middle; }
    </style>
    """
    html_parts = [
        "<html><head>", css, "</head><body><div class='scrollable'>",
        "<p class='legend'><b>Legend:</b> <span class='line-match'></span> Same as reference line, <span class='line-mismatch'></span> Different from reference or missing line</p>",
        "<p><b>Reference file:</b> " + html.escape(reference_file) + "</p>",
        "<table><tr><th>Line #</th>",
        "".join(f"<th>{html.escape(fname)}</th>" for fname in files),
        "</tr>",
    ]

    for i in range(max_lines):
        html_parts.append(f"<tr><td class='line-number'>{i + 1}</td>")
        reference_line = file_lines[reference_file][i] if i < len(file_lines[reference_file]) else ""
        for fname in files:
            line_text = file_lines[fname][i] if i < len(file_lines[fname]) else ""
            if line_text == reference_line and line_text != "":
                cell_html = f"<span class='line-match'>{html.escape(line_text)}</span>"
            elif line_text == reference_line == "":
                cell_html = "&nbsp;"
            else:
                cell_html = f"<span class='line-mismatch'>{html.escape(line_text)}</span>"
            html_parts.append(f"<td>{cell_html}</td>")
        html_parts.append("</tr>")

    html_parts.append("</table></div></body></html>")
    return "".join(html_parts)


def highlight_multi_file_differences(file_texts, comparison_mode="Exact inline word diff", reference_file=None):
    if comparison_mode == "Side-by-side line diff":
        return highlight_side_by_side_differences_cached(
            tuple((fname, str(text)) for fname, text in file_texts.items()),
            reference_file=reference_file
        )
    return highlight_multi_file_differences_cached(
        tuple((fname, str(text)) for fname, text in file_texts.items()),
        comparison_mode=comparison_mode,
        reference_file=reference_file
    )


def build_semantic_diff_explanation(file_texts):
    """Explain meaning-level changes beside the visual word/line diff."""
    if not file_texts or len(file_texts) < 2:
        return "Select at least two files to generate a semantic difference explanation."

    file_names = list(file_texts.keys())
    per_file = {}
    for file_name, text in file_texts.items():
        text = str(text or "")
        per_file[file_name] = {
            "themes": set(extract_key_themes(text, limit=14)),
            "entities": set(extract_entities(text, limit=20)),
            "risks": set(extract_risk_signals(text, limit=10)),
            "length": len(text),
        }

    baseline = file_names[0]
    base = per_file[baseline]
    sections = [f"### Semantic Diff Explanation\nBaseline: **{html.escape(baseline)}**"]
    for file_name in file_names[1:]:
        current = per_file[file_name]
        added_themes = sorted(current["themes"] - base["themes"])[:8]
        removed_themes = sorted(base["themes"] - current["themes"])[:8]
        added_entities = sorted(current["entities"] - base["entities"])[:8]
        removed_entities = sorted(base["entities"] - current["entities"])[:8]
        added_risks = sorted(current["risks"] - base["risks"])[:5]
        delta = current["length"] - base["length"]
        delta_label = "expanded" if delta > 0 else "contracted" if delta < 0 else "unchanged in size"

        sections.append(f"#### {html.escape(file_name)}")
        sections.append(f"- Structural signal: content {delta_label} by {abs(delta):,} extracted characters.")
        sections.append(f"- New themes: {html.escape(', '.join(added_themes) if added_themes else 'None detected')}.")
        sections.append(f"- Missing themes: {html.escape(', '.join(removed_themes) if removed_themes else 'None detected')}.")
        sections.append(f"- New entities: {html.escape(', '.join(added_entities) if added_entities else 'None detected')}.")
        sections.append(f"- Missing entities: {html.escape(', '.join(removed_entities) if removed_entities else 'None detected')}.")
        if added_risks:
            sections.append("- Risk/signals introduced:")
            sections.extend(f"  - {html.escape(str(risk)[:220])}" for risk in added_risks)
        else:
            sections.append("- Risk/signals introduced: None detected.")

    return "\n".join(sections)


@st.cache_data(show_spinner=False)
def generate_word_level_comparison_excel_cached(file_items):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Comparison"

    files = [fname for fname, _ in file_items]
    file_texts = {fname: text for fname, text in file_items}
    ws.append(["Line #"] + files)
    file_lines = {f: [l.split() for l in t.splitlines()] for f, t in file_texts.items()}

    red_fill = PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid")
    green_fill = PatternFill(start_color="CCFFCC", end_color="CCFFCC", fill_type="solid")

    max_lines = max(len(l) for l in file_lines.values())

    for i in range(max_lines):
        max_words = max(len(file_lines[f][i]) if i < len(file_lines[f]) else 0 for f in files)
        for w_idx in range(max_words):
            row_values = [i + 1 if w_idx == 0 else ""]
            for f in files:
                line_words = file_lines[f][i] if i < len(file_lines[f]) else []
                word = line_words[w_idx] if w_idx < len(line_words) else ""
                row_values.append(word)
            ws.append(row_values)

            # Highlight exact matches in green and missing/mismatched content in red
            all_words_set = set()
            for f in files:
                if i < len(file_lines[f]):
                    all_words_set.update(file_lines[f][i])
            for col_idx, f in enumerate(files, start=2):
                cell = ws.cell(row=ws.max_row, column=col_idx)
                line_words = file_lines[f][i] if i < len(file_lines[f]) else []
                if w_idx >= len(line_words):
                    cell.fill = red_fill
                elif all(word == line_words[w_idx] for other_file in files
                         for word in ([file_lines[other_file][i][w_idx]]
                                      if i < len(file_lines[other_file]) and w_idx < len(file_lines[other_file][i])
                                      else ["__missing__"])):
                    cell.fill = green_fill
                else:
                    cell.fill = red_fill

    excel_io = BytesIO()
    wb.save(excel_io)
    return excel_io.getvalue()


def generate_word_level_comparison_excel(file_texts):
    excel_io = BytesIO(generate_word_level_comparison_excel_cached(tuple((fname, str(text)) for fname, text in file_texts.items())))
    excel_io.seek(0)
    return excel_io


@st.cache_data(show_spinner=False)
def analyze_capl_code_with_suggestions_cached(code):
    issues = []

    brace_stack = []
    declared_vars = []
    used_vars = []

    lines = code.splitlines()

    for i, line in enumerate(lines, 1):
        stripped = line.strip()

        if not stripped or stripped.startswith("//"):
            continue  # Skip empty lines/comments

        # Track braces
        for c in stripped:
            if c == "{":
                brace_stack.append(i)
            elif c == "}":
                if brace_stack:
                    brace_stack.pop()
                else:
                    issues.append({
                        "line": i,
                        "error": "Unmatched closing brace",
                        "suggestion": "Remove or match with an opening '{'"
                    })

        # Detect variable declarations
        var_match = re.match(r'\b(int|float|byte|char|mstimer|timer|enum)\b\s+(\w+)', stripped)
        if var_match:
            declared_vars.append(var_match.group(2))

        # Track all used variable names
        used_vars += re.findall(r'\b([a-zA-Z_]\w*)\b', stripped)

        # Check for case sensitivity in keywords
        if re.search(r'\b(If|Else|For|While|Switch|Case|Return|On|Variables|Includes|Enum|Mstimer|Timer)\b', stripped):
            issues.append({
                "line": i,
                "error": "CAPL keywords should be lowercase",
                "suggestion": "Use lowercase keywords like 'if', 'else', 'on', etc."
            })

        # Check for incomplete if conditions
        if re.match(r'^\s*(if|else if)\s*\(', stripped) and not re.search(r'\)\s*(\{)?\s*$', stripped):
            issues.append({
                "line": i,
                "error": "Incomplete if condition",
                "suggestion": "Add closing parenthesis ')' and possibly opening brace '{'"
            })

        # Check for missing opening brace after control statements
        if re.match(r'^\s*(if|else if|else|for|while|switch)\b', stripped) and not stripped.endswith(
                '{') and not re.search(r'\)\s*\{', stripped):
            # Check if next line starts with '{'
            if i < len(lines) and not lines[i].strip().startswith('{'):
                issues.append({
                    "line": i,
                    "error": "Missing opening brace after control statement",
                    "suggestion": "Add '{' after the condition or on the next line"
                })

        # Detect missing semicolon
        if not stripped.endswith(";") and not stripped.endswith("{") and not stripped.endswith("}"):
            if not re.match(r'^(on|variables|includes|enum|mstimer|timer|if|else|switch|case|for|while|return)\b',
                            stripped):
                issues.append({
                    "line": i,
                    "error": "Missing semicolon",
                    "suggestion": "Add ';' at the end of this line"
                })

    # Check unmatched opening braces
    for open_line in brace_stack:
        issues.append({
            "line": open_line,
            "error": "Unmatched opening brace",
            "suggestion": "Add closing '}' to match this '{'"
        })

    # Check for 'on message' presence
    if "on message" not in code.lower():
        issues.append({
            "line": None,
            "error": "No 'on message' handler found",
            "suggestion": "Add an 'on message' event handler as required"
        })

    # Check for unused declared variables
    for var in declared_vars:
        if var not in used_vars:
            issues.append({
                "line": None,
                "error": f"Unused variable: {var}",
                "suggestion": "Consider removing this variable or using it in the code"
            })

    # Detect undeclared variables starting with PT4_ or $PT4_ used in code
    for i, line in enumerate(lines, 1):
        pt4_vars = re.findall(r'\b(PT4_[a-zA-Z_]\w*|\$PT4_[a-zA-Z_]\w*)\b', line)
        for var in pt4_vars:
            if var not in declared_vars and not var.startswith("$"):
                issues.append({
                    "line": i,
                    "error": f"Undeclared variable used: {var}",
                    "suggestion": f"Declare '{var}' in the variables section before using it"
                })

    return issues


def analyze_capl_code_with_suggestions(code):
    return analyze_capl_code_with_suggestions_cached(code)


def is_capl_code(text):
    """Check if the given text contains CAPL-specific keywords or syntax."""
    capl_keywords = [
        "on message", "variables", "includes", "mstimer", "timer", "byte", "char", "int", "float",
        "enum", "if", "else", "switch", "case", "for", "while", "return", "write", "output",
        "setTimer", "cancelTimer", "getTimer", "putValue", "getValue", "testcase", "teststep"
    ]
    text_lower = text.lower()
    return any(keyword in text_lower for keyword in capl_keywords)


@st.cache_data(show_spinner=False)
def render_capl_code_with_highlights_cached(code, issues_key):
    """Render CAPL code with IDE-like line highlighting for detected issues."""
    issues = [
        {"line": line, "error": error, "suggestion": suggestion}
        for line, error, suggestion in issues_key
    ]
    issue_lines = defaultdict(list)

    for issue in issues:
        line_no = issue.get("line")
        if isinstance(line_no, int):
            issue_lines[line_no].append(issue.get("error", "Issue detected"))

    code_lines = code.splitlines() or [""]
    rendered_lines = []

    for line_no, line in enumerate(code_lines, 1):
        escaped_line = html.escape(line) if line else "&nbsp;"
        line_classes = ["capl-line"]
        if line_no in issue_lines:
            line_classes.append("capl-line-error")

        tooltip = " | ".join(issue_lines[line_no]) if line_no in issue_lines else ""
        title_attr = f' title="{html.escape(tooltip)}"' if tooltip else ""

        rendered_lines.append(
            f"<div class=\"{' '.join(line_classes)}\"{title_attr}>"
            f"<span class=\"capl-gutter\">{line_no:>4}</span>"
            f"<span class=\"capl-code-text\">{escaped_line}</span>"
            f"</div>"
        )

    code_html = """
    <style>
        .capl-code-block {
            background: #0f172a;
            border: 1px solid #cbd5e1;
            border-radius: 10px;
            font-family: Consolas, "Courier New", monospace;
            font-size: 14px;
            line-height: 1.5;
            max-height: 420px;
            overflow: auto;
            padding: 12px 0;
        }
        .capl-line {
            color: #e2e8f0;
            display: flex;
            white-space: pre;
        }
        .capl-line-error {
            background: rgba(239, 68, 68, 0.22);
            border-left: 4px solid #ef4444;
        }
        .capl-gutter {
            color: #94a3b8;
            display: inline-block;
            min-width: 52px;
            padding: 0 12px;
            text-align: right;
            user-select: none;
        }
        .capl-code-text {
            display: inline-block;
            padding: 0 16px 0 0;
            width: 100%;
        }
    </style>
    """
    return code_html + f"<div class='capl-code-block'>{''.join(rendered_lines)}</div>"


def render_capl_code_with_highlights(code, issues=None):
    issues_key = tuple(
        (
            issue.get("line"),
            issue.get("error", "Issue detected"),
            issue.get("suggestion", "")
        )
        for issue in (issues or [])
    )
    return render_capl_code_with_highlights_cached(code, issues_key)


def render_capl_issue_table(issues):
    if not issues:
        st.success("✅ No issues detected!")
        return

    df_issues = pd.DataFrame(issues).fillna("-")
    st.dataframe(df_issues, use_container_width=True, hide_index=True)


def get_combined_vector_store(file_names):
    """Get vector store with intelligent caching to avoid redundant processing"""
    ensure_files_processed(file_names)
    selection_key = get_selection_signature(file_names)
    
    # Check cache first
    cached_vs = VECTOR_STORE_CACHE.get(selection_key)
    if cached_vs is not None:
        st.session_state.vector_stores[selection_key] = cached_vs
        return cached_vs
    
    # Create vector store if not cached
    if selection_key not in st.session_state.vector_stores:
        combined_text = "\n".join(st.session_state.file_texts.get(file_name, "") for file_name in file_names)
        vs = create_vector_store(combined_text)
        st.session_state.vector_stores[selection_key] = vs
        VECTOR_STORE_CACHE.set(selection_key, vs)
    return st.session_state.vector_stores[selection_key]


def show_current_sidebar_selection():
    selected = st.session_state.get("selected_files", [])
    if selected:
        st.info("Sidebar selected files: " + ", ".join(selected))
    else:
        st.info("No sidebar files selected yet. Upload and select files from the sidebar first.")


def render_file_context_card(title, available_files, active_files=None):
    active_files = active_files or []
    chips_html = "".join(
        f"<span class='file-chip'>{html.escape(file_name)}</span>"
        for file_name in active_files[:12]
    )
    if len(active_files) > 12:
        chips_html += f"<span class='file-chip'>+{len(active_files) - 12} more</span>"

    st.markdown(
        f"""
        <div class="app-card">
            <h4>{html.escape(title)}</h4>
            <p class="app-muted">Available from sidebar: {len(available_files)} file(s)</p>
            <p class="app-muted">Selected in this tab: {len(active_files)} file(s)</p>
            <div class="file-chip-wrap">
                {chips_html if chips_html else "<span class='file-chip'>No tab files selected yet</span>"}
            </div>
        </div>
        """,
        unsafe_allow_html=True
    )


def render_autonomous_workspace_shell():
    """Render the connected AI operating system layer above module controls."""
    memory = normalize_workspace_memory(st.session_state.workspace_memory)
    indexed_count = len(memory.get("indexed_files", []))
    chat_count = len(memory.get("chat", []))
    agent_count = len(memory.get("agent_runs", []))
    event_count = len(memory.get("memory_events", []))
    st.markdown(
        f"""
        <div class="ai-os-shell">
            <div class="ai-os-kicker">Autonomous AI Operating System</div>
            <div class="ai-os-title">One shared AI brain across Chat, Upload, Dashboard, Compare, and CAPL.</div>
            <div class="ai-os-loop">
                <span>Chat -> Memory</span>
                <span>Upload -> FAISS</span>
                <span>Dashboard -> Insights</span>
                <span>Compare -> Semantic Diff</span>
                <span>CAPL -> Agents</span>
            </div>
            <div class="ai-os-metrics">
                <span>{indexed_count} indexed files</span>
                <span>{chat_count} chat memories</span>
                <span>{agent_count} agent runs</span>
                <span>{event_count} memory events</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def track_user_behavior(tab_name):
    """Tracks user actions to detect skill level progression."""
    if "behavior_tracker" not in st.session_state:
        st.session_state.behavior_tracker = {
            "chat": {"queries": 0, "actions": []},
            "dashboard": {"queries": 0, "actions": []},
            "compare": {"queries": 0, "actions": []},
            "capl": {"queries": 0, "actions": []}
        }
    
    if tab_name not in st.session_state.behavior_tracker:
        st.session_state.behavior_tracker[tab_name] = {"queries": 0, "actions": []}
    
    st.session_state.behavior_tracker[tab_name]["queries"] += 1
    tracker = st.session_state.behavior_tracker[tab_name]
    return tracker


def infer_user_workflow():
    """Auto-detects user skill level from query patterns."""
    if "behavior_tracker" not in st.session_state:
        return "beginner"
    
    total_queries = sum(t.get("queries", 0) for t in st.session_state.behavior_tracker.values())
    
    if total_queries > 15:
        return "advanced"
    elif total_queries > 5:
        return "intermediate"
    else:
        return "beginner"


def get_dynamic_suggestions(tab_name, skill_level):
    """Returns context-aware suggestions based on skill level."""
    suggestions_by_skill = {
        "chat": {
            "beginner": ["Analyze", "Summary", "Ask document question", "Review confidence"],
            "intermediate": ["Use agent routing", "Ask follow-up with memory", "Find exact evidence", "Compare selected docs"],
            "advanced": ["Validate citations", "Use cached summaries", "Audit confidence", "Cross-document reasoning"]
        },
        "dashboard": {
            "beginner": ["Select HTML/XLSX report", "Choose chart type", "Review metrics"],
            "intermediate": ["Inspect report tables", "Switch orientation", "Review selected source"],
            "advanced": ["Validate dashboard source", "Correlate report metrics", "Reset and re-scope"]
        },
        "compare": {
            "beginner": ["Exact inline diff", "Side-by-side line diff", "Select two files"],
            "intermediate": ["Word presence summary", "Download Excel diff", "Review semantic summary"],
            "advanced": ["Multi-file comparison", "Change impact analysis", "Validate changed sections", "Use chat comparison"]
        },
        "capl": {
            "beginner": ["Analyze CAPL syntax", "Review issue table", "Select a .can file"],
            "intermediate": ["Generate AI fix", "Run CAPL goal", "Inspect unused variables"],
            "advanced": ["Goal-driven CAPL run", "Review run history", "Validate generated fixes", "Coordinate final output"]
        }
    }
    
    return suggestions_by_skill.get(tab_name, {}).get(skill_level, [])


def get_next_best_action(tab_name, skill_level):
    """Intelligently recommends the next workflow step."""
    workflow_paths = {
        "chat": {
            "beginner": "Pro Tip: Use Analyze for full-document analysis, Summary for a concise executive summary, and normal questions for grounded answers.",
            "intermediate": "Next: Ask naturally and the chat assistant will choose the right style for analysis, summary, tables, components, diagrams, comparison, or exact search.",
            "advanced": "Next: Validate confidence, citations, cached summaries, and exact find/count evidence before using the answer downstream."
        },
        "dashboard": {
            "beginner": "Pro Tip: Select an HTML/HTM/XLSX report from Uploaded files, then choose it in Dashboard.",
            "intermediate": "Next: Use chart type and orientation controls to inspect structured report distributions.",
            "advanced": "Next: Cross-check dashboard findings with Chat or Compare when you need evidence, explanation, or file-to-file validation."
        },
        "compare": {
            "beginner": "Pro Tip: Select at least two files, then start with exact inline word diff.",
            "intermediate": "Next: Switch modes to line diff or word presence summary, then export the Excel workbook.",
            "advanced": "Next: Use the semantic summary to capture comparison findings into shared memory."
        },
        "capl": {
            "beginner": "Pro Tip: Select a CAPL file, run analysis, and review line-level issues first.",
            "intermediate": "Next: Use AI fix suggestions or run a focused autonomous CAPL goal.",
            "advanced": "Next: Let the planning, retrieval, execution, reasoning, and coordination agents work across selected files."
        }
    }
    
    return workflow_paths.get(tab_name, {}).get(skill_level, "Keep exploring the features available.")


def show_help_popup(tab_name, selected_files):
    state_key = ensure_help_popup_state(tab_name)

    if not st.session_state[state_key]:
        return

    tracker = track_user_behavior(tab_name)
    skill_level = infer_user_workflow()
    selected_types = {os.path.splitext(f)[1].lower() for f in (selected_files or [])}
    selected_types_text = ", ".join(sorted(selected_types)) if selected_types else "No files selected"
    selected_type_names = ", ".join(sorted(ext.lstrip('.') for ext in selected_types if ext)) or "none"
    selected_file_count = len(selected_files or [])

    support_hint_msg = "Select files from the sidebar to use this feature."
    if selected_file_count == 0:
        support_hint_msg = "No files selected yet. Choose files from the sidebar to get started."
    elif tab_name == "dashboard":
        if not selected_types.intersection({".html", ".htm", ".xlsx"}):
            support_hint_msg = "Dashboard works only with HTML, HTM, or XLSX report files. Select one of these formats."
        else:
            support_hint_msg = f"Ready to analyze {selected_file_count} selected file(s) ({selected_type_names})."
    elif tab_name == "compare":
        if selected_file_count < 2:
            support_hint_msg = "Select at least two files from the sidebar to compare them."
        else:
            support_hint_msg = f"Ready to compare {selected_file_count} selected files ({selected_type_names})."
    elif tab_name == "capl":
        if not selected_types.intersection({".can", ".txt"}):
            support_hint_msg = "CAPL analysis requires a .can or .txt CAPL source file. Select one of those file types."
        else:
            support_hint_msg = f"Ready to analyze CAPL files ({selected_type_names})."
    elif tab_name == "chat":
        support_hint_msg = f"Ready for grounded multi-format document chat over {selected_file_count} selected file(s) ({selected_type_names})."

    helper_defs = {
        "chat": {
            "title": "Chat Agent Helper",
            "text": "Use the document chat assistant to get natural summaries, deeper analysis, overviews, comparisons, table answers, diagram explanations, exact search, and grounded Q&A.",
            "hint": "Ask for Summary, Analyze, or Overview when you want a reader-friendly synthesis. Every answer includes sources and a confidence label so you can judge how complete the context was.",
            "workflow": [
                "Upload files in the sidebar, then explicitly select only the files you want available in Chat.",
                "Ask naturally for analysis, summary, overview, question answering, tables, components, diagrams, comparison, or search.",
                "Use Analyze for deeper reasoning about purpose, structure, strengths, risks, design patterns, and implications.",
                "Use Summary for a concise executive summary.",
                "Ask specific questions when you need targeted answers from the selected documents.",
                "Ask for tables, components, diagrams, comparisons, or searches when you need those specific outputs.",
                "Review Sources and Confidence. Sources show the files, pages, slides, sheets, or sections used for the answer.",
                "Follow-up questions reuse memory for the same user and document selection.",
                "Use find \"phrase\" or count \"phrase\" when you need exact text evidence."
            ],
            "outputs": ["Full analysis", "Executive summary", "Grounded answers", "Citations", "Confidence", "Per-document memory"],
            "shortcuts": [
                ("Analyze", "Create a deeper expert analysis."),
                ("Summary", "Return a concise executive summary."),
                ("Overview", "Show a high-level explanation of the document."),
                ("table data", "Extract clean table-like content."),
                ("components", "Extract modules/components and roles."),
                ("find \"phrase\"", "Locate exact occurrences in the selected files."),
                ("count \"phrase\"", "Count exact matches in selected document text.")
            ]
        },
        "dashboard": {
            "title": "Dashboard Helper",
            "text": "Use Dashboard for structured HTML/HTM and XLSX report inspection with selectable chart types and focused source-file context.",
            "hint": "Only files selected in the sidebar appear here. Pick a dashboard-compatible file, then choose the chart type, orientation, sheet/column, or report section you want to inspect.",
            "workflow": [
                "Select an HTML, HTM, or XLSX file from the sidebar Uploaded files list.",
                "Choose the active dashboard file in this tab so Dashboard state stays independent from Chat, Compare, and CAPL.",
                "For XLSX files, select a column and chart type to inspect value distributions.",
                "For HTML/HTM reports, review extracted login/statistics/test-result sections and generated charts.",
                "Use Reset to clear Dashboard-specific file and chart selections without clearing uploaded documents."
            ],
            "outputs": ["Excel distributions", "HTML report metrics", "Interactive charts", "Fixture summaries", "Dashboard file context"],
            "shortcuts": [
                ("Select report", "Load the active HTML/HTM/XLSX file for analysis."),
                ("Choose chart", "Switch chart types or orientations."),
                ("Reset", "Clear dashboard filters and selections.")
            ]
        },
        "compare": {
            "title": "Compare Helper",
            "text": "Use Compare for file-to-file review: exact inline diffs, side-by-side line changes, word presence differences, semantic summaries, and Excel export.",
            "hint": "Select at least two sidebar files, choose only the files you want in Compare, run the diff, then use Chat for explanation if you need a document-grounded comparison answer.",
            "workflow": [
                "Select two or more files from the sidebar and confirm them in this tab.",
                "Choose a comparison mode: exact inline diff, side-by-side line diff, or word presence summary.",
                "Run the compare action, inspect exact and semantic changes, and download the generated Excel workbook.",
                "For natural-language comparison questions, switch to Chat and ask a comparison query so the agent routes it to COMPARISON."
            ],
            "outputs": ["Inline diffs", "Side-by-side line comparison", "Word presence summaries", "Semantic summary", "Excel export"],
            "shortcuts": [
                ("Select 2+ files", "Add at least two files for comparison."),
                ("Compare mode", "Choose the best diff style for your review."),
                ("Download Excel", "Save comparison results for offline review.")
            ]
        },
        "capl": {
            "title": "CAPL Helper",
            "text": "Use CAPL to analyze CANoe/CANalyzer scripts, detect syntax and structure issues, review highlighted code, generate fixes, and run focused CAPL goals.",
            "hint": "Select a .can or .txt CAPL source file. CAPL state is independent from Chat, but selected CAPL files can still be discussed in Chat when you need document-grounded explanation.",
            "workflow": [
                "Select CAPL source files in the sidebar and choose one for analysis.",
                "Run CAPL analysis to detect issues, review live preview output, and inspect the generated issue table.",
                "Use the AI fix button for code suggestions or enter a goal to run the autonomous CAPL workflow.",
                "Validate generated fixes in CANoe or the official CAPL toolchain before production use."
            ],
            "outputs": ["Issue diagnosis", "Live code preview", "Suggested fixes", "Agent run summary", "Run history"],
            "shortcuts": [
                ("Analyze CAPL", "Scan the selected CAPL file for issues."),
                ("Suggest Fix", "Generate AI-based corrections for the code."),
                ("Run goal", "Execute an autonomous CAPL workflow.")
            ]
        }
    }

    helper_def = helper_defs.get(tab_name, helper_defs["chat"])
    suggestions = get_dynamic_suggestions(tab_name, skill_level)[:4]
    if not suggestions:
        suggestions = ["Review the documents", "Ask a question", "Request an overview", "Search for exact phrases"]
    next_action = get_next_best_action(tab_name, skill_level)
    modal_key = f"helper_modal_{tab_name}"
    helper_close_key = f"helper_close_{tab_name}"

    # Build shortcuts HTML
    shortcuts_html = ""
    for cmd, desc in helper_def["shortcuts"]:
        shortcuts_html += f'<div class="helper-shortcut"><code>{html.escape(cmd)}</code><span>{html.escape(desc)}</span></div>'

    st.markdown(
        f"""
        <style>
        .helper-modal-backdrop {{
            position: fixed;
            right: 18px;
            bottom: 18px;
            width: 0;
            height: 0;
            background: transparent;
            z-index: 99998;
            pointer-events: none;
        }}
        .st-key-{modal_key} {{
            position: fixed !important;
            right: clamp(12px, 2vw, 24px) !important;
            bottom: clamp(12px, 2vw, 24px) !important;
            top: auto !important;
            left: auto !important;
            transform: none !important;
            width: min(460px, calc(100vw - 24px)) !important;
            max-height: min(76vh, 620px) !important;
            overflow-y: auto !important;
            z-index: 99999 !important;
            padding: 0 !important;
            background: #ffffff !important;
            border: 1px solid rgba(15, 23, 42, 0.10) !important;
            border-radius: 14px !important;
            box-shadow: 0 18px 46px rgba(15, 23, 42, 0.22) !important;
        }}
        .st-key-{modal_key} > div {{
            padding: 0 !important;
        }}
        .st-key-{modal_key} [data-testid="stHorizontalBlock"]:first-of-type {{
            position: sticky !important;
            top: 0 !important;
            z-index: 1 !important;
            align-items: center !important;
            padding: 16px 18px 12px !important;
            background: rgba(255, 255, 255, 0.96) !important;
            border-bottom: 1px solid rgba(15, 23, 42, 0.08) !important;
            backdrop-filter: blur(10px);
        }}
        .st-key-{modal_key} h3 {{
            margin: 0 !important;
            color: #111827 !important;
            font-size: 1.08rem !important;
            letter-spacing: 0 !important;
        }}
        .st-key-{modal_key} h4 {{
            margin: 14px 0 6px !important;
            color: #1f2937 !important;
            font-size: 0.94rem !important;
        }}
        .st-key-{modal_key} p,
        .st-key-{modal_key} li {{
            color: #4b5563 !important;
            font-size: 0.92rem !important;
            line-height: 1.45 !important;
        }}
        .st-key-{modal_key} .helper-modal-body {{
            padding: 14px 18px 18px;
        }}
        .st-key-{modal_key} .helper-info-row {{
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin: 4px 0 10px;
        }}
        .st-key-{modal_key} .helper-chip {{
            display: inline-flex;
            align-items: center;
            border-radius: 999px;
            padding: 5px 9px;
            background: #f3f4f6;
            color: #374151;
            border: 1px solid rgba(15, 23, 42, 0.06);
            font-size: 0.78rem;
            font-weight: 650;
        }}
        .st-key-{modal_key} .helper-callout {{
            padding: 10px 12px;
            border-radius: 10px;
            background: #f8fbff;
            border: 1px solid #dbeafe;
            color: #1f3b57;
            font-size: 0.9rem;
            line-height: 1.45;
            margin: 8px 0 12px;
        }}
        .st-key-{modal_key} .helper-shortcut {{
            display: grid;
            grid-template-columns: 88px 1fr;
            gap: 8px;
            padding: 7px 0;
            border-bottom: 1px solid rgba(15, 23, 42, 0.06);
            color: #4b5563;
            font-size: 0.9rem;
        }}
        .st-key-{modal_key} .helper-shortcut code {{
            color: #111827;
            background: #f3f4f6;
            border-radius: 6px;
            padding: 2px 6px;
            font-size: 0.82rem;
        }}
        .st-key-{modal_key} .helper-suggestions {{
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 8px;
        }}
        .st-key-{modal_key} .helper-suggestions span {{
            background: #fff7ed;
            color: #9a3412;
            border: 1px solid #fed7aa;
            border-radius: 999px;
            padding: 6px 10px;
            font-size: 0.82rem;
            line-height: 1.3;
        }}
        .st-key-{helper_close_key} {{
            display: flex !important;
            justify-content: flex-end !important;
        }}
        .st-key-{helper_close_key} button {{
            width: 36px !important;
            height: 36px !important;
            min-height: 36px !important;
            padding: 0 !important;
            border: 1px solid rgba(15, 23, 42, 0.08) !important;
            border-radius: 50% !important;
            background: #f9fafb !important;
            color: #374151 !important;
            font-size: 1rem !important;
            line-height: 1 !important;
            box-shadow: none !important;
        }}
        .st-key-{helper_close_key} button:hover {{
            background: #eef2ff !important;
            color: #1d4ed8 !important;
        }}
        @media (max-width: 640px) {{
            .st-key-{modal_key} {{
                right: 10px !important;
                bottom: 10px !important;
                width: calc(100vw - 20px) !important;
                max-height: 72vh !important;
                border-radius: 12px !important;
            }}
            .st-key-{modal_key} [data-testid="stHorizontalBlock"]:first-of-type {{
                padding: 12px 14px 10px !important;
            }}
        }}
        </style>
        <div class="helper-modal-backdrop"></div>
        """,
        unsafe_allow_html=True,
    )

    with st.container(key=modal_key):
        header_col, close_col = st.columns([8, 1], vertical_alignment="center")
        with header_col:
            st.markdown(f"### 🧠 {html.escape(helper_def['title'])}")
        with close_col:
            if st.button("✕", key=helper_close_key, help="Close helper"):
                set_help_popup_state(tab_name, False)
                st.rerun()

        suggestion_tags = "".join(f"<span>{html.escape(s)}</span>" for s in suggestions)
        workflow_items = "".join(f"<li>{html.escape(item)}</li>" for item in helper_def.get("workflow", []))
        output_tags = "".join(f"<span>{html.escape(item)}</span>" for item in helper_def.get("outputs", []))
        render_html_frame(
            f"""
            <style>
                body {{
                    margin: 0;
                    background: transparent;
                    font-family: "Segoe UI", Tahoma, sans-serif;
                    color: #374151;
                }}
                .helper-modal-body {{
                    padding: 14px 18px 18px;
                }}
                .helper-info-row {{
                    display: flex;
                    flex-wrap: wrap;
                    gap: 8px;
                    margin: 4px 0 10px;
                }}
                .helper-chip {{
                    display: inline-flex;
                    align-items: center;
                    border-radius: 999px;
                    padding: 5px 9px;
                    background: #f3f4f6;
                    color: #374151;
                    border: 1px solid rgba(15, 23, 42, 0.06);
                    font-size: 0.78rem;
                    font-weight: 650;
                }}
                h4 {{
                    margin: 14px 0 6px;
                    color: #1f2937;
                    font-size: 0.94rem;
                }}
                li {{
                    color: #4b5563;
                    font-size: 0.92rem;
                    line-height: 1.45;
                    margin-bottom: 4px;
                }}
                .helper-callout {{
                    padding: 10px 12px;
                    border-radius: 10px;
                    background: #f8fbff;
                    border: 1px solid #dbeafe;
                    color: #1f3b57;
                    font-size: 0.9rem;
                    line-height: 1.45;
                    margin: 8px 0 12px;
                }}
                .helper-shortcut {{
                    display: grid;
                    grid-template-columns: 88px 1fr;
                    gap: 8px;
                    padding: 7px 0;
                    border-bottom: 1px solid rgba(15, 23, 42, 0.06);
                    color: #4b5563;
                    font-size: 0.9rem;
                }}
                .helper-shortcut code {{
                    color: #111827;
                    background: #f3f4f6;
                    border-radius: 6px;
                    padding: 2px 6px;
                    font-size: 0.82rem;
                }}
                .helper-suggestions,
                .helper-outputs {{
                    display: flex;
                    flex-wrap: wrap;
                    gap: 8px;
                    margin-top: 8px;
                }}
                .helper-suggestions span,
                .helper-outputs span {{
                    background: #fff7ed;
                    color: #9a3412;
                    border: 1px solid #fed7aa;
                    border-radius: 999px;
                    padding: 6px 10px;
                    font-size: 0.82rem;
                    line-height: 1.3;
                }}
                .helper-outputs span {{
                    background: #eef2ff;
                    color: #3730a3;
                    border-color: #c7d2fe;
                }}
            </style>
            <div class="helper-modal-body">
                <div class="helper-info-row">
                    <span class="helper-chip">Skill: {html.escape(skill_level.title())}</span>
                    <span class="helper-chip">Queries: {tracker.get('queries', 0)}</span>
                    <span class="helper-chip">Files: {html.escape(selected_types_text)}</span>
                </div>

                <h4>💡 Quick Guide</h4>
                <div class="helper-callout">{html.escape(helper_def['text'])}</div>

                <h4>Workflow</h4>
                <ul>
                    {workflow_items}
                </ul>

                <h4>⚡ Shortcuts</h4>
                {shortcuts_html}

                <h4>📌 Tips</h4>
                <ul>
                    <li>{html.escape(helper_def['hint'])}</li>
                    <li>{html.escape(next_action)}</li>
                    <li>{html.escape(support_hint_msg)}</li>
                    <li>Click suggestions for faster input when available.</li>
                </ul>

                <h4>Outputs</h4>
                <div class="helper-outputs">{output_tags}</div>

                <h4>Suggestions</h4>
                <div class="helper-suggestions">{suggestion_tags}</div>
            </div>
            """,
            height=390,
        )


def hex_to_rgb_values(hex_color):
    clean = str(hex_color or "#38bdf8").lstrip("#")
    if len(clean) != 6:
        clean = "38bdf8"
    try:
        return tuple(int(clean[index:index + 2], 16) for index in (0, 2, 4))
    except ValueError:
        return (56, 189, 248)


def ensure_tab_glow_colors(tab_options):
    """Assign each tab one random neon identity color and keep it across reruns."""
    neon_palette = [
        "#00E5FF", "#7C4DFF", "#FF4081", "#69F0AE", "#FFEA00", "#FF6D00",
        "#18FFFF", "#B388FF", "#F50057", "#64FFDA", "#40C4FF", "#EEFF41",
    ]
    existing_colors = st.session_state.get("tab_colors")
    if not isinstance(existing_colors, dict):
        existing_colors = {}

    assigned_colors = {
        tab_name: existing_colors[tab_name]
        for tab_name in tab_options
        if tab_name in existing_colors and existing_colors[tab_name]
    }
    used_colors = set(assigned_colors.values())
    available_colors = [color for color in neon_palette if color not in used_colors]
    random.shuffle(available_colors)

    for tab_name in tab_options:
        if tab_name in assigned_colors:
            continue
        if available_colors:
            assigned_colors[tab_name] = available_colors.pop()
        else:
            while True:
                generated_color = "#{:06X}".format(random.randint(0x3030A0, 0xFFFFFF))
                if generated_color not in used_colors:
                    assigned_colors[tab_name] = generated_color
                    break
        used_colors.add(assigned_colors[tab_name])

    st.session_state.tab_colors = assigned_colors
    return assigned_colors


# ==============================
# CONTEXT-AWARE TAB SUGGESTION ENGINE
# Moved into functions.py so app.py remains an orchestrator.
# Reads chat, recent behavior, and CAPL context, then suggests the best tab.
# ==============================
def ensure_context_memory():
    """Initialize lightweight context memory used by automatic tab suggestion."""
    default_memory = {
        "recent_messages": [],
        "recent_actions": [],
        "capl_context": {},
        "last_signature": "",
        "suggested_tab": None,
    }
    memory = st.session_state.get("context_memory")
    if not isinstance(memory, dict):
        memory = {}
    normalized = {**default_memory, **memory}
    if not isinstance(normalized.get("recent_messages"), list):
        normalized["recent_messages"] = []
    if not isinstance(normalized.get("recent_actions"), list):
        normalized["recent_actions"] = []
    if not isinstance(normalized.get("capl_context"), dict):
        normalized["capl_context"] = {}
    st.session_state.context_memory = normalized
    return normalized


def build_context_memory_snapshot():
    """Collect chat messages, recent actions, and CAPL state without heavy work."""
    memory = ensure_context_memory()

    message_texts = []
    for message in st.session_state.get("messages", [])[-8:]:
        if isinstance(message, dict):
            message_texts.append(str(message.get("content") or message.get("user") or message.get("assistant") or ""))
        else:
            message_texts.append(str(message))

    action_texts = []
    tracker = st.session_state.get("behavior_tracker", {})
    if isinstance(tracker, dict):
        for tab_data in tracker.values():
            if isinstance(tab_data, dict):
                action_texts.extend(str(action) for action in tab_data.get("actions", [])[-5:])

    capl_issues = st.session_state.get("capl_last_issues") or []
    selected_capl_file = st.session_state.get("selected_capl_file", "")
    if str(selected_capl_file).strip() == "--Select CAPL file--":
        selected_capl_file = ""
    capl_context = {
        "selected_file": selected_capl_file,
        "last_file": st.session_state.get("capl_last_analyzed_file", ""),
        "issue_count": f"{len(capl_issues)} CAPL issue(s)" if isinstance(capl_issues, list) and capl_issues else "",
        "goal": st.session_state.get("capl_autonomous_goal", ""),
        "agent_result": st.session_state.get("capl_agent_result", ""),
    }

    memory["recent_messages"] = message_texts[-8:]
    memory["recent_actions"] = action_texts[-12:]
    memory["capl_context"] = capl_context
    st.session_state.context_memory = memory
    return memory


def suggest_tab_from_context(tab_options):
    """Return the best tab label from simple keyword rules."""
    memory = build_context_memory_snapshot()
    context_text = " ".join(
        memory.get("recent_messages", [])
        + memory.get("recent_actions", [])
        + [str(value) for value in memory.get("capl_context", {}).values()]
    ).lower()

    tab_lookup = {tab_name.lower(): tab_name for tab_name in tab_options}
    chat_tab = next((tab for key, tab in tab_lookup.items() if "chat" in key), tab_options[0])
    dashboard_tab = next((tab for key, tab in tab_lookup.items() if "dashboard" in key), chat_tab)
    compare_tab = next((tab for key, tab in tab_lookup.items() if "compare" in key), chat_tab)
    capl_tab = next((tab for key, tab in tab_lookup.items() if "capl" in key), chat_tab)

    if any(keyword in context_text for keyword in ["error", "fix", "capl", "syntax", "compile", "debug"]):
        return capl_tab, context_text
    if any(keyword in context_text for keyword in ["compare", "difference", "differences", "diff", "changes"]):
        return compare_tab, context_text
    if any(keyword in context_text for keyword in ["overview", "summary", "dashboard", "chart", "metric", "statistics"]):
        return dashboard_tab, context_text
    return chat_tab, context_text


def apply_auto_tab_suggestion(tab_options):
    """Switch active_main_tab only when the context snapshot changes."""
    memory = ensure_context_memory()
    suggested_tab, context_text = suggest_tab_from_context(tab_options)
    signature = hashlib.md5(context_text.encode("utf-8", errors="ignore")).hexdigest() if context_text else ""

    if signature and signature != memory.get("last_signature"):
        memory["last_signature"] = signature
        memory["suggested_tab"] = suggested_tab
        st.session_state.context_memory = memory
        if suggested_tab in tab_options:
            st.session_state.active_main_tab = suggested_tab

    return st.session_state.get("active_main_tab", suggested_tab)
